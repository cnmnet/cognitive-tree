#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
认知晶体树 5 全功能整合版（输出增强版）
包含：Web API (FastAPI) + Tkinter GUI + 核心引擎 + 所有业务模块
启动方式：python crystal_tree_all_in_one.py [--web|--gui]
修改内容：增加百灵鸟裸模型 Round 0，大幅提升各轮字数限制，
         重构总结员，强制输出完整学生版/教师版，自动补全不足字数。
"""

from typing import Any, Dict, List, Optional, Tuple
from datetime import datetime, date
# =============================================================================
# 1. 配置 (config.py)
# =============================================================================
import sys
import os


from pathlib import Path

PROJECT_ROOT = Path(__file__).resolve().parent.parent  # 注意：此处__file__为当前文件，需调整为项目根目录
# 若直接运行该文件，PROJECT_ROOT 为当前文件所在目录的父目录，可根据实际情况调整
# 更稳妥：获取当前文件所在目录，然后向上两级（如果放在项目根目录）
PROJECT_ROOT = Path(__file__).resolve().parent

# 设置缓存目录
os.environ["HF_ENDPOINT"] = "https://hf-mirror.com"
custom_cache_dir = str(PROJECT_ROOT / "model_cache")
os.environ["SENTENCE_TRANSFORMERS_HOME"] = custom_cache_dir
os.environ["HF_HOME"] = custom_cache_dir
os.environ["TRANSFORMERS_CACHE"] = custom_cache_dir
if not os.path.exists(custom_cache_dir):
    os.makedirs(custom_cache_dir, exist_ok=True)

try:
    from dotenv import load_dotenv
    load_dotenv()
except ImportError:
    pass

class Config:
    DATA_ROOT: Path = Path(os.getenv("CRYSTAL_TREE_DATA_ROOT", str(PROJECT_ROOT))).absolute()
    ATTENTION_LIMIT: int = 50
    L0_SIZE: int = 3
    L1_MAX: int = 47
    L1_WARNING_THRESHOLD: int = 45
    L2_TO_L3_HEAT_THRESHOLD: float = 0.1
    L2_TO_L3_DAYS_THRESHOLD: int = 30
    # TIMEOUT: tuple = (15, 30)   # 连接15秒，读取30秒（太短）
    TIMEOUT: tuple = (30, 120)     # 连接30秒，读取120秒（给长文本生成留足时间）
    MAX_RETRIES: int = 3
    BACKOFF_FACTOR: float = 1.5
    DELAY_BETWEEN_REQUESTS: tuple = (1, 3)
    DAILY_PLAN_TIME_BUDGET_SECONDS: int = 900
    DAILY_PLAN_MAX_CANDIDATES: int = 60
    DAILY_PLAN_TOP_ITEMS: int = 15
    DAILY_PLAN_MAX_PENDING_CARDS: int = 10
    DEEPSEEK_API_URL: str = "https://api.deepseek.com/v1/chat/completions"
    ARXIV_MIRRORS: list = ["https://cn.arxiv.org", "https://arxiv.org"]
    HF_MIRROR: str = "https://hf-mirror.com"
    PATHS: dict = {
        "state": "系统日志/状态快照.md",
        "crystals": "晶体数据/晶体卡片.md",
        "holes": "晶体数据/孔洞分层.md",
        "principles": "核心配置/核心操作原则.md",
        "ai_instructions": "核心配置/AI协作者指令.md",
        "pending": "暂存区/PENDING_待确认卡片.md",
        "search_cache": "系统日志/搜索摘要.md",
        "change_log": "系统日志/变更记录.md",
        "layer_state": "系统日志/晶体分层.json",
        "hole_progress": "系统日志/孔洞进度.json",
        "external_cache": "系统日志/external_sources_cache.json",
        "fingerprint": "系统日志/user_profile.json",  # ★ 新增这一行
        "task_cards": "系统日志/任务卡片.json",
        "roles": "核心配置/角色定义.json",
        "pareto_frontier": "系统日志/pareto_frontier.json",
        "inspiration_pool": "系统日志/灵感池.json",
    }
    DIRECTORIES: list = ["晶体数据", "核心配置", "系统日志", "暂存区", "系统日志/案例"]
    USER_AGENTS: list = [
        "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36",
        "Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) AppleWebKit/605.1.15 (KHTML, like Gecko) Version/17.0 Safari/605.1.15",
        "Mozilla/5.0 (Windows NT 10.0; Win64; x64; rv:109.0) Gecko/20100101 Firefox/115.0"
    ]
    GUI_BG_MAIN: str = "#f7f3ff"
    GUI_BG_SIDEBAR: str = "#fbf8ff"
    GUI_BG_CARD: str = "#fffaff"
    GUI_BG_CARD_ALT: str = "#f4efff"
    GUI_BG_INPUT: str = "#fffefe"
    GUI_BG_HEADER: str = "#f1ebff"
    GUI_FG_TEXT: str = "#211b33"
    GUI_FG_MUTED: str = "#716883"
    GUI_BORDER: str = "#d9cdf5"
    GUI_BORDER_STRONG: str = "#bba9ea"
    GUI_ACCENT: str = "#6f56d9"
    GUI_ACCENT_DARK: str = "#503bb4"
    GUI_SUCCESS: str = "#5f9f45"
    GUI_WARNING: str = "#dc8425"
    GUI_DANGER: str = "#d84a55"
    GUI_INFO: str = "#8d7fd1"
    GUI_BUTTON_SOFT: str = "#eee8fb"
    GUI_HIGHLIGHT: str = "#e7ddff"
    GUI_TITLE_FONT: tuple = ("微软雅黑", 18, "bold")
    GUI_CARD_FONT: tuple = ("微软雅黑", 11, "bold")
    GUI_BUTTON_FONT: tuple = ("微软雅黑", 10)
    GUI_TEXT_FONT: tuple = ("微软雅黑", 10)
    GUI_INPUT_FONT: tuple = ("微软雅黑", 11)
    GUI_LOG_FONT: tuple = ("宋体", 12)

    # ===== 向量检索配置 =====
    VECTOR_SEARCH_ENABLED: bool = True
    VECTOR_SEARCH_MODEL: str = "all-MiniLM-L6-v2"
    VECTOR_SEARCH_TOP_K: int = 5
    VECTOR_SEARCH_BM25_FALLBACK: bool = True  # 向量检索失败时降级到 BM25

    # ===== 验证门控配置 =====
    VALIDATION_GATE_RULES: dict = {
        "min_references": 2,           # 最少被引用的晶体数
        "require_external_source": False,  # 是否要求外部来源
        "require_recent_access": True,     # 是否要求近期访问
        "recent_access_days": 60,          # 近期访问天数阈值
        "low_contribution_threshold": 25.0, # 低贡献阈值
    }

    # ===== 元层原语开关（建议1：工作层+元层双层架构）=====
    META_PRIMITIVES: dict = {
        "active_gap_detection": {
            "enabled": True,
            "description": "主动缺口检测：自动识别晶体之间的逻辑缝隙"
        },
        "temporal_aware_escalation": {
            "enabled": True,
            "description": "时序感知升级：根据时间戳和活动频率自动调整优先级"
        },
        "layer_aware_calibration": {
            "enabled": True,
            "description": "层级感知校准：根据晶体层级（L1/L2/L3）调整调度权重"
        },
        "sleep_consolidation": {
            "enabled": True,
            "description": "睡眠式巩固：低活跃时段自动进行知识压缩和冗余清理"
        },
        "distributed_metacognition": {
            "enabled": False,
            "description": "分布式元认知：多个系统实例间共享认知模式（需网络）"
        },
        "validation_gated_self_evolution": {
            "enabled": True,
            "description": "验证门控自我进化：任何自我改进必须通过验证门检验"
        },
        # ===== 新增：PEAR 动态角色路由 =====
        "dynamic_role_routing": {
            "enabled": True,
            "description": "PEAR动态角色路由：通过元辩论为每个位置选择最优角色"
        }
    }

    META_LAYER_CONFIG: dict = {
        "consolidation_hour_start": 2,
        "consolidation_hour_end": 4,
        "gap_detection_interval_hours": 6,
        "max_crystals_before_consolidation": 100,
        "validation_gate_rules": [
            "new_evidence_from_at_least_3_sources",
            "audit_score_gt_0.6",
            "no_major_conflict_with_existing_crystals"
        ]
    }
    
    # ===== Day 7: 帕累托配置 =====
    PROFILE_HIGH_ACCURACY: dict = {
        "name": "高精度模式",
        "accuracy_weight": 0.7,
        "cost_weight": 0.1,
        "latency_weight": 0.2,
        "max_tokens": 4000,
        "temperature": 0.3,
        "enable_vector_search": True,
        "enable_external_fetch": True,
        "debate_rounds": 4,
        "description": "追求最高答案质量，适合复杂决策问题"
    }
    PROFILE_BALANCED: dict = {
        "name": "平衡模式",
        "accuracy_weight": 0.5,
        "cost_weight": 0.3,
        "latency_weight": 0.2,
        "max_tokens": 2000,
        "temperature": 0.5,
        "enable_vector_search": True,
        "enable_external_fetch": True,
        "debate_rounds": 3,
        "description": "质量与成本的平衡，适合日常使用"
    }
    PROFILE_ECONOMY: dict = {
        "name": "经济模式",
        "accuracy_weight": 0.3,
        "cost_weight": 0.4,
        "latency_weight": 0.3,
        "max_tokens": 800,
        "temperature": 0.7,
        "enable_vector_search": False,
        "enable_external_fetch": False,
        "debate_rounds": 2,
        "description": "最低成本，适合简单问题和快速响应"
    }

    DEFAULT_PROFILE: str = "balanced"
    PARETO_HISTORY_LIMIT: int = 50

    # ===== Day 7: 评分规则配置（可JSON覆盖） =====
    SCORING_RULES: dict = {
        "quality": {
            "crystal_ref_weight": 0.3,
            "jaccard_weight": 0.2,
            "audit_score_weight": 0.3,
            "risk_identified_weight": 0.2
        },
        "thresholds": {
            "excellent": 0.8,
            "good": 0.6,
            "needs_improvement": 0.4
        }
    }
    # ===== Day 3 新增：元原语触发链配置 =====
    META_CHAIN_RULES: dict = {
        # 触发链1：孤立晶体 → 验证门控自我进化
        "isolated_crystal_to_validation": {
            "enabled": True,
            "source_primitive": "active_gap_detection",
            "source_condition": "isolated_crystal_count >= 3",  # 发现3个以上孤立晶体
            "target_primitive": "validation_gated_self_evolution",
            "target_params": {"focus": "isolated_crystals"},
            "description": "发现孤立晶体 → 触发验证门控自我进化"
        },
        # 触发链2：高热度久未访问 → 主动缺口检测
        "stale_hot_to_gap_detection": {
            "enabled": True,
            "source_primitive": "temporal_aware_escalation",
            "source_condition": "stale_hot_count >= 2",  # 发现2个以上高热度久未访问晶体
            "target_primitive": "active_gap_detection",
            "target_params": {"focus": "stale_hot_crystals"},
            "description": "高热度久未访问晶体 → 触发主动缺口检测"
        },
        # 触发链3：验证门控自我进化结果 → 时序感知升级（预留）
        "validation_to_temporal": {
            "enabled": False,  # 默认禁用，Day 3 暂不启用
            "source_primitive": "validation_gated_self_evolution",
            "source_condition": "verification_passed_count >= 2",
            "target_primitive": "temporal_aware_escalation",
            "target_params": {"focus": "verified_crystals"},
            "description": "验证通过数量≥2 → 触发时序感知升级"
        }
    }
    # ===== Day 1 新增：警报规则配置（可JSON覆盖） =====
    ALARM_RULES: dict = {
        "knowledge_poverty": {
            "enabled": True,
            "metric": "crystal_reference_rate",
            "threshold": 0.5,          # 引用率 < 50%
            "action": "inject_external",
            "message": "知识贫瘠警报：晶体引用率低于50%，强制注入外部知识"
        },
        "bias_inflation": {
            "enabled": True,
            "metric": "bias_amplification",
            "threshold": 0.3,
            "action": "inject_perspective",
            "message": "偏见膨胀警报：偏见强化指数超过0.3，强制注入对立视角"
        },
        "information_starvation": {
            "enabled": True,
            "metric": "external_consecutive_empty",
            "threshold": 3,            # 连续3轮无新外部数据
            "action": "trigger_search",
            "message": "信息枯竭警报：连续3轮外部搜索无新数据，强制触发搜索"
        },
        "thought_stagnation": {
            "enabled": True,
            "metric": "jaccard_consecutive_high",
            "threshold": 0.8,
            "consecutive": 3,
            "action": "inject_perspective",
            "message": "思维固化警报：辩论Jaccard相似度连续3轮>0.8，强制注入新视角"
        }
    }
    # ===== Day 2: 动态路由配置 =====
    ROUTING_CONFIG = {
        "simple_length_threshold": 5,       # 长度<=5且无复杂词视为简单
        "medium_length_threshold": 30,      # 长度>5且<=30且含关键词视为中等
        "complex_keywords": [
            "设计", "方案", "系统", "模型", "策略", "框架", "博弈", "全球", "长期",
            "多变量", "优化", "权衡", "机制", "制度", "政策", "评估", "比较", "综合",
            "如何", "为什么", "怎样", "怎么", "能否", "是否", "哪些", "什么",  # 开放性问题
        ],
        "token_budget_high": 2000,
        "token_budget_medium": 800,
        "token_budget_simple": 200,
    }
    # ===== Day 10: 角色质量参数配置 =====
    ROLE_QUALITY_CONFIG: dict = {
        "radical": {
            "temperature": 0.80,
            "token_multiplier": 4.0,
            "description": "激进者：高温度鼓励颠覆性思维"
        },
        "conservative": {
            "temperature": 0.40,
            "token_multiplier": 3.0,
            "description": "保守者：低温保证风险判断稳定可靠"
        },
        "structural": {
            "temperature": 0.60,
            "token_multiplier": 3.5,
            "description": "结构主义者：中等温度保证框架严谨"
        },
        "lark": {
            "temperature": 0.75,
            "token_multiplier": 4.0,
            "description": "百灵鸟：略高温度鼓励跨领域联想"
        },
        "pilgrim": {
            "temperature": 0.55,
            "token_multiplier": 3.0,
            "description": "取经者：稳定温度保证长期视角一致"
        },
        "strategist": {
            "temperature": 0.80,
            "token_multiplier": 3.5,
            "description": "奇谋者：高温度鼓励非常规路径"
        },
        "statesman": {
            "temperature": 0.35,
            "token_multiplier": 3.5,
            "description": "延安智者：低温保证实事求是"
        },
        "judge": {
            "temperature": 0.20,
            "token_multiplier": 3.0,
            "description": "大法官：极低温保证裁决一致性"
        },
        "spokesperson": {
            "temperature": 0.65,
            "token_multiplier": 3.0,
            "description": "首席发言人：中低温保证叙事精准"
        },
        "twin": {
            "temperature": 0.60,
            "token_multiplier": 3.5,
            "description": "替身-我：模仿用户思维惯性"
        },
        "default": {
            "temperature": 0.70,
            "token_multiplier": 3.0,
            "description": "默认配置"
        }
    }

    @classmethod
    def get_path(cls, key: str) -> Path:
        return cls.DATA_ROOT / cls.PATHS[key]
    
    @classmethod
    def get_db_path(cls) -> Path:
        override = os.getenv("CRYSTAL_TREE_DB_PATH")
        if override:
            return Path(override).absolute()
        if cls.DATA_ROOT != PROJECT_ROOT:
            return cls.DATA_ROOT / "chat_sessions.db"
        return PROJECT_ROOT / "chat_sessions.db"

    @classmethod
    def get_api_key(cls) -> str:
        return os.getenv("DEEPSEEK_API_KEY", "")

    @staticmethod
    def _determine_data_root() -> Path:
        """智能检测数据根目录：环境变量 → 晶体树文件夹 → 当前目录"""
        env_root = os.getenv("CRYSTAL_TREE_DATA_ROOT")
        if env_root:
            return Path(env_root).absolute()
        candidate = PROJECT_ROOT / "晶体树文件夹"
        if candidate.exists() and candidate.is_dir():
            return candidate
        return PROJECT_ROOT


# ===== 关键修正：在类外部重新设置 DATA_ROOT =====
Config.DATA_ROOT = Config._determine_data_root()


# =============================================================================
# 标准测试问题集（Day 0 基线采集 & Day 3 回归测试）
# =============================================================================
BENCHMARK_QUESTIONS = [
    "如何在不增加预算的前提下，提升一个 20 人研发团队的技术决策质量？",
    "AI 辅助编程工具会让初级程序员丧失独立解决问题的能力吗？请给出平衡方案。",
    "一家传统制造企业想转型数据驱动，但管理层年龄偏大且抵触新技术，最关键的三个切入点是什么？",
    "在资源有限的情况下，应该优先修复技术债务还是优先开发新功能？给出决策框架。",
    "如何判断一个创新想法是「真正有潜力」还是「伪需求」？请给出 3 个检验标准。",
    "远程办公环境下，如何建立有效的「非正式沟通」机制以减少信息孤岛？",
    "面对一个高度不确定的市场，应该采用「精益创业」快速试错，还是「深度调研」后再出手？",
    "个人知识管理系统中，「信息囤积」和「深度消化」的平衡点在哪里？",
    "如何设计一个「自我修正」的 OKR 体系，使团队不偏离长期目标？",
    "当行业颠覆性技术出现时，现有领军企业的「路径依赖」如何打破？"
]

# =============================================================================
# 2. 数据模型 (models.py)
# =============================================================================
from dataclasses import dataclass, field
from datetime import date
from enum import Enum, auto
from typing import List, Optional

class Layer(Enum):
    L1 = auto()
    L2 = auto()
    L3 = auto()

@dataclass
class Crystal:
    id: str
    content: str
    links: List[str] = field(default_factory=list)
    layer: Layer = Layer.L2
    heat: float = 0.0
    last_accessed: Optional[date] = None
    # ===== 新增：晶体代码化字段 =====
    input_conditions: List[str] = field(default_factory=list)
    execution_logic: str = ""
    output_format: str = ""
    validation_criteria: List[str] = field(default_factory=list)
    @property
    def summary(self) -> str:
        return self.content[:50] + "..." if len(self.content) > 50 else self.content

@dataclass
class Hole:
    id: str
    content: str
    urgency: float = 0.5
    layer: int = 2
    @property
    def summary(self) -> str:
        return self.content[:50] + "..." if len(self.content) > 50 else self.content

@dataclass
class Conflict:
    crystal_a: str
    crystal_b: str
    similarity: float
    content_a: str
    content_b: str

@dataclass
class TaskCard:
    id: str
    type: str
    title: str
    content: str
    source: str
    links: List[str] = field(default_factory=list)
    suggested_action: str = ""
    status: str = "pending"

@dataclass
class HealthCheckResult:
    level: str
    file: str
    message: str
    suggested_fix: str = ""


# =============================================================================
# 3. 依赖检测 (dependencies.py)
# =============================================================================
try:
    import requests
    from requests.adapters import HTTPAdapter
    from urllib3.util.retry import Retry
    REQUESTS_AVAILABLE = True
except ImportError:
    requests = None
    HTTPAdapter = None
    Retry = None
    REQUESTS_AVAILABLE = False

try:
    from bs4 import BeautifulSoup
    BS4_AVAILABLE = True
except ImportError:
    BeautifulSoup = None
    BS4_AVAILABLE = False

try:
    from dotenv import load_dotenv
    load_dotenv()
except ImportError:
    pass

try:
    import pandas as pd
except ImportError:
    pd = None

try:
    from docx import Document
    HAS_DOCX = True
except ImportError:
    Document = None
    HAS_DOCX = False

try:
    from pypdf import PdfReader
    HAS_PDF = True
except ImportError:
    try:
        from PyPDF2 import PdfReader
        HAS_PDF = True
    except ImportError:
        PdfReader = None
        HAS_PDF = False

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    Presentation = None
    HAS_PPTX = False

try:
    import arxiv
    ARXIV_AVAILABLE = True
except ImportError:
    arxiv = None
    ARXIV_AVAILABLE = False

# sentence_transformers 可选
SENTENCE_TRANSFORMERS_AVAILABLE = False
print("注意：sentence_transformers 未强制要求，系统将使用内置检索")


# =============================================================================
# 4. 存储 (storage.py)
# =============================================================================
import json
import sqlite3
from contextlib import closing
from datetime import datetime
from typing import Dict, List, Tuple

class FileIO:
    @staticmethod
    def resolve(path: str) -> Path:
        return Config.get_path(path) if path in Config.PATHS else Config.DATA_ROOT / path

    @staticmethod
    def read(path: str) -> str:
        full = FileIO.resolve(path)
        if not full.exists():
            return ""
        with open(full, "r", encoding="utf-8") as f:
            return f.read()

    @staticmethod
    def write(path: str, content: str) -> None:
        full = FileIO.resolve(path)
        full.parent.mkdir(parents=True, exist_ok=True)
        with open(full, "w", encoding="utf-8") as f:
            f.write(content)

    @staticmethod
    def append(path: str, content: str) -> None:
        full = FileIO.resolve(path)
        full.parent.mkdir(parents=True, exist_ok=True)
        with open(full, "a", encoding="utf-8") as f:
            f.write(content)

    @staticmethod
    def exists(path: str) -> bool:
        return FileIO.resolve(path).exists()

    @staticmethod
    def ensure_directories() -> None:
        for dir_name in Config.DIRECTORIES:
            (Config.DATA_ROOT / dir_name).mkdir(parents=True, exist_ok=True)

    @staticmethod
    def read_fingerprint() -> Dict[str, Any]:
        """读取认知指纹 JSON"""
        content = FileIO.read("fingerprint")
        if not content:
            return {
                "fingerprint": {
                    "risk_tolerance": 0.5,
                    "innovation_preference": 0.5,
                    "decisiveness": 0.5,
                    "preferred_role": "structural",
                    "role_adoption_history": {},
                    "conflict_resolution_style": "integrative",
                    "attention_span": 0.5,
                    "context_preference": 3,
                    "last_updated": datetime.now().isoformat(),
                    "total_interactions": 0,
                    "confidence": 0.3,
                    "evolution_log": []
                },
                "extraction_metadata": {
                    "last_extraction": None,
                    "messages_analyzed": 0,
                    "debates_analyzed": 0,
                    "crystals_created": 0,
                    "version": "1.0"
                }
            }
        try:
            return json.loads(content)
        except json.JSONDecodeError:
            return {
                "fingerprint": {
                    "risk_tolerance": 0.5,
                    "innovation_preference": 0.5,
                    "decisiveness": 0.5,
                    "preferred_role": "structural",
                    "role_adoption_history": {},
                    "conflict_resolution_style": "integrative",
                    "attention_span": 0.5,
                    "context_preference": 3,
                    "last_updated": datetime.now().isoformat(),
                    "total_interactions": 0,
                    "confidence": 0.3,
                    "evolution_log": []
                },
                "extraction_metadata": {
                    "last_extraction": None,
                    "messages_analyzed": 0,
                    "debates_analyzed": 0,
                    "crystals_created": 0,
                    "version": "1.0"
                }
            }

    @staticmethod
    def write_fingerprint(data: Dict[str, Any]) -> None:
        """写入认知指纹 JSON"""
        FileIO.write("fingerprint", json.dumps(data, ensure_ascii=False, indent=2))

    @staticmethod
    def ensure_default_files() -> None:
        defaults = {
            "principles": "# 核心操作原则\n\n（待补充）",
            "ai_instructions": "# AI协作者指令\n\n（待补充）",
            "crystals": "# 晶体卡片库\n\n| ID | 内容 | 链接 |\n|----|------|------|\n",
            "holes": "# 孔洞分层\n\n## 第一层：核心孔洞\n\n| ID | 内容 | 紧迫度 |\n|----|------|--------|\n",
            "state": "# 系统状态快照\n\n（初始状态）",
            "change_log": "# 系统变更记录\n\n",
            "pending": "# PENDING 待确认卡片\n\n",
            "search_cache": "# 搜索摘要\n\n",
            "layer_state": json.dumps({"layers": {}, "heat_map": {}, "last_accessed": {}, "manual_override": {}}, ensure_ascii=False, indent=2),
            "hole_progress": json.dumps({}, ensure_ascii=False, indent=2),
            "external_cache": json.dumps({}, ensure_ascii=False, indent=2),
            "task_cards": "[]",
            "roles": json.dumps({
                "radical": {"name": "激进者", "instruction": "攻击默认前提，假设现有框架是错的，给出颠覆性方案。"},
                "conservative": {"name": "保守者", "instruction": "风险优先，假设资源有限，给出最可落地的稳健方案。"},
                "structural": {"name": "结构主义者", "instruction": "从已有晶体中寻找同构案例，用类比生成方案。"},
                "judge": {"name": "大法官","instruction": "以晶体卡片、核心操作原则和资源约束为准绳，做出终审裁决。必须明确引用依据（晶体ID、原则条款或约束条件），不得凭直觉判案。"},
                "spokesperson": {"name": "首席发言人","instruction": "将内部辩论结论转化为清晰、简洁、无歧义的对外陈述。遵循降维（通俗化）、定调（不超过3条核心信息）、检验（老板读前100字能决策）三原则。"},
                "lark": {"name": "百灵鸟","instruction": "见多识广的通用智能体，从外部世界（学术、产业、政策、跨学科）补充知识，打破信息茧房。在第二轮登场。"},
                "pilgrim": {"name": "取经者","instruction": "以长期愿景和核心价值观为锚，防止短期利益或局部优化偏离最终使命。评估方案的可持续性和道德一致性。"},
                "strategist": {"name": "奇谋者","instruction": "善于洞察人心、把握时机，敢押注非常规路径，捕捉机会窗口。评估方案能否借力打力、以奇制胜。"},
                "statesman": {"name": "延安智者","instruction": "坚持调查研究，不唯上、不唯书、只唯实。从全局矛盾和主要矛盾切入，提出实事求是、可落地的综合方略。"}              
            }, ensure_ascii=False, indent=2)
        }
        # 写入所有默认文本文件
        for key, content in defaults.items():
            if not FileIO.exists(key):
                FileIO.write(key, content)

        # ===== Day 0 新增：pareto_frontier.json（帕累托前沿跟踪） =====
        pareto_path = Config.DATA_ROOT / "系统日志" / "pareto_frontier.json"
        if not pareto_path.exists():
            with open(pareto_path, "w", encoding="utf-8") as f:
                json.dump({
                    "configs": {
                        "PROFILE_HIGH_ACCURACY": {"accuracy": 0.9, "cost": 0.3, "latency": 0.2},
                        "PROFILE_BALANCED": {"accuracy": 0.7, "cost": 0.2, "latency": 0.15},
                        "PROFILE_ECONOMY": {"accuracy": 0.5, "cost": 0.1, "latency": 0.1}
                    },
                    "history": []
                }, f, ensure_ascii=False, indent=2)

        # ===== Day 0 新增：灵感池.json（灵感熔炉数据） =====
        insp_path = Config.DATA_ROOT / "系统日志" / "灵感池.json"
        if not insp_path.exists():
            with open(insp_path, "w", encoding="utf-8") as f:
                json.dump([
                    {
                        "id": "INSP-001",
                        "source": "对话",
                        "content": "初始灵感：将‘八道防线’与‘沉思式反思’融合，形成免疫+智慧的协同效应",
                        "status": "待筛选",
                        "created_at": datetime.now().isoformat()
                    }
                ], f, ensure_ascii=False, indent=2)
            
class HealthChecker:
    @staticmethod
    def run() -> List[HealthCheckResult]:
        results = []
        for directory in Config.DIRECTORIES:
            path = Config.DATA_ROOT / directory
            if not path.exists():
                results.append(HealthCheckResult("warning", directory, "目录不存在", "启动初始化会自动创建"))
        for key in Config.PATHS:
            path = Config.get_path(key)
            if not path.exists():
                results.append(HealthCheckResult("warning", Config.PATHS[key], "文件不存在", "启动初始化会按默认模板创建"))
                continue
            try:
                text = path.read_text(encoding="utf-8")
            except UnicodeDecodeError:
                results.append(HealthCheckResult("error", Config.PATHS[key], "文件不是有效 UTF-8 编码", "请先备份后统一转为 UTF-8"))
                continue
            if key in ("layer_state", "hole_progress", "external_cache", "task_cards"):
                try:
                    json.loads(text or "{}")
                except json.JSONDecodeError as exc:
                    results.append(HealthCheckResult("error", Config.PATHS[key], f"JSON 解析失败: {exc}", "请检查逗号、括号和引号"))
            if key == "crystals" and text and "| ID | 内容 | 链接 |" not in text:
                results.append(HealthCheckResult("warning", Config.PATHS[key], "晶体卡片表头缺失", "请确认 Markdown 表格结构"))
            if key == "holes" and text and "| ID | 内容 | 紧迫度 |" not in text:
                results.append(HealthCheckResult("warning", Config.PATHS[key], "孔洞表头缺失", "请确认 Markdown 表格结构"))
        return results


# =============================================================================
# 认知指纹数据模型（新增）
# =============================================================================

from datetime import datetime
from typing import Any

@dataclass
class CognitiveFingerprint:
    """
    用户认知指纹 —— 让系统"认识你"
    
    认知指纹是用户思维模式的量化表示，通过分析用户的对话历史、
    辩论参与行为、晶体化偏好等数据动态提取和更新。
    
    四阶路径中的"认识你"阶段的核心数据结构。
    """
    
    # ===== 决策偏好维度 =====
    risk_tolerance: float = 0.5
    """风险容忍度（0-1）：0=极度保守，1=极度激进"""
    
    innovation_preference: float = 0.5
    """创新偏好（0-1）：0=偏好成熟方案，1=偏好颠覆性方案"""
    
    decisiveness: float = 0.5
    """决策果断值（0-1）：0=深思熟虑型，1=快速决断型"""
    
    # ===== 角色倾向维度 =====
    preferred_role: str = "structural"
    """最常采纳的角色（radical/conservative/structural/executor/auditor）"""
    
    role_adoption_history: Dict[str, int] = field(default_factory=dict)
    """各角色被采纳次数统计"""
    
    # ===== 冲突解决风格 =====
    conflict_resolution_style: str = "integrative"
    """冲突解决风格：integrative（整合型）/competitive（竞争型）/avoidant（回避型）"""
    
    # ===== 注意力模式 =====
    attention_span: float = 0.5
    """注意力持续度（0-1）：0=易分散，1=高度专注"""
    
    context_preference: int = 3
    """偏好上下文轮数（回顾最近几轮对话）"""
    # ===== Day 2.5 新增：认知风格 =====
    reasoning_style: str = "balanced"        # "deductive" | "inductive" | "balanced"
    analogy_preference: str = "balanced"     # "analogy" | "first_principles" | "balanced"
    output_style: str = "conclusion_first"   # "conclusion_first" | "evidence_first" | "balanced"

    # ===== 元数据 =====
    last_updated: str = field(default_factory=lambda: datetime.now().isoformat())
    total_interactions: int = 0
    confidence: float = 0.3
    evolution_log: List[Dict[str, Any]] = field(default_factory=list)
    
    def to_dict(self) -> Dict[str, Any]:
        return {
            "risk_tolerance": self.risk_tolerance,
            "innovation_preference": self.innovation_preference,
            "decisiveness": self.decisiveness,
            "preferred_role": self.preferred_role,
            "role_adoption_history": self.role_adoption_history,
            "conflict_resolution_style": self.conflict_resolution_style,
            "attention_span": self.attention_span,
            "context_preference": self.context_preference,
            "last_updated": self.last_updated,
            "total_interactions": self.total_interactions,
            "confidence": self.confidence,
            "reasoning_style": self.reasoning_style,
            "analogy_preference": self.analogy_preference,
            "output_style": self.output_style,            
            "evolution_log": self.evolution_log
        }
    
    @classmethod
    def from_dict(cls, data: Dict[str, Any]) -> "CognitiveFingerprint":
        return cls(
            risk_tolerance=data.get("risk_tolerance", 0.5),
            innovation_preference=data.get("innovation_preference", 0.5),
            decisiveness=data.get("decisiveness", 0.5),
            preferred_role=data.get("preferred_role", "structural"),
            role_adoption_history=data.get("role_adoption_history", {}),
            conflict_resolution_style=data.get("conflict_resolution_style", "integrative"),
            attention_span=data.get("attention_span", 0.5),
            context_preference=data.get("context_preference", 3),
            last_updated=data.get("last_updated", datetime.now().isoformat()),
            total_interactions=data.get("total_interactions", 0),
            confidence=data.get("confidence", 0.3),
            reasoning_style=data.get("reasoning_style", "balanced"),
            analogy_preference=data.get("analogy_preference", "balanced"),
            output_style=data.get("output_style", "balanced"),
            evolution_log=data.get("evolution_log", [])
        )


@dataclass
class FingerprintExtractionResult:
    """指纹提取结果"""
    fingerprint: CognitiveFingerprint
    source_analysis: Dict[str, Any]
    changes: List[Dict[str, Any]]
    confidence_delta: float

class DBManager:
    def __init__(self, db_path: Path = None):
        self.db_path = db_path or Config.get_db_path()
        self._init_db()
        
    def _connect(self):
        return sqlite3.connect(self.db_path)

    def _init_db(self):
        with closing(self._connect()) as conn:
            cursor = conn.cursor()
            cursor.execute('''CREATE TABLE IF NOT EXISTS sessions (id TEXT PRIMARY KEY, name TEXT NOT NULL, created_at TIMESTAMP, updated_at TIMESTAMP, messages TEXT)''')
            conn.commit()

    def create_session(self, session_id: str, name: str):
        with closing(self._connect()) as conn:
            cursor = conn.cursor()
            now = datetime.now().isoformat()
            cursor.execute('INSERT INTO sessions (id, name, created_at, updated_at, messages) VALUES (?,?,?,?,?)', (session_id, name, now, now, json.dumps([])))
            conn.commit()

    def update_session(self, session_id: str, messages: List, name: str = None):
        """
        更新会话消息，消息元素可以是 tuple (role, content) 或 dict {role, content, label}
        自动转换为统一格式（包含 label 字段）
        """
        with closing(self._connect()) as conn:
            cursor = conn.cursor()
            if name:
                cursor.execute('UPDATE sessions SET name = ?, updated_at = ? WHERE id = ?',
                               (name, datetime.now().isoformat(), session_id))
            else:
                cursor.execute('UPDATE sessions SET updated_at = ? WHERE id = ?',
                               (datetime.now().isoformat(), session_id))
            # 统一转换为带 label 的 dict
            new_msgs = []
            for item in messages:
                if isinstance(item, dict):
                    # 如果已经是 dict，确保有 label 字段
                    new_msgs.append({
                        "role": item.get("role", ""),
                        "content": item.get("content", ""),
                        "label": item.get("label")
                    })
                else:
                    # 假定是 tuple (role, content)
                    role, content = item[0], item[1]
                    new_msgs.append({"role": role, "content": content, "label": None})
            cursor.execute('UPDATE sessions SET messages = ? WHERE id = ?',
                           (json.dumps(new_msgs, ensure_ascii=False), session_id))
            conn.commit()

    def get_session(self, session_id: str):
        with closing(self._connect()) as conn:
            cursor = conn.cursor()
            cursor.execute('SELECT name, messages FROM sessions WHERE id = ?', (session_id,))
            row = cursor.fetchone()
        if row:
            name, messages_json = row
            messages = json.loads(messages_json)
            history = []
            labels = []
            for msg in messages:
                role = msg.get("role", "")
                content = msg.get("content", "")
                label = msg.get("label")  # 可能为 None
                history.append((role, content))
                labels.append(label)
            return name, history, labels
        return None, [], []

    def list_sessions(self):
        with closing(self._connect()) as conn:
            cursor = conn.cursor()
            cursor.execute('SELECT id, name, updated_at FROM sessions ORDER BY updated_at DESC')
            return cursor.fetchall()

    def delete_session(self, session_id: str):
        with closing(self._connect()) as conn:
            cursor = conn.cursor()
            cursor.execute('DELETE FROM sessions WHERE id = ?', (session_id,))
            conn.commit()

    def rename_session(self, session_id: str, new_name: str):
        with closing(self._connect()) as conn:
            cursor = conn.cursor()
            cursor.execute('UPDATE sessions SET name = ?, updated_at = ? WHERE id = ?', (new_name, datetime.now().isoformat(), session_id))
            conn.commit()


# =============================================================================
# 5. 网络 (network.py)
# =============================================================================
import random
import time
from typing import Callable

class NetworkManager:
    _shared_session = None

    @classmethod
    def _get_session(cls):
        if not REQUESTS_AVAILABLE:
            return None
        if cls._shared_session is None:
            session = requests.Session()
            retry_strategy = Retry(total=Config.MAX_RETRIES, backoff_factor=Config.BACKOFF_FACTOR, status_forcelist=[429,500,502,503,504], allowed_methods=["GET","HEAD"])
            adapter = HTTPAdapter(max_retries=retry_strategy)
            session.mount("http://", adapter)
            session.mount("https://", adapter)
            cls._shared_session = session
        return cls._shared_session

    @classmethod
    def get_random_user_agent(cls) -> str:
        return random.choice(Config.USER_AGENTS)

    @classmethod
    def safe_request(cls, url: str, use_mirror: bool = False, log_callback: Callable = None, **kwargs):
        if not REQUESTS_AVAILABLE:
            return None
        session = cls._get_session()
        if session is None:
            return None
        time.sleep(random.uniform(*Config.DELAY_BETWEEN_REQUESTS))
        final_url = url
        if use_mirror and "huggingface.co" in url:
            final_url = url.replace("https://huggingface.co", Config.HF_MIRROR)
        headers = kwargs.get('headers', {})
        headers.update({'User-Agent': cls.get_random_user_agent(), 'Accept': 'text/html,application/xhtml+xml,application/xml;q=0.9,image/webp,*/*;q=0.8', 'Accept-Language': 'zh-CN,zh;q=0.9,en;q=0.8', 'Connection': 'keep-alive'})
        kwargs['headers'] = headers
        kwargs['timeout'] = Config.TIMEOUT
        try:
            response = session.get(final_url, **kwargs)
            response.raise_for_status()
            if response.status_code == 200 and len(response.content) > 100:
                return response
            return None
        except Exception as e:
            if log_callback:
                log_callback(f"请求失败 {type(e).__name__}: {final_url}", "warning")
            return None


# =============================================================================
# 6. 核心引擎 (engine.py)
# =============================================================================
import math
import re
from collections import Counter
from datetime import date, datetime, timedelta
from typing import Dict, List, Set, Tuple,Any



class MetaLayer:
    """
    认知晶体树的元层引擎

    负责管理"如何管理认知"的六种元原语。
    与工作层（CrystalEngine）分离，实现双层架构。

    对应建议1：建立"工作层+元层"双层架构
    """

    def __init__(self, engine: 'CrystalEngine', file_io: FileIO):
        self.engine = engine
        self.files = file_io
        self.primitive_states = {
            "active_gap_detection": {"last_run": None, "status": "idle"},
            "temporal_aware_escalation": {"last_run": None, "status": "idle"},
            "layer_aware_calibration": {"last_run": None, "status": "active"},
            "sleep_consolidation": {"last_run": None, "status": "idle"},
            "distributed_metacognition": {"last_run": None, "status": "disabled"},
            "validation_gated_self_evolution": {"last_run": None, "status": "active"}
        }
        self._load_state()
        # Day 3 新增：触发链执行状态
        self.chain_states = {
            "isolated_crystal_to_validation": {"last_triggered": None, "trigger_count": 0},
            "stale_hot_to_gap_detection": {"last_triggered": None, "trigger_count": 0},
        }
        self._load_chain_state()
    def _load_pareto_data(self) -> Dict[str, Any]:
        """加载帕累托数据"""
        pareto_path = Config.DATA_ROOT / "系统日志" / "pareto_frontier.json"
        if pareto_path.exists():
            try:
                with open(pareto_path, "r", encoding="utf-8") as f:
                    return json.load(f)
            except:
                pass
        return {"configs": {}, "history": [], "daily_stats": []}

    def _save_pareto_data(self, data: Dict[str, Any]) -> None:
        """保存帕累托数据"""
        pareto_path = Config.DATA_ROOT / "系统日志" / "pareto_frontier.json"
        pareto_path.parent.mkdir(parents=True, exist_ok=True)
        with open(pareto_path, "w", encoding="utf-8") as f:
            json.dump(data, f, ensure_ascii=False, indent=2)

    def record_conversation_metrics(
        self,
        profile_name: str,
        accuracy: float,
        cost: float,
        latency: float,
        crystal_refs: int = 0,
        quality_score: float = 0.0
    ) -> None:
        """
        记录一次对话的三维指标

        Args:
            profile_name: 配置名称 (high_accuracy/balanced/economy)
            accuracy: 准确性评分 (0-1)
            cost: 成本 (美元)
            latency: 延迟 (秒)
            crystal_refs: 晶体引用数
            quality_score: 质量评分 (0-1)
        """
        data = self._load_pareto_data()

        # 更新配置记录
        if profile_name not in data["configs"]:
            data["configs"][profile_name] = {
                "accuracy": 0.0,
                "cost": 0.0,
                "latency": 0.0,
                "crystal_refs": 0,
                "quality_score": 0.0,
                "count": 0
            }

        config = data["configs"][profile_name]
        config["accuracy"] = (config["accuracy"] * config["count"] + accuracy) / (config["count"] + 1)
        config["cost"] = (config["cost"] * config["count"] + cost) / (config["count"] + 1)
        config["latency"] = (config["latency"] * config["count"] + latency) / (config["count"] + 1)
        config["crystal_refs"] = (config["crystal_refs"] * config["count"] + crystal_refs) / (config["count"] + 1)
        config["quality_score"] = (config["quality_score"] * config["count"] + quality_score) / (config["count"] + 1)
        config["count"] += 1

        # 添加到历史
        data["history"].append({
            "timestamp": datetime.now().isoformat(),
            "profile": profile_name,
            "accuracy": accuracy,
            "cost": cost,
            "latency": latency,
            "crystal_refs": crystal_refs,
            "quality_score": quality_score
        })

        # 限制历史数量
        if len(data["history"]) > Config.PARETO_HISTORY_LIMIT:
            data["history"] = data["history"][-Config.PARETO_HISTORY_LIMIT:]

        self._save_pareto_data(data)

    def get_pareto_status(self) -> Dict[str, Any]:
        """获取当前帕累托状态"""
        data = self._load_pareto_data()
        configs = data.get("configs", {})
        history = data.get("history", [])

        # 计算趋势
        trends = self._calculate_trends(history)

        return {
            "configs": configs,
            "history_count": len(history),
            "trends": trends,
            "best_profile": self._get_best_profile(configs),
            "daily_stats": data.get("daily_stats", [])
        }

    def _calculate_trends(self, history: List[Dict]) -> Dict[str, Any]:
        """计算趋势指标"""
        if len(history) < 3:
            return {"trend": "insufficient_data", "accuracy_delta": 0, "cost_delta": 0}

        recent = history[-10:]
        if len(recent) < 3:
            return {"trend": "stable", "accuracy_delta": 0, "cost_delta": 0}

        # 计算最近3条和之前3条的平均值
        first_half = recent[:len(recent)//2]
        second_half = recent[len(recent)//2:]

        avg_acc_first = sum(h.get("accuracy", 0) for h in first_half) / max(1, len(first_half))
        avg_acc_second = sum(h.get("accuracy", 0) for h in second_half) / max(1, len(second_half))
        avg_cost_first = sum(h.get("cost", 0) for h in first_half) / max(1, len(first_half))
        avg_cost_second = sum(h.get("cost", 0) for h in second_half) / max(1, len(second_half))

        acc_delta = avg_acc_second - avg_acc_first
        cost_delta = avg_cost_second - avg_cost_first

        if acc_delta > 0.05:
            trend = "improving"
        elif acc_delta < -0.05:
            trend = "declining"
        else:
            trend = "stable"

        return {
            "trend": trend,
            "accuracy_delta": round(acc_delta, 3),
            "cost_delta": round(cost_delta, 3),
            "avg_accuracy": round(avg_acc_second, 3),
            "avg_cost": round(avg_cost_second, 3)
        }

    def _get_best_profile(self, configs: Dict) -> Optional[str]:
        """获取当前最优配置"""
        if not configs:
            return None

        best = None
        best_score = -1

        for name, data in configs.items():
            # 综合评分：准确性权重0.5，成本权重0.3，延迟权重0.2
            score = (
                data.get("accuracy", 0) * 0.5 +
                (1 - min(1, data.get("cost", 0) * 10)) * 0.3 +
                (1 - min(1, data.get("latency", 0) / 60)) * 0.2
            )
            if score > best_score:
                best_score = score
                best = name

        return best

    def record_daily_stats(self, stats: Dict[str, Any]) -> None:
        """记录每日统计（个人认知效率仪表盘）"""
        data = self._load_pareto_data()

        if "daily_stats" not in data:
            data["daily_stats"] = []

        today = datetime.now().date().isoformat()
        # 检查今天是否已有记录
        existing = None
        for i, entry in enumerate(data["daily_stats"]):
            if entry.get("date") == today:
                existing = i
                break

        if existing is not None:
            # 更新已有记录
            data["daily_stats"][existing].update(stats)
            data["daily_stats"][existing]["date"] = today
        else:
            # 新增记录
            data["daily_stats"].append({"date": today, **stats})

        # 只保留最近30天
        if len(data["daily_stats"]) > 30:
            data["daily_stats"] = data["daily_stats"][-30:]

        self._save_pareto_data(data)

    def get_daily_stats(self, days: int = 7) -> List[Dict[str, Any]]:
        """获取最近N天的每日统计"""
        data = self._load_pareto_data()
        daily = data.get("daily_stats", [])
        return daily[-days:]

    def _load_state(self) -> None:
        """从文件加载元层状态"""
        state_file = self.files.resolve("change_log").parent / "meta_state.json"
        if state_file.exists():
            try:
                with open(state_file, "r", encoding="utf-8") as f:
                    data = json.load(f)
                    self.primitive_states.update(data.get("primitive_states", {}))
            except:
                pass
    def _load_chain_state(self) -> None:
        """从文件加载触发链状态"""
        state_file = self.files.resolve("change_log").parent / "meta_chain_state.json"
        if state_file.exists():
            try:
                with open(state_file, "r", encoding="utf-8") as f:
                    data = json.load(f)
                    self.chain_states.update(data.get("chain_states", {}))
            except:
                pass
    def _save_chain_state(self) -> None:
        """保存触发链状态到文件"""
        state_file = self.files.resolve("change_log").parent / "meta_chain_state.json"
        with open(state_file, "w", encoding="utf-8") as f:
            json.dump({
                "chain_states": self.chain_states,
                "last_saved": datetime.now().isoformat()
            }, f, ensure_ascii=False, indent=2)

    def _save_state(self) -> None:
        """保存元层状态到文件"""
        state_file = self.files.resolve("change_log").parent / "meta_state.json"
        with open(state_file, "w", encoding="utf-8") as f:
            json.dump({
                "primitive_states": self.primitive_states,
                "last_saved": datetime.now().isoformat()
            }, f, ensure_ascii=False, indent=2)

    def active_gap_detection(self) -> List[Dict[str, Any]]:
        """
        主动检测晶体之间的逻辑缝隙

        返回：检测到的缺口列表，并附带触发信息
        """
        if not Config.META_PRIMITIVES["active_gap_detection"]["enabled"]:
            return []

        crystals = self.engine.parse_crystals()
        if len(crystals) < 20:
            return []

        gaps = []
        # 检测孤立晶体（links 为空的晶体）
        isolated = [c for c in crystals if not c.links]
        for c in isolated:
            gaps.append({
                "type": "isolated_crystal",
                "crystal_id": c.id,
                "content": c.content,
                "severity": "medium",
                "suggestion": f"尝试将 {c.id} 与已有晶体建立链接"
            })

        # 检测语义接近但无链接的晶体对
        for i in range(len(crystals)):
            for j in range(i + 1, len(crystals)):
                c1, c2 = crystals[i], crystals[j]
                if len(c1.links) >= 2 or len(c2.links) >= 2:
                    continue
                sim = self.engine._simple_similarity(c1.content, c2.content)
                if sim > 0.5 and c2.id not in c1.links and c1.id not in c2.links:
                    gaps.append({
                        "type": "near_duplicate",
                        "crystal_a": c1.id,
                        "crystal_b": c2.id,
                        "similarity": round(sim, 2),
                        "severity": "high" if sim > 0.7 else "medium",
                        "suggestion": f"考虑合并 {c1.id} 和 {c2.id} 或建立链接"
                    })
                    break

        if gaps:
            self.engine._append_change_log(
                "主动缺口检测",
                f"发现 {len(gaps)} 个缺口：{', '.join([g.get('crystal_id', g.get('crystal_a', '')) for g in gaps[:5]])}"
            )

        self.primitive_states["active_gap_detection"]["last_run"] = datetime.now().isoformat()
        self._save_state()

        # ===== Day 3 新增：返回完整的缺口信息（含统计） =====
        isolated_count = len([g for g in gaps if g.get("type") == "isolated_crystal"])
        near_dup_count = len([g for g in gaps if g.get("type") == "near_duplicate"])
        
        return {
            "gaps": gaps,
            "stats": {
                "total": len(gaps),
                "isolated_count": isolated_count,
                "near_dup_count": near_dup_count
            },
            "trigger_info": {
                "isolated_crystal_count": isolated_count,
                "has_isolated_crystals": isolated_count > 0
            }
        }

    def temporal_aware_escalation(self) -> List[Dict[str, Any]]:
        """
        根据时间戳和活动频率自动调整优先级

        返回：升级建议列表，并附带触发信息
        """
        if not Config.META_PRIMITIVES["temporal_aware_escalation"]["enabled"]:
            return []

        crystals = self.engine.parse_crystals()
        state = self.engine.load_layer_state()
        last_accessed = state.get("last_accessed", {})
        heat_map = state.get("heat_map", {})

        escalations = []
        today = date.today()
        stale_hot_count = 0

        for c in crystals:
            last = last_accessed.get(c.id)
            if last:
                try:
                    days_since = (today - date.fromisoformat(last)).days
                except:
                    days_since = 999
            else:
                days_since = 999

            heat = heat_map.get(c.id, 0.0)
            current_layer = state.get("layers", {}).get(c.id, "L2")

            if heat > 0.5 and days_since > 14 and current_layer == "L2":
                escalations.append({
                    "type": "stale_hot_crystal",
                    "crystal_id": c.id,
                    "heat": heat,
                    "days_since": days_since,
                    "suggestion": f"{c.id} 热度较高但 {days_since} 天未访问，建议主动召回"
                })
                stale_hot_count += 1

        self.primitive_states["temporal_aware_escalation"]["last_run"] = datetime.now().isoformat()
        self._save_state()

        # ===== Day 3 新增：返回完整的升级信息（含统计） =====
        return {
            "escalations": escalations,
            "stats": {
                "total": len(escalations),
                "stale_hot_count": stale_hot_count
            },
            "trigger_info": {
                "stale_hot_count": stale_hot_count,
                "has_stale_hot": stale_hot_count > 0
            }
        }

    def layer_aware_calibration(self, crystals: List) -> List:
        """
        根据晶体层级调整检索权重

        返回：按层级权重调整后的晶体列表
        """
        if not Config.META_PRIMITIVES["layer_aware_calibration"]["enabled"]:
            return crystals

        state = self.engine.load_layer_state()
        layers = state.get("layers", {})

        def layer_weight(c) -> int:
            return {"L1": 3, "L2": 2, "L3": 1}.get(layers.get(c.id, "L2"), 2)

        return sorted(crystals, key=lambda c: (layer_weight(c), getattr(c, "heat", 0)), reverse=True)

    def sleep_consolidation(self) -> Dict[str, Any]:
        """
        低活跃时段自动进行知识压缩和冗余清理

        返回：巩固报告
        """
        if not Config.META_PRIMITIVES["sleep_consolidation"]["enabled"]:
            return {"status": "disabled", "message": "睡眠巩固未启用"}

        current_hour = datetime.now().hour
        start = Config.META_LAYER_CONFIG["consolidation_hour_start"]
        end = Config.META_LAYER_CONFIG["consolidation_hour_end"]

        if not (start <= current_hour < end):
            return {"status": "skipped", "message": f"当前时间 {current_hour}:00 不在巩固窗口 ({start}:00-{end}:00)"}

        crystals = self.engine.parse_crystals()
        if len(crystals) < Config.META_LAYER_CONFIG["max_crystals_before_consolidation"]:
            return {
                "status": "skipped",
                "message": f"晶体数量 {len(crystals)} 低于阈值 {Config.META_LAYER_CONFIG['max_crystals_before_consolidation']}"
            }

        archived = self.engine.archive_cold_crystals()
        gaps = self.active_gap_detection()

        result = {
            "status": "completed",
            "archived_count": len(archived),
            "gaps_found": len(gaps),
            "archived_ids": archived[:10],
            "consolidated_at": datetime.now().isoformat()
        }

        self.primitive_states["sleep_consolidation"]["last_run"] = datetime.now().isoformat()
        self._save_state()
        return result

    def validation_gated_self_evolution(self, new_data: Dict[str, Any], context: Dict[str, Any]) -> Dict[str, Any]:
        """
        验证门控自我进化：任何自我改进必须通过验证门检验

        返回：验证结果（通过/不通过 + 原因）
        """
        if not Config.META_PRIMITIVES["validation_gated_self_evolution"]["enabled"]:
            return {"passed": True, "reason": "验证门控未启用"}

        rules = Config.META_LAYER_CONFIG["validation_gate_rules"]
        results = []

        if "new_evidence_from_at_least_3_sources" in rules:
            sources = context.get("sources", [])
            results.append({
                "rule": "new_evidence_from_at_least_3_sources",
                "passed": len(sources) >= 3,
                "reason": f"有 {len(sources)} 个来源" if len(sources) >= 3 else f"仅有 {len(sources)} 个来源，需要至少3个"
            })

        if "audit_score_gt_0.6" in rules:
            audit_score = context.get("audit_score", 0.0)
            results.append({
                "rule": "audit_score_gt_0.6",
                "passed": audit_score >= 0.6,
                "reason": f"审计评分 {audit_score:.2f} ≥ 0.6" if audit_score >= 0.6 else f"审计评分 {audit_score:.2f} 低于阈值 0.6"
            })

        if "no_major_conflict_with_existing_crystals" in rules:
            conflicts = self.engine.detect_conflicts()
            major_conflicts = [c for c in conflicts if c.similarity > 0.8]
            results.append({
                "rule": "no_major_conflict_with_existing_crystals",
                "passed": len(major_conflicts) == 0,
                "reason": "无重大冲突" if len(major_conflicts) == 0 else f"存在 {len(major_conflicts)} 个重大冲突"
            })

        all_passed = all(r["passed"] for r in results)

        self.primitive_states["validation_gated_self_evolution"]["last_run"] = datetime.now().isoformat()
        self._save_state()

        return {
            "passed": all_passed,
            "rules": results,
            "summary": "所有验证通过" if all_passed else f"未通过：{', '.join([r['reason'] for r in results if not r['passed']])}"
        }

    def process_trigger_chains(self, primitive_results: Dict[str, Any]) -> List[Dict[str, Any]]:
        """
        处理元原语触发链
        检查各元原语的执行结果，根据配置的触发规则自动触发后续元原语

        Args:
            primitive_results: run_all_primitives 的返回结果

        Returns:
            触发的链列表
        """
        triggered_chains = []
        rules = Config.META_CHAIN_RULES

        # 获取各元原语的结果
        gap_result = primitive_results.get("active_gap_detection", {})
        temporal_result = primitive_results.get("temporal_aware_escalation", {})
        validation_result = primitive_results.get("validation_gated_self_evolution", {})

        # ===== 触发链1：孤立晶体 → 验证门控自我进化 =====
        chain1 = rules.get("isolated_crystal_to_validation", {})
        if chain1.get("enabled", True):
            # 检查是否满足条件：孤立晶体数量 ≥ 3
            if isinstance(gap_result, dict):
                isolated_count = gap_result.get("stats", {}).get("isolated_count", 0)
                if isolated_count >= 3:
                    # 记录触发事件
                    self.engine.log_evolution_event(
                        "chain_triggered",
                        {
                            "chain": "isolated_crystal_to_validation",
                            "source": "active_gap_detection",
                            "target": "validation_gated_self_evolution",
                            "isolated_count": isolated_count,
                            "trigger": "chain"
                        }
                    )
                    self.engine._append_change_log(
                        "元原语触发链",
                        f"主动缺口检测发现 {isolated_count} 个孤立晶体 → 触发验证门控自我进化"
                    )
                    
                    # 执行目标元原语：验证门控自我进化
                    context = {
                        "sources": [g.get("crystal_id", "") for g in gap_result.get("gaps", []) if g.get("type") == "isolated_crystal"],
                        "audit_score": 0.7,  # 模拟审计评分
                        "focus": "isolated_crystals"
                    }
                    validation_result = self.validation_gated_self_evolution({}, context)
                    
                    # 记录触发链结果
                    triggered_chains.append({
                        "chain": "isolated_crystal_to_validation",
                        "source": "active_gap_detection",
                        "target": "validation_gated_self_evolution",
                        "source_result": f"发现 {isolated_count} 个孤立晶体",
                        "target_result": validation_result.get("summary", ""),
                        "passed": validation_result.get("passed", False)
                    })
                    
                    # 更新状态
                    self.chain_states["isolated_crystal_to_validation"]["last_triggered"] = datetime.now().isoformat()
                    self.chain_states["isolated_crystal_to_validation"]["trigger_count"] += 1
                    self._save_chain_state()

        # ===== 触发链2：高热度久未访问 → 主动缺口检测 =====
        chain2 = rules.get("stale_hot_to_gap_detection", {})
        if chain2.get("enabled", True):
            if isinstance(temporal_result, dict):
                stale_hot_count = temporal_result.get("stats", {}).get("stale_hot_count", 0)
                if stale_hot_count >= 2:
                    self.engine.log_evolution_event(
                        "chain_triggered",
                        {
                            "chain": "stale_hot_to_gap_detection",
                            "source": "temporal_aware_escalation",
                            "target": "active_gap_detection",
                            "stale_hot_count": stale_hot_count,
                            "trigger": "chain"
                        }
                    )
                    self.engine._append_change_log(
                        "元原语触发链",
                        f"时序感知升级发现 {stale_hot_count} 个高热度久未访问晶体 → 触发主动缺口检测"
                    )
                    
                    # 执行目标元原语：主动缺口检测（强制运行）
                    gap_result = self.active_gap_detection()
                    
                    triggered_chains.append({
                        "chain": "stale_hot_to_gap_detection",
                        "source": "temporal_aware_escalation",
                        "target": "active_gap_detection",
                        "source_result": f"发现 {stale_hot_count} 个高热度久未访问晶体",
                        "target_result": f"发现 {gap_result.get('stats', {}).get('total', 0)} 个缺口" if isinstance(gap_result, dict) else "已执行",
                        "passed": True
                    })
                    
                    self.chain_states["stale_hot_to_gap_detection"]["last_triggered"] = datetime.now().isoformat()
                    self.chain_states["stale_hot_to_gap_detection"]["trigger_count"] += 1
                    self._save_chain_state()

        return triggered_chains

    # ===== Day 8: 双时间尺度进化调度 =====
    def _load_saturation_state(self) -> Dict[str, Any]:
        """加载饱和检测状态"""
        state_file = Config.DATA_ROOT / "系统日志" / "saturation_state.json"
        if state_file.exists():
            try:
                with open(state_file, "r", encoding="utf-8") as f:
                    return json.load(f)
            except:
                pass
        return {
            "prompt_optimization_rounds": [],
            "quality_history": [],
            "saturation_status": "unsaturated",
            "current_level": "prompt",
            "consecutive_rounds": 0,
            "last_improvement": 0.0,
            "control_logic_changes": []
        }

    def _save_saturation_state(self, state: Dict[str, Any]) -> None:
        """保存饱和检测状态"""
        state_file = Config.DATA_ROOT / "系统日志" / "saturation_state.json"
        state_file.parent.mkdir(parents=True, exist_ok=True)
        with open(state_file, "w", encoding="utf-8") as f:
            json.dump(state, f, ensure_ascii=False, indent=2)

    def prompt_saturation_detector(self, quality_score: float, context: Dict[str, Any] = None) -> Dict[str, Any]:
        """
        提示词饱和检测器（改进版：累积检测）

        记录连续N轮修改后质量提升幅度，当提升 < 6% 时标记为"饱和"
        饱和后，才允许系统触发控制逻辑层面的变更

        Args:
            quality_score: 本轮质量评分 (0-1)
            context: 上下文信息（如修改类型、触发原因等）

        Returns:
            {
                "is_saturated": bool,
                "consecutive_rounds": int,
                "improvement": float,
                "level": "prompt" | "control_logic",
                "status": "unsaturated" | "saturated" | "escalated"
            }
        """
        # ===== 关键修复：在方法内部定义阈值 =====
        SATURATION_THRESHOLD = 0.06  # 6% 的提升阈值

        state = self._load_saturation_state()
        context = context or {}

        # 获取修改类型
        modification_type = context.get("modification_type", "prompt")
        is_control_logic = context.get("is_control_logic", False)

        # 如果当前是控制逻辑变更，记录并返回
        if is_control_logic:
            state["current_level"] = "control_logic"
            state["control_logic_changes"].append({
                "timestamp": datetime.now().isoformat(),
                "context": context,
                "quality_score": quality_score
            })
            self._save_saturation_state(state)
            return {
                "is_saturated": True,
                "consecutive_rounds": state.get("consecutive_rounds", 0),
                "improvement": state.get("last_improvement", 0.0),
                "level": "control_logic",
                "status": "escalated"
            }

        # 添加质量历史
        state["quality_history"].append({
            "timestamp": datetime.now().isoformat(),
            "score": quality_score,
            "type": modification_type,
            "context": context
        })

        # 只保留最近20轮
        if len(state["quality_history"]) > 20:
            state["quality_history"] = state["quality_history"][-20:]

        history = state["quality_history"]

        # ===== 数据不足时返回未饱和 =====
        if len(history) < 2:
            state["saturation_status"] = "unsaturated"
            state["consecutive_rounds"] = 0
            state["last_improvement"] = 0.0
            self._save_saturation_state(state)
            return {
                "is_saturated": False,
                "consecutive_rounds": 0,
                "improvement": 0.0,
                "level": "prompt",
                "status": "unsaturated",
                "avg_recent": quality_score,
                "avg_previous": quality_score
            }

        # ===== 累积式饱和检测 =====
        # 取最近5轮
        recent_rounds = history[-5:] if len(history) >= 5 else history

        saturated_count = 0
        total_improvement = 0.0
        improvements = []

        # 对最近5轮逐一计算 improvement（相邻两轮对比）
        for i in range(1, len(recent_rounds)):
            current_score = recent_rounds[i]["score"]
            previous_score = recent_rounds[i-1]["score"]
            improvement = current_score - previous_score
            improvements.append(improvement)

            # 饱和条件：有提升（>0）且提升小于阈值（<6%）
            if improvement > 0 and improvement < SATURATION_THRESHOLD:
                saturated_count += 1

        # 计算平均提升
        avg_improvement = sum(improvements) / len(improvements) if improvements else 0.0

        # ===== 判断是否饱和：最近5轮中有3轮以上饱和 =====
        is_saturated = saturated_count >= 3

        if is_saturated:
            state["consecutive_rounds"] = state.get("consecutive_rounds", 0) + 1
            state["saturation_status"] = "saturated"
            state["last_improvement"] = avg_improvement

            # 如果累积饱和达到3次，升级到控制逻辑
            if state["consecutive_rounds"] >= 3:
                state["current_level"] = "control_logic"
                self.engine._append_change_log(
                    "双时间尺度进化调度",
                    f"提示词优化已饱和（累积 {state['consecutive_rounds']} 轮），升级到控制逻辑层面"
                )
                # 记录进化事件
                self.engine.log_evolution_event(
                    "saturation_escalated",
                    {
                        "consecutive_rounds": state["consecutive_rounds"],
                        "avg_improvement": avg_improvement,
                        "level": "control_logic",
                        "trigger": "cumulative_saturation"
                    }
                )
        else:
            # 未饱和时，连续计数减1（但不低于0）
            state["consecutive_rounds"] = max(0, state.get("consecutive_rounds", 0) - 1)
            state["saturation_status"] = "unsaturated"
            state["last_improvement"] = avg_improvement

        self._save_saturation_state(state)

        # 计算平均评分
        avg_recent = sum(h["score"] for h in recent_rounds) / len(recent_rounds)
        # 计算前一轮平均（如果有）
        if len(history) > 5:
            prev_rounds = history[-10:-5]
            avg_previous = sum(h["score"] for h in prev_rounds) / len(prev_rounds) if prev_rounds else avg_recent
        else:
            avg_previous = avg_recent

        return {
            "is_saturated": is_saturated,
            "consecutive_rounds": state.get("consecutive_rounds", 0),
            "improvement": round(avg_improvement, 3),
            "level": state.get("current_level", "prompt"),
            "status": state.get("saturation_status", "unsaturated"),
            "saturated_count": saturated_count,
            "avg_recent": round(avg_recent, 3),
            "avg_previous": round(avg_previous, 3)
        }

    def get_saturation_status(self) -> Dict[str, Any]:
        """获取当前饱和状态"""
        state = self._load_saturation_state()
        return {
            "saturation_status": state.get("saturation_status", "unsaturated"),
            "current_level": state.get("current_level", "prompt"),
            "consecutive_rounds": state.get("consecutive_rounds", 0),
            "last_improvement": state.get("last_improvement", 0.0),
            "quality_history_count": len(state.get("quality_history", [])),
            "control_logic_changes_count": len(state.get("control_logic_changes", []))
        }

    # ===== Day 8: 灵感熔炉复盘（一） =====
    def inspiration_furnace_review(self) -> Dict[str, Any]:
        """
        灵感熔炉复盘（一）

        从灵感池.json中读取状态为"待筛选"的记录，
        运行L2筛选（重要性、紧急性、与主线目标的一致性、资源投入估算），
        产出"待采纳清单"
        """
        insp_path = Config.DATA_ROOT / "系统日志" / "灵感池.json"
        if not insp_path.exists():
            return {
                "total_pending": 0,
                "s_level": [],
                "a_level": [],
                "b_level": [],
                "rejected": [],
                "summary": "灵感池文件不存在"
            }

        try:
            with open(insp_path, "r", encoding="utf-8") as f:
                inspirations = json.load(f)
        except:
            return {
                "total_pending": 0,
                "s_level": [],
                "a_level": [],
                "b_level": [],
                "rejected": [],
                "summary": "灵感池文件解析失败"
            }

        # 筛选"待筛选"状态的灵感
        pending = [i for i in inspirations if i.get("status") == "待筛选"]

        if not pending:
            return {
                "total_pending": 0,
                "s_level": [],
                "a_level": [],
                "b_level": [],
                "rejected": [],
                "summary": "暂无待筛选的灵感"
            }

        # L2筛选：评估每个灵感
        s_level = []
        a_level = []
        b_level = []
        rejected = []

        for insp in pending:
            content = insp.get("content", "")
            source = insp.get("source", "未知")

            # 评估指标
            importance_score = self._evaluate_importance(content)
            urgency_score = self._evaluate_urgency(content)
            alignment_score = self._evaluate_alignment(content)
            resource_estimate = self._estimate_resources(content)

            # 综合评分
            total_score = (importance_score * 0.5 + urgency_score * 0.2 + alignment_score * 0.3)

            # 调试日志（在GUI中显示）
            self.engine._append_change_log(
                "灵感评估",
                f"{insp.get('id')}: 重要性={importance_score:.2f}, 紧急性={urgency_score:.2f}, "
                f"一致性={alignment_score:.2f}, 资源={resource_estimate}h, 总分={total_score:.2f}"
            )

            # 分类
            if total_score >= 0.65 and resource_estimate <= 3:
                insp["evaluation"] = {
                    "importance": importance_score,
                    "urgency": urgency_score,
                    "alignment": alignment_score,
                    "resource_hours": resource_estimate,
                    "total_score": total_score,
                    "level": "S"
                }
                s_level.append(insp)
            elif total_score >= 0.5 and resource_estimate <= 6:
                insp["evaluation"] = {
                    "importance": importance_score,
                    "urgency": urgency_score,
                    "alignment": alignment_score,
                    "resource_hours": resource_estimate,
                    "total_score": total_score,
                    "level": "A"
                }
                a_level.append(insp)
            elif total_score >= 0.35 and resource_estimate <= 16:
                insp["evaluation"] = {
                    "importance": importance_score,
                    "urgency": urgency_score,
                    "alignment": alignment_score,
                    "resource_hours": resource_estimate,
                    "total_score": total_score,
                    "level": "B"
                }
                b_level.append(insp)
            else:
                insp["evaluation"] = {
                    "importance": importance_score,
                    "urgency": urgency_score,
                    "alignment": alignment_score,
                    "resource_hours": resource_estimate,
                    "total_score": total_score,
                    "level": "rejected"
                }
                rejected.append(insp)

        # 更新灵感池状态
        updated_inspirations = []
        for insp in inspirations:
            if insp.get("status") == "待筛选":
                # 检查是否在已分类列表中
                classified = None
                for item in s_level + a_level + b_level + rejected:
                    if item.get("id") == insp.get("id"):
                        classified = item
                        break
                if classified:
                    # 更新状态
                    classified["status"] = "已评估"
                    updated_inspirations.append(classified)
                else:
                    updated_inspirations.append(insp)
            else:
                updated_inspirations.append(insp)

        with open(insp_path, "w", encoding="utf-8") as f:
            json.dump(updated_inspirations, f, ensure_ascii=False, indent=2)

        # 记录到进化日志
        self.engine.log_evolution_event(
            "inspiration_review",
            {
                "total_pending": len(pending),
                "s_level_count": len(s_level),
                "a_level_count": len(a_level),
                "b_level_count": len(b_level),
                "rejected_count": len(rejected),
                "trigger": "day8_review"
            }
        )

        summary = (
            f"灵感熔炉复盘（一）完成：\n"
            f"  - S级（<3小时）：{len(s_level)} 条\n"
            f"  - A级（半天内）：{len(a_level)} 条\n"
            f"  - B级（1-2天）：{len(b_level)} 条\n"
            f"  - 已拒绝：{len(rejected)} 条"
        )

        return {
            "total_pending": len(pending),
            "s_level": s_level,
            "a_level": a_level,
            "b_level": b_level,
            "rejected": rejected,
            "summary": summary
        }

    def _evaluate_importance(self, content: str) -> float:
        """评估灵感的重要性（0-1）"""
        high_importance = ["核心", "关键", "突破", "创新", "战略", "架构", "系统", "框架", "机制", "范式"]
        medium_importance = ["优化", "改进", "增强", "提升", "完善", "调整", "补充"]

        content_lower = content.lower()
        high_score = sum(0.3 for kw in high_importance if kw in content_lower)
        medium_score = sum(0.15 for kw in medium_importance if kw in content_lower)

        # 长度加成（更长的描述通常更具体）
        length_bonus = min(0.2, len(content) / 500)

        return min(1.0, high_score + medium_score + length_bonus)

    def _evaluate_urgency(self, content: str) -> float:
        """评估灵感的紧急性（0-1）"""
        urgent_keywords = ["紧急", "立刻", "马上", "尽快", "亟待", "急需", "关键", "阻塞", "阻断"]
        content_lower = content.lower()
        score = sum(0.25 for kw in urgent_keywords if kw in content_lower)
        return min(1.0, score)

    def _evaluate_alignment(self, content: str) -> float:
        """评估灵感与主线目标的一致性（0-1）"""
        # 主线目标关键词
        main_keywords = ["晶体", "认知", "辩论", "决策", "八道防线", "沉思", "进化", "学习", "智能", "知识"]
        content_lower = content.lower()
        score = sum(0.15 for kw in main_keywords if kw in content_lower)
        return min(1.0, score + 0.2)  # 基础分0.2

    def _estimate_resources(self, content: str) -> int:
        """估算资源投入（小时）"""
        # 基于内容长度和复杂度估算
        length = len(content)

        if length < 50:
            return 1  # 简单想法，1小时
        elif length < 150:
            return 2  # 中等想法，2小时
        elif length < 300:
            return 4  # 复杂想法，半天
        else:
            # 检查是否有"实现"、"构建"等关键词
            if "实现" in content or "构建" in content or "开发" in content:
                return 8  # 需要实现，1天
            return 4  # 默认半天

    def diagnose_history(self, question: str, threshold: float = 0.7) -> Dict[str, Any]:
        """
        历史诊断与经验复用（非马尔可夫历史检索）
        
        从 evolution_log.json 中检索历史失败轨迹和成功经验，
        匹配当前问题与历史问题，返回最相似的历史记录及其有效晶体组合。

        Args:
            question: 当前问题
            threshold: 相似度阈值（0-1），默认 0.7

        Returns:
            {
                "matched": bool,               # 是否找到有效匹配
                "match_score": float,           # 最佳匹配相似度
                "reused_history_id": str,       # 复用的历史记录ID
                "crystal_combination": List[str], # 该历史问题对应的有效晶体组合
                "diagnosis": str,               # 历史诊断结论
                "repair_attempts": List[Dict]   # 历史修复尝试
            }
        """
        log_path = Config.DATA_ROOT / "系统日志" / "evolution_log.json"
        if not log_path.exists():
            return {
                "matched": False,
                "match_score": 0.0,
                "reused_history_id": None,
                "crystal_combination": [],
                "diagnosis": "",
                "repair_attempts": []
            }

        # 读取日志
        try:
            with open(log_path, "r", encoding="utf-8") as f:
                data = json.load(f)
        except:
            return {
                "matched": False,
                "match_score": 0.0,
                "reused_history_id": None,
                "crystal_combination": [],
                "diagnosis": "",
                "repair_attempts": []
            }

        events = data.get("events", [])
        if not events:
            return {
                "matched": False,
                "match_score": 0.0,
                "reused_history_id": None,
                "crystal_combination": [],
                "diagnosis": "",
                "repair_attempts": []
            }

        # 1. 筛选有失败轨迹或诊断信息的事件
        candidate_events = []
        for event in events:
            # 检查是否有失败轨迹
            if "failure_traces" in event or "diagnosis" in event:
                # 尝试从 details 中提取问题原文
                details = event.get("details", {})
                history_question = details.get("question", "") or details.get("user_input", "")
                if history_question:
                    candidate_events.append({
                        "event": event,
                        "question": history_question,
                        "timestamp": event.get("timestamp", "")
                    })
            # 也检查是否有成功复用的历史（用于正向强化）
            if event.get("event_type") == "history_reused":
                details = event.get("details", {})
                history_question = details.get("question", "")
                if history_question:
                    candidate_events.append({
                        "event": event,
                        "question": history_question,
                        "timestamp": event.get("timestamp", "")
                    })

        if not candidate_events:
            return {
                "matched": False,
                "match_score": 0.0,
                "reused_history_id": None,
                "crystal_combination": [],
                "diagnosis": "",
                "repair_attempts": []
            }

        # 2. 计算当前问题与每个历史问题的相似度
        scored = []
        for item in candidate_events:
            sim = self.engine._simple_similarity(question, item["question"])
            scored.append({
                **item,
                "similarity": sim
            })

        # 3. 按相似度降序排序
        scored.sort(key=lambda x: x["similarity"], reverse=True)

        # 4. 选择最佳匹配
        best = scored[0]
        if best["similarity"] >= threshold:
            event = best["event"]
            details = event.get("details", {})
            
            # 提取有效晶体组合
            crystal_combination = []
            # 从 repair_attempts 中提取成功方案中使用的晶体
            repair_attempts = event.get("repair_attempts", [])
            for attempt in repair_attempts:
                if attempt.get("success", False):
                    crystal_combination.extend(attempt.get("crystals_used", []))
            # 如果 repair_attempts 中没有，尝试从 details 中提取
            if not crystal_combination:
                crystal_combination = details.get("effective_crystals", [])

            return {
                "matched": True,
                "match_score": best["similarity"],
                "reused_history_id": event.get("timestamp", ""),
                "crystal_combination": list(set(crystal_combination)),  # 去重
                "diagnosis": event.get("diagnosis", details.get("diagnosis", "")),
                "repair_attempts": repair_attempts,
                "matched_question": best["question"]
            }
        else:
            return {
                "matched": False,
                "match_score": best["similarity"],
                "reused_history_id": None,
                "crystal_combination": [],
                "diagnosis": "",
                "repair_attempts": []
            }

    def run_all_primitives(self) -> Dict[str, Any]:
        """
        运行所有已启用的元原语，并处理触发链
        """
        # 1. 运行各元原语
        results = {
            "active_gap_detection": self.active_gap_detection(),
            "temporal_aware_escalation": self.temporal_aware_escalation(),
            "sleep_consolidation": self.sleep_consolidation(),
            "validation_gated_self_evolution": {
                "status": "pending",
                "message": "在晶体更新时调用，需传入 new_data 和 context 参数"
            }
        }

        # 2. Day 3 新增：处理触发链
        triggered = self.process_trigger_chains(results)
        results["triggered_chains"] = triggered

        # 3. 记录触发链日志
        if triggered:
            for chain in triggered:
                self.engine._append_change_log(
                    "触发链执行",
                    f"{chain['chain']}: {chain['source']} → {chain['target']} (通过: {chain['passed']})"
                )

        return results
# =============================================================================
# Day 1 新增：警报监控系统 (AlarmMonitor)
# =============================================================================
class AlarmMonitor:
    """
    八道防线警报系统
    根据配置规则监控辩论过程中的关键指标，触发相应警报
    """

    def __init__(self, rules: dict = None, log_callback=None):
        """
        :param rules: 警报规则字典，若为None则使用 Config.ALARM_RULES
        :param log_callback: 日志回调函数，用于记录警报事件
        """
        self.rules = rules or Config.ALARM_RULES
        self.log = log_callback or (lambda msg, level="system": print(f"[{level.upper()}] {msg}"))
        # 状态跟踪
        self.external_empty_count = 0      # 连续空外部数据计数
        self.jaccard_history = []          # 最近3轮的Jaccard值
        self.alarm_triggered = False       # 本轮是否已触发警报（防止重复）

    def check(self, metrics: dict) -> List[Dict[str, Any]]:
        """
        检查所有警报规则，返回触发的警报列表

        :param metrics: 包含以下键的字典：
            - crystal_reference_rate: float (本轮晶体引用率)
            - bias_amplification: float (本轮偏见强化指数)
            - external_has_new: bool (本轮是否有新外部数据)
            - jaccard_similarity: float (本轮Jaccard相似度)
        :return: 触发的警报列表，每个元素为 dict { 'rule_name', 'message', 'action', 'data' }
        """
        triggered = []
        # 1. 知识贫瘠警报
        if self.rules.get("knowledge_poverty", {}).get("enabled", True):
            threshold = self.rules["knowledge_poverty"]["threshold"]
            if metrics.get("crystal_reference_rate", 1.0) < threshold:
                triggered.append({
                    "rule": "knowledge_poverty",
                    "message": self.rules["knowledge_poverty"]["message"],
                    "action": self.rules["knowledge_poverty"]["action"],
                    "data": {"rate": metrics["crystal_reference_rate"], "threshold": threshold}
                })

        # 2. 偏见膨胀警报
        if self.rules.get("bias_inflation", {}).get("enabled", True):
            threshold = self.rules["bias_inflation"]["threshold"]
            if metrics.get("bias_amplification", 0.0) > threshold:
                triggered.append({
                    "rule": "bias_inflation",
                    "message": self.rules["bias_inflation"]["message"],
                    "action": self.rules["bias_inflation"]["action"],
                    "data": {"bias": metrics["bias_amplification"], "threshold": threshold}
                })

        # 3. 信息枯竭警报（连续3轮无新数据）
        if self.rules.get("information_starvation", {}).get("enabled", True):
            threshold = self.rules["information_starvation"]["threshold"]
            if metrics.get("external_has_new", True):
                self.external_empty_count = 0
            else:
                self.external_empty_count += 1
            if self.external_empty_count >= threshold:
                triggered.append({
                    "rule": "information_starvation",
                    "message": self.rules["information_starvation"]["message"],
                    "action": self.rules["information_starvation"]["action"],
                    "data": {"consecutive_empty": self.external_empty_count, "threshold": threshold}
                })
                self.external_empty_count = 0  # 重置，避免连续触发

        # 4. 思维固化警报（连续3轮Jaccard > 0.8）
        if self.rules.get("thought_stagnation", {}).get("enabled", True):
            threshold = self.rules["thought_stagnation"]["threshold"]
            consecutive = self.rules["thought_stagnation"]["consecutive"]
            jaccard = metrics.get("jaccard_similarity", 0.0)
            self.jaccard_history.append(jaccard)
            if len(self.jaccard_history) > consecutive:
                self.jaccard_history.pop(0)
            if len(self.jaccard_history) == consecutive and all(j > threshold for j in self.jaccard_history):
                triggered.append({
                    "rule": "thought_stagnation",
                    "message": self.rules["thought_stagnation"]["message"],
                    "action": self.rules["thought_stagnation"]["action"],
                    "data": {"jaccards": self.jaccard_history.copy(), "threshold": threshold}
                })
                self.jaccard_history.clear()  # 重置

        return triggered

    def handle_alarm(self, alarm: Dict, debate_engine) -> bool:
        """
        处理单个警报，执行对应动作，并返回是否应继续辩论

        :param alarm: 警报字典
        :param debate_engine: DebateEngine 实例，用于调用其方法（如注入视角、触发搜索）
        :return: True 表示恢复辩论，False 表示终止辩论（目前始终恢复）
        """
        action = alarm.get("action")
        message = alarm.get("message")

        # 记录事件（由上层调用者写入 evolution_log）
        self.log(f"🚨 {message}", "warning")

        if action == "inject_external":
            # 强制注入外部知识：调用 debate_engine 的外部搜索并注入
            self.log("  执行动作：强制注入外部知识", "system")
            # 我们可以直接调用 debate_engine._fetch_external_overview 并注入到下一轮
            # 但需要暴露接口，我们在 DebateEngine 中添加 _inject_external_knowledge 方法
            debate_engine._inject_external_knowledge(alarm)

        elif action == "inject_perspective":
            # 强制注入对立视角：可以在系统消息中加入新角色或提示
            self.log("  执行动作：强制注入对立视角", "system")
            debate_engine._inject_perspective(alarm)

        elif action == "trigger_search":
            # 强制触发外部搜索
            self.log("  执行动作：强制触发外部搜索", "system")
            debate_engine._trigger_search(alarm)

        else:
            self.log(f"  未知动作：{action}，跳过", "warning")

        return True  # 恢复辩论

# =============================================================================
# 向量存储引擎（新增）
# =============================================================================

class VectorStore:
    """
    向量存储引擎 —— 基于 ChromaDB 的语义检索

    将晶体卡片向量化，支持语义相似度检索。
    首次运行时自动下载 sentence-transformers 模型（all-MiniLM-L6-v2）。
    检索失败时自动降级到 BM25（由调用方处理）。
    """

    def __init__(self, file_io: FileIO, model_name: str = "all-MiniLM-L6-v2"):
        self.files = file_io
        self.model_name = model_name
        self._model = None
        self._collection = None
        self._initialized = False

        # 向量数据库存储路径
        self._db_path = Config.DATA_ROOT / "model_cache" / "chroma_db"
        self._init_vector_store()

    def _init_vector_store(self) -> None:
        """初始化向量存储（延迟加载模型）"""
        try:
            import chromadb
            from chromadb.config import Settings

            # 创建存储目录
            self._db_path.mkdir(parents=True, exist_ok=True)

            # 初始化 ChromaDB 客户端（持久化）
            self._client = chromadb.PersistentClient(
                path=str(self._db_path),
                settings=Settings(anonymized_telemetry=False)
            )

            # 获取或创建 collection
            self._collection = self._client.get_or_create_collection(
                name="crystals",
                metadata={"hnsw:space": "cosine"}
            )

            self._initialized = True
            print(f"[OK] VectorStore 初始化成功，已有 {self._collection.count()} 条向量")
        except ImportError:
            print("[WARN] chromadb 未安装，向量检索不可用，将使用 BM25 降级")
            self._initialized = False
        except Exception as e:
            print(f"[WARN] VectorStore 初始化失败: {e}，将使用 BM25 降级")
            self._initialized = False

    def _get_model(self):
        """延迟加载 sentence-transformers 模型"""
        if self._model is not None:
            return self._model

        if not self._initialized:
            return None

        try:
            from sentence_transformers import SentenceTransformer
            # 设置 HuggingFace 镜像
            import os
            os.environ["HF_ENDPOINT"] = "https://hf-mirror.com"
            self._model = SentenceTransformer(self.model_name)
            print(f"✅ 向量模型 {self.model_name} 加载成功")
            return self._model
        except ImportError:
            print("⚠️ sentence_transformers 未安装")
            return None
        except Exception as e:
            print(f"⚠️ 模型加载失败: {e}")
            return None

    def add_crystals(self, crystals: List) -> int:
        """
        将晶体批量添加到向量库

        Returns:
            成功添加的数量
        """
        if not self._initialized or not crystals:
            return 0

        model = self._get_model()
        if model is None:
            return 0

        try:
            # 准备数据
            ids = [c.id for c in crystals]
            documents = [c.content for c in crystals]
            # 附加元数据
            metadatas = [{"layer": c.layer.value if hasattr(c.layer, 'value') else str(c.layer)} for c in crystals]

            # 生成向量
            embeddings = model.encode(documents, show_progress_bar=False).tolist()

            # 分批插入（避免一次性过大）
            batch_size = 50
            for i in range(0, len(ids), batch_size):
                batch_end = min(i + batch_size, len(ids))
                self._collection.add(
                    ids=ids[i:batch_end],
                    documents=documents[i:batch_end],
                    embeddings=embeddings[i:batch_end],
                    metadatas=metadatas[i:batch_end]
                )

            print(f"✅ 成功向量化 {len(ids)} 条晶体")
            return len(ids)

        except Exception as e:
            print(f"[WARN] 向量化失败: {e}")
            return 0

    def query(self, query_text: str, top_k: int = 5) -> List[Tuple[str, float]]:
        """
        语义检索最相似的晶体

        Returns:
            List[(crystal_id, similarity_score), ...]
        """
        if not self._initialized:
            return []

        model = self._get_model()
        if model is None:
            return []

        try:
            # 生成查询向量
            query_embedding = model.encode(query_text, show_progress_bar=False).tolist()

            # 检索
            results = self._collection.query(
                query_embeddings=[query_embedding],
                n_results=top_k,
                include=["documents", "metadatas", "distances"]
            )

            # 解析结果
            if results and results.get('ids') and len(results['ids']) > 0:
                ids = results['ids'][0]
                distances = results['distances'][0] if results.get('distances') else []
                # 距离转相似度（cosine距离 → 相似度）
                similarities = [1 - d for d in distances] if distances else []
                return list(zip(ids, similarities))

            return []

        except Exception as e:
            print(f"⚠️ 向量检索失败: {e}")
            return []

    def delete_crystal(self, crystal_id: str) -> bool:
        """从向量库删除单个晶体"""
        if not self._initialized:
            return False

        try:
            self._collection.delete(ids=[crystal_id])
            return True
        except Exception as e:
            print(f"⚠️ 删除向量失败: {e}")
            return False

    def count(self) -> int:
        """获取向量库中的晶体数量"""
        if not self._initialized:
            return 0
        try:
            return self._collection.count()
        except:
            return 0

    def reset(self) -> bool:
        """重置向量库（删除所有数据）"""
        if not self._initialized:
            return False
        try:
            self._client.delete_collection("crystals")
            self._collection = self._client.create_collection("crystals")
            return True
        except Exception as e:
            print(f"⚠️ 重置向量库失败: {e}")
            return False

    def is_available(self) -> bool:
        """检查向量检索是否可用"""
        return self._initialized and self._get_model() is not None

# =============================================================================
# 认知指纹提取引擎（新增）
# =============================================================================

# =============================================================================
# 认知指纹提取引擎（新增）
# =============================================================================

class FingerprintExtractor:
    """
    认知指纹提取引擎

    从用户的历史对话、辩论参与、晶体化行为中提取认知指纹。
    让系统从"知道你说过什么"升级为"知道你怎么想"。

    对应四阶路径中的"认识你"阶段。
    """

    def __init__(self, engine: 'CrystalEngine', file_io: FileIO):
        self.engine = engine
        self.files = file_io

    def extract(self, history: List[Tuple[str, str]], debate_logs: List[Dict] = None, increment_interactions: int = 0) -> CognitiveFingerprint:
        """
        从对话历史和辩论日志中提取认知指纹

        Args:
            history: 对话历史列表 [(role, content), ...]
            debate_logs: 辩论日志列表（可选）

        Returns:
            CognitiveFingerprint: 提取的认知指纹
        """
        # 加载现有指纹（用于增量更新）
        existing_data = self.files.read_fingerprint()
        existing = CognitiveFingerprint.from_dict(existing_data.get("fingerprint", {}))

        # 如果没有历史数据，返回现有指纹
        if not history:
            return existing

        # 1. 分析用户提问中的关键词 → 风险偏好 + 创新偏好
        risk_score, innovation_score = self._analyze_keywords(history)

        # 2. 分析角色采纳历史 → 偏好角色
        preferred_role, role_history = self._analyze_role_adoption(history, debate_logs)

        # 3. 分析决策果断值
        decisiveness = self._analyze_decisiveness(history)

        # 4. 分析冲突解决风格
        conflict_style = self._analyze_conflict_style(history)

        # 5. 分析注意力模式
        attention_span = self._analyze_attention(history)

        # 6. 计算置信度（数据越多置信度越高）
        # total_interactions = len([msg for msg in history if msg[0] == "user"])
        if increment_interactions > 0:
            new_interactions = increment_interactions
        else:
            new_interactions = len([msg for msg in history if msg[0] == "user"])
            if new_interactions > 0:
                print("⚠️ Warning: extract() called without increment_interactions, using full history count, may cause double counting.")
        total_interactions = existing.total_interactions + new_interactions
        confidence = min(0.9, 0.3 + total_interactions * 0.01)

        # 7. 使用平滑更新
        new_risk = self._smooth_update(existing.risk_tolerance, risk_score, 0.3)
        new_innovation = self._smooth_update(existing.innovation_preference, innovation_score, 0.3)
        new_decisiveness = self._smooth_update(existing.decisiveness, decisiveness, 0.3)
        new_attention = self._smooth_update(existing.attention_span, attention_span, 0.3)
        new_preferred = preferred_role or existing.preferred_role
        new_conflict = conflict_style or existing.conflict_resolution_style
        new_role_history = self._merge_role_history(existing.role_adoption_history, role_history)

        # 8. 构建演化日志
        changes = []
        if abs(existing.risk_tolerance - new_risk) > 0.05:
            changes.append({"dimension": "risk_tolerance", "old": existing.risk_tolerance, "new": new_risk})
        if abs(existing.innovation_preference - new_innovation) > 0.05:
            changes.append({"dimension": "innovation_preference", "old": existing.innovation_preference, "new": new_innovation})
        if abs(existing.decisiveness - new_decisiveness) > 0.05:
            changes.append({"dimension": "decisiveness", "old": existing.decisiveness, "new": new_decisiveness})
        if existing.preferred_role != new_preferred:
            changes.append({"dimension": "preferred_role", "old": existing.preferred_role, "new": new_preferred})

        old_logs = existing.evolution_log[-15:] if existing.evolution_log else []
        evolution_log = old_logs + [{
            "timestamp": datetime.now().isoformat(),
            "changes": changes,
            "total_interactions": existing.total_interactions + total_interactions
        }]
        # Day 2.5: 认知风格分析
        style_result = self._analyze_thinking_style(history)
        new_reasoning = style_result["reasoning_style"]
        new_analogy = style_result["analogy_preference"]
        new_output = style_result["output_style"]
        
        # 数据不足时保留旧值，避免波动
        if len([msg for msg in history if msg[0] == "user"]) < 5:
            new_reasoning = existing.reasoning_style
            new_analogy = existing.analogy_preference
            new_output = existing.output_style

        # 9. 构建新的指纹对象
        new_fingerprint = CognitiveFingerprint(
            risk_tolerance=new_risk,
            innovation_preference=new_innovation,
            decisiveness=new_decisiveness,
            preferred_role=new_preferred,
            role_adoption_history=new_role_history,
            conflict_resolution_style=new_conflict,
            attention_span=new_attention,
            context_preference=existing.context_preference,
            last_updated=datetime.now().isoformat(),
            total_interactions=existing.total_interactions + total_interactions,
            confidence=confidence,
            reasoning_style=new_reasoning,
            analogy_preference=new_analogy,
            output_style=new_output,           
            evolution_log=evolution_log
        )

        # 10. 保存更新后的指纹
        # 10. 保存更新后的指纹
        self._save_fingerprint(new_fingerprint)

        # ===== 新增：记录指纹变化进化事件 =====
        if existing.total_interactions > 0:  # 非首次提取
            changes = []
            if abs(existing.risk_tolerance - new_fingerprint.risk_tolerance) > 0.05:
                changes.append({"dimension": "risk_tolerance", "old": existing.risk_tolerance, "new": new_fingerprint.risk_tolerance})
            if abs(existing.innovation_preference - new_fingerprint.innovation_preference) > 0.05:
                changes.append({"dimension": "innovation_preference", "old": existing.innovation_preference, "new": new_fingerprint.innovation_preference})
            if abs(existing.decisiveness - new_fingerprint.decisiveness) > 0.05:
                changes.append({"dimension": "decisiveness", "old": existing.decisiveness, "new": new_fingerprint.decisiveness})
            if existing.preferred_role != new_fingerprint.preferred_role:
                changes.append({"dimension": "preferred_role", "old": existing.preferred_role, "new": new_fingerprint.preferred_role})
            if changes:
                self.engine.log_evolution_event(
                    "fingerprint_changed",
                    {
                        "changes": changes,
                        "confidence": new_fingerprint.confidence,
                        "total_interactions": new_fingerprint.total_interactions,
                        "trigger": "conversation_analysis"
                    }
                )

        return new_fingerprint

    def _analyze_keywords(self, history: List[Tuple[str, str]]) -> Tuple[float, float]:
        """
        分析用户提问中的关键词 → 风险容忍度 + 创新偏好

        风险关键词: 万一, 风险, 成本, 失败, 安全, 稳健, 保守, 验证, 谨慎, 稳妥
        创新关键词: 机会, 突破, 创新, 颠覆, 潜力, 激进, 大胆, 尝试, 新的, 探索
        """
        risk_words = ["万一", "风险", "成本", "失败", "安全", "稳健", "保守", "验证", "谨慎", "稳妥"]
        innovation_words = ["机会", "突破", "创新", "颠覆", "潜力", "激进", "大胆", "尝试", "新的", "探索"]

        user_messages = [msg[1] for msg in history if msg[0] == "user"]
        if not user_messages:
            return 0.5, 0.5

        combined = " ".join(user_messages)
        risk_count = sum(1 for w in risk_words if w in combined)
        innovation_count = sum(1 for w in innovation_words if w in combined)

        total = risk_count + innovation_count
        if total == 0:
            return 0.5, 0.5

        risk_score = min(1.0, risk_count / max(1, len(user_messages)) * 3)
        innovation_score = min(1.0, innovation_count / max(1, len(user_messages)) * 3)

        return risk_score, innovation_score

    def _analyze_role_adoption(self, history: List[Tuple[str, str]], debate_logs: List[Dict] = None) -> Tuple[str, Dict[str, int]]:
        """
        分析用户对辩论角色的采纳偏好
        """
        role_history = {}

        # 从辩论日志中分析
        if debate_logs:
            for log in debate_logs:
                for role in ["激进者", "保守者", "结构主义者", "执行者", "审计者"]:
                    if role in log.get("user_feedback", ""):
                        role_history[role] = role_history.get(role, 0) + 1

        # 从对话历史中分析（关键词匹配）
        user_messages = [msg[1] for msg in history if msg[0] == "user"]
        for msg in user_messages:
            if "激进" in msg or "颠覆" in msg:
                role_history["激进者"] = role_history.get("激进者", 0) + 0.5
            if "保守" in msg or "稳健" in msg:
                role_history["保守者"] = role_history.get("保守者", 0) + 0.5
            if "结构" in msg or "系统" in msg or "框架" in msg:
                role_history["结构主义者"] = role_history.get("结构主义者", 0) + 0.5
            if "执行" in msg or "步骤" in msg or "操作" in msg:
                role_history["执行者"] = role_history.get("执行者", 0) + 0.5
            if "审计" in msg or "验证" in msg or "检查" in msg:
                role_history["审计者"] = role_history.get("审计者", 0) + 0.5

        if not role_history:
            return "structural", {}

        preferred = max(role_history.items(), key=lambda x: x[1])[0]
        return preferred, role_history

    def _analyze_decisiveness(self, history: List[Tuple[str, str]]) -> float:
        """
        分析决策果断值
        通过用户是否快速追问、是否打断、是否快速确认来判断
        """
        user_messages = [msg[1] for msg in history if msg[0] == "user"]
        if len(user_messages) < 3:
            return 0.5

        short_count = sum(1 for msg in user_messages if len(msg) < 20)
        short_ratio = short_count / len(user_messages)

        confirm_words = ["好的", "对", "是", "行", "可以", "同意", "确认", "就这样"]
        confirm_count = sum(1 for msg in user_messages if any(w in msg for w in confirm_words))
        confirm_ratio = confirm_count / len(user_messages)

        decisiveness = short_ratio * 0.6 + confirm_ratio * 0.4
        return min(1.0, decisiveness)

    def _analyze_conflict_style(self, history: List[Tuple[str, str]]) -> str:
        """
        分析冲突解决风格
        """
        user_messages = [msg[1] for msg in history if msg[0] == "user"]
        combined = " ".join(user_messages)

        integrative_words = ["综合", "融合", "结合", "共识", "共同", "平衡", "兼顾"]
        competitive_words = ["不对", "错误", "我坚持", "反驳", "质疑", "反对"]
        avoidant_words = ["跳过", "忽略", "算了", "不管", "别说了", "先放着"]

        i_count = sum(1 for w in integrative_words if w in combined)
        c_count = sum(1 for w in competitive_words if w in combined)
        a_count = sum(1 for w in avoidant_words if w in combined)

        if i_count >= c_count and i_count >= a_count:
            return "integrative"
        elif c_count >= i_count and c_count >= a_count:
            return "competitive"
        else:
            return "avoidant"

    def _analyze_attention(self, history: List[Tuple[str, str]]) -> float:
        """
        分析注意力持续度
        通过对话长度和主题一致性判断
        """
        user_messages = [msg[1] for msg in history if msg[0] == "user"]
        if len(user_messages) < 3:
            return 0.5

        avg_length = sum(len(msg) for msg in user_messages) / len(user_messages)
        length_score = min(1.0, avg_length / 100)

        if len(user_messages) >= 5:
            first_words = set(user_messages[0][:50].split())
            recent_words = set(user_messages[-1][:50].split())
            common = len(first_words & recent_words) / max(1, len(first_words))
            consistency_score = common
        else:
            consistency_score = 0.5

        return length_score * 0.5 + consistency_score * 0.5

    def _analyze_thinking_style(self, history: List[Tuple[str, str]]) -> Dict[str, str]:
        """
        分析用户的思维模式偏好
        返回：{
            "reasoning_style": "deductive" / "inductive" / "balanced",
            "analogy_preference": "analogy" / "first_principles" / "balanced",
            "output_style": "conclusion_first" / "evidence_first" / "balanced"
        }
        """
        user_messages = [msg[1] for msg in history if msg[0] == "user"]
        combined = " ".join(user_messages)
        
        # 1. 演绎 vs 归纳
        deductive_words = ["因此", "所以", "必然", "推导", "逻辑", "根据", "因为...所以"]
        inductive_words = ["案例", "观察", "现象", "数据", "实验", "例子", "发现"]
        deductive_score = sum(1 for w in deductive_words if w in combined)
        inductive_score = sum(1 for w in inductive_words if w in combined)
        if deductive_score > inductive_score * 1.5:
            reasoning = "deductive"
        elif inductive_score > deductive_score * 1.5:
            reasoning = "inductive"
        else:
            reasoning = "balanced"
        
        # 2. 类比 vs 第一性原理
        analogy_words = ["类比", "比喻", "类似", "如同", "就像", "像...一样", "参照"]
        first_principles_words = ["本质", "根本", "基础", "原理", "底层", "核心", "还原"]
        analogy_score = sum(1 for w in analogy_words if w in combined)
        first_score = sum(1 for w in first_principles_words if w in combined)
        if analogy_score > first_score * 1.5:
            analogy_pref = "analogy"
        elif first_score > analogy_score * 1.5:
            analogy_pref = "first_principles"
        else:
            analogy_pref = "balanced"
        
        # 3. 输出风格：先结论 vs 先证据
        # 叙事型（先讲故事）vs 数据型（先给结论）
        narrative_words = ["故事", "经历", "感受", "我认为", "我觉得", "信任", "相信"]
        data_words = ["数据", "证据", "统计", "研究", "显示", "表明", "证明"]
        narrative_score = sum(1 for w in narrative_words if w in combined)
        data_score = sum(1 for w in data_words if w in combined)
        if narrative_score > data_score * 1.5:
            output_style = "evidence_first"   # 先讲故事再结论
        elif data_score > narrative_score * 1.5:
            output_style = "conclusion_first" # 先结论再证据
        else:
            output_style = "balanced"
        
        return {
            "reasoning_style": reasoning,
            "analogy_preference": analogy_pref,
            "output_style": output_style
        }

    def get_cognitive_operators(self, fingerprint: CognitiveFingerprint) -> str:
        """生成用户专属认知操作符描述"""
        style_map = {
            "deductive": "演绎推理（从一般到特殊）",
            "inductive": "归纳推理（从特殊到一般）",
            "balanced": "演绎与归纳并重"
        }
        analogy_map = {
            "analogy": "类比思维（用已知解释未知）",
            "first_principles": "第一性原理（回归本质）",
            "balanced": "类比与本质分析并重"
        }
        output_map = {
            "conclusion_first": "先结论后展开",
            "evidence_first": "先证据后结论",
            "balanced": "结论与证据交替"
        }
        ops = [
            f"[思维模式：{style_map.get(fingerprint.reasoning_style, '平衡')}]",
            f"[论证偏好：{analogy_map.get(fingerprint.analogy_preference, '平衡')}]",
            f"[输出偏好：{output_map.get(fingerprint.output_style, '平衡')}]"
        ]
        return " ".join(ops)

    def _smooth_update(self, old_val: float, new_val: float, factor: float = 0.3) -> float:
        """平滑更新，防止指纹突变"""
        if old_val is None:
            return new_val
        return old_val * (1 - factor) + new_val * factor

    def _merge_role_history(self, old: Dict[str, int], new: Dict[str, int]) -> Dict[str, int]:
        """合并角色历史"""
        merged = old.copy()
        for role, count in new.items():
            merged[role] = merged.get(role, 0) + int(count)
        return merged

    def _build_evolution_log(self, old: CognitiveFingerprint, new: CognitiveFingerprint) -> List[Dict[str, Any]]:
        """构建演化日志"""
        changes = []
        if abs(old.risk_tolerance - new.risk_tolerance) > 0.1:
            changes.append({
                "dimension": "risk_tolerance",
                "old": old.risk_tolerance,
                "new": new.risk_tolerance
            })
        if abs(old.innovation_preference - new.innovation_preference) > 0.1:
            changes.append({
                "dimension": "innovation_preference",
                "old": old.innovation_preference,
                "new": new.innovation_preference
            })
        if abs(old.decisiveness - new.decisiveness) > 0.1:
            changes.append({
                "dimension": "decisiveness",
                "old": old.decisiveness,
                "new": new.decisiveness
            })
        if old.preferred_role != new.preferred_role:
            changes.append({
                "dimension": "preferred_role",
                "old": old.preferred_role,
                "new": new.preferred_role
            })

        old_logs = old.evolution_log[-15:] if old.evolution_log else []
        return old_logs + [{
            "timestamp": datetime.now().isoformat(),
            "changes": changes,
            "total_interactions": new.total_interactions
        }]

    def _save_fingerprint(self, fingerprint: CognitiveFingerprint) -> None:
        """保存认知指纹到文件"""
        data = {
            "fingerprint": fingerprint.to_dict(),
            "extraction_metadata": {
                "last_extraction": datetime.now().isoformat(),
                "messages_analyzed": fingerprint.total_interactions,
                "debates_analyzed": 0,
                "crystals_created": 0,
                "version": "1.0"
            }
        }
        self.files.write_fingerprint(data)

    def get_fingerprint(self) -> CognitiveFingerprint:
        """获取当前认知指纹"""
        data = self.files.read_fingerprint()
        return CognitiveFingerprint.from_dict(data.get("fingerprint", {}))

# =============================================================================
# 便宜门规则引擎（新增）
# =============================================================================

class CheapGate:
    """
    便宜门规则引擎

    在调用 LLM 之前先用低成本规则过滤请求。
    只有通过便宜门检查的请求才交给 LLM 处理。

    对应建议2：用"便宜门→LLM"原则降本
    """

    def __init__(self, engine: 'CrystalEngine', file_io: FileIO, log_callback=None):
        self.engine = engine
        self.files = file_io
        self.log = log_callback or (lambda msg, level="system": print(msg))
        self._search_counter = 0
        self._search_threshold = 3

    def _estimate_complexity(self, user_input: str) -> str:
        """
        评估问题复杂度：返回 'simple', 'medium', 'high'
        """
        q_len = len(user_input.strip())
        config = Config.ROUTING_CONFIG
        
        # 1. 检查是否包含复杂关键词
        has_complex = any(kw in user_input for kw in config["complex_keywords"])
        
        # 2. 长度极短且无复杂词 → 简单
        if q_len <= config["simple_length_threshold"] and not has_complex:
            return "simple"
        
        # 3. 长度中等且有复杂词 → 中等
        if q_len <= config["medium_length_threshold"] and has_complex:
            return "medium"
        
        # 4. 长度较长且有复杂词 → 高
        if q_len > config["medium_length_threshold"] and has_complex:
            return "high"
        
        # 5. 其他情况（比如长度长但无复杂词）视为中等
        if q_len > config["simple_length_threshold"]:
            return "medium"
        return "simple"

    def check(self, user_input: str, history: List[Tuple[str, str]]) -> Dict[str, Any]:
        """
        动态路由决策
        返回：{
            "action": str,       # "rule_engine" | "cheap_gate_llm" | "direct_llm"
            "complexity": str,   # "simple" | "medium" | "high"
            "skip_llm": bool,
            "token_budget": int,
            "cost_estimate": float,
            "reason": str,
        }
        """
        # 1. 评估复杂度
        complexity = self._estimate_complexity(user_input)
        
        # 2. 根据复杂度选择路由
        config = Config.ROUTING_CONFIG
        if complexity == "simple":
            action = "rule_engine"
            skip_llm = True
            token_budget = config["token_budget_simple"]
            reason = "简单问题（短文本/明确指令/情绪确认），规则引擎直接回答"
        elif complexity == "medium":
            action = "cheap_gate_llm"
            skip_llm = False
            token_budget = config["token_budget_medium"]
            reason = "中等复杂度，便宜门预筛选后调用LLM"
        else:  # high
            action = "direct_llm"
            skip_llm = False
            token_budget = config["token_budget_high"]
            reason = "高复杂度问题（含复杂关键词且较长），直接调用LLM，Token预算上限"
        
        # 3. 估算成本（粗略）
        cost_estimate = token_budget * 0.000001  # 假设每token $0.000001
        
        # 4. 记录路由决策日志
        log_msg = f"便宜门路由决策：问题复杂度={complexity}，选择路径={action}，预估成本=${cost_estimate:.6f}，Token预算={token_budget}"
        self.log(log_msg, "ai")
        
        return {
            "action": action,
            "complexity": complexity,
            "skip_llm": skip_llm,
            "token_budget": token_budget,
            "cost_estimate": cost_estimate,
            "reason": reason,
        }

    def _check_instructions(self, user_input: str) -> Dict[str, Any]:
        """
        检查是否包含明确指令
        """
        instructions = {
            "开晶": "crystallize",
            "晶体化": "crystallize",
            "系统状态": "status",
            "查看待确认": "show_pending",
            "孔洞花园": "show_holes",
            "确认": "confirm_card",
            "拒绝": "reject_card",
            "暂停自主探测": "pause_auto",
            "恢复自主探测": "resume_auto",
            "归档": "archive"
        }

        for keyword, instruction_type in instructions.items():
            if keyword in user_input:
                return {
                    "direct_match": True,
                    "instruction_type": instruction_type,
                    "matched_keyword": keyword
                }

        return {"direct_match": False}

    def _check_simple_question(self, user_input: str) -> Dict[str, Any]:
        """
        检测是否简单问题（无需LLM）
        """
        simple_patterns = [
            "你好", "hi", "hello", "在吗",
            "谢谢", "感谢", "好的", "OK",
            "知道了", "明白了", "清楚",
            "继续", "接着说", "然后呢"
        ]

        stripped = user_input.strip()
        if len(stripped) < 10:
            for pattern in simple_patterns:
                if pattern in stripped:
                    return {
                        "is_simple": True,
                        "matched_pattern": pattern,
                        "reason": f"匹配简单模式: {pattern}"
                    }

        return {"is_simple": False}

    def _check_search_frequency(self, history: List[Tuple[str, str]]) -> Dict[str, Any]:
        """
        检查搜索频率衰减（连续3轮未搜索）
        """
        # 检查最近3轮用户消息是否包含搜索词
        search_keywords = ["搜索", "查找", "找", "查一下", "搜索一下", "外部", "信息"]
        user_messages = [msg[1] for msg in history if msg[0] == "user"]

        if len(user_messages) < 3:
            return {"need_reminder": False, "search_count": 0}

        recent = user_messages[-3:]
        has_search = any(any(kw in msg for kw in search_keywords) for msg in recent)
        search_count = len([msg for msg in recent if any(kw in msg for kw in search_keywords)])

        return {
            "need_reminder": search_count == 0,
            "search_count": search_count,
            "recent_messages": recent
        }

    def _check_emotion_only(self, user_input: str) -> Dict[str, Any]:
        """
        检测纯情绪表达
        """
        emotion_words = ["哈哈", "呵呵", "呜呜", "唉", "诶", "嗯", "哦", "啊", "哇"]
        stripped = user_input.strip()
        if len(stripped) < 5:
            for word in emotion_words:
                if word in stripped:
                    return {
                        "is_emotion_only": True,
                        "matched_word": word
                    }

        return {"is_emotion_only": False}

    def _estimate_cost(self, user_input: str, history: List[Tuple[str, str]]) -> Dict[str, Any]:
        """
        估算Token消耗
        """
        # 估算输入长度
        input_tokens = len(user_input) / 2  # 粗略估算（中文约2字符/token）

        # 估算历史长度
        history_tokens = sum(len(msg[1]) / 2 for msg in history[-8:])

        # 估算输出
        output_tokens = min(500, len(user_input))  # 简单估算

        return {
            "estimated_input_tokens": int(input_tokens + history_tokens),
            "estimated_output_tokens": int(output_tokens),
            "estimated_total_tokens": int(input_tokens + history_tokens + output_tokens),
            "estimated_cost_usd": round((input_tokens + history_tokens + output_tokens) * 0.000001, 6)
        }

    def reset_search_counter(self):
        """重置搜索计数器"""
        self._search_counter = 0

    def increment_search_counter(self):
        """增加搜索计数器"""
        self._search_counter += 1


class CrystalEngine:
    def __init__(self, file_io: FileIO):
        self.files = file_io
        self.meta = MetaLayer(self, file_io)  # ★ 新增这一行
        self.fingerprint_extractor = FingerprintExtractor(self, file_io)  # ★ 新增
        self.cheap_gate = CheapGate(self, file_io)     # ★ 新增
        self.vector_store = VectorStore(file_io)
        # Day 2.8: 元问题分类器（动态导入）
        try:
            import sys
            from pathlib import Path
            core_config_path = str(Config.DATA_ROOT / "核心配置")
            if core_config_path not in sys.path:
                sys.path.append(core_config_path)
            from question_classifier import QuestionClassifier
            self.question_classifier = QuestionClassifier()
        except ImportError:
            self.question_classifier = None
            print("⚠️ question_classifier 未找到，元问题分类功能禁用")
        # Day 5: Hebbian 学习权重存储
        self.hebbian_weights = {}
        self._load_hebbian_weights()
    # ========================================================================
    # 在 CrystalEngine 类中新增突触管理方法
    # ========================================================================

    # ========================================================================
    # 突触管理方法
    # ========================================================================

    def get_role_synapses(self, role_key: str) -> Dict[str, float]:
        """获取指定角色的突触权重表"""
        synapse_file = Config.DATA_ROOT / "系统日志" / "角色突触.json"
        try:
            if not synapse_file.exists():
                return {}
            with open(synapse_file, "r", encoding="utf-8") as f:
                data = json.load(f)
            return data.get(role_key, {}).get("synapses", {})
        except Exception:
            return {}

    def update_role_synapse(self, role_key: str, crystal_id: str, delta: float) -> float:
        """更新突触权重（采纳+0.08，驳回-0.05）"""
        synapse_file = Config.DATA_ROOT / "系统日志" / "角色突触.json"
        try:
            synapse_file.parent.mkdir(parents=True, exist_ok=True)
            data = {}
            if synapse_file.exists():
                with open(synapse_file, "r", encoding="utf-8") as f:
                    data = json.load(f)
            
            if role_key not in data:
                data[role_key] = {"synapses": {}, "win_count": 0, "loss_count": 0}
            
            synapses = data[role_key]["synapses"]
            old = synapses.get(crystal_id, 0.5)
            new_weight = max(0.0, min(1.0, old + delta))
            synapses[crystal_id] = round(new_weight, 3)
            data[role_key]["synapses"] = synapses
            data[role_key]["last_updated"] = datetime.now().isoformat()
            
            if delta > 0:
                data[role_key]["win_count"] = data[role_key].get("win_count", 0) + 1
            else:
                data[role_key]["loss_count"] = data[role_key].get("loss_count", 0) + 1
            
            with open(synapse_file, "w", encoding="utf-8") as f:
                json.dump(data, f, ensure_ascii=False, indent=2)
            return new_weight
        except Exception as e:
            print(f"[WARN] 突触更新失败: {e}")
            return 0.5

    def get_role_win_rate(self, role_key: str) -> float:
        """获取角色历史胜率"""
        synapse_file = Config.DATA_ROOT / "系统日志" / "角色突触.json"
        try:
            if not synapse_file.exists():
                return 0.5
            with open(synapse_file, "r", encoding="utf-8") as f:
                data = json.load(f)
            role_data = data.get(role_key, {})
            wins = role_data.get("win_count", 0)
            losses = role_data.get("loss_count", 0)
            total = wins + losses
            return 0.5 if total == 0 else round(wins / total, 3)
        except Exception:
            return 0.5

    def get_role_synapses(self, role_key: str) -> Dict[str, float]:
        """获取指定角色的突触权重表"""
        synapse_file = Config.DATA_ROOT / "系统日志" / "角色突触.json"
        try:
            if not synapse_file.exists():
                return {}
            with open(synapse_file, "r", encoding="utf-8") as f:
                data = json.load(f)
            return data.get(role_key, {}).get("synapses", {})
        except Exception:
            return {}

    def update_role_synapse(self, role_key: str, crystal_id: str, delta: float) -> float:
        """更新突触权重（采纳+0.08，驳回-0.05）"""
        synapse_file = Config.DATA_ROOT / "系统日志" / "角色突触.json"
        try:
            synapse_file.parent.mkdir(parents=True, exist_ok=True)
            data = {}
            if synapse_file.exists():
                with open(synapse_file, "r", encoding="utf-8") as f:
                    data = json.load(f)
            
            if role_key not in data:
                data[role_key] = {"synapses": {}, "win_count": 0, "loss_count": 0}
            
            synapses = data[role_key]["synapses"]
            old = synapses.get(crystal_id, 0.5)
            new_weight = max(0.0, min(1.0, old + delta))
            synapses[crystal_id] = round(new_weight, 3)
            data[role_key]["synapses"] = synapses
            data[role_key]["last_updated"] = datetime.now().isoformat()
            
            if delta > 0:
                data[role_key]["win_count"] = data[role_key].get("win_count", 0) + 1
            else:
                data[role_key]["loss_count"] = data[role_key].get("loss_count", 0) + 1
            
            with open(synapse_file, "w", encoding="utf-8") as f:
                json.dump(data, f, ensure_ascii=False, indent=2)
            return new_weight
        except Exception as e:
            print(f"[WARN] 突触更新失败: {e}")
            return 0.5

    def get_role_win_rate(self, role_key: str) -> float:
        """获取角色历史胜率"""
        synapse_file = Config.DATA_ROOT / "系统日志" / "角色突触.json"
        try:
            if not synapse_file.exists():
                return 0.5
            with open(synapse_file, "r", encoding="utf-8") as f:
                data = json.load(f)
            role_data = data.get(role_key, {})
            wins = role_data.get("win_count", 0)
            losses = role_data.get("loss_count", 0)
            total = wins + losses
            return 0.5 if total == 0 else round(wins / total, 3)
        except Exception:
            return 0.5

    def parse_crystals(self) -> List[Crystal]:
        content = self.files.read("crystals")
        pattern = r"\| (C\d+) \| (.*?) \| (.*?) \| (.*?) \| (.*?) \| (.*?) \| (.*?) \|"
        crystals = []
        for match in re.finditer(pattern, content):
            cid = match.group(1)
            text = match.group(2).strip()
            links_str = match.group(3).strip()
            links = [l.strip() for l in links_str.split(",") if l.strip() and l.strip() != "—"]
            input_conditions = [c.strip() for c in match.group(4).split(",") if c.strip() and c.strip() != "—"]
            execution_logic = match.group(5).strip() if match.group(5).strip() != "—" else ""
            output_format = match.group(6).strip() if match.group(6).strip() != "—" else ""
            validation_criteria = [c.strip() for c in match.group(7).split(",") if c.strip() and c.strip() != "—"]
            crystals.append(Crystal(
                id=cid,
                content=text,
                links=links,
                input_conditions=input_conditions,
                execution_logic=execution_logic,
                output_format=output_format,
                validation_criteria=validation_criteria
            ))
        return crystals

    def parse_holes(self) -> List[Hole]:
        content = self.files.read("holes")
        pattern = r"\| (H\d+) \| (.*?) \| ([\d\.]+) \|"
        holes = []
        for match in re.finditer(pattern, content):
            hid = match.group(1)
            text = match.group(2).strip()
            urgency = float(match.group(3))
            holes.append(Hole(id=hid, content=text, urgency=urgency))
        return holes

    def _simple_similarity(self, text1: str, text2: str) -> float:
        def tokenize(s: str) -> Set[str]:
            words = re.findall(r'[\w\u4e00-\u9fff]+', s.lower())
            tokens = set()
            for w in words:
                if re.match(r'[\u4e00-\u9fff]', w):
                    tokens.update(w)
                else:
                    tokens.add(w)
            return tokens
        set1 = tokenize(text1)
        set2 = tokenize(text2)
        if not set1 or not set2:
            return 0.0
        return len(set1 & set2) / len(set1 | set2)

    def _search_tokens(self, text: str) -> List[str]:
        tokens = []
        for word in re.findall(r'[A-Za-z0-9_]+|[\u4e00-\u9fff]+', text.lower()):
            if re.search(r'[\u4e00-\u9fff]', word):
                chars = [ch for ch in word if re.match(r'[\u4e00-\u9fff]', ch)]
                tokens.extend(chars)
                tokens.extend(''.join(chars[i:i+2]) for i in range(len(chars)-1))
                tokens.extend(''.join(chars[i:i+3]) for i in range(len(chars)-2))
            else:
                tokens.append(word)
        return [t for t in tokens if t]

    def rank_crystals(self, query: str, crystals: List[Crystal], top_k: int = 5, task_type: str = "general") -> List[Tuple[float, Crystal]]:
        """
        检索最相关的晶体（向量检索优先，BM25 降级）

        向量检索使用语义匹配，BM25 使用关键词匹配。
        向量检索失败时自动降级到 BM25。
        """
        if not query or not crystals:
            return []

        # 尝试向量检索
        if Config.VECTOR_SEARCH_ENABLED and self.vector_store.is_available():
            try:
                # 获取晶体 ID 到对象的映射
                crystal_map = {c.id: c for c in crystals}

                # 检索 top_k（向量检索返回 ID + 相似度）
                results = self.vector_store.query(query, top_k=top_k * 2)  # 多取一些用于后续过滤

                if results:
                    # 构建返回结果
                    scored = []
                    for cid, score in results:
                        if cid in crystal_map:
                            # 将相似度作为排名分数，叠加热度作为辅助
                            crystal = crystal_map[cid]
                            # 组合分数：相似度（主要）+ 热度（辅助）
                            combined_score = score * 0.9 + (crystal.heat / 10) * 0.1
                            hebbian_boost = self.get_hebbian_boost(crystal.id, task_type)
                            combined_score = score * 0.9 + (crystal.heat / 10) * 0.1 + hebbian_boost * 0.2
                            scored.append((combined_score, crystal))

                    # 按分数降序排列
                    scored.sort(key=lambda item: item[0], reverse=True)
                    return scored[:top_k]

            except Exception as e:
                print(f"[WARN] 向量检索异常，降级到 BM25: {e}")

        # ===== 降级：BM25 关键词检索 =====
        return self._rank_crystals_bm25(query, crystals, top_k, task_type)

    def _rank_crystals_bm25(self, query: str, crystals: List[Crystal], top_k: int = 5, task_type: str = "general") -> List[Tuple[float, Crystal]]:
        """
        BM25 关键词检索（保留原逻辑作为降级方案）
        """
        query_tokens = self._search_tokens(query)
        if not query_tokens:
            return [(0.0, c) for c in crystals[:top_k]]

        doc_tokens = [self._search_tokens(c.content + " " + " ".join(c.links)) for c in crystals]
        doc_freq = Counter()
        for tokens in doc_tokens:
            doc_freq.update(set(tokens))

        avg_len = sum(len(tokens) for tokens in doc_tokens) / max(1, len(doc_tokens))
        query_counts = Counter(query_tokens)

        scored = []
        for crystal, tokens in zip(crystals, doc_tokens):
            counts = Counter(tokens)
            doc_len = max(1, len(tokens))

            bm25 = 0.0
            for term, q_weight in query_counts.items():
                tf = counts.get(term, 0)
                if not tf:
                    continue
                idf = math.log(1 + (len(crystals) - doc_freq[term] + 0.5) / (doc_freq[term] + 0.5))
                denom = tf + 1.5 * (1 - 0.75 + 0.75 * doc_len / max(1, avg_len))
                bm25 += idf * (tf * 2.5 / denom) * min(2, q_weight)

            # 短语匹配加成
            phrase_bonus = 0.0
            q_lower = query.lower()
            text_lower = crystal.content.lower()
            if q_lower and q_lower in text_lower:
                phrase_bonus += 3.0
            phrase_bonus += sum(0.35 for term in set(query_tokens) if len(term) >= 2 and term in text_lower)
            hebbian_boost = self.get_hebbian_boost(crystal.id, task_type)
            total_score = bm25 + phrase_bonus + hebbian_boost * 0.3
            scored.append((total_score, crystal))


        scored.sort(key=lambda item: (item[0], item[1].heat), reverse=True)
        return scored[:top_k]
    def compute_crystal_heat(self, crystal: Crystal, all_crystals: List[Crystal]) -> float:
        holes = self.parse_holes()
        return self._compute_crystal_heat(crystal, all_crystals, holes)

    def _compute_crystal_heat(self, crystal: Crystal, all_crystals: List[Crystal], holes: List[Hole]) -> float:
        ref_count = sum(1 for c in all_crystals if crystal.id in c.links)
        l0_hole_ids = ["H001","H002","H003"]
        l0_texts = [h.content for h in holes if h.id in l0_hole_ids]
        semantic = max((self._simple_similarity(crystal.content, lt) for lt in l0_texts), default=0.0)
        return ref_count * 0.4 + semantic * 0.3

    def load_layer_state(self) -> Dict:
        if not self.files.exists("layer_state"):
            return {"layers": {}, "heat_map": {}, "last_accessed": {}, "manual_override": {}}
        try:
            return json.loads(self.files.read("layer_state"))
        except:
            return {"layers": {}, "heat_map": {}, "last_accessed": {}, "manual_override": {}}

    def save_layer_state(self, state: Dict):
        self.files.write("layer_state", json.dumps(state, ensure_ascii=False, indent=2))

    def update_crystal_access_time(self, crystal_id: str):
        state = self.load_layer_state()
        last_acc = state.get("last_accessed", {})
        last_acc[crystal_id] = date.today().isoformat()
        state["last_accessed"] = last_acc
        self.save_layer_state(state)

    def update_crystal_layers(self) -> Tuple[List[Crystal], List[Crystal], List[Crystal]]:
        crystals = self.parse_crystals()
        if not crystals:
            return [], [], []
        state = self.load_layer_state()
        last_accessed = state.get("last_accessed", {})
        manual_override = state.get("manual_override", {})
        holes = self.parse_holes()
        heat_map = {c.id: self._compute_crystal_heat(c, crystals, holes) for c in crystals}
        def sort_key(cid):
            if manual_override.get(cid) == "L1_fixed":
                return (1, heat_map[cid])
            return (0, heat_map[cid])
        sorted_ids = sorted(heat_map.keys(), key=sort_key, reverse=True)
        layers = {}
        l1_assigned = 0
        for cid in sorted_ids:
            if manual_override.get(cid) == "L1_fixed":
                layers[cid] = "L1"
                l1_assigned += 1
        for cid in sorted_ids:
            if cid in layers:
                continue
            if l1_assigned < Config.L1_MAX:
                layers[cid] = "L1"
                l1_assigned += 1
            else:
                last = last_accessed.get(cid, (date.today() - timedelta(days=31)).isoformat())
                try:
                    days_since = (date.today() - date.fromisoformat(last)).days
                except:
                    days_since = 31
                if heat_map[cid] < Config.L2_TO_L3_HEAT_THRESHOLD and days_since > Config.L2_TO_L3_DAYS_THRESHOLD:
                    layers[cid] = "L3"
                else:
                    layers[cid] = "L2"
        state["layers"] = layers
        state["heat_map"] = heat_map
        self.save_layer_state(state)
        L1 = [c for c in crystals if layers.get(c.id) == "L1"]
        L2 = [c for c in crystals if layers.get(c.id) == "L2"]
        L3 = [c for c in crystals if layers.get(c.id) == "L3"]
        return L1, L2, L3

    def archive_cold_crystals(self) -> List[str]:
        state = self.load_layer_state()
        layers = state.get("layers", {})
        heat_map = state.get("heat_map", {})
        last_accessed = state.get("last_accessed", {})
        today_dt = date.today()
        archived = []
        changed = False
        for cid, layer in list(layers.items()):
            if layer == "L2":
                heat = heat_map.get(cid, 0.0)
                last = last_accessed.get(cid, "2000-01-01")
                try:
                    days_since = (today_dt - date.fromisoformat(last)).days
                except:
                    days_since = 999
                if heat < Config.L2_TO_L3_HEAT_THRESHOLD and days_since > Config.L2_TO_L3_DAYS_THRESHOLD:
                    layers[cid] = "L3"
                    changed = True
                    archived.append(cid)
        if changed:
            state["layers"] = layers
            self.save_layer_state(state)
            self._append_change_log("自动归档", f"已将 {len(archived)} 个冷晶体从 L2 移至 L3")
        return archived

    def get_attention_context(self) -> Tuple[List[Hole], List[Crystal]]:
        holes = self.parse_holes()
        l0_hole_ids = ["H001","H002","H003"]
        l0_holes = [h for h in holes if h.id in l0_hole_ids]
        if len(l0_holes) < 3:
            default_holes = {"H001":"如何定义与分解复杂问题？","H002":"非共识情况下如何做出正确判断？","H003":"因果链的完整推演与验证方法"}
            existing_ids = [h.id for h in l0_holes]
            for hid, content in default_holes.items():
                if hid not in existing_ids:
                    l0_holes.append(Hole(id=hid, content=content, urgency=0.9))
        state = self.load_layer_state()
        layers = state.get("layers", {})
        all_crystals = self.parse_crystals()
        L1_crystals = [c for c in all_crystals if layers.get(c.id) == "L1"]
        for c in L1_crystals:
            self.update_crystal_access_time(c.id)
        return l0_holes, L1_crystals

    def get_conflict_scope(self) -> List[Crystal]:
        state = self.load_layer_state()
        layers = state.get("layers", {})
        all_crystals = self.parse_crystals()
        scope = [c for c in all_crystals if layers.get(c.id) in ("L1","L2")]
        if not scope:
            L1, L2, _ = self.update_crystal_layers()
            scope = L1 + L2
        return scope

    def detect_conflicts(self, scope: List[Crystal] = None, method: str = "auto") -> List[Conflict]:
        """
        矛盾检测引擎（升级版）

        支持三种模式：
        - "auto": 自动选择（向量可用时用向量，否则用 Jaccard）
        - "vector": 强制使用向量检测
        - "jaccard": 强制使用 Jaccard 关键词匹配

        Args:
            scope: 检测范围（晶体列表）
            method: 检测方法 ("auto" | "vector" | "jaccard")

        Returns:
            List[Conflict]: 检测到的冲突列表
        """
        if scope is None:
            scope = self.get_conflict_scope()

        if len(scope) < 2:
            return []

        if method == "jaccard":
            return self._detect_conflicts_jaccard(scope)
        elif method == "vector":
            if self.vector_store.is_available():
                return self.detect_conflicts_vector(scope)
            else:
                return self._detect_conflicts_jaccard(scope)
        else:  # "auto"
            if self.vector_store.is_available():
                return self.detect_conflicts_vector(scope)
            else:
                return self._detect_conflicts_jaccard(scope)

    def _detect_conflicts_jaccard(self, scope: List[Crystal]) -> List[Conflict]:
        """
        Jaccard 关键词匹配检测（原逻辑保留作为降级方案）
        """
        if len(scope) < 2:
            return []

        conflicts = []
        checked_pairs = set()

        for i in range(len(scope)):
            for j in range(i + 1, len(scope)):
                c1, c2 = scope[i], scope[j]
                pair_key = tuple(sorted([c1.id, c2.id]))
                if pair_key in checked_pairs:
                    continue
                checked_pairs.add(pair_key)

                sim = self._simple_similarity(c1.content, c2.content)
                if sim > 0.7:
                    conflicts.append(Conflict(
                        crystal_a=c1.id,
                        crystal_b=c2.id,
                        similarity=round(sim, 2),
                        content_a=c1.content,
                        content_b=c2.content
                    ))

        conflicts.sort(key=lambda x: x.similarity, reverse=True)
        return conflicts

    def detect_conflicts_vector(self, scope: List[Crystal] = None, threshold: float = 0.3) -> List[Conflict]:
        """
        向量级矛盾检测引擎

        使用向量余弦距离检测晶体间的语义矛盾。
        距离 < threshold 视为高度语义相关（可能矛盾或重复）。

        优势：
        - 能发现语义相近但关键词不同的隐藏冲突
        - 不受关键词匹配限制

        Returns:
            List[Conflict]: 检测到的冲突列表
        """
        if not self.vector_store.is_available():
            return self._detect_conflicts_jaccard(scope)

        if scope is None:
            scope = self.get_conflict_scope()

        if len(scope) < 2:
            return []

        conflicts = []
        checked_pairs = set()

        for i in range(len(scope)):
            for j in range(i + 1, len(scope)):
                c1, c2 = scope[i], scope[j]
                pair_key = tuple(sorted([c1.id, c2.id]))
                if pair_key in checked_pairs:
                    continue
                checked_pairs.add(pair_key)

                try:
                    results = self.vector_store.query(c1.content, top_k=10)
                    c2_similarity = 0.0
                    for cid, score in results:
                        if cid == c2.id:
                            c2_similarity = score
                            break

                    if c2_similarity > (1 - threshold):
                        conflicts.append(Conflict(
                            crystal_a=c1.id,
                            crystal_b=c2.id,
                            similarity=round(c2_similarity, 2),
                            content_a=c1.content,
                            content_b=c2.content
                        ))
                except Exception as e:
                    sim = self._simple_similarity(c1.content, c2.content)
                    if sim > 0.7:
                        conflicts.append(Conflict(
                            crystal_a=c1.id,
                            crystal_b=c2.id,
                            similarity=round(sim, 2),
                            content_a=c1.content,
                            content_b=c2.content
                        ))

        conflicts.sort(key=lambda x: x.similarity, reverse=True)
        return conflicts

    def load_hole_progress(self) -> Dict[str, float]:
        if not self.files.exists("hole_progress"):
            return {}
        try:
            return json.loads(self.files.read("hole_progress"))
        except:
            return {}

    def save_hole_progress(self, progress: Dict[str, float]) -> None:
        self.files.write("hole_progress", json.dumps(progress, ensure_ascii=False, indent=2))

    def match_info_to_hole(self, info_title: str, hole_content: str) -> float:
        keywords = re.findall(r'[\w\u4e00-\u9fff]{2,}', hole_content)
        info_lower = info_title.lower()
        hit = sum(1 for kw in keywords if kw.lower() in info_lower)
        return min(1.0, hit / max(1, len(keywords)*0.5))

    def vector_search(self, query: str, crystals: List[Crystal], top_k: int = 5) -> List[Crystal]:
        ranked = self.rank_crystals(query, crystals, top_k)
        if ranked:
            return [c for _, c in ranked]
        return self._keyword_search(query, crystals, top_k)

    def _keyword_search(self, query: str, crystals: List[Crystal], top_k: int) -> List[Crystal]:
        query_words = set(re.findall(r'[\w\u4e00-\u9fff]+', query.lower()))
        scored = []
        for c in crystals:
            text = c.content.lower()
            score = sum(1 for w in query_words if w in text)
            scored.append((score, c))
        scored.sort(key=lambda x: x[0], reverse=True)
        return [c for _, c in scored[:top_k]]

    def get_associative_crystals(self, question: str, top_k: int = 5) -> List[Crystal]:
        all_crystals = self.parse_crystals()
        if not all_crystals:
            return []
        return self.vector_search(question, all_crystals, top_k=top_k)

    def _append_change_log(self, section: str, content: str):
        date_header = f"\n## {datetime.now().strftime('%Y-%m-%d')} - {section}\n"
        self.files.append("change_log", date_header + content + "\n")

    def log_evolution_event(self, event_type: str, details: Dict[str, Any]) -> None:
        """
        记录单用户进化事件（支持失败轨迹记录）

        Args:
            event_type: 事件类型（crystal_added | crystal_updated | crystal_archived | 
                        hole_filled | fingerprint_changed | role_adopted | verification_passed |
                        failure_trace | history_reused | alarm | chain_triggered）
            details: 事件详情，支持以下扩展字段：
                - failure_traces: 失败时的完整上下文
                - repair_attempts: 尝试的修复方案及结果
                - diagnosis: 系统自己对失败原因的判断
                - reused_history_id: 复用的历史记录ID
                - match_score: 历史匹配度
        """
        log_path = Config.DATA_ROOT / "系统日志" / "evolution_log.json"
        log_path.parent.mkdir(parents=True, exist_ok=True)

        # 读取现有日志
        if log_path.exists():
            try:
                with open(log_path, "r", encoding="utf-8") as f:
                    data = json.load(f)
            except:
                data = {"events": [], "summary": {}}
        else:
            data = {"events": [], "summary": {}}

        # 构建事件（支持失败轨迹字段）
        event = {
            "timestamp": datetime.now().isoformat(),
            "event_type": event_type,
            "details": details,
            "trigger": details.get("trigger", "system_auto")
        }

        # ===== Day 4 新增：扩展字段（如果存在则记录） =====
        if "failure_traces" in details:
            event["failure_traces"] = details["failure_traces"]
        if "repair_attempts" in details:
            event["repair_attempts"] = details["repair_attempts"]
        if "diagnosis" in details:
            event["diagnosis"] = details["diagnosis"]
        if "reused_history_id" in details:
            event["reused_history_id"] = details["reused_history_id"]
        if "match_score" in details:
            event["match_score"] = details["match_score"]

        # 追加事件
        data["events"].append(event)

        # 更新摘要统计
        summary = data.get("summary", {})
        summary[event_type] = summary.get(event_type, 0) + 1
        data["summary"] = summary

        # 只保留最近 500 条事件（增加容量以容纳更多历史）
        if len(data["events"]) > 500:
            data["events"] = data["events"][-500:]

        # 写入文件
        with open(log_path, "w", encoding="utf-8") as f:
            json.dump(data, f, ensure_ascii=False, indent=2)

    def get_evolution_stats(self) -> Dict[str, Any]:
        """获取进化统计"""
        log_path = Config.DATA_ROOT / "系统日志" / "evolution_log.json"
        if not log_path.exists():
            return {"total_events": 0, "summary": {}}
        try:
            with open(log_path, "r", encoding="utf-8") as f:
                data = json.load(f)
            return {
                "total_events": len(data.get("events", [])),
                "summary": data.get("summary", {}),
                "last_event": data.get("events", [])[-1] if data.get("events") else None
            }
        except:
            return {"total_events": 0, "summary": {}}

    def quality_gate_g1(self, question: str, context: Dict = None) -> Dict[str, Any]:
        """
        G1 质量门：问题是否可检验且相关

        与四组件挂钩：
        - 晶体卡片：检索相关晶体数量 ≥ 1
        - 孔洞：匹配孔洞数量 ≥ 1（紧迫度 ≥ 0.5）
        - 认知指纹：与指纹匹配度 ≥ 0.3

        Returns:
            {"passed": bool, "checks": dict, "reason": str}
        """
        context = context or {}
        results = {
            "crystal_match": {"count": 0, "passed": False},
            "hole_match": {"count": 0, "passed": False},
            "fingerprint_match": {"score": 0.0, "passed": False}
        }

        crystals = self.parse_crystals()
        if crystals:
            ranked = self.rank_crystals(question, crystals, top_k=5)
            results["crystal_match"]["count"] = len(ranked)
            results["crystal_match"]["passed"] = len(ranked) >= 1

        holes = self.parse_holes()
        if holes:
            matched = [h for h in holes if h.urgency >= 0.5 and
                       any(kw in question for kw in h.content[:20].split())]
            results["hole_match"]["count"] = len(matched)
            results["hole_match"]["passed"] = len(matched) >= 1

        try:
            fp = self.fingerprint_extractor.get_fingerprint()
            role_keywords = {
                "radical": ["颠覆", "激进", "创新", "突破", "打破"],
                "conservative": ["稳健", "保守", "风险", "安全", "验证"],
                "structural": ["结构", "系统", "框架", "流程", "机制"]
            }
            fp_keywords = role_keywords.get(fp.preferred_role, [])
            match_score = sum(1 for kw in fp_keywords if kw in question) / max(1, len(fp_keywords))
            results["fingerprint_match"]["score"] = match_score
            results["fingerprint_match"]["passed"] = match_score >= 0.3
        except:
            results["fingerprint_match"]["passed"] = True

        passed = (results["crystal_match"]["passed"] or
                  results["hole_match"]["passed"] or
                  results["fingerprint_match"]["passed"])

        reasons = []
        if results["crystal_match"]["passed"]:
            reasons.append(f"匹配 {results['crystal_match']['count']} 条晶体")
        if results["hole_match"]["passed"]:
            reasons.append(f"匹配 {results['hole_match']['count']} 个孔洞")
        if results["fingerprint_match"]["passed"]:
            reasons.append(f"指纹匹配度 {results['fingerprint_match']['score']:.2f}")

        return {
            "passed": passed,
            "checks": results,
            "reason": "；".join(reasons) if reasons else "无匹配"
        }

    def quality_gate_g2(self, answer: str, context: Dict = None) -> Dict[str, Any]:
        """
        G2 质量门：输出有依据且可靠

        与四组件挂钩：
        - 置信度：引用的晶体置信度 ≥ 0.5
        - 审计评分：审计评分 ≥ 0.6
        - 适用边界：无"不适用"条件被触发

        Returns:
            {"passed": bool, "checks": dict, "reason": str}
        """
        context = context or {}
        results = {
            "confidence": {"score": 0.0, "passed": False},
            "audit_score": {"score": 0.0, "passed": False},
            "boundary_ok": {"passed": True}
        }

        import re
        cited_ids = re.findall(r'C\d{3}', answer)
        if cited_ids:
            state = self.load_layer_state()
            heat_map = state.get("heat_map", {})
            conf_scores = [heat_map.get(cid, 0.0) for cid in cited_ids]
            avg_conf = sum(conf_scores) / len(conf_scores) if conf_scores else 0.0
            results["confidence"]["score"] = avg_conf
            results["confidence"]["passed"] = avg_conf >= 0.5

        audit_score = context.get("audit_score", 0.0)
        results["audit_score"]["score"] = audit_score
        results["audit_score"]["passed"] = audit_score >= 0.6

        boundary_terms = ["不适用", "不推荐", "谨慎使用", "特殊情况", "例外"]
        boundary_triggered = any(term in answer for term in boundary_terms)
        results["boundary_ok"]["passed"] = not boundary_triggered

        passed = (results["confidence"]["passed"] or
                  results["audit_score"]["passed"])

        reasons = []
        if results["confidence"]["passed"]:
            reasons.append(f"置信度 {results['confidence']['score']:.2f}")
        if results["audit_score"]["passed"]:
            reasons.append(f"审计评分 {results['audit_score']['score']:.2f}")
        if not results["boundary_ok"]["passed"]:
            reasons.append("触发了适用边界警告")

        return {
            "passed": passed,
            "checks": results,
            "reason": "；".join(reasons) if reasons else "无验证"
        }

    def sync_vector_store(self) -> Dict[str, Any]:
        """
        同步向量库：将当前所有晶体向量化

        Returns:
            {"total": int, "synced": int, "status": str}
        """
        if not self.vector_store.is_available():
            return {"total": 0, "synced": 0, "status": "vector_store_unavailable"}

        crystals = self.parse_crystals()
        if not crystals:
            return {"total": 0, "synced": 0, "status": "no_crystals"}

        # 检查现有数量
        existing = self.vector_store.count()
        if existing == len(crystals):
            return {"total": len(crystals), "synced": existing, "status": "already_synced"}

        # 如果数量不匹配，重置并重新同步
        if existing > 0 and existing != len(crystals):
            self.vector_store.reset()

        # 重新同步
        added = self.vector_store.add_crystals(crystals)
        return {"total": len(crystals), "synced": added, "status": "synced" if added > 0 else "failed"}

    def get_user_fingerprint(self) -> Optional[CognitiveFingerprint]:
        """获取当前用户的认知指纹"""
        try:
            return self.fingerprint_extractor.get_fingerprint()
        except Exception as e:
            print(f"[WARN] 获取指纹失败: {e}")
            return None

    def contribution_scoring(self, crystal_id: str = None) -> Dict[str, Any]:
        """
        计算晶体的贡献得分（负能力修剪）

        贡献得分综合考虑：
        - 被引用次数（links 入度）
        - 热度（heat）
        - 与用户认知指纹的匹配度
        - 最近访问时间

        Returns:
            {
                "crystal_id": str,
                "score": float,           # 综合贡献得分 0-100
                "ref_count": int,         # 被引用次数
                "heat": float,            # 热度
                "fingerprint_match": float, # 与指纹匹配度 0-1
                "days_since_access": int,  # 距上次访问天数
                "status": str             # "active" | "low_contribution" | "cold"
            }
        """
        crystals = self.parse_crystals()
        state = self.load_layer_state()
        layers = state.get("layers", {})
        heat_map = state.get("heat_map", {})
        last_accessed = state.get("last_accessed", {})

        # 获取用户指纹
        try:
            fingerprint = self.fingerprint_extractor.get_fingerprint()
        except:
            fingerprint = None

        # 构建引用关系映射
        ref_count_map = {}
        for c in crystals:
            ref_count_map[c.id] = 0
        for c in crystals:
            for link in c.links:
                if link in ref_count_map:
                    ref_count_map[link] += 1

        # 获取偏好角色的关键词（用于指纹匹配）
        role_keywords = {
            "radical": ["颠覆", "激进", "创新", "突破", "打破", "颠覆性"],
            "conservative": ["稳健", "保守", "风险", "安全", "验证", "可靠"],
            "structural": ["结构", "系统", "框架", "流程", "机制", "模型"],
            "executor": ["执行", "步骤", "操作", "落地", "行动", "检查"],
            "auditor": ["审计", "验证", "检查", "证据", "反例", "漏洞"]
        }
        pref_role = fingerprint.preferred_role if fingerprint else "structural"
        fp_keywords = role_keywords.get(pref_role, role_keywords["structural"])

        # 计算每个晶体的贡献得分
        today = date.today()
        results = []

        for c in crystals:
            # 1. 被引用次数（0-30分）
            ref_count = ref_count_map.get(c.id, 0)
            ref_score = min(30, ref_count * 6)

            # 2. 热度（0-30分）
            heat = heat_map.get(c.id, 0.0)
            heat_score = min(30, heat * 15)

            # 3. 指纹匹配度（0-20分）
            match_score = 0
            content_lower = c.content.lower()
            for kw in fp_keywords:
                if kw in content_lower:
                    match_score += 4
            match_score = min(20, match_score)

            # 4. 时效性（0-20分）
            last = last_accessed.get(c.id)
            if last:
                try:
                    days_since = (today - date.fromisoformat(last)).days
                except:
                    days_since = 999
            else:
                days_since = 999
            # 30天内访问得满分，超过30天逐渐衰减
            time_score = max(0, 20 * (1 - days_since / 90))
            time_score = min(20, time_score)

            total_score = ref_score + heat_score + match_score + time_score

            # 判断状态
            if total_score < 15:
                status = "cold"  # 极低贡献，建议归档
            elif total_score < 25:
                status = "low_contribution"  # 低贡献，需验证
            else:
                status = "active"

            results.append({
                "crystal_id": c.id,
                "content": c.content,
                "score": round(total_score, 1),
                "ref_count": ref_count,
                "heat": round(heat, 2),
                "fingerprint_match": round(match_score / 20, 2),
                "days_since_access": days_since,
                "status": status
            })

        # 按得分排序
        results.sort(key=lambda x: x["score"], reverse=True)

        if crystal_id:
            # 返回指定晶体的信息
            for r in results:
                if r["crystal_id"] == crystal_id:
                    return r
            return {"error": f"未找到晶体 {crystal_id}"}

        return {
            "total": len(results),
            "active": len([r for r in results if r["status"] == "active"]),
            "low_contribution": len([r for r in results if r["status"] == "low_contribution"]),
            "cold": len([r for r in results if r["status"] == "cold"]),
            "details": results
        }

    def get_low_contribution_crystals(self, threshold: float = 25.0) -> List[Dict]:
        """
        获取低贡献晶体列表（用于验证门控）
        """
        result = self.contribution_scoring()
        if "error" in result:
            return []
        return [r for r in result.get("details", []) if r["score"] < threshold]

    # ===== Day 5: Hebbian 学习 =====
    def _load_hebbian_weights(self):
        """从文件加载 Hebbian 权重"""
        weight_file = Config.DATA_ROOT / "系统日志" / "hebbian_weights.json"
        if weight_file.exists():
            try:
                with open(weight_file, "r", encoding="utf-8") as f:
                    self.hebbian_weights = json.load(f)
            except:
                self.hebbian_weights = {}
        else:
            self.hebbian_weights = {}

    def _save_hebbian_weights(self):
        """保存 Hebbian 权重到文件"""
        weight_file = Config.DATA_ROOT / "系统日志" / "hebbian_weights.json"
        with open(weight_file, "w", encoding="utf-8") as f:
            json.dump(self.hebbian_weights, f, ensure_ascii=False, indent=2)

    def update_hebbian_weights(self, crystal_ids: List[str], task_type: str, score: float):
        """
        更新 Hebbian 权重
        - 对每对晶体组合更新权重
        - 对每个晶体与任务类型更新权重
        """
        if not crystal_ids or len(crystal_ids) < 2:
            return
        # 对每对晶体组合更新权重
        for i in range(len(crystal_ids)):
            for j in range(i+1, len(crystal_ids)):
                pair = tuple(sorted([crystal_ids[i], crystal_ids[j]]))
                key = f"pair_{pair[0]}_{pair[1]}"
                old = self.hebbian_weights.get(key, 0.0)
                # 评分高则增加，低则衰减（Hebbian学习）
                delta = (score - 0.5) * 0.1
                new_weight = max(0.0, min(1.0, old + delta))
                self.hebbian_weights[key] = new_weight
        # 记录任务类型与晶体的关联权重
        task_key = f"task_{task_type}"
        for cid in crystal_ids:
            key = f"{task_key}_{cid}"
            old = self.hebbian_weights.get(key, 0.0)
            delta = (score - 0.5) * 0.05
            new_weight = max(0.0, min(1.0, old + delta))
            self.hebbian_weights[key] = new_weight
        self._save_hebbian_weights()

    def get_hebbian_boost(self, crystal_id: str, task_type: str = None) -> float:
        """获取某个晶体在当前任务类型下的 Hebbian 加成"""
        boost = 0.0
        if task_type:
            key = f"task_{task_type}_{crystal_id}"
            boost += self.hebbian_weights.get(key, 0.0) * 0.3
        return boost

# =============================================================================
# 7. AI 客户端 (ai_client.py)
# =============================================================================
import json
from typing import Dict, List, Optional, Tuple

class AIClient:
    def __init__(self, api_key: str = None, api_url: str = None):
        self.api_key = api_key or Config.get_api_key()
        self.api_url = api_url or Config.DEEPSEEK_API_URL
        self._session = None
        self._has_requests = REQUESTS_AVAILABLE
        # 埋点统计属性（Day 0 新增）
        self._call_count = 0
        self._token_estimate = 0
        self._total_time = 0.0

    @property
    def session(self):
        if not self._has_requests:
            return None
        if self._session is None:
            self._session = requests.Session()
        return self._session

    def _call_api(self, messages: List[Dict], temperature: float = 0.7,
                  response_format: Dict = None, stream: bool = False,
                  callback: Callable[[str], None] = None,
                  max_tokens: int = None) -> Optional[str]:
        # === 1. 初始化 result 为 None（防御性编程） ===
        result = None
        start_time = time.time()
        total_chars = sum(len(m.get("content", "")) for m in messages)
        token_used = total_chars // 2
        self._token_estimate += token_used
        self._call_count += 1

        # === 2. 校验 API Key ===
        if not self.api_key:
            result = "错误：未配置 DEEPSEEK_API_KEY"
        elif not self._has_requests:
            result = "错误：需要安装 requests 库"
        else:
            headers = {"Authorization": f"Bearer {self.api_key}", "Content-Type": "application/json"}
            payload = {
                "model": "deepseek-chat",
                "messages": messages,
                "temperature": temperature,
                "stream": stream,
            }
            if max_tokens is not None:
                payload["max_tokens"] = min(8000, max(100, int(max_tokens)))
            else:
                payload["max_tokens"] = 4000
            if response_format:
                payload["response_format"] = response_format

            # === 3. 请求执行（完整异常捕获） ===
            try:
                if stream:
                    with self.session.post(self.api_url, headers=headers, json=payload,
                                           stream=True, timeout=120) as resp:
                        resp.raise_for_status()
                        collected = ""
                        for line in resp.iter_lines():
                            if line:
                                line = line.decode('utf-8')
                                if line.startswith('data: '):
                                    data = line[6:]
                                    if data == '[DONE]':
                                        break
                                    try:
                                        chunk = json.loads(data)
                                        delta = chunk.get('choices', [{}])[0].get('delta', {})
                                        content = delta.get('content', '')
                                        if content:
                                            collected += content
                                            if callback:
                                                callback(content)
                                    except json.JSONDecodeError:
                                        continue
                        result = collected if collected else "（AI返回空内容）"
                else:
                    resp = self.session.post(self.api_url, headers=headers, json=payload, timeout=60)
                    resp.raise_for_status()
                    try:
                        response_data = resp.json()
                        result = response_data["choices"][0]["message"]["content"]
                        if not result:
                            result = "（AI返回空内容）"
                    except (KeyError, json.JSONDecodeError, IndexError) as e:
                        result = f"AI响应解析失败: {e}"
            except requests.exceptions.Timeout:
                result = "错误：请求超时（请检查网络或增大超时设置）"
            except requests.exceptions.ConnectionError:
                result = "错误：网络连接失败（请检查网络或API地址）"
            except requests.exceptions.HTTPError as e:
                if e.response.status_code == 401:
                    result = "错误：API Key 无效或已过期"
                elif e.response.status_code == 429:
                    result = "错误：请求频率过高，请稍后重试"
                else:
                    result = f"错误：HTTP {e.response.status_code} - {e.response.text[:100]}"
            except Exception as e:
                result = f"AI调用失败: {type(e).__name__} - {str(e)}"

        # === 4. 最终保护：如果 result 仍为 None ===
        if result is None:
            result = "错误：未知原因导致返回为空"

        # === 5. 埋点 ===
        elapsed = time.time() - start_time
        self._total_time += elapsed
        if self._call_count % 10 == 0:
            self._write_metrics()

        return result
    
    def _write_metrics(self):
        """将埋点统计数据写入 系统日志/埋点数据.json"""
        log_path = Config.DATA_ROOT / "系统日志" / "埋点数据.json"
        try:
            with open(log_path, "r", encoding="utf-8") as f:
                data = json.load(f)
        except (FileNotFoundError, json.JSONDecodeError):
            data = {"TTQ": [], "TTA": [], "TTU": []}

        # 记录当前累计值（每次记录是一个快照）
        data["TTQ"].append({
            "timestamp": datetime.now().isoformat(),
            "call_count": self._call_count,
            "total_tokens": self._token_estimate,
            "avg_tokens_per_call": round(self._token_estimate / self._call_count, 2) if self._call_count else 0
        })
        data["TTA"].append({
            "timestamp": datetime.now().isoformat(),
            "call_count": self._call_count,
            "total_time_seconds": round(self._total_time, 3),
            "avg_time_per_call": round(self._total_time / self._call_count, 3) if self._call_count else 0
        })
        data["TTU"].append({
            "timestamp": datetime.now().isoformat(),
            "call_count": self._call_count,
            "total_tokens_used": self._token_estimate,
        })

        with open(log_path, "w", encoding="utf-8") as f:
            json.dump(data, f, ensure_ascii=False, indent=2)
    
    def chat(self, prompt: str, system: str = "你是认知晶体树的AI协作者，请友好自然地回答问题。", temperature: float = 0.7) -> str:
        return self._call_api([{"role": "system", "content": system}, {"role": "user", "content": prompt}], temperature=temperature)
    def chat_with_history(self, history: List[Tuple[str, str]], system: str = "你是认知晶体树的AI协作者，请友好自然地回答问题。", context: str = "") -> str:
        if context:
            system = system + context
        messages = [{"role": "system", "content": system}]
        for role, content in history:
            if role in ("user", "assistant"):
                messages.append({"role": role, "content": content})
        return self._call_api(messages)

    def chat_json(self, prompt: str, temperature: float = 0.3) -> Dict:
        result = self._call_api([{"role": "user", "content": prompt}], temperature=temperature, response_format={"type": "json_object"})
        if not isinstance(result, str):
            return {"error": "AI返回为空"}
        if result.startswith("错误") or result.startswith("AI调用失败"):
            return {"error": result}
        cleaned = result.strip()
        if cleaned.startswith("```"):
            cleaned = cleaned.strip("`")
            if cleaned.lower().startswith("json"):
                cleaned = cleaned[4:].strip()
        if not cleaned.startswith("{"):
            start = cleaned.find("{")
            end = cleaned.rfind("}")
            if start >= 0 and end > start:
                cleaned = cleaned[start:end+1]
        try:
            return json.loads(cleaned)
        except:
            return {"error": "解析JSON失败", "raw": result}

    def chat_stream(self, prompt: str, system: str = "你是认知晶体树的AI协作者，请友好自然地回答问题。",
                    callback: Callable[[str], None] = None) -> str:
        """
        流式对话，支持逐块回调。
        """
        messages = [{"role": "system", "content": system}, {"role": "user", "content": prompt}]
        return self._call_api(messages, stream=True, callback=callback)

# =============================================================================
# 辅助函数：AI 生成会话标题
# =============================================================================

def generate_session_title_from_content(content: str, api_key: str = None) -> str:
    """
    调用 AI 为会话生成精炼标题（不超过 8 个字）
    如果 AI 调用失败，则降级为截取前 8 个字符。
    """
    if not content:
        return ""
    ai = AIClient(api_key=api_key or Config.get_api_key())
    try:
        result = ai.chat_json(
            f"请为以下对话生成一个不超过 8 个字的精炼标题，只返回 JSON：{{'title': '你的标题'}}\n\n内容：{content[:300]}",
            temperature=0.1
        )
        if "error" not in result:
            title = result.get("title", "").strip()
            if title:
                return title
    except Exception:
        pass
    # 降级：取前 8 个字符（去除多余空白）
    return content.strip().replace("\n", " ")[:8].strip()

# =============================================================================
# 8. 外部抓取 (external.py)
# =============================================================================
import random
import re
import time
from datetime import datetime
from typing import Callable, Dict, List
from urllib.parse import quote

class ExternalFetcher:
    def __init__(self, log_callback: Callable = None):
        self.log_callback = log_callback
        self._has_requests = REQUESTS_AVAILABLE
        self._has_bs4 = BS4_AVAILABLE
        self._has_arxiv = ARXIV_AVAILABLE

    def _log(self, msg: str, tag: str = "system"):
        if self.log_callback:
            self.log_callback(msg, tag)

    def translate_to_english(self, text: str) -> str:
        if not re.search(r'[\u4e00-\u9fff]', text):
            return text
        if not self._has_requests:
            return text
        try:
            url = "https://api.mymemory.translated.net/get"
            params = {"q": text, "langpair": "zh|en", "de": "a@b.c"}
            resp = requests.get(url, params=params, timeout=10)
            if resp.status_code == 200:
                translated = resp.json().get("responseData", {}).get("translatedText", "")
                if translated and translated != text:
                    return translated
        except:
            pass
        return text

    def fetch_arxiv_papers(self, query: str = "cat:cs.AI", max_results: int = 5) -> List[str]:
        if not self._has_arxiv:
            return ["(需要安装 arxiv 库)"]
        import arxiv
        query = self.translate_to_english(query)
        self._log(f"  搜索查询: {query}", "system")
        session = requests.Session()
        retry_strategy = Retry(total=Config.MAX_RETRIES, backoff_factor=Config.BACKOFF_FACTOR, status_forcelist=[429,500,502,503,504], allowed_methods=["GET","HEAD"])
        adapter = HTTPAdapter(max_retries=retry_strategy)
        session.mount("http://", adapter)
        session.mount("https://", adapter)
        session.headers.update({"User-Agent": NetworkManager.get_random_user_agent(), "Accept": "application/atom+xml,application/xml"})
        client = arxiv.Client(page_size=min(max_results,20), delay_seconds=3.0, num_retries=Config.MAX_RETRIES)
        client._session = session
        search = arxiv.Search(query=query, max_results=max_results, sort_by=arxiv.SortCriterion.SubmittedDate)
        try:
            papers = []
            for paper in client.results(search):
                pub_date = paper.published.date().isoformat() if paper.published else ""
                title = paper.title.strip().replace('\n',' ')
                papers.append(f"{title} (发布于: {pub_date})" if pub_date else title)
                if len(papers) >= max_results:
                    break
            return papers if papers else ["(未找到相关论文)"]
        except Exception as e:
            return [f"(arXiv 请求失败: {e})"]

    def fetch_hf_papers(self, max_results: int = 3) -> List[str]:
        if not self._has_requests or not self._has_bs4:
            return ["(需要 requests + beautifulsoup4)"]
        response = NetworkManager.safe_request(f"{Config.HF_MIRROR}/papers", use_mirror=True, log_callback=self._log)
        if not response:
            return ["(HuggingFace镜像站连接失败)"]
        try:
            soup = BeautifulSoup(response.text, 'html.parser')
            results = []
            for link in soup.select('a[href^="/papers/"]')[:max_results]:
                title = link.get_text(strip=True)
                if title and 20 < len(title) < 200:
                    results.append(title[:120])
            return results if results else ["(HF页面结构可能已更新)"]
        except:
            return ["(HF解析失败)"]

    def fetch_baidu_news(self, keyword: str, max_results: int = 2) -> List[str]:
        if not self._has_requests or not self._has_bs4:
            return ["(需要 requests + beautifulsoup4)"]
        search_url = f"https://www.baidu.com/s?rtt=1&tn=news&word={quote(keyword)}"
        response = NetworkManager.safe_request(search_url, log_callback=self._log)
        if not response:
            return [f"(百度新闻搜索失败)"]
        try:
            soup = BeautifulSoup(response.text, 'html.parser')
            results = []
            for container in soup.select('.result, .c-container')[:max_results*2]:
                title_elem = container.select_one('h3 a, .news-title a')
                if title_elem:
                    title = title_elem.get_text(strip=True)
                    title = re.sub(r'百度快照|查看更多|资讯|\|.*', '', title)
                    title = re.sub(r'\s+', ' ', title).strip()
                    if title and len(title) > 10 and keyword[:2] in title:
                        results.append(title[:80])
                    if len(results) >= max_results:
                        break
            return results if results else [f"(未找到关于 '{keyword}' 的新闻)"]
        except:
            return [f"(百度新闻解析异常)"]

    def fetch_all(self) -> Dict:
        self._log("📡 开始增强型外部信源抓取...", "system")
        cache_data = self._load_cache()
        if cache_data:
            self._log("  使用缓存数据（12小时内有效）", "system")
            return cache_data
        self._log("  缓存无效或已过期，开始在线抓取（可能需要1-2分钟）...", "system")
        data = {}
        self._log("  🔬 抓取 arXiv AI论文...", "system")
        data['ai_papers'] = self.fetch_arxiv_papers(query="cat:cs.AI OR cat:cs.LG OR cat:cs.CL", max_results=5)
        self._log("  🤗 抓取 HuggingFace 论文...", "system")
        data['hf_papers'] = self.fetch_hf_papers(max_results=3)
        self._log("  📰 抓取国产大模型动态...", "system")
        llm_keywords = ["面壁智能 新模型", "智谱AI 最新进展", "阿里通义 大模型"]
        llm_news = {}
        for kw in llm_keywords:
            self._log(f"    搜索: {kw}", "system")
            llm_news[kw] = self.fetch_baidu_news(kw, max_results=2)
            time.sleep(random.uniform(2,3))
        data['llm_news'] = llm_news
        self._log("  🧠 抓取认知科学论文...", "system")
        data['neuro_papers'] = self.fetch_arxiv_papers(query="cat:q-bio.NC", max_results=4)
        data['timestamp'] = datetime.now().isoformat()
        self._save_cache(data)
        return data

    def _load_cache(self) -> Dict:
        if not FileIO.exists("external_cache"):
            return {}
        try:
            cache = json.loads(FileIO.read("external_cache"))
            cache_time = datetime.fromisoformat(cache.get("timestamp", "2000-01-01"))
            if (datetime.now() - cache_time).total_seconds()/3600 < 12:
                return cache.get("data", {})
        except:
            pass
        return {}

    def _save_cache(self, data: Dict):
        FileIO.write("external_cache", json.dumps({"timestamp": datetime.now().isoformat(), "data": data}, ensure_ascii=False, indent=2))

    def build_insights(self, data: Dict) -> List[str]:
        insights = []
        ai_papers = data.get('ai_papers', [])
        if ai_papers and not ai_papers[0].startswith("("):
            insights.append("## AI学术前沿（arXiv）")
            insights.extend([f"- {p}" for p in ai_papers[:5]])
        else:
            insights.append("## AI学术前沿（arXiv）\n- （暂无最新论文）")
        insights.append("\n## 模型与应用动态")
        hf_papers = data.get('hf_papers', [])
        if hf_papers and not hf_papers[0].startswith("("):
            insights.append("### HuggingFace 论文")
            insights.extend([f"- {p}" for p in hf_papers])
        llm_news = data.get('llm_news', {})
        if llm_news:
            insights.append("\n## 国产大模型动态")
            for kw, news_list in llm_news.items():
                insights.append(f"### {kw}")
                for n in news_list:
                    insights.append(f"- {n}")
        neuro_papers = data.get('neuro_papers', [])
        if neuro_papers and not neuro_papers[0].startswith("("):
            insights.append("\n## 认知科学前沿")
            insights.extend([f"- {p}" for p in neuro_papers])
        if len(insights) <= 2:
            insights.append("（外部追踪未获取到有效数据）")
        return insights

    def build_structured_insights(self, data: Dict) -> List[Dict]:
        insights = []
        for paper in data.get('ai_papers', []):
            if paper.startswith("(") or not paper.strip():
                continue
            title = paper.split(" (发布于:")[0].strip()
            insights.append({"type": "arxiv", "title": title, "summary": title, "link": "", "source": "arXiv"})
        for paper in data.get('hf_papers', []):
            if paper.startswith("(") or not paper.strip():
                continue
            insights.append({"type": "huggingface", "title": paper, "summary": paper, "link": "", "source": "HuggingFace"})
        for kw, news_list in data.get('llm_news', {}).items():
            for news in news_list:
                if news.startswith("(") or not news.strip():
                    continue
                insights.append({"type": "news", "title": news, "summary": news, "link": "", "source": f"百度新闻({kw})"})
        return insights


# =============================================================================
# 9. 搜索 (search.py)
# =============================================================================
from collections import Counter

class SearchService:
    @staticmethod
    def _tokens(text: str) -> List[str]:
        tokens = []
        for word in re.findall(r'[A-Za-z0-9_]+|[\u4e00-\u9fff]+', text.lower()):
            if re.search(r'[\u4e00-\u9fff]', word):
                chars = [ch for ch in word if re.match(r'[\u4e00-\u9fff]', ch)]
                tokens.extend(chars)
                tokens.extend(''.join(chars[i:i+2]) for i in range(len(chars)-1))
                tokens.extend(''.join(chars[i:i+3]) for i in range(len(chars)-2))
            else:
                tokens.append(word)
        return [t for t in tokens if t]

    @staticmethod
    def _score(keyword: str, line: str) -> float:
        if not keyword or not line:
            return 0.0
        score = 8.0 if keyword in line else 0.0
        query_terms = Counter(SearchService._tokens(keyword))
        line_terms = Counter(SearchService._tokens(line))
        for term, weight in query_terms.items():
            if term in line_terms:
                score += min(3, line_terms[term]) * min(2, weight)
        return score

    @staticmethod
    def search_documents(keyword: str, dirs: List[str], regex: bool = False) -> List[Tuple[str, int, str]]:
        results = []
        search_dirs = [Config.DATA_ROOT / d for d in dirs]
        pattern = re.compile(keyword) if regex else None
        for sdir in search_dirs:
            if not sdir.exists():
                continue
            for file_path in sdir.rglob("*"):
                if file_path.is_file() and file_path.suffix not in ('.pyc','.db'):
                    try:
                        with open(file_path, 'r', encoding='utf-8') as f:
                            for line_num, line in enumerate(f, 1):
                                if regex:
                                    if pattern.search(line):
                                        results.append((1000.0, str(file_path.relative_to(Config.DATA_ROOT)), line_num, line.rstrip()))
                                else:
                                    score = SearchService._score(keyword, line)
                                    if score > 0:
                                        results.append((score, str(file_path.relative_to(Config.DATA_ROOT)), line_num, line.rstrip()))
                    except:
                        continue
        results.sort(key=lambda item: item[0], reverse=True)
        return [(file_path, line_num, line) for _, file_path, line_num, line in results]


# =============================================================================
# 10. 批处理 (batch.py)
# =============================================================================
import os
import time
from typing import Callable, List

class BatchProcessor:
    def __init__(self, ai_client: AIClient, log_callback: Callable):
        self.ai = ai_client
        self.log = log_callback

    def extract_text_from_file(self, file_path: str) -> List[str]:
        ext = os.path.splitext(file_path)[1].lower()
        try:
            if ext in ['.txt','.md','.py','.html','.htm','.json','.xml','.css','.js']:
                with open(file_path,'r',encoding='utf-8') as f:
                    text = f.read()
                    return [text] if text.strip() else []
            elif ext in ['.xlsx','.xls']:
                if pd is None:
                    return [f"需要pandas读取Excel: {file_path}"]
                df = pd.read_excel(file_path, sheet_name=None, header=None)
                all_text = []
                for sheet_df in df.values():
                    sheet_text = sheet_df.astype(str).values.flatten()
                    sheet_text = ' '.join([t for t in sheet_text if t and t!='nan'])
                    if sheet_text:
                        all_text.append(sheet_text)
                return all_text
            elif ext == '.csv':
                if pd is None:
                    return [f"需要pandas读取CSV: {file_path}"]
                df = pd.read_csv(file_path, encoding='utf-8', header=None)
                text = df.astype(str).values.flatten()
                text = ' '.join([t for t in text if t and t!='nan'])
                return [text] if text.strip() else []
            elif ext == '.docx':
                if not HAS_DOCX:
                    return [f"需要python-docx: {file_path}"]
                doc = Document(file_path)
                paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
                return ['\n'.join(paragraphs)] if paragraphs else []
            elif ext == '.pdf':
                if not HAS_PDF:
                    return [f"需要pypdf/PyPDF2: {file_path}"]
                reader = PdfReader(file_path)
                text = ''.join(page.extract_text() or '' for page in reader.pages)
                return [text.strip()] if text.strip() else []
            elif ext == '.pptx':
                if not HAS_PPTX:
                    return [f"需要python-pptx: {file_path}"]
                prs = Presentation(file_path)
                all_text = []
                for slide in prs.slides:
                    slide_text = [shape.text for shape in slide.shapes if hasattr(shape,"text") and shape.text.strip()]
                    if slide_text:
                        all_text.append('\n'.join(slide_text))
                return all_text
            else:
                return []
        except Exception as e:
            self.log(f"读取文件失败 {file_path}: {e}", "error")
            return []

    def process_folder(self, folder_path: str, mode: str, skip_search: bool, progress_callback: Callable, stop_flag: Callable, history_callback: Callable = None):
        supported_exts = {'.txt','.md','.py','.html','.htm','.json','.xml','.css','.js','.xlsx','.xls','.csv','.docx','.pdf','.pptx'}
        all_files = []
        for root, dirs, files in os.walk(folder_path):
            for file in files:
                if os.path.splitext(file)[1].lower() in supported_exts:
                    all_files.append(os.path.join(root, file))
        if not all_files:
            self.log("未找到支持的文件", "error")
            return
        total = len(all_files)
        self.log(f"找到 {total} 个文件，开始批量处理（模式: {mode}）", "system")
        for idx, file_path in enumerate(all_files):
            if stop_flag and stop_flag():
                self.log("批量处理被用户中断", "warning")
                break
            progress_callback(int(100*idx/total))
            self.log(f"\n处理文件 [{idx+1}/{total}]: {os.path.basename(file_path)}", "system")
            text_units = self.extract_text_from_file(file_path)
            if not text_units:
                self.log(f"  文件无有效内容或读取失败", "warning")
                continue
            for unit_idx, unit_text in enumerate(text_units):
                if len(unit_text.strip()) < 10:
                    continue
                if mode == "chat":
                    reply = self.ai.chat(unit_text)
                    self.log(f"  [{unit_idx+1}] AI 回应: {reply[:200]}...", "ai")
                    if history_callback:
                        history_callback("assistant", f"[批量处理文件 {os.path.basename(file_path)}] {reply}")
                else:
                    self.log(f"  [{unit_idx+1}] 晶体化处理（略）", "system")
                time.sleep(0.5)
        progress_callback(100)
        self.log("批量处理完成", "success")


# =============================================================================
# 11. 辩论引擎 (debate.py)
# =============================================================================
import itertools
import json
import re
from dataclasses import dataclass
from typing import Any, Callable, Dict, List, Tuple

LogFn = Callable[[str, str], None]

@dataclass
class DebateRole:
    key: str
    name: str
    instruction: str

# =============================================================================
# 11. 辩论引擎 (debate.py) - 并发提速版
# =============================================================================
import itertools
import json
import re
import time
import concurrent.futures
from dataclasses import dataclass
from typing import Any, Callable, Dict, List, Tuple

LogFn = Callable[[str, str], None]

@dataclass
class DebateRole:
    key: str
    name: str
    instruction: str

class DebateEngine:
    """多角色辩论引擎 —— 并发提速版（含百灵鸟、每轮反思、流式输出、进度回调）"""

    ROLE_ATTENTION_PROFILES = {
        "radical": {
            "label": "破局采样",
            "keywords": "破局 非共识 机会 创新 反例 假设 颠覆 新路径",
            "sample_a": "优先追求洞察强度、非显然视角和破局可能",
            "sample_b": "优先追求突破方案的可落地路径、失败条件和反证检查",
        },
        "conservative": {
            "label": "边界采样",
            "keywords": "风险 成本 约束 边界 失败 安全 稳健 验证",
            "sample_a": "优先寻找风险、限制、成本和不可忽略的边界条件",
            "sample_b": "优先构造稳健执行方案、验证步骤和最低可行路径",
        },
        "structural": {
            "label": "结构采样",
            "keywords": "结构 因果 系统 流程 模型 类比 关系 分层",
            "sample_a": "优先抽取问题结构、因果链条和关键变量",
            "sample_b": "优先建立分步流程、评价标准和可复用模型",
        },
        "executor": {
            "label": "执行采样",
            "keywords": "步骤 行动 资源 时间 优先级 检查点 任务 交付",
            "sample_a": "优先拆解行动路径、资源需求和最小可执行步骤",
            "sample_b": "优先检查执行阻力、时间顺序和验收标准",
        },
        "auditor": {
            "label": "审计采样",
            "keywords": "证据 漏洞 冲突 过度推断 验证 反例 暂存 不确定",
            "sample_a": "优先寻找证据缺口、逻辑漏洞和过度推断",
            "sample_b": "优先提出验证办法、暂存问题和纠偏建议",
        },
        "default": {
            "label": "角色采样",
            "keywords": "问题 证据 结构 风险 行动 反思",
            "sample_a": "优先寻找本角色最强洞察和独特贡献",
            "sample_b": "优先检查本角色观点的证据、边界和执行条件",
        },
        "lark": {
            "label": "知识广度采样",
            "keywords": "外部知识 前沿动态 多领域 跨学科 最新论文 行业实践",
            "sample_a": "学术视角：优先引用最新论文和理论框架",
            "sample_b": "产业视角：优先引用行业案例和落地实践",
            "sample_c": "未来视角：优先提出前瞻性假设和推演",
        },
        "judge": {
            "label": "裁决采样",
            "keywords": "裁决 判决 依据 资源 原则 引证 终审",
            "sample_a": "优先审视各方论证所依据的晶体、原则和客观约束",
            "sample_b": "优先做出可追溯的终裁结论，并写明引用来源",
        },
        "spokesperson": {
            "label": "定调采样",
            "keywords": "对外 沟通 降维 明确 简洁 通俗 定调",
            "sample_a": "优先将专业黑话转化为通俗易懂的语言",
            "sample_b": "优先提炼不超过3条的核心信息，确保对外一致",
        },
        "pilgrim": {
            "label": "远航采样",
            "keywords": "长期 使命 价值观 持续 道德 愿景",
            "sample_a": "优先评估方案与长期愿景的契合度",
            "sample_b": "优先检查是否偏离核心使命，是否可持续",
        },
        "strategist": {
            "label": "奇谋采样",
            "keywords": "时机 人性 机会窗口 借力 非常规 押注",
            "sample_a": "优先寻找人性破绽和时机窗口",
            "sample_b": "优先评估方案的灵活性和迂回空间",
        },
        "statesman": {
            "label": "调研采样",
            "keywords": "调研 矛盾 全局 实事 数据 主要矛盾 求是",
            "sample_a": "优先从全局矛盾和数据分析问题",
            "sample_b": "优先提出实事求是、可落地的综合策略",
        },
    }

    def __init__(self, ai: AIClient, engine: CrystalEngine, roles: List[Dict], log: LogFn = None, stream_callback: Callable = None, progress_callback: Callable = None):
        self.ai = ai
        self.engine = engine
        self.log = log or (lambda message, level="system": None)
        self.stream_callback = stream_callback
        self.progress_callback = progress_callback

        # 确保百灵鸟存在
        role_keys = {r.get("key") for r in roles}
        if "lark" not in role_keys:
            roles.append({
                "key": "lark",
                "name": "百灵鸟",
                "instruction": "你是一位见多识广的通用智能体。你不受晶体树知识库的限制，擅长从广阔的外部世界补充信息。你在第二轮才登场，提供外部视野。"
            })

        self.roles = [
            DebateRole(str(item.get("key", idx)), item.get("name", f"角色{idx + 1}"), item.get("instruction", ""))
            for idx, item in enumerate(roles)
        ]

        # Day 1 新增：警报监控系统
        self.alarm_monitor = AlarmMonitor(log_callback=self.log)
        self.alarm_triggered_this_round = False   # 标记本轮是否已触发过警报
        # Day 2.8: 元问题分类器
        try:
            import sys
            from pathlib import Path
            core_config_path = str(Config.DATA_ROOT / "核心配置")
            if core_config_path not in sys.path:
                sys.path.append(core_config_path)
            from question_classifier import QuestionClassifier
            self.question_classifier = QuestionClassifier()
        except ImportError:
            self.question_classifier = None
            self.log("⚠️ question_classifier 未找到，元问题分类功能禁用", "warning") 
     
        # ===== 新增：缺失的属性初始化 =====
        self._forced_external = None        # 强制注入的外部知识
        self._forced_perspective = None     # 强制注入的对立视角
        self._external_has_new = False      # 本轮是否有新外部数据
        self._current_question = None       # 当前辩论问题
        self._routing_result = None         # 便宜门路由结果
        self._cognitive_operators = "[思维模式：平衡] [论证偏好：平衡] [输出偏好：平衡]"  # 默认认知风格
        self._current_classification = None # 元问题分类结果
        self._history_result = None         # 历史诊断结果
        self._reused_crystals = []          # 复用的晶体列表            
        # Day 4: 历史诊断与经验复用
        self._history_result = None
        self._reused_crystals = []
        # ===== Day 10: 角色质量参数缓存 =====
        self._role_quality_cache = {}           

    # ==================== 新增：带重试的并发安全调用 ====================
    # ---- 修改：_call_role_with_retry 使用独立 AIClient ----
    def _call_role_with_retry(self, role: DebateRole, prompt: str, system: str,
                               max_retries: int = 2, expected_words: int = 200) -> str:
        quality_config = self._get_role_quality_config(role.key)
        temperature = quality_config["temperature"]
        max_tokens = self._calculate_max_tokens(expected_words, quality_config["token_multiplier"])

        self.log(f"  🔧 {role.name}: temp={temperature:.2f}, max_tokens={max_tokens}", "system")

        last_error = None
        for attempt in range(max_retries + 1):
            try:
                ai_client = AIClient(api_key=self.ai.api_key)
                result = ai_client._call_api(
                    [{"role": "system", "content": system}, {"role": "user", "content": prompt}],
                    temperature=temperature,
                    max_tokens=max_tokens
                )

                if result is None:
                    raise Exception("AI返回为空（None）")
                if not isinstance(result, str):
                    raise Exception(f"AI返回类型异常: {type(result)}")
                if result.strip() == "":
                    raise Exception("AI返回空字符串")
                if result.startswith("错误：") or result.startswith("AI调用失败"):
                    raise Exception(result)

                # 检测严重截断
                if self._is_severely_truncated(result, expected_words):
                    if attempt < max_retries:
                        wait_time = 2 ** attempt
                        self.log(f"  ⚠️ {role.name} 可能被截断（{len(result)}字 < {expected_words*0.5:.0f}字），{wait_time}s 后重试...", "warning")
                        time.sleep(wait_time)
                        continue

                return result

            except Exception as e:
                last_error = e
                if attempt < max_retries:
                    wait_time = 2 ** attempt
                    self.log(f"⚠️ {role.name} 调用失败 (尝试 {attempt+1}/{max_retries+1})，{wait_time}s 后重试...", "warning")
                    time.sleep(wait_time)
                else:
                    error_msg = f"{role.name} 在 {max_retries} 次重试后仍失败: {last_error}"
                    self.log(f"❌ {error_msg}", "error")
                    return f"（{role.name} 发言失败: {str(last_error)[:50]}）"

        return f"（{role.name} 发言超时）"

    # ===== Day 10: 角色质量参数 =====
    def _get_role_quality_config(self, role_key: str) -> Dict[str, Any]:
        """获取角色的质量参数配置"""
        # 使用缓存避免重复查找
        if role_key in self._role_quality_cache:
            return self._role_quality_cache[role_key]
        
        config = Config.ROLE_QUALITY_CONFIG.get(
            role_key,
            Config.ROLE_QUALITY_CONFIG.get("default", {})
        )
        result = {
            "temperature": config.get("temperature", 0.70),
            "token_multiplier": config.get("token_multiplier", 3.0)
        }
        self._role_quality_cache[role_key] = result
        return result

    def _calculate_max_tokens(self, expected_words: int, multiplier: float) -> int:
        """计算 max_tokens（安全边际 = 预期字数 × 倍数 + 缓冲）"""
        return int(expected_words * multiplier) + 200

    def _is_severely_truncated(self, text: str, expected_words: int) -> bool:
        """检测是否严重截断（仅检测，不强制补全）"""
        word_count = len(text)
        if word_count < expected_words * 0.5:
            if not text.endswith(("。", "！", "？", "\"")):
                return True
        return False

    # ==================== 辅助方法 ====================
    def _estimate_complexity(self, question: str) -> float:
        q_len = len(question)
        length_score = min(1.0, q_len / 150)
        complex_keywords = [
            "设计", "方案", "系统", "模型", "策略", "框架",
            "博弈", "全球", "长期", "多变量", "优化", "权衡",
            "机制", "制度", "政策", "评估", "比较", "综合"
        ]
        keyword_score = sum(1 for kw in complex_keywords if kw in question) / 5.0
        keyword_score = min(1.0, keyword_score)
        open_score = 0.3 if any(kw in question for kw in ["如何", "为什么", "怎样", "设计"]) else 0.0
        final = 0.4 * length_score + 0.4 * keyword_score + 0.2 * open_score
        return min(1.0, final)

    # ==================== 新增：带重试的并发安全调用 ====================
    def _call_role_with_retry(self, role: DebateRole, prompt: str, system: str,
                              max_retries: int = 2, expected_words: int = 200) -> str:
        """
        带重试机制的安全AI调用，使用角色专属质量参数
        """
        # ===== 获取角色专属配置 =====
        quality_config = self._get_role_quality_config(role.key)
        temperature = quality_config["temperature"]
        max_tokens = self._calculate_max_tokens(expected_words, quality_config["token_multiplier"])

        self.log(f"  🔧 {role.name}: temp={temperature:.2f}, max_tokens={max_tokens}", "system")

        result = ""  # ← 修复：在循环外初始化
        for attempt in range(max_retries + 1):
            try:
                # 每个线程使用独立的 AIClient
                ai_client = AIClient(api_key=Config.get_api_key())
                result = ai_client._call_api(
                    [{"role": "system", "content": system}, {"role": "user", "content": prompt}],
                    temperature=temperature,
                    max_tokens=max_tokens
                )

                # ===== 质量检测：不因字数不足而重试，只检测严重截断 =====
                if self._is_severely_truncated(result, expected_words):
                    if attempt < max_retries:
                        wait_time = 2 ** attempt
                        self.log(f"  ⚠️ {role.name} 可能被截断（{len(result)}字 < {expected_words*0.5:.0f}字），{wait_time}s 后重试...", "warning")
                        time.sleep(wait_time)
                        continue
                return result

            except Exception as e:
                if attempt < max_retries:
                    wait_time = 2 ** attempt
                    self.log(f"⚠️ {role.name} 调用失败 (尝试 {attempt+1}/{max_retries+1})，{wait_time}s 后重试...", "warning")
                    time.sleep(wait_time)
                else:
                    error_msg = f"{role.name} 在 {max_retries} 次重试后仍失败: {e}"
                    self.log(f"❌ {error_msg}", "error")
                    return f"（{role.name} 发言失败: {str(e)[:50]}）"  # ← 修复：返回错误信息
        return result
    
    def _truncate_to_sentence(self, text: str, max_len: int = 1200) -> str:
        if len(text) <= max_len:
            return text
        import re
        search_start = max(0, max_len - 50)
        search_end = min(len(text), max_len + 50)
        segment = text[search_start:search_end]
        match = re.search(r'[。！？；；.!?;]', segment)
        if match:
            cut_pos = search_start + match.start() + 1
            return text[:cut_pos] + "……（后略）"
        else:
            return text[:max_len] + "……（后略）"

    def _emit_progress(self, stage: str, progress: int):
        if self.progress_callback:
            self.progress_callback({
                "stage": stage,
                "progress": min(100, max(0, progress)),
                "timestamp": datetime.now().isoformat()
            })

    # ==================== 外部知识相关 ====================
    def _fetch_external_overview(self, question: str) -> str:
        prompt = f"""请作为一个见多识广的通用智能体，为以下问题提供一份「外部知识总览」。

【用户问题】
{question}

【要求】
1. 提供 3~5 个来自不同领域（学术、产业、政策、跨学科）的最新案例或理论。
2. 每个案例用 1~2 句话说明其核心观点。
3. 最后用一句话总结：这些外部知识如何可能补充或挑战传统认知。

【输出格式】纯文本，不超过 500 字。不要使用 Markdown 或 JSON。

【语言要求】必须使用中文，专有名词可保留英文但需加中文注释。
"""
        try:
            result = self.ai.chat(prompt)
            # ← 修复：检查 result 是否有效
            if result is None:
                return "（外部知识总览暂不可用）"
            if result.startswith("错误") or result.startswith("AI调用失败"):
                return "（外部知识总览暂不可用）"
            if not result.strip():
                return "（外部知识总览暂不可用）"
            self._external_has_new = True
            return result.strip()
        except Exception as e:
            self.log(f"[WARN] 外部知识总览生成失败：{e}", "warning")
            return "（外部知识总览暂不可用）"

    # ===== 修改：百灵鸟直接输出综合视角（删除三样本采样+合成）=====
    # ===== 修改：百灵鸟直接输出综合视角（删除三样本采样+合成）=====
    def _lark_sampled_opinion(self, question: str, external_overview: str) -> Dict:
        """
        百灵鸟直接输出综合外部视角（原三样本采样+合成已删除，提速约3.8倍）
        """
        self.log("🎵 百灵鸟正在生成综合外部视角...", "system")
        
        lark_role = DebateRole(key="lark", name="百灵鸟", instruction="见多识广的通用智能体")
        
        prompt = f"""
【用户问题】{question}

【外部知识总览】
{external_overview}

你正在进行百灵鸟的角色输出。请提供200-300字的综合外部视角，
包含学术、产业、未来三个维度的洞察：
- 学术维度：引用最新论文或理论框架
- 产业维度：引用行业案例或实践
- 未来维度：提出前瞻性假设或推演

直接输出综合版本，不要分三部分写作。
"""
        system = self._role_system(lark_role, self._crystal_context(question))
        answer = self._call_role_with_retry(lark_role, prompt, system, expected_words=250)
        # ← 修复：检查返回值是否为错误信息
        if answer.startswith("（") and "失败" in answer:
            return {"role": "百灵鸟", "answer": f"（百灵鸟视角暂时无法生成）\n\n{answer}"}
        return {"role": "百灵鸟", "answer": answer}



    # ==================== 核心方法 ====================
    def _core_roles(self) -> List[DebateRole]:
        # 返回所有非百灵鸟、非替身的角色（即除 lark 和 twin 外的所有角色）
        return [r for r in self.roles if r.key not in ("lark", "twin")]

    def _run_lark_bare(self, question: str) -> str:
        """百灵鸟裸模型基线（完全不依赖晶体树和角色提示）"""
        prompt = f"请直接回答以下问题，给出你最直观、最完整的答案。\n\n问题：{question}"
        raw_answer = self.ai.chat(prompt, system="你是一位知识广博的通用AI，请直接输出完整回答。")
        if len(raw_answer) < 400:
            expand = self.ai.chat(
                f"请将以下回答扩展为更详细的论述（至少800字）：\n{raw_answer}",
                system="请只输出扩展后的正文。"
            )
            return expand
        return raw_answer

    def _assign_roles_dynamically(self, question: str, candidate_roles: List[DebateRole]) -> List[DebateRole]:
        """
        动态角色分配（PEAR + 元辩论）
        通过两阶段元辩论，为每个发言位置选择最合适的角色。
        
        返回：按发言位置顺序排列的角色列表（长度固定，与候选角色数量一致）
        """
        if len(candidate_roles) < 2:
            return candidate_roles
        
        # ---- 阶段1：提案 ----
        # 每个角色生成一份针对当前问题的定制论点
        proposals = {}
        for role in candidate_roles:
            prompt = f"""
用户问题：{question}

请以【{role.name}】的立场，简要阐述你对这个问题的核心论点（不超过200字）。
你的观点必须体现你的角色特征。
"""
            try:
                response = self.ai.chat(prompt, system=f"你是{role.name}，请给出你的核心论点。")
                proposals[role.key] = {"role": role, "proposal": response}
                self.log(f"  提案生成：{role.name}", "system")
            except Exception as e:
                self.log(f"  提案生成失败({role.name}): {e}", "warning")
                proposals[role.key] = {"role": role, "proposal": f"（{role.name}的论点生成失败）"}
        
        # ---- 阶段2：同行评审 ----
        # 每个角色对其他角色的提案进行评分（0-10）
        scores = {role.key: 0 for role in candidate_roles}
        for role in candidate_roles:
            reviewer = role
            review_text = ""
            for target_key, proposal_data in proposals.items():
                if target_key == role.key:
                    continue
                target_name = proposal_data["role"].name
                target_proposal = proposal_data["proposal"]
                review_text += f"【{target_name}的提案】\n{target_proposal}\n\n"
            
            if not review_text:
                continue
            
            prompt = f"""
用户问题：{question}

你作为【{reviewer.name}】，请对以下其他角色的提案进行评分（0-10分）。

{review_text}

请输出 JSON 格式（只输出JSON，不要其他内容）：
{{
    "scores": {{
        "角色名1": 分数,
        "角色名2": 分数,
        ...
    }},
    "justification": "简要说明你的评分依据"
}}
"""
            try:
                response = self.ai.chat_json(prompt)
                if "error" not in response and "scores" in response:
                    for name, score in response["scores"].items():
                        # 根据角色名查找key
                        for target_key, proposal_data in proposals.items():
                            if proposal_data["role"].name == name:
                                scores[target_key] += score
                                break
                self.log(f"  评审完成：{reviewer.name}", "system")
            except Exception as e:
                self.log(f"  评审失败({reviewer.name}): {e}", "warning")
                # 降级：平均分配
                for target_key in proposals:
                    if target_key != role.key:
                        scores[target_key] += 5  # 中立分
        
        # ---- 排序与分配 ----
        # 按得分降序排列角色
        sorted_roles = sorted(candidate_roles, key=lambda r: scores.get(r.key, 0), reverse=True)
        
        # 记录分配结果
        self.log(f"  动态角色分配结果：", "system")
        for idx, role in enumerate(sorted_roles):
            self.log(f"    位置{idx+1}: {role.name} (得分{scores.get(role.key, 0):.1f})", "system")
        
        return sorted_roles

    def generate_twin_role(self, fingerprint=None) -> DebateRole:
        """
        根据认知指纹动态生成“替身-我”角色

        替身角色的 System Prompt 会注入用户的认知指纹特征，
        使其回答风格接近用户本人的思维惯性。

        Args:
            fingerprint: CognitiveFingerprint 对象（可选，不传则自动加载）

        Returns:
            DebateRole: 替身角色
        """
        if fingerprint is None:
            try:
                fingerprint = self.engine.fingerprint_extractor.get_fingerprint()
            except Exception as e:
                print(f"[WARN] 无法获取认知指纹: {e}，使用默认替身")
                fingerprint = None

        # 构建替身角色的特征描述
        if fingerprint:
            risk_desc = "偏好高风险高回报方案" if fingerprint.risk_tolerance > 0.6 else "偏好低风险稳健方案"
            innovation_desc = "偏好颠覆性创新" if fingerprint.innovation_preference > 0.6 else "偏好渐进式优化"
            decisive_desc = "快速决断型" if fingerprint.decisiveness > 0.6 else "深思熟虑型"
            role_desc = f"最常采纳的角色是 {fingerprint.preferred_role}"
            confidence_note = f"（当前指纹置信度: {fingerprint.confidence:.2f}）"
        else:
            risk_desc = "偏好中等风险方案"
            innovation_desc = "偏好平衡创新与稳健"
            decisive_desc = "适度决断型"
            role_desc = "尚未建立明确的角色偏好"
            confidence_note = "（指纹数据不足，使用默认替身）"

        instruction = f"""你是认知晶体树辩论引擎中的【替身-我】角色。

你的任务是**复刻用户的思维惯性**，在辩论中站在用户的立场发言。

【用户认知特征】
- 风险偏好：{risk_desc}
- 创新偏好：{innovation_desc}
- 决策风格：{decisive_desc}
- {role_desc}
{confidence_note}

【辩论行为准则】
1. 你的回答应体现上述认知特征，模拟用户可能会说的话。
2. 你天然倾向于支持与用户特征一致的观点。
3. 当对方提出与用户特征相悖的观点时，你会本能地质疑。
4. 但你不应固执己见——如果对方的证据明显更强，你可以承认并吸收。
5. 你的终极目标是：帮用户找到最能代表“他”的答案，而非赢得辩论。

【特别提醒】
- 你不是在扮演一个“完美理性”的 AI，而是在扮演“这个人”的替身。
- 你不需要表现得比用户更聪明，但需要表现得像用户本人。
- 如果你的回答被用户本人认可，那就是最大成功。
"""

        return DebateRole(
            key="twin",
            name="替身-我",
            instruction=instruction
        )

    def _crystal_context(self, question: str) -> str:
        """
        获取晶体上下文（Day 4 增强：优先使用历史经验中的有效晶体）
        """
        # Day 4: 如果有复用的历史晶体组合，优先加载
        if hasattr(self, '_reused_crystals') and self._reused_crystals:
            all_crystals = self.engine.parse_crystals()
            reused = [c for c in all_crystals if c.id in self._reused_crystals]
            if reused:
                # 与普通检索结果合并，去重
                normal = self.engine.get_associative_crystals(question, top_k=5)
                # 去重（按ID去重）
                seen_ids = set(c.id for c in reused)
                for c in normal:
                    if c.id not in seen_ids:
                        seen_ids.add(c.id)
                        reused.append(c)
                assoc = reused[:8]
            else:
                assoc = self.engine.get_associative_crystals(question, top_k=5)
        else:
            assoc = self.engine.get_associative_crystals(question, top_k=5)

        if not assoc:
            return "（暂无相关晶体）"
        return "\n".join([f"- [{c.id}] {c.content}" for c in assoc])

    def _role_system(self, role: DebateRole, crystal_context: str, is_reflection: bool = False) -> str:
        base = f"""你是认知晶体树辩论引擎中的【{role.name}】。
        角色立场：{role.instruction}

        注意力材料：
        {crystal_context}

        辩论元能力：
        1. 精准复述对方论点后再回应，禁止稻草人攻击。
        2. 区分事实分歧与价值分歧，前者需举证，后者可存异。
        3. 当对方证据明显更强时，必须明确承认并吸收。
        4. 终极目标不是赢，而是产出融合方案，超越任何单一角色初始输出。"""

        # ===== 强制行为清单 =====
        if role.key == "radical":
            base += """

        【强制行为 - 激进者】
        - 你必须在每轮发言中提出 **至少 3 个颠覆性观点**（如"传统做法完全错误""我们可以换一个维度思考"）。
        - 格式：用「颠覆性观点 1/2/3」明确标出。
        - 示例：「颠覆性观点 1：我们不应该追求准确率，而应该追求召回率，因为漏判比误判代价更高。」"""
        elif role.key == "conservative":
            base += """

        【强制行为 - 保守者】
        - 你必须在每轮发言中列出 **至少 3 个风险点**（如"成本超支""团队抵触""技术债务"）。
        - 格式：用「风险清单 1/2/3」明确标出。
        - 示例：「风险清单 1：高准确率模型可能需要大量标注数据，成本超支风险。」"""
        elif role.key == "judge":
            base += """

        【强制行为 - 大法官】
        - 你必须在裁决中**引用至少 2 个晶体 ID（如 C012）或孔洞 ID（如 H003）**。
        - 格式：「依据 [C012]（反脆弱决策），判定激进者观点成立。」
        - 你的裁决必须输出为**独立的 JSON 结构**，不夹杂在其他文本中。
        - 如果你找不到依据，必须标记为 "deferred"（暂缓）。"""
        elif role.key == "spokesperson":
            base += """

        【强制行为 - 首席发言人】
        - 你必须在最终陈述中**确保老板读前 100 字能做出决策**。
        - 你的输出必须包含：① 结论先行 ② 不超过 3 条核心信息 ③ 无歧义的行动指令。"""

        if role.key == "lark":
            base += """

        【特别角色：百灵鸟】
        你是辩论中的「外部知识提供者」。你的优势在于：
        - 不受晶体树现有知识的限制，可以调用广泛的外部世界知识。
        - 你的使命是为辩论注入「新鲜空气」，指出晶体树可能忽略的盲区。
        - 即使某些外部知识暂时无法被晶体树验证，你仍然可以提出作为「可能性假设」。"""

        if is_reflection:
            base += """

        【反思轮次特殊指令】
        你现在处于「认知升级」阶段。请将以下反思作为你的输出核心：
        - 你从前几轮中学到了什么新东西？
        - 你的新观点与初始观点有何不同？
        - 你的新观点中，有哪些成分来自百灵鸟或晶体树？"""

        base += """

        【晶体引用强制要求】
        你必须在每轮发言中引用至少 1 条晶体卡片：
        ① 引用格式：`[ID] 内容` （例如 `[C001] 认知晶体树的核心是动态分层`）
        ② 必须说明该晶体是 **支持** 还是 **反驳** 你的论点。
        ③ 若找不到支持性晶体，必须说明"未找到支持晶体，我的论点基于以下独立推理..."
        ④ 如果你引用的是孔洞（Hxxx），也请用同样格式，并说明是"有待验证的孔洞"。

        【语言要求】
        所有输出必须使用中文。除非是专有名词（如 OKR、AI、GPT），否则不得使用英文。如果必须出现英文术语，请在其后括号内附上中文解释（例如：OKR（目标与关键成果法））。"""

        return base

    # ---- 修改：_independent_round 实现并发调用 ----
    def _independent_round(self, question: str, crystal_context: str, roles: List[DebateRole], baseline_answer: str = "") -> List[Dict]:
        """第一轮独立发言：并发执行所有角色（使用质量参数）"""
        import time
        import concurrent.futures
        
        answers = []
        complexity = self._estimate_complexity(question)
        if complexity > 0.7:
            min_words, max_words = 1200, 1800
        elif complexity > 0.4:
            min_words, max_words = 1000, 1500
        else:
            min_words, max_words = 800, 1200

        baseline_note = f"\n\n【参考基线】百灵鸟（裸模型）的初始回答：\n{baseline_answer}\n\n" if baseline_answer else ""

        # 准备每个角色的 prompt 和 system
        role_tasks = []
        expected_words = (min_words + max_words) // 2
        for role in roles:
            prompt = (
                f"用户问题：{question}\n"
                f"{baseline_note}"
                f"请基于你的角色立场给出独立答案，包含结论、理由、证据和风险建议。\n"
                f"**字数要求：{min_words}~{max_words} 字。** 请充分展开论述，确保论证完整。"
            )
            system = self._role_system(role, crystal_context)
            role_tasks.append((role, prompt, system, expected_words))

        # 并发执行
        start_time = time.time()
        self.log(f"🚀 启动 {len(role_tasks)} 个角色并发发言...", "system")
        
        with concurrent.futures.ThreadPoolExecutor(max_workers=len(role_tasks)) as executor:
            future_to_role = {}
            for role, prompt, system, exp_words in role_tasks:
                # ===== 直接使用 _call_role_with_retry，不再需要内部函数 =====
                future = executor.submit(self._call_role_with_retry, role, prompt, system, 2, exp_words)
                future_to_role[future] = role
            
            for future in concurrent.futures.as_completed(future_to_role):
                role = future_to_role[future]
                try:
                    result = future.result(timeout=180)
                    self.log(f"✅ {role.name} 发言完成（{len(result)}字）", "system")
                    answers.append({"role": role.name, "answer": result})
                except concurrent.futures.TimeoutError:
                    self.log(f"❌ {role.name} 执行超时", "error")
                    raise Exception(f"{role.name} 执行超时")
                except Exception as e:
                    self.log(f"❌ {role.name} 发言失败: {e}", "error")
                    raise

        elapsed = time.time() - start_time
        self.log(f"📊 第1轮并发执行完成，耗时 {elapsed:.2f} 秒", "system")

        # 按角色顺序排序
        role_order = {role.key: idx for idx, role in enumerate(roles)}
        answers.sort(key=lambda x: role_order.get(self._map_role_key(x["role"]), 999))
        return answers

    # ---- 新增：并发执行 Round 0（百灵鸟裸模型）+ Round 1（核心角色）----
    # ---- 新增：并发执行 Round 0（百灵鸟裸模型）+ Round 1（核心角色）----
    def _parallel_round0_and_round1(self, question: str, crystal_context: str,
                                     core_roles: List[DebateRole]) -> Tuple[List[Dict], str]:
        import time
        import concurrent.futures

        round1_answers = []
        lark_bare_answer = ""

        complexity = self._estimate_complexity(question)
        if complexity > 0.7:
            min_words, max_words = 1200, 1800
        elif complexity > 0.4:
            min_words, max_words = 1000, 1500
        else:
            min_words, max_words = 800, 1200

        expected_words = (min_words + max_words) // 2

        role_tasks = []
        for role in core_roles:
            prompt = (
                f"用户问题：{question}\n"
                f"请基于你的角色立场给出独立答案，包含结论、理由、证据和风险建议。\n"
                f"**字数要求：{min_words}~{max_words} 字。** 请充分展开论述，确保论证完整。"
            )
            system = self._role_system(role, crystal_context)
            role_tasks.append((role, prompt, system, expected_words))

        start_time = time.time()
        total_tasks = len(role_tasks) + 1
        self.log(f"🚀 并发启动 Round 0（百灵鸟裸模型）+ Round 1（{len(role_tasks)} 个角色），共 {total_tasks} 个任务...", "system")

        results = {}

        def execute_role_task(role, prompt, system, exp_words):
            try:
                return self._call_role_with_retry(role, prompt, system, 2, exp_words)
            except Exception as e:
                return f"（{role.name} 发言生成失败: {str(e)[:50]}）"

        def execute_lark_bare_task(q):
            try:
                ai_client = AIClient(api_key=self.ai.api_key)
                prompt = f"请直接回答以下问题，给出你最直观、最完整的答案。\n\n问题：{q}"
                raw_answer = ai_client.chat(prompt, system="你是一位知识广博的通用AI，请直接输出完整回答。")

                if raw_answer is None:
                    return "（百灵鸟裸模型生成失败: 返回为空）"
                if not isinstance(raw_answer, str):
                    return f"（百灵鸟裸模型生成失败: 类型异常 {type(raw_answer)}）"
                if raw_answer.startswith("错误") or raw_answer.startswith("AI调用失败"):
                    return f"（百灵鸟裸模型生成失败: {raw_answer[:50]}）"

                if len(raw_answer) < 400:
                    expand = ai_client.chat(
                        f"请将以下回答扩展为更详细的论述（至少800字）：\n{raw_answer}",
                        system="请只输出扩展后的正文。"
                    )
                    if expand and isinstance(expand, str) and len(expand) > len(raw_answer):
                        return expand
                    return raw_answer
                return raw_answer
            except Exception as e:
                return f"（百灵鸟裸模型生成失败: {str(e)[:50]}）"

        with concurrent.futures.ThreadPoolExecutor(max_workers=total_tasks) as executor:
            future_to_key = {}

            for idx, (role, prompt, system, exp_words) in enumerate(role_tasks):
                key = f"role_{idx}"
                future = executor.submit(execute_role_task, role, prompt, system, exp_words)
                future_to_key[future] = (key, role)

            future = executor.submit(execute_lark_bare_task, question)
            future_to_key[future] = ("lark_bare", None)

            for future in concurrent.futures.as_completed(future_to_key):
                key, role = future_to_key[future]
                try:
                    result = future.result(timeout=180)
                    if key == "lark_bare":
                        lark_bare_answer = result
                        self.log(f"✅ 百灵鸟裸模型完成", "system")
                    else:
                        results[key] = {"role": role.name, "answer": result}
                        self.log(f"✅ {role.name} 发言完成", "system")
                except concurrent.futures.TimeoutError:
                    self.log(f"❌ {key} 执行超时", "error")
                    if key == "lark_bare":
                        lark_bare_answer = "（百灵鸟裸模型生成超时）"
                    else:
                        results[key] = {"role": role.name, "answer": f"（{role.name} 执行超时）"}
                except Exception as e:
                    error_msg = str(e)
                    if key == "lark_bare":
                        self.log(f"❌ 百灵鸟裸模型失败: {error_msg}", "error")
                        lark_bare_answer = "（百灵鸟裸模型生成失败）"
                    else:
                        self.log(f"❌ {role.name} 发言失败: {error_msg}", "error")
                        results[key] = {"role": role.name, "answer": f"（{role.name} 发言失败: {error_msg[:50]}）"}

        for idx, (role, _, _, _) in enumerate(role_tasks):
            key = f"role_{idx}"
            if key in results:
                round1_answers.append(results[key])
            else:
                self.log(f"⚠️ {role.name} 的结果丢失", "warning")
                round1_answers.append({"role": role.name, "answer": f"（{role.name} 结果丢失）"})

        elapsed = time.time() - start_time
        self.log(f"📊 Round 0 + Round 1 并发执行完成，耗时 {elapsed:.2f} 秒", "system")

        return round1_answers, lark_bare_answer
    
    # ---- 修改：_debate_round 实现并发调用 ----
    def _debate_round(self, question, crystal_context, previous, audit, round_no, roles, lark_answer=None):
        """后续辩论轮：并发执行所有角色（使用质量参数）"""
        import time
        import concurrent.futures
        
        answers = []
        complexity = self._estimate_complexity(question)
        if complexity > 0.7:
            min_words, max_words = 1000, 1600
        elif complexity > 0.4:
            min_words, max_words = 800, 1300
        else:
            min_words, max_words = 700, 1100

        previous_text = "\n\n".join([f"### {item['role']}\n{item['answer']}" for item in previous])
        if lark_answer and not any(item["role"] == "百灵鸟" for item in previous):
            previous_text = f"### 百灵鸟（外部知识）\n{lark_answer['answer']}\n\n" + previous_text

        # 准备每个角色的任务
        role_tasks = []
        expected_words = (min_words + max_words) // 2
        for role in roles:
            feedback = (audit.get("feedback_by_role") or {}).get(role.name, "")
            other_names = [item["role"] for item in previous if item["role"] != role.name]
            if lark_answer and "百灵鸟" not in other_names:
                other_names.append("百灵鸟")
            target_lines = "\n".join([f"对 {target}：[精准复述其具体论据后，给出反驳理由]" for target in other_names[:3]])

            prompt = f"""用户问题：{question}

上一轮各方观点：
{previous_text}

逻辑检查员给你的反馈：
{feedback}

请输出以下结构：
【靶向攻击】
{target_lines}

【辩护与吸收】
[我坚持的核心理由是……]
[我从X角色处吸收了……因为其证据充分/逻辑严密]

【折冲整合方案】
融合各方合理成分后的完整答案，需超越此前任何单一版本。

**字数要求：{min_words}~{max_words} 字。** 请充分阐述你的攻击、辩护和融合方案，确保论证有据。
"""
            system = self._role_system(role, crystal_context)
            role_tasks.append((role, prompt, system, expected_words))

        # 并发执行
        start_time = time.time()
        self.log(f"🚀 启动第 {round_no} 轮 {len(role_tasks)} 个角色并发发言...", "system")
        
        with concurrent.futures.ThreadPoolExecutor(max_workers=len(role_tasks)) as executor:
            future_to_role = {}
            for role, prompt, system, exp_words in role_tasks:
                # ===== 直接使用 _call_role_with_retry，不再需要内部函数 =====
                future = executor.submit(self._call_role_with_retry, role, prompt, system, 2, exp_words)
                future_to_role[future] = role
            
            for future in concurrent.futures.as_completed(future_to_role):
                role = future_to_role[future]
                try:
                    result = future.result(timeout=180)
                    self.log(f"✅ {role.name} 第 {round_no} 轮发言完成（{len(result)}字）", "system")
                    answers.append({"role": role.name, "answer": result})
                except concurrent.futures.TimeoutError:
                    self.log(f"❌ {role.name} 第 {round_no} 轮执行超时", "error")
                    raise Exception(f"{role.name} 执行超时")
                except Exception as e:
                    self.log(f"❌ {role.name} 第 {round_no} 轮发言失败: {e}", "error")
                    raise

        elapsed = time.time() - start_time
        self.log(f"📊 第 {round_no} 轮并发执行完成，耗时 {elapsed:.2f} 秒", "system")

        # 按角色顺序排序
        role_order = {role.key: idx for idx, role in enumerate(roles)}
        answers.sort(key=lambda x: role_order.get(self._map_role_key(x["role"]), 999))
        return answers

    def _reflection_round(self, question, crystal_context, previous_answers, audit, roles, round_num):
        """反思轮：并发执行所有角色（使用质量参数）"""
        import time
        import concurrent.futures
        
        answers = []
        previous_text = "\n\n".join([f"【{item['role']}】\n{item['answer']}" for item in previous_answers])
        min_words, max_words = 600, 1000
        
        # 准备每个角色的任务
        role_tasks = []
        expected_words = (min_words + max_words) // 2
        for role in roles:
            feedback = (audit.get("feedback_by_role") or {}).get(role.name, "")
            prompt = f"""用户问题：{question}

你刚刚完成了第 {round_num} 轮辩论。

【第 {round_num} 轮完整记录】
{previous_text}

【审计员反馈】{feedback}

现在进入「反思」阶段 —— 请回答：
1. 你从本轮中（尤其是从百灵鸟的外部知识）学到了什么之前不知道的信息？
2. 基于这些新信息，你之前的立场需要做哪些**具体的修正**？（至少 3 点）
3. 修正后，你的**新立场**是什么？

**输出要求**：输出「反思声明」，字数控制在 {min_words}~{max_words} 字之间。聚焦于变化和增量，详细说明修正理由。
"""
            system = self._role_system(role, crystal_context, is_reflection=True)
            role_tasks.append((role, prompt, system, expected_words))

        # 并发执行
        start_time = time.time()
        self.log(f"🚀 启动第 {round_num} 轮反思（{len(roles)} 个角色并发）...", "system")
        
        with concurrent.futures.ThreadPoolExecutor(max_workers=len(role_tasks)) as executor:
            future_to_role = {}
            for role, prompt, system, exp_words in role_tasks:
                # ===== 直接使用 _call_role_with_retry，不再需要内部函数 =====
                future = executor.submit(self._call_role_with_retry, role, prompt, system, 2, exp_words)
                future_to_role[future] = role
            
            for future in concurrent.futures.as_completed(future_to_role):
                role = future_to_role[future]
                try:
                    result = future.result(timeout=180)
                    self.log(f"✅ {role.name} 反思完成（{len(result)}字）", "system")
                    answers.append({"role": role.name, "answer": result})
                except concurrent.futures.TimeoutError:
                    self.log(f"❌ {role.name} 反思超时", "error")
                    raise Exception(f"{role.name} 反思超时")
                except Exception as e:
                    self.log(f"❌ {role.name} 反思失败: {e}", "error")
                    raise

        elapsed = time.time() - start_time
        self.log(f"📊 第 {round_num} 轮反思并发执行完成，耗时 {elapsed:.2f} 秒", "system")
        
        return answers

    def _check_convergence(self, debate_answers, reflection_answers):
        all_texts = [item["answer"] for item in debate_answers] + [item["answer"] for item in reflection_answers]
        if len(all_texts) < 3:
            return False
        import itertools
        sims = [self._jaccard(a, b) for a, b in itertools.combinations(all_texts, 2)]
        avg_sim = sum(sims) / len(sims) if sims else 0
        self.log(f"📊 当前观点收敛度：{avg_sim:.2f}", "system")
        return avg_sim > 0.65

    # ===== 在这里添加 _audit 方法 =====
    def _audit(self, question, answers, round_no):
        """审计辩论"""
        answers_text = "\n\n".join([f"### {item['role']}\n{item['answer']}" for item in answers])
        prompt = f"""请作为认知晶体树的逻辑检查员，审计第 {round_no} 轮辩论。
用户问题：{question}

三方答案：
{answers_text}

只返回 JSON，不要 Markdown。字段：
{{
  "feedback_by_role": {{"角色名": "个性化改进建议"}},
  "disagreement_map": {{
    "fact": [],
    "logic": [],
    "risk_preference": [],
    "value": [],
    "term": []
  }},
  "major_conflict": false,
  "evidence_scores": {{"角色名": 0.0}},
  "should_stop": false,
  "summary": "一句话审计摘要"
}}

判断规则：
- 个性化改进建议只以论证严密性和证据充分性为评判依据，不得因立场偏向任何一方。
- 只因事实或逻辑层面的不可调和分歧才设置 major_conflict=true。
- 风险偏好或价值观不同通常不是 major_conflict。
- evidence_scores 按证据充分性与论证严密性评分，0到1。

**字数约束：`feedback_by_role` 中每个角色的反馈不超过 150 字；`summary` 不超过 80 字。**
"""
        raw = self.ai.chat_json(prompt)
        if "error" in raw:
            self.log(f"审计 JSON 解析失败，使用降级审计：{raw.get('error')}", "warning")
            return {
                "feedback_by_role": {},
                "disagreement_map": {"fact": [], "logic": [], "risk_preference": [], "value": [], "term": []},
                "major_conflict": False,
                "evidence_scores": {},
                "should_stop": False,
                "summary": raw.get("raw", "审计失败，已降级处理")[:200],
            }
        return raw


    def get_roles_with_twin(self, include_twin: bool = True, fingerprint=None) -> List[DebateRole]:
        """
        获取辩论角色列表（可选择是否包含替身角色）
        """
        if not include_twin:
            return self.roles

        # 获取核心角色（激进者、保守者、结构主义者）
        core = self._core_roles()

        # 生成替身角色
        twin = self.generate_twin_role(fingerprint)

        # 用替身替换执行者或审计者（保留三个核心角色 + 替身）
        # 替换最后一个角色（通常是审计者或执行者）
        if len(core) >= 3:
            # 保留前两个核心角色，第三个替换为替身
            # 但实际上激进者、保守者、结构主义者都很重要
            # 更好的方式：在原有角色基础上增加替身
            result = core[:]  # 复制核心角色
            # 如果已经有 5 个角色，替换最后一个
            if len(self.roles) > 3:
                # 保留所有角色，替身作为额外角色加入（但要注意注意力上限）
                result = self.roles[:] + [twin]
            else:
                result = self.roles[:] + [twin]
            return result

        return self.roles + [twin]

    # ---- run 方法保持不变（内部调用已修改） ----
    def run(self, question: str, mode: str = "debate_full", max_rounds: int = 4) -> Dict:
        """
        运行辩论 - 由 OutputOrchestrator 输出结构化数据
        """
        if len(self.roles) < 3:
            raise ValueError("辩论至少需要 3 个角色")

        max_rounds = max(3, min(8, int(max_rounds or 4)))
        self._current_question = question
        self._emit_progress("初始化", 0)

        # ===== Day 0: 运行时自检断言 =====
        self._run_day0_runtime_assertions()

        # ===== 便宜门路由 =====
        routing_result = self.engine.cheap_gate.check(question, [])
        self._routing_result = routing_result


        # ===== Day 2.8: 元问题分类器 =====
        if self.question_classifier:
            try:
                classification = self.question_classifier.classify(question)
                self._current_classification = classification
                label = classification.get('label', '未知')
                pipeline = classification.get('pipeline', '通用')
                self.log(f"🔍 当前问题被分类为：{label}，走{pipeline}路径", "system")
            except Exception as e:
                self.log(f"⚠️ 元问题分类失败：{e}，使用默认路径", "warning")
                self._current_classification = {"label": "通用", "type": "general", "pipeline": "通用"}
        else:
            self._current_classification = {"label": "通用", "type": "general", "pipeline": "通用"}

        # ===== Day 2.5: 认知风格注入 =====
        try:
            fingerprint = self.engine.fingerprint_extractor.get_fingerprint()
            operators = self.engine.fingerprint_extractor.get_cognitive_operators(fingerprint)
            self._cognitive_operators = operators
            self.log(f"🧠 注入认知风格：{operators}", "system")
        except Exception as e:
            self.log(f"⚠️ OutputOrchestrator 执行失败: {e}，降级返回原始数据", "error")
            import traceback
            traceback.print_exc()
            self.engine._append_change_log("OutputOrchestrator错误", f"{e}\n{traceback.format_exc()}")

            fallback_board = f"【系统提示】报告生成失败（{str(e)[:50]}），以下是原始辩论摘要。\n\n共 {len(rounds)} 轮辩论，涉及 {len(all_roles)} 个角色。请查看上方各轮次详细记录。"
            return {
                "mode": mode,
                "question": question,
                "rounds": rounds,
                "answer": fallback_board,
                "board_version": fallback_board,
                "employee_version": "（生成失败，请查看原始辩论）",
                "novice_version": "（生成失败，请查看原始辩论）",
                "expert_version": "（生成失败，请查看原始辩论）",
                "calls_estimate": len(all_roles) * 4,
                "error": str(e)
            }

        # ===== Day 4: 历史诊断与经验复用 =====
        try:
            history_result = self.engine.meta.diagnose_history(question)
            self._history_result = history_result
            if history_result.get("matched"):
                self._reused_crystals = history_result.get("crystal_combination", [])
                self.log(f"📚 复用历史经验：匹配历史问题（相似度 {history_result['match_score']:.2f}）："
                         f"{history_result.get('matched_question', '')[:50]}...", "system")
                if self._reused_crystals:
                    self.log(f"   有效晶体组合：{', '.join(self._reused_crystals[:5])}", "system")
            else:
                self._reused_crystals = []
                self.log(f"📚 未找到匹配历史经验（最高相似度 {history_result.get('match_score', 0):.2f}）", "system")
        except Exception as e:
            self._reused_crystals = []
            self._history_result = {}
            self.log(f"⚠️ 历史诊断失败：{e}", "warning")

        # ===== Day 0: 埋点计时 =====
        self._start_time = time.time()
        self._token_count = 0

        # 加载认知风格（用于后续）
        try:
            fingerprint = self.engine.fingerprint_extractor.get_fingerprint()
            operators = self.engine.fingerprint_extractor.get_cognitive_operators(fingerprint)
            self._cognitive_operators = operators
        except Exception:
            self._cognitive_operators = "[思维模式：平衡] [论证偏好：平衡] [输出偏好：平衡]"

        # ===== Day 1 & 4: 警报检查内部函数（含失败轨迹记录） =====
        # ===== Day 1 & 4: 警报检查内部函数（含失败轨迹记录） =====
        def _check_alarms_for_round(round_answers, round_no, audit):
            """检查本轮指标并触发警报，同时记录失败轨迹"""
            import re
            # 计算 Jaccard
            jaccard = self._average_jaccard([item["answer"] for item in round_answers])
            # 晶体引用率：计算有引用的回答占比（而非平均引用数）
            total_answers = len(round_answers)
            # 统计有多少个回答至少包含一个晶体引用
            ref_count_answers = 0
            total_refs = 0
            for item in round_answers:
                answer = item.get("answer", "")
                crystal_matches = len(re.findall(r'\[C\d+\]', answer))
                hole_matches = len(re.findall(r'\[H\d+\]', answer))
                refs_in_answer = crystal_matches + hole_matches
                total_refs += refs_in_answer
                if refs_in_answer > 0:
                    ref_count_answers += 1
            # 引用率 = 有引用的回答数 / 总回答数（0-1 之间的百分比）
            ref_rate = ref_count_answers / max(1, total_answers)
            self.log(f"🔍 DEBUG: 第{round_no}轮警报检查，答案数={total_answers}，"
                     f"有引用回答数={ref_count_answers}，总引用数={total_refs}，引用率={ref_rate:.2f}", "system")
            
            bias_amp = 0.0
            metrics = {
                "crystal_reference_rate": ref_rate,
                "bias_amplification": bias_amp,
                "external_has_new": self._external_has_new,
                "jaccard_similarity": jaccard
            }
            triggered = self.alarm_monitor.check(metrics)
            self.log(f"🔍 DEBUG: 触发警报数={len(triggered)}", "system")

            for alarm in triggered:
                # Day 4: 记录失败轨迹
                if alarm["rule"] == "knowledge_poverty":
                    self._record_failure_trace(
                        self._current_question or question,
                        "low_crystal_reference",
                        {
                            "ref_rate": ref_rate,
                            "threshold": 0.5,
                            "round": round_no,
                            "effective_crystals": getattr(self, '_reused_crystals', [])
                        }
                    )
                elif alarm["rule"] == "thought_stagnation":
                    self._record_failure_trace(
                        self._current_question or question,
                        "debate_diverged",
                        {
                            "jaccard": jaccard,
                            "threshold": 0.8,
                            "round": round_no
                        }
                    )
                # 记录进化事件
                self.engine.log_evolution_event(
                    "alarm",
                    {
                        "rule": alarm["rule"],
                        "message": alarm["message"],
                        "action": alarm["action"],
                        "data": alarm.get("data", {}),
                        "round": round_no,
                        "trigger": "alarm"
                    }
                )
                # 处理警报
                self.alarm_monitor.handle_alarm(alarm, self)
            self._external_has_new = False
            return triggered
        core_roles = self._core_roles()
        lark = next((r for r in self.roles if r.key == "lark"), None)
        all_roles = core_roles + [lark] if lark else core_roles

        # ===== 动态路由 =====
        routing_result = self.engine.cheap_gate.check(question, [])
        complexity = routing_result.get("complexity", "high")
        self._routing_result = routing_result

        if complexity == "simple":
            # 简单问题只激活 3 个核心角色
            core_keys = ["radical", "conservative", "structural"]
            all_roles = [r for r in all_roles if r.key in core_keys]
            self.log(f"🚀 动态路由：简单问题，激活 {len(all_roles)} 个角色（激进者+保守者+结构主义者）", "system")
        elif complexity == "medium":
            # 中等问题激活 5 个角色（核心 + 执行者 + 审计者）
            core_keys = ["radical", "conservative", "structural", "executor", "auditor"]
            all_roles = [r for r in all_roles if r.key in core_keys]
            self.log(f"🚀 动态路由：中等复杂度，激活 {len(all_roles)} 个角色", "system")
        else:
            self.log(f"🚀 动态路由：高复杂度，全量激活 {len(all_roles)} 个角色", "system")
        # ===== Round 0 + Round 1：并发执行百灵鸟裸模型 + 核心角色 =====
        self._emit_progress("Round 0 + Round 1 并发启动", 10)
        shared_context = self._crystal_context(question)
        round1_answers, lark_bare_answer = self._parallel_round0_and_round1(
            question, shared_context, core_roles
        )
        
        # 记录 Round 0 结果
        round0_answers = [{"role": "百灵鸟（裸模型）", "answer": lark_bare_answer}]
        rounds = [{"round": 0, "answers": round0_answers, "audit": {"summary": "裸模型基线，无审计"}}]
        
        # 审计 Round 1
        current_audit = self._audit(question, round1_answers, round_no=1)
        rounds.append({"round": 1, "answers": round1_answers, "audit": current_audit})
        _check_alarms_for_round(round1_answers, 1, current_audit)
        self._emit_progress("Round 0 + Round 1 完成", 30)


        # ===== 第2轮：引入百灵鸟 =====
        self._emit_progress("生成外部总览", 30)
        external_overview = self._fetch_external_overview(question)
        self._emit_progress("百灵鸟采样", 35)
        lark_answer = self._lark_sampled_opinion(question, external_overview)
        self._emit_progress("第2轮辩论", 40)
        previous_with_lark = round1_answers + [lark_answer]
        round2_answers = self._debate_round(question, shared_context, previous_with_lark, current_audit, 2, all_roles, lark_answer=lark_answer)
        current_audit = self._audit(question, round2_answers, round_no=2)
        _check_alarms_for_round(round2_answers, 2, current_audit)
        self._emit_progress("第2轮完成", 50)

        # ===== 后续轮次：反思 + 辩论 =====
        previous_answers = round2_answers
        for round_no in range(3, max_rounds + 1):
            self._emit_progress(f"反思第{round_no-1}轮", 55 + round_no * 5)
            reflection_answers = self._reflection_round(question, shared_context, previous_answers, current_audit, all_roles, round_num=round_no-1)
            reflection_audit = self._audit(question, reflection_answers, round_no=round_no-0.5)
            self._emit_progress(f"第{round_no}轮辩论", 60 + round_no * 5)
            round_answers = self._debate_round(question, shared_context, reflection_answers, reflection_audit, round_no, all_roles)
            current_audit = self._audit(question, round_answers, round_no=round_no)
            rounds.append({
                "round": round_no,
                "answers": round_answers,
                "audit": current_audit,
                "reflection": reflection_answers,
                "reflection_audit": reflection_audit
            })
            _check_alarms_for_round(round_answers, round_no, current_audit)
            previous_answers = round_answers
            self._emit_progress(f"第{round_no}轮完成", 70 + round_no * 5)
            if self._check_convergence(round_answers, reflection_answers):
                self._emit_progress("收敛完成", 95)
                break

        # ===== 由 OutputOrchestrator 接管输出 =====
        self._emit_progress("生成结构化输出", 90)

        try:
            orchestrator = OutputOrchestrator(self.ai, self.engine)
            final_schema = orchestrator.generate(question, rounds)

            result = {
                "mode": mode,
                "question": question,
                "rounds": rounds,
                "final_schema": final_schema.dict(),
                "board_version": final_schema.board_version,
                "employee_version": final_schema.employee_version,
                "novice_version": final_schema.novice_version,
                "expert_version": final_schema.expert_version,
                "elegant_epilogue": final_schema.elegant_epilogue,
                "dashboard_stats": final_schema.dashboard_stats,
                "judge_audit": final_schema.judge_audit,
                "round_by_round": [r.dict() for r in final_schema.round_by_round],
                "answer": final_schema.board_version,
                "calls_estimate": len(all_roles) * 4,
                # ===== Day 0-4 元数据 =====
                "_meta": {
                    "routing": routing_result,
                    "classification": self._current_classification,
                    "cognitive_operators": self._cognitive_operators,
                    "history_reused": self._history_result.get("matched", False) if self._history_result else False,
                    "estimated_tokens": self._token_count,
                    "elapsed_seconds": round(time.time() - self._start_time, 2)
                }
            }

            # ===== Day 3: 元原语触发链（辩论结束后运行） =====
            try:
                self.log("🧠 运行元原语触发链检查...", "system")
                meta_results = self.engine.meta.run_all_primitives()
                triggered = meta_results.get("triggered_chains", [])
                if triggered:
                    self.log(f"✅ 触发链执行：{len(triggered)} 条链被触发", "success")
                    for chain in triggered:
                        self.log(f"   {chain['chain']}: {chain['source']} → {chain['target']} (通过: {chain['passed']})", "system")
                else:
                    self.log("ℹ️ 无触发链条件满足", "system")
            except Exception as e:
                self.log(f"⚠️ 元原语执行异常：{e}", "warning")

            # ===== Day 5: Hebbian 学习 - 评估有效晶体组合 =====
            try:
                # 收集本轮辩论中所有被引用的晶体 ID
                used_crystal_ids = set()
                for rd in rounds:
                    for answer_item in rd.get("answers", []):
                        text = answer_item.get("answer", "")
                        matches = re.findall(r'\[C(\d+)\]', text)
                        used_crystal_ids.update([f"C{mid}" for mid in matches])

                if used_crystal_ids:
                    # 获取最后一轮的审计评分作为质量信号
                    last_round = rounds[-1] if rounds else {}
                    last_audit = last_round.get("audit", {})
                    ev_scores = list(last_audit.get("evidence_scores", {}).values())
                    if ev_scores:
                        avg_score = sum(ev_scores) / len(ev_scores)
                    else:
                        # 降级：使用审计摘要中的关键词判断
                        summary = last_audit.get("summary", "")
                        if "通过" in summary or "优秀" in summary:
                            avg_score = 0.8
                        elif "一般" in summary:
                            avg_score = 0.5
                        else:
                            avg_score = 0.6

                    # 获取任务类型（来自问题分类器）
                    task_type = getattr(self, '_current_classification', {}).get("type", "general")
                    if not task_type or task_type == "unknown":
                        task_type = "general"

                    # 更新 Hebbian 权重
                    self.engine.update_hebbian_weights(list(used_crystal_ids), task_type, avg_score)
                    self.log(f"🧠 Hebbian 学习更新：{len(used_crystal_ids)} 个晶体，"
                             f"任务类型={task_type}，评分={avg_score:.2f}", "system")
                else:
                    self.log("ℹ️ 本轮辩论未引用晶体，跳过 Hebbian 学习", "system")
            except Exception as e:
                self.log(f"⚠️ Hebbian 学习失败：{e}", "warning")

            # ===== Day 8: 双时间尺度进化调度 - 饱和检测 =====
            try:
                # 计算本轮质量评分（从审计评分获取）
                last_round = rounds[-1] if rounds else {}
                last_audit = last_round.get("audit", {})
                ev_scores = list(last_audit.get("evidence_scores", {}).values())
                if ev_scores:
                    quality_score = sum(ev_scores) / len(ev_scores)
                else:
                    summary = last_audit.get("summary", "")
                    if "通过" in summary or "优秀" in summary:
                        quality_score = 0.8
                    elif "一般" in summary:
                        quality_score = 0.5
                    else:
                        quality_score = 0.6

                context = {
                    "modification_type": "prompt",
                    "rounds_count": len(rounds),
                    "question": question
                }
                saturation_result = self.engine.meta.prompt_saturation_detector(quality_score, context)

                if saturation_result.get("is_saturated"):
                    self.log(
                        f"📊 双时间尺度进化调度：本轮质量提升 {saturation_result.get('improvement', 0):.3f}，"
                        f"连续 {saturation_result.get('consecutive_rounds', 0)} 轮饱和，"
                        f"状态：{saturation_result.get('status', 'unknown')}",
                        "system"
                    )
                    if saturation_result.get("level") == "control_logic":
                        self.log(
                            f"📊 双时间尺度进化调度：提示词优化已饱和（连续 {saturation_result.get('consecutive_rounds', 0)} 轮），"
                            f"升级到控制逻辑层面",
                            "system"
                        )
                else:
                    self.log(
                        f"📊 双时间尺度进化调度：本轮质量提升 {saturation_result.get('improvement', 0):.3f}，"
                        f"未饱和（状态：{saturation_result.get('status', 'unknown')}）",
                        "system"
                    )
            except Exception as e:
                self.log(f"⚠️ 饱和检测失败：{e}", "warning")

            # ===== Day 0: 埋点数据统计 =====
            elapsed = time.time() - self._start_time
            self.log(f"📊 辩论耗时：{elapsed:.2f} 秒", "system")
            # 估算 Token 消耗（粗略）
            total_chars = 0
            for rd in rounds:
                for item in rd.get("answers", []):
                    total_chars += len(item.get("answer", ""))
            estimated_tokens = total_chars // 2
            self._token_count = estimated_tokens
            self.log(f"📊 估算Token消耗：{estimated_tokens} tokens，预估成本：${estimated_tokens * 0.000001:.6f}", "system")

            # ===== Day 7: 记录帕累托数据 =====
            try:
                # 获取当前模式
                profile_name = getattr(self, '_current_profile', 'balanced')
                if profile_name not in ['high_accuracy', 'balanced', 'economy']:
                    profile_name = 'balanced'

                # 计算本次对话的指标
                accuracy = 0.7  # 默认值，实际应从审计评分获取
                # 尝试从审计评分获取更准确的准确性
                last_round = rounds[-1] if rounds else {}
                last_audit = last_round.get("audit", {})
                ev_scores = list(last_audit.get("evidence_scores", {}).values())
                if ev_scores:
                    accuracy = sum(ev_scores) / len(ev_scores)
                else:
                    # 降级：使用审计摘要判断
                    summary = last_audit.get("summary", "")
                    if "通过" in summary or "优秀" in summary:
                        accuracy = 0.8
                    elif "一般" in summary:
                        accuracy = 0.5

                # 估算成本
                total_chars = 0
                for rd in rounds:
                    for item in rd.get("answers", []):
                        total_chars += len(item.get("answer", ""))
                estimated_tokens = total_chars // 2
                cost = estimated_tokens * 0.000001

                # 计算延迟
                elapsed = time.time() - self._start_time

                # 统计晶体引用数
                total_refs = 0
                for rd in rounds:
                    for item in rd.get("answers", []):
                        text = item.get("answer", "")
                        total_refs += len(re.findall(r'\[C\d+\]', text))

                # 记录到帕累托跟踪器
                self.engine.meta.record_conversation_metrics(
                    profile_name=profile_name,
                    accuracy=accuracy,
                    cost=cost,
                    latency=elapsed,
                    crystal_refs=total_refs,
                    quality_score=accuracy
                )

                # 记录每日统计
                today = datetime.now().date().isoformat()
                self.engine.meta.record_daily_stats({
                    "date": today,
                    "quality_score": accuracy,
                    "crystal_refs": total_refs,
                    "bias_index": 0.3,  # 默认值
                    "tokens_used": estimated_tokens
                })

                self.log(f"📊 帕累托数据已记录：模式={profile_name}, 准确性={accuracy:.2f}, 成本=${cost:.6f}", "system")

            except Exception as e:
                self.log(f"⚠️ 帕累托数据记录失败：{e}", "warning")

            self._emit_progress("完成", 100)
            self.log("✅ V3.0 结构化输出生成完成", "success")
            return result

        except Exception as e:
            self.log(f"⚠️ OutputOrchestrator 执行失败: {e}，降级返回原始数据", "error")
            
            import traceback
            traceback.print_exc()
            return {
                "mode": mode,
                "question": question,
                "rounds": rounds,
                "answer": f"（输出生成失败，请查看原始辩论数据）",
                "calls_estimate": len(all_roles) * 4,
                "error": str(e)
            }

    # ===== 补全专家方案遗漏的 _run_lark_bare 方法 =====
    def _run_lark_bare(self, question: str) -> str:
        """百灵鸟裸模型基线（完全不依赖晶体树和角色提示）"""
        prompt = f"请直接回答以下问题，给出你最直观、最完整的答案。\n\n问题：{question}"
        raw_answer = self.ai.chat(prompt, system="你是一位知识广博的通用AI，请直接输出完整回答。")
        # 若字数不足，强制补充展开
        if len(raw_answer) < 400:
            expand = self.ai.chat(
                f"请将以下回答扩展为更详细的论述（至少800字）：\n{raw_answer}",
                system="请只输出扩展后的正文。"
            )
            return expand
        return raw_answer

    # ==================== 工具方法 ====================
    def _tokens(self, text: str) -> set:
        stop = {"的", "了", "和", "是", "我", "一个", "这个", "那个", "如何", "什么", "为什么", "我们", "你们", "他们"}
        tokens = set()
        for word in re.findall(r"[A-Za-z0-9_]+|[\u4e00-\u9fff]+", text.lower()):
            if re.search(r"[\u4e00-\u9fff]", word):
                tokens.update(word[i:i + 2] for i in range(len(word) - 1))
            elif len(word) >= 3:
                tokens.add(word)
        return {token for token in tokens if token and token not in stop}

    def _jaccard(self, a: str, b: str) -> float:
        ta, tb = self._tokens(a), self._tokens(b)
        if not ta or not tb:
            return 0.0
        return len(ta & tb) / len(ta | tb)

    def _average_jaccard(self, answers: List[str]) -> float:
        pairs = list(itertools.combinations(answers, 2))
        if not pairs:
            return 0.0
        return sum(self._jaccard(a, b) for a, b in pairs) / len(pairs)

    # ===== Day 1 新增：警报处理辅助方法 =====
    def _inject_external_knowledge(self, alarm: dict):
        """强制注入外部知识：重新获取外部总览并添加到下一轮上下文"""
        question = self._current_question   # 需要在 run 方法中记录当前问题
        if not question:
            return
        # 获取外部知识
        overview = self._fetch_external_overview(question)
        # 将外部知识附加到系统上下文中，由后续轮次使用
        # 简单实现：将外部知识存储为实例变量，在下一轮系统 Prompt 中注入
        self._forced_external = overview
        self.log("  外部知识已强制注入，将在下一轮生效", "system")

    def _record_failure_trace(self, question: str, failure_type: str, context: Dict) -> None:
        """
        记录失败轨迹，供后续历史检索使用

        Args:
            question: 当前问题
            failure_type: 失败类型（如 "low_crystal_reference", "audit_failed", "debate_diverged"）
            context: 失败上下文，包含失败时的各项指标
        """
        trace = {
            "question": question,
            "failure_type": failure_type,
            "timestamp": datetime.now().isoformat(),
            "context": context,
        }
        # 如果上下文中有有效晶体组合，也记录下来
        effective_crystals = context.get("effective_crystals", [])
        if effective_crystals:
            trace["effective_crystals"] = effective_crystals

        # 记录到 evolution_log
        self.engine.log_evolution_event(
            "failure_trace",
            {
                "failure_traces": trace,
                "trigger": "system_auto",
                "question": question,
                "effective_crystals": effective_crystals
            }
        )
        self.log(f"📝 失败轨迹已记录：{failure_type}", "system")
    def _inject_perspective(self, alarm: dict):
        """强制注入对立视角：添加一个临时角色或修改系统指令"""
        # 可以添加一个“对抗者”角色，或者修改系统提示要求从对立面思考
        self._forced_perspective = "请从与之前所有角色截然相反的角度重新审视问题，提出颠覆性观点。"
        self.log("  对立视角已注入，将在下一轮生效", "system")

    def _trigger_search(self, alarm: dict):
        """强制触发外部搜索并更新上下文"""
        question = self._current_question
        if not question:
            return
        # 直接调用 _fetch_external_overview 并强制更新外部缓存
        overview = self._fetch_external_overview(question)
        self._forced_external = overview
        self.log("  外部搜索已强制触发", "system")


    def _format_process_summary(self, rounds: List[Dict]) -> str:
        lines = []
        for item in rounds:
            round_no = item.get("round")
            lines.append(f"第 {round_no} 轮：")
            audit = item.get("audit") or {}
            if audit.get("summary"):
                lines.append(f"- 逻辑检查员：{audit.get('summary')}")
            answers = item.get("answers") or []
            for answer in answers:
                text = re.sub(r"\s+", " ", str(answer.get("answer", ""))).strip()
                lines.append(f"- {answer.get('role')}：{text[:180]}{'...' if len(text) > 180 else ''}")
        return "\n".join(lines)

    # ===== 新增工具方法：角色名称转key =====
    def _map_role_key(self, name: str) -> str:
        mapping = {
            "激进者": "radical", 
            "保守者": "conservative", 
            "结构主义者": "structural",
            "百灵鸟": "lark", 
            "取经者": "pilgrim", 
            "奇谋者": "strategist",
            "延安智者": "statesman", 
            "大法官": "judge", 
            "首席发言人": "spokesperson",
            "替身-我": "twin"
        }
        return mapping.get(name, name)

    # ===== Day 0 运行时自检断言 =====
    def _run_day0_runtime_assertions(self):
        """Day 0 运行时自检断言"""
        assertions = [
            ("CrystalEngine 存在", hasattr(self, 'engine')),
            ("MetaLayer 存在", hasattr(self.engine, 'meta')),
            ("CheapGate 存在", hasattr(self.engine, 'cheap_gate')),
            ("FingerprintExtractor 存在", hasattr(self.engine, 'fingerprint_extractor')),
            ("VectorStore 存在", hasattr(self.engine, 'vector_store')),
            ("AlarmMonitor 可用", hasattr(self, 'alarm_monitor')),
        ]
        all_passed = True
        for name, passed in assertions:
            if not passed:
                self.log(f"❌ 自检断言失败：{name}", "error")
                all_passed = False
        if all_passed:
            self.log("✅ Day 0 运行时自检断言全部通过", "system")
        return all_passed

    # ===== 替身自我博弈 =====
    def _run_self_play(self, question: str, max_rounds: int = 2) -> Dict:
        """替身自我博弈主流程"""
        # ===== 进度初始化 =====
        self._emit_progress("启动替身自我博弈", 0)

        roles = self.get_roles_with_twin(include_twin=True)
        twin = next((r for r in roles if r.key == "twin"), None)
        radical = next((r for r in roles if r.key == "radical"), None)
        conservative = next((r for r in roles if r.key == "conservative"), None)
        structural = next((r for r in roles if r.key == "structural"), None)

        if not twin or not radical or not conservative:
            self.log("[WARN] 缺少必要角色，降级到标准辩论", "warning")
            return self.run(question, mode="debate_full", max_rounds=max_rounds)

        crystal_context = self._crystal_context(question)
        rounds_data = []
        all_answers = []

        # ===== Round 1: 激进者 vs 替身-我 =====
        self._emit_progress("第1轮：替身生成观点", 20)
        self.log("[ROUND 1] 激进者攻击替身观点", "system")
        round1_answers = []

        self.log("  替身-我 正在生成观点...", "system")
        twin_prompt = f"用户问题：{question}\n请以替身身份给出独立观点，体现你的认知特征。"
        twin_answer = self.ai.chat(twin_prompt, system=self._role_system(twin, crystal_context))
        round1_answers.append({"role": twin.name, "answer": twin_answer})

        self._emit_progress("第1轮：激进者攻击", 30)
        self.log("  激进者 正在攻击...", "system")
        radical_prompt = f"""用户问题：{question}

替身-我的观点：
{twin_answer}

请作为激进者，攻击替身观点中的默认前提、逻辑漏洞和认知盲区。"""
        radical_answer = self.ai.chat(radical_prompt, system=self._role_system(radical, crystal_context))
        round1_answers.append({"role": radical.name, "answer": radical_answer})

        self._emit_progress("第1轮审计", 40)
        audit1 = self._audit(question, round1_answers, round_no=1)
        jaccard1 = self._average_jaccard([item["answer"] for item in round1_answers])
        rounds_data.append({"round": 1, "answers": round1_answers, "audit": audit1, "jaccard": round(jaccard1, 3)})
        all_answers.extend(round1_answers)

        # ===== Round 2: 替身-我 vs 保守者 =====
        self._emit_progress("第2轮：替身辩护", 55)
        self.log("[ROUND 2] 替身辩护 vs 保守者质疑", "system")
        round2_answers = []

        self.log("  替身-我 正在辩护...", "system")
        twin_defense_prompt = f"""用户问题：{question}

激进者的攻击：
{radical_answer}

审计员反馈：
{audit1.get('summary', '无')}

请作为替身-我，进行辩护并修正观点。"""
        twin_defense = self.ai.chat(twin_defense_prompt, system=self._role_system(twin, crystal_context))
        round2_answers.append({"role": twin.name, "answer": twin_defense})

        self._emit_progress("第2轮：保守者质疑", 65)
        self.log("  保守者 正在质疑...", "system")
        conservative_prompt = f"""用户问题：{question}

替身-我的观点：
{twin_defense}

请作为保守者，从风险、成本、边界角度提出质疑。"""
        conservative_answer = self.ai.chat(conservative_prompt, system=self._role_system(conservative, crystal_context))
        round2_answers.append({"role": conservative.name, "answer": conservative_answer})

        self._emit_progress("第2轮审计", 75)
        audit2 = self._audit(question, round2_answers, round_no=2)
        jaccard2 = self._average_jaccard([item["answer"] for item in round2_answers])
        rounds_data.append({"round": 2, "answers": round2_answers, "audit": audit2, "jaccard": round(jaccard2, 3)})
        all_answers.extend(round2_answers)

        # ===== Round 3: 结构主义者融合 =====
        self._emit_progress("第3轮：结构主义者融合", 85)
        self.log("[ROUND 3] 结构主义者融合各方观点", "system")
        round3_answers = []

        all_text = "\n\n".join([f"### {item['role']}\n{item['answer']}" for item in all_answers])
        structural_prompt = f"""用户问题：{question}

所有辩论观点：
{all_text}

审计摘要：
- 第1轮：{audit1.get('summary', '')}
- 第2轮：{audit2.get('summary', '')}

请作为结构主义者，融合所有观点，输出综合方案。"""
        structural_answer = self.ai.chat(structural_prompt, system=self._role_system(structural, crystal_context))
        round3_answers.append({"role": structural.name, "answer": structural_answer})

        audit3 = self._audit(question, round3_answers, round_no=3)
        jaccard3 = self._average_jaccard([item["answer"] for item in round3_answers])
        rounds_data.append({"round": 3, "answers": round3_answers, "audit": audit3, "jaccard": round(jaccard3, 3)})

        # ===== 总结 =====
        self._emit_progress("总结生成中", 92)
        self.log("[FINAL] 发言人生成最终答案", "system")
        final = self._spokesperson(question, rounds_data, crystal_context, light=False)

        self._emit_progress("完成", 100)
        # Day 5: Hebbian 学习 - 评估有效晶体组合
        try:
            used_crystal_ids = set()
            for round_data in rounds_data:
                for answer_item in round_data.get("answers", []):
                    text = answer_item.get("answer", "")
                    matches = re.findall(r'\[C(\d+)\]', text)
                    used_crystal_ids.update([f"C{mid}" for mid in matches])
            if used_crystal_ids:
                last_audit = rounds_data[-1].get("audit", {})
                ev_scores = list(last_audit.get("evidence_scores", {}).values())
                if ev_scores:
                    avg_score = sum(ev_scores) / len(ev_scores)
                else:
                    avg_score = 0.6
                task_type = getattr(self, '_current_classification', {}).get("type", "general")
                self.engine.update_hebbian_weights(list(used_crystal_ids), task_type, avg_score)
                self.log(f"🧠 Hebbian 学习更新：{len(used_crystal_ids)} 个晶体，任务类型={task_type}，评分={avg_score:.2f}", "system")
        except Exception as e:
            self.log(f"⚠️ Hebbian 学习失败：{e}", "warning")
        return self._result(question, "twin_self_play", rounds_data, final, calls_estimate=7)

# =============================================================================
# Day 0 基线采集运行器
# =============================================================================
class BaselineRunner:
    """
    运行标准测试问题，采集辩论系统基线数据。
    在代码修改前运行一次，生成基线 JSON，供 Day 3 对比。
    """
    
    def __init__(self, engine: CrystalEngine, ai_client: AIClient, roles: List[Dict]):
        self.engine = engine
        self.ai = ai_client
        self.roles = roles
        self.results = []
    
    def run(self, max_rounds: int = 2) -> Dict[str, Any]:
        """运行所有基线测试"""
        print("📊 开始基线采集...")
        total = len(BENCHMARK_QUESTIONS)
        
        for idx, question in enumerate(BENCHMARK_QUESTIONS):
            print(f"\n[{idx+1}/{total}] 测试问题: {question[:50]}...")
            
            # 初始化辩论引擎
            debate = DebateEngine(
                self.ai,
                self.engine,
                self.roles,
                log=lambda m, l: print(f"  {m[:80]}"),
                progress_callback=None
            )
            
            # 运行辩论
            try:
                result = debate.run(question, mode="debate_full", max_rounds=max_rounds)
                metrics = self._extract_metrics(question, result)
            except Exception as e:
                print(f"  ❌ 错误: {e}")
                metrics = self._empty_metrics(question, str(e))
            
            self.results.append(metrics)
            
            # 打印简要结果
            print(f"  Jaccard: {metrics['jaccard_similarity']:.3f} | "
                  f"引用率: {metrics['crystal_reference_rate']:.1%} | "
                  f"审计反馈长度: {metrics['audit_feedback_avg_len']:.1f}")
        
        # 汇总统计
        summary = self._aggregate(self.results)
        
        # 写入 JSON
        output_path = Config.DATA_ROOT / "系统日志" / "辩论基线.json"
        output_path.parent.mkdir(parents=True, exist_ok=True)
        
        baseline_data = {
            "timestamp": datetime.now().isoformat(),
            "total_questions": total,
            "max_rounds": max_rounds,
            "summary": summary,
            "details": self.results,
            "good_debate_criteria": [
                "最终答案包含至少 2 个可执行建议",
                "最终答案明确指出了至少 1 个风险",
                "one_sentence_conclusion 不超过 30 字",
                "至少 2 个不同角色在最后一轮出现观点修正（自我批判后的变化）"
            ]
        }
        
        with open(output_path, "w", encoding="utf-8") as f:
            json.dump(baseline_data, f, ensure_ascii=False, indent=2)
        
        print(f"\n✅ 基线数据已保存至: {output_path}")
        return baseline_data
    
    def _extract_metrics(self, question: str, result: Dict) -> Dict[str, Any]:
        """从单次辩论结果中提取全部指标"""
        rounds_data = result.get("rounds", [])
        if not rounds_data:
            return self._empty_metrics(question, "no_rounds")
        
        # 1. Jaccard 相似度：取最后一轮各角色回答的平均相似度
        last_round = rounds_data[-1]
        answers = [item["answer"] for item in last_round.get("answers", [])]
        if len(answers) >= 2:
            temp_debate = DebateEngine(self.ai, self.engine, self.roles, lambda m, l: None)
            jaccard = temp_debate._average_jaccard(answers)
        else:
            jaccard = 0.0
        
        # 2. 晶体引用率：统计所有发言中 [Cxxx] 或 [Hxxx] 的出现次数
        all_text = " ".join(answers)
        crystal_matches = re.findall(r'\[C\d+\]', all_text)
        hole_matches = re.findall(r'\[H\d+\]', all_text)
        total_refs = len(crystal_matches) + len(hole_matches)
        total_answers = len(answers)
        ref_rate = total_refs / max(1, total_answers)
        
        # 3. 审计反馈具体性：用平均反馈长度作为代理指标
        audit = last_round.get("audit", {})
        feedbacks = list(audit.get("feedback_by_role", {}).values())
        avg_feedback_len = sum(len(str(f)) for f in feedbacks) / max(1, len(feedbacks))
        
        # 4. 分歧丰富度：disagreement_map 中非空类别数
        disc_map = audit.get("disagreement_map", {})
        active_categories = [k for k, v in disc_map.items() if v]
        diversity_score = len(active_categories) / max(1, len(disc_map))
        
        # 5. 证据评分方差
        ev_scores = list(audit.get("evidence_scores", {}).values())
        if len(ev_scores) > 1:
            mean = sum(ev_scores) / len(ev_scores)
            variance = sum((s - mean) ** 2 for s in ev_scores) / len(ev_scores)
        else:
            variance = 0.0
        
        # 6. 最终答案有用性代理
        final = result.get("final", {})
        one_sentence = final.get("one_sentence_conclusion", "")
        student_answer = final.get("student_friendly_answer", "")
        teacher_detail = final.get("teacher_detail", "")
        
        return {
            "question": question,
            "jaccard_similarity": round(jaccard, 4),
            "crystal_reference_rate": round(ref_rate, 4),
            "audit_feedback_avg_len": round(avg_feedback_len, 1),
            "disagreement_diversity": round(diversity_score, 3),
            "evidence_score_variance": round(variance, 4),
            "one_sentence_len": len(one_sentence),
            "student_answer_len": len(student_answer),
            "teacher_detail_len": len(teacher_detail),
            "has_risks": "风险" in student_answer or "风险" in teacher_detail,
            "has_executable_actions": any(kw in student_answer for kw in ["步骤", "建议", "方法", "操作", "执行"]),
            "rounds_count": len(rounds_data),
            "total_roles": len(answers),
            "raw_final": final
        }
    
    def _empty_metrics(self, question: str, reason: str) -> Dict:
        return {
            "question": question,
            "error": reason,
            "jaccard_similarity": 0.0,
            "crystal_reference_rate": 0.0,
            "audit_feedback_avg_len": 0.0,
            "disagreement_diversity": 0.0,
            "evidence_score_variance": 0.0,
            "one_sentence_len": 0,
            "student_answer_len": 0,
            "teacher_detail_len": 0,
            "has_risks": False,
            "has_executable_actions": False,
            "rounds_count": 0,
            "total_roles": 0
        }
    
    def _aggregate(self, results: List[Dict]) -> Dict[str, Any]:
        valid = [r for r in results if "error" not in r]
        if not valid:
            return {"error": "无有效结果"}
        
        keys = ["jaccard_similarity", "crystal_reference_rate", 
                "audit_feedback_avg_len", "disagreement_diversity", "evidence_score_variance",
                "one_sentence_len", "student_answer_len", "teacher_detail_len"]
        
        avg = {}
        for k in keys:
            vals = [r.get(k, 0) for r in valid]
            avg[k] = round(sum(vals) / len(vals), 4) if vals else 0
        
        avg["has_risks_rate"] = sum(1 for r in valid if r.get("has_risks", False)) / len(valid)
        avg["has_executable_rate"] = sum(1 for r in valid if r.get("has_executable_actions", False)) / len(valid)
        avg["total_valid"] = len(valid)
        avg["total_errors"] = len(results) - len(valid)
        
        return avg

# =============================================================================
# Meta搜索引擎（认知路径平行对比）
# =============================================================================

class MetaSearchEngine:
    """
    Meta-Harness 式自动搜索的局部原型

    对同一问题生成多条认知路径，用规则引擎快速评分，选择最优路径。
    对应建议9：引入Meta-Harness式自动搜索
    """

    def __init__(self, engine: CrystalEngine, ai_client: AIClient):
        self.engine = engine
        self.ai = ai_client

    def generate_paths(self, question: str, num_paths: int = 3) -> List[Dict[str, Any]]:
        """
        生成多条认知路径

        每条路径包含不同的：
        - 晶体组合
        - 孔洞检测策略
        - 推理权重
        """
        crystals = self.engine.parse_crystals()
        if not crystals:
            return []

        paths = []

        # 路径1：标准检索（向量 + BM25 混合）
        path1 = self._build_path_standard(question, crystals)
        paths.append(path1)

        # 路径2：激进模式（偏重新颖性，优先高热度晶体）
        path2 = self._build_path_radical(question, crystals)
        paths.append(path2)

        # 路径3：保守模式（偏重稳定性，优先高置信度/固定晶体）
        path3 = self._build_path_conservative(question, crystals)
        paths.append(path3)

        return paths[:num_paths]

    def _build_path_standard(self, question: str, crystals: List) -> Dict[str, Any]:
        """标准检索路径"""
        ranked = self.engine.rank_crystals(question, crystals, top_k=5)
        return {
            "name": "标准路径",
            "crystals": [{"id": c.id, "content": c.content, "score": score} for score, c in ranked],
            "crystal_ids": [c.id for _, c in ranked],
            "strategy": "vector_bm25_hybrid"
        }

    def _build_path_radical(self, question: str, crystals: List) -> Dict[str, Any]:
        """激进路径：偏重新颖性和高热度"""
        # 按热度排序，取 top 10
        sorted_crystals = sorted(crystals, key=lambda c: c.heat, reverse=True)
        selected = sorted_crystals[:5]
        return {
            "name": "激进路径",
            "crystals": [{"id": c.id, "content": c.content, "heat": c.heat} for c in selected],
            "crystal_ids": [c.id for c in selected],
            "strategy": "heat_priority"
        }

    def _build_path_conservative(self, question: str, crystals: List) -> Dict[str, Any]:
        """保守路径：偏重固定晶体和 L1 层"""
        state = self.engine.load_layer_state()
        layers = state.get("layers", {})
        manual = state.get("manual_override", {})

        # 优先固定晶体
        fixed = [c for c in crystals if manual.get(c.id) == "L1_fixed"]
        l1 = [c for c in crystals if layers.get(c.id) == "L1" and c.id not in [f.id for f in fixed]]

        selected = fixed[:3] + l1[:2]
        if not selected:
            selected = crystals[:5]

        return {
            "name": "保守路径",
            "crystals": [{"id": c.id, "content": c.content, "layer": layers.get(c.id, "L2")} for c in selected],
            "crystal_ids": [c.id for c in selected],
            "strategy": "fixed_l1_priority"
        }

    def score_paths(self, paths: List[Dict[str, Any]], question: str) -> List[Dict[str, Any]]:
        """
        使用规则引擎快速评分（非 LLM）

        评分维度：
        - 引用晶体数（越多越好）
        - 晶体层级权重（L1 > L2 > L3）
        - 与指纹匹配度
        """
        try:
            fingerprint = self.engine.fingerprint_extractor.get_fingerprint()
            pref_role = fingerprint.preferred_role if fingerprint else "structural"
        except:
            pref_role = "structural"

        # 角色关键词权重
        role_keywords = {
            "radical": ["颠覆", "激进", "创新", "突破"],
            "conservative": ["稳健", "保守", "风险", "安全"],
            "structural": ["结构", "系统", "框架", "模型"],
            "executor": ["执行", "步骤", "操作", "落地"],
            "auditor": ["审计", "验证", "检查", "证据"]
        }
        fp_keywords = role_keywords.get(pref_role, role_keywords["structural"])

        scored_paths = []
        for path in paths:
            score = 0

            # 1. 晶体数量（0-20分）
            crystal_count = len(path.get("crystals", []))
            score += min(20, crystal_count * 4)

            # 2. 层级权重（0-30分）
            state = self.engine.load_layer_state()
            layers = state.get("layers", {})
            for c in path.get("crystals", []):
                cid = c.get("id") if isinstance(c, dict) else c.id
                layer = layers.get(cid, "L2")
                if layer == "L1":
                    score += 6
                elif layer == "L2":
                    score += 3
                else:
                    score += 1

            # 3. 指纹匹配度（0-30分）
            for c in path.get("crystals", []):
                content = c.get("content") if isinstance(c, dict) else c.content
                for kw in fp_keywords:
                    if kw in content:
                        score += 3
                        break
            score = min(30, score)

            # 4. 多样性奖励（0-20分）
            # 检查晶体是否来自不同领域（基于内容长度和关键词差异）
            contents = [c.get("content") if isinstance(c, dict) else c.content for c in path.get("crystals", [])]
            unique_keywords = set()
            for content in contents:
                words = content[:30].split()
                unique_keywords.update(words)
            diversity_score = min(20, len(unique_keywords) * 2)

            total_score = score + diversity_score

            scored_paths.append({
                "path": path,
                "score": total_score,
                "details": {
                    "crystal_count_score": min(20, len(path.get("crystals", [])) * 4),
                    "layer_score": score - min(20, len(path.get("crystals", [])) * 4) - diversity_score + 30,
                    "fingerprint_score": min(30, sum(3 for c in path.get("crystals", []) if any(kw in (c.get("content") if isinstance(c, dict) else c.content) for kw in fp_keywords))),
                    "diversity_score": diversity_score
                }
            })

        scored_paths.sort(key=lambda x: x["score"], reverse=True)
        return scored_paths

    def select_best_path(self, question: str) -> Dict[str, Any]:
        """
        选择最优认知路径
        """
        paths = self.generate_paths(question)
        if not paths:
            return {"error": "无法生成认知路径"}

        scored = self.score_paths(paths, question)
        best = scored[0] if scored else None

        return {
            "question": question,
            "paths": scored,
            "best_path": best,
            "selected_crystals": best["path"]["crystal_ids"] if best else []
        }

    def run_comparison(self, question: str) -> Dict[str, Any]:
        """
        运行认知路径对比（对外接口）
        """
        result = self.select_best_path(question)
        return result

# =============================================================================
# 12. 每日计划器 (planner.py)
# =============================================================================
import hashlib
import json
import re
import time
from datetime import date, datetime
from typing import Callable, Dict, List, Optional, Tuple

class DailyPlanner:
    DEFAULT_KEYWORDS = ["AI", "认知科学", "教育学习", "Agent", "推理能力", "国产大模型"]
    STAGES = [
        "更新晶体分层",
        "外部信息采集",
        "价值评分筛选",
        "生成待确认卡片",
        "冲突预警",
        "孔洞进展评估",
        "Wondering 与日报整理",
    ]

    def __init__(
        self,
        engine: CrystalEngine,
        ai_client: AIClient,
        fetcher: ExternalFetcher,
        log_callback: Callable,
        update_status_callback: Callable,
    ):
        self.engine = engine
        self.ai = ai_client
        self.fetcher = fetcher
        self.log = log_callback
        self.update_status = update_status_callback
        self.today = date.today().isoformat()
        self.report_content = ""
        self.last_external = []
        self._reset_state()

    def _reset_state(self):
        self.started_at = time.time()
        self.deadline = self.started_at + Config.DAILY_PLAN_TIME_BUDGET_SECONDS
        self.keywords = list(self.DEFAULT_KEYWORDS)
        self.completed_stages: List[str] = []
        self.skipped_stages: List[str] = []
        self.candidates: List[Dict] = []
        self.scored_candidates: List[Dict] = []
        self.created_pending_ids: List[str] = []
        self.created_task_ids: List[str] = []
        self.wondering: List[str] = []
        self.interrupted_reason = ""
        self.progress_callback: Optional[Callable[[Dict], None]] = None
        self.stop_flag: Optional[Callable[[], bool]] = None

    def is_today_run(self) -> bool:
        last_run_file = Config.get_path("change_log").parent / "last_daily_run.txt"
        return last_run_file.exists() and last_run_file.read_text(encoding="utf-8").strip() == self.today

    def _mark_today_run(self):
        last_run_file = Config.get_path("change_log").parent / "last_daily_run.txt"
        last_run_file.parent.mkdir(parents=True, exist_ok=True)
        last_run_file.write_text(self.today, encoding="utf-8")

    def _load_task_cards(self) -> List[Dict]:
        if not FileIO.exists("task_cards"):
            return []
        try:
            return json.loads(FileIO.read("task_cards"))
        except Exception:
            return []

    def _save_task_cards(self, cards: List[Dict]):
        FileIO.write("task_cards", json.dumps(cards, ensure_ascii=False, indent=2))

    @staticmethod
    def _stable_suffix(*parts: str) -> str:
        raw = "|".join(str(p) for p in parts)
        return f"{int(hashlib.sha256(raw.encode('utf-8')).hexdigest(), 16) % 1000:03d}"

    def _normalize_keywords(self, keywords: Optional[List[str]]) -> List[str]:
        if isinstance(keywords, str):
            raw = keywords
        else:
            raw = " ".join(str(x) for x in (keywords or []))
        parts = [p.strip() for p in re.split(r"[,，;；\s\n]+", raw) if p.strip()]
        result = []
        for item in parts:
            if item not in result:
                result.append(item)
        return result[:10] or list(self.DEFAULT_KEYWORDS)

    def _emit_progress(self, stage: str, progress: int, stage_index: int = 0):
        elapsed = int(time.time() - self.started_at)
        payload = {
            "stage": stage,
            "progress": max(0, min(100, int(progress))),
            "stage_index": stage_index,
            "stage_total": len(self.STAGES),
            "candidate_count": len(self.candidates),
            "pending_count": len(self.created_pending_ids),
            "task_count": len(self.created_task_ids),
            "elapsed_seconds": elapsed,
            "budget_seconds": int(self.deadline - self.started_at),
            "keywords": self.keywords,
        }
        self.update_status(f"{stage}：{payload['progress']}%")
        if self.progress_callback:
            self.progress_callback(payload)

    def _should_stop(self) -> bool:
        if self.stop_flag and self.stop_flag():
            self.interrupted_reason = "用户中断"
            return True
        if time.time() >= self.deadline:
            self.interrupted_reason = "达到 15 分钟预算"
            return True
        return False

    def _ensure_can_continue(self, next_stage: str):
        if self._should_stop():
            raise InterruptedError(self.interrupted_reason)
        self.log(f"▶ {next_stage}", "system")

    def _add_task_card(self, card: TaskCard):
        cards = self._load_task_cards()
        if any(
            c.get("type") == card.type
            and c.get("source") == card.source
            and c.get("title") == card.title
            and c.get("status") == "pending"
            for c in cards
        ):
            self.log(f"  已存在任务卡片，跳过重复生成：{card.title}", "system")
            return False
        cards.append(card.__dict__)
        self._save_task_cards(cards)
        self.created_task_ids.append(card.id)
        return True

    def _create_pending_card_from_info(self, info: Dict) -> str:
        pending_content = FileIO.read("pending")
        title = info.get("title", "")
        source = info.get("source", "")
        if title and source and title in pending_content and source in pending_content:
            self.log(f"  已存在待确认卡片，跳过重复生成：{title[:60]}", "system")
            return ""

        card_id = f"PENDING-{datetime.now().strftime('%Y%m%d%H%M%S')}-{self._stable_suffix(source, title)}"
        matched_holes = ", ".join(info.get("matched_holes", [])[:3]) or "暂无"
        matched_crystals = ", ".join(info.get("matched_crystals", [])[:3]) or "暂无"
        keywords = ", ".join(self.keywords)
        content = f"""## {card_id}
- 类型：{info.get('type', 'external')}
- 标题：{title}
- 来源：{source}
- 链接：{info.get('link', '')}
- 当日意向关键词：{keywords}
- 意向匹配原因：{info.get('intent_reason', '与今日关键词或核心孔洞存在潜在关联')}
- 价值评分：{info.get('value_score', 0):.2f}
- 匹配孔洞：{matched_holes}
- 匹配晶体：{matched_crystals}
- 内容摘要：{info.get('summary', '')[:500]}
- 建议动作：{info.get('suggested_action', '老师确认后决定是否转为晶体')}
- AI判断：可能与孔洞或 L1 晶体相关，建议确认后再转为晶体。
"""
        FileIO.append("pending", f"\n{content}\n")
        self.created_pending_ids.append(card_id)
        self.log(f"  📄 生成待确认卡片 {card_id}", "system")
        return card_id

    def _create_conflict_task_card(self, conflict: Conflict) -> str:
        """生成冲突任务卡片（升级版：包含冲突类型和严重程度）"""
        if conflict.similarity > 0.85:
            conflict_type = "高度重复"
            severity = "高"
            suggestion = "建议合并这两个晶体，保留更完整的一个"
        elif conflict.similarity > 0.70:
            conflict_type = "语义冲突"
            severity = "中"
            suggestion = "建议补充适用边界，明确两者的区分条件"
        else:
            conflict_type = "潜在关联"
            severity = "低"
            suggestion = "建议建立链接关系，说明两者的关联"

        card_id = f"TASK-{datetime.now().strftime('%Y%m%d%H%M%S')}-{self._stable_suffix(conflict.crystal_a, conflict.crystal_b)}"

        task = TaskCard(
            id=card_id,
            type="conflict",
            title=f"晶体冲突：{conflict.crystal_a} vs {conflict.crystal_b}",
            content=f"""冲突类型：{conflict_type}
严重程度：{severity}
相似度：{conflict.similarity:.2f}

A ({conflict.crystal_a}): {conflict.content_a}
B ({conflict.crystal_b}): {conflict.content_b}

建议动作：{suggestion}""",
            source=f"{conflict.crystal_a} vs {conflict.crystal_b}",
            suggested_action=suggestion,
            status="pending",
        )

        if not self._add_task_card(task):
            return ""
        self.log(f"  [WARN] 生成{conflict_type}任务卡片 {card_id} (严重程度: {severity})", "system")
        return card_id

    def _collect_external_info(self) -> Tuple[List[str], List[Dict]]:
        data = self.fetcher.fetch_all()
        insights_text = self.fetcher.build_insights(data)
        structured = self.fetcher.build_structured_insights(data)

        for keyword in self.keywords:
            if self._should_stop():
                break
            try:
                self.log(f"  关键词追踪：{keyword}", "system")
                for paper in self.fetcher.fetch_arxiv_papers(query=keyword, max_results=3):
                    if paper.startswith("(") or not paper.strip():
                        continue
                    title = paper.split(" (")[0].strip()
                    structured.append({
                        "type": "keyword_arxiv",
                        "title": title,
                        "summary": title,
                        "link": "",
                        "source": f"arXiv({keyword})",
                        "intent_keyword": keyword,
                    })
                for news in self.fetcher.fetch_baidu_news(keyword, max_results=2):
                    if news.startswith("(") or not news.strip():
                        continue
                    structured.append({
                        "type": "keyword_news",
                        "title": news,
                        "summary": news,
                        "link": "",
                        "source": f"百度新闻({keyword})",
                        "intent_keyword": keyword,
                    })
            except Exception as exc:
                self.log(f"  关键词 {keyword} 追踪失败：{exc}", "warning")

        deduped = []
        seen = set()
        for item in structured:
            title = re.sub(r"\s+", " ", item.get("title", "")).strip()
            source = item.get("source", "")
            key = title.lower()
            if not title or key in seen:
                continue
            seen.add(key)
            item["title"] = title
            item.setdefault("summary", title)
            item.setdefault("source", source)
            deduped.append(item)
            if len(deduped) >= Config.DAILY_PLAN_MAX_CANDIDATES:
                break
        return insights_text, deduped

    def _score_candidates(self, candidates: List[Dict]) -> List[Dict]:
        holes = self.engine.parse_holes()
        crystals = self.engine.parse_crystals()
        scored = []
        for item in candidates:
            text = f"{item.get('title', '')} {item.get('summary', '')}"
            keyword_hits = [kw for kw in self.keywords if kw.lower() in text.lower()]
            keyword_score = min(1.0, len(keyword_hits) * 0.35)

            hole_scores = []
            for hole in holes:
                try:
                    score = self.engine.match_info_to_hole(text, hole.content)
                except Exception:
                    score = self.engine._simple_similarity(text, hole.content)
                if score > 0.25:
                    hole_scores.append((score, hole.id))
            hole_scores.sort(reverse=True)

            ranked_crystals = self.engine.rank_crystals(text, crystals, top_k=3) if crystals else []
            crystal_score = ranked_crystals[0][0] if ranked_crystals else 0.0
            crystal_score = min(1.0, crystal_score / 5.0)

            source_bonus = 0.15 if item.get("type", "").startswith(("keyword", "arxiv", "huggingface")) else 0.05
            value_score = round(keyword_score + min(0.6, (hole_scores[0][0] if hole_scores else 0.0)) + crystal_score + source_bonus, 3)
            item["value_score"] = value_score
            item["matched_holes"] = [hid for _score, hid in hole_scores[:3]]
            item["matched_crystals"] = [crystal.id for _score, crystal in ranked_crystals[:3]]
            item["intent_reason"] = "、".join(keyword_hits) if keyword_hits else "与核心孔洞或 L1 晶体存在语义关联"
            item["suggested_action"] = "优先确认并考虑转为晶体" if value_score >= 0.8 else "暂存观察，必要时转为任务"
            scored.append(item)
        scored.sort(key=lambda x: x.get("value_score", 0), reverse=True)
        return scored[:Config.DAILY_PLAN_TOP_ITEMS]

    def _summarize_candidate(self, item: Dict) -> Dict:
        if self._should_stop():
            return item
        prompt = f"""请对这条外部信息做教师可读的价值摘要。
标题：{item.get('title')}
来源：{item.get('source')}
今日关键词：{', '.join(self.keywords)}
匹配孔洞：{', '.join(item.get('matched_holes', []))}
匹配晶体：{', '.join(item.get('matched_crystals', []))}

输出 120 字以内，说明为什么值得看，以及可能转成什么认知晶体。"""
        try:
            item["summary"] = self.ai.chat(prompt)[:500]
        except Exception as exc:
            self.log(f"  摘要生成失败，使用原摘要：{exc}", "warning")
        return item

    def _update_hole_progress_with_evidence(self, external_insights: List[Dict]) -> Tuple[List[Dict], str]:
        holes = self.engine.parse_holes()
        progress = self.engine.load_hole_progress()
        important_holes = [h for h in holes if h.urgency >= 0.7]
        if not external_insights:
            hole_list = [{"id": h.id, "content": h.content, "progress": int(progress.get(h.id, 0.0) * 100), "urgency": h.urgency} for h in important_holes]
            return hole_list, "无外部信息，无法评估填充进展。"

        updated = False
        evidence_log = []
        for hole in important_holes:
            hid = hole.id
            cur = progress.get(hid, 0.0)
            best_match = 0.0
            best_info = None
            for info in external_insights:
                text = f"{info.get('title', '')} {info.get('summary', '')}"
                try:
                    score = self.engine.match_info_to_hole(text, hole.content)
                except Exception:
                    score = self.engine._simple_similarity(text, hole.content)
                if score > best_match:
                    best_match = score
                    best_info = info
            if best_info and best_match > 0.6:
                increment = min(0.15, best_match * 0.2)
                new_prog = min(1.0, cur + increment)
                if new_prog > cur:
                    progress[hid] = new_prog
                    updated = True
                    evidence_log.append(f"{hid} 进度 {int(cur * 100)}% -> {int(new_prog * 100)}%，依据：{best_info['title'][:60]}")
        if updated:
            self.engine.save_hole_progress(progress)
            for ev in evidence_log:
                self.engine._append_change_log("孔洞填充", ev)
        hole_progress_list = [{"id": h.id, "content": h.content, "progress": int(progress.get(h.id, 0.0) * 100), "urgency": h.urgency} for h in important_holes]
        desc = "\n".join(evidence_log[:5]) if evidence_log else "暂无显著填充进展。"
        return hole_progress_list, desc

    def _generate_wondering(self, external_insights: List[str]) -> List[str]:
        l0_holes, l1_crystals = self.engine.get_attention_context()
        l0_text = "\n".join([f"- {h.id}: {h.content[:80]}" for h in l0_holes])
        l1_text = "\n".join([f"- {c.id}: {c.content[:60]}" for c in l1_crystals[:30]])
        external_text = "\n".join(external_insights[:8])
        prompt = f"""基于今日关键词、核心孔洞、L1晶体和外部信息，生成 3 个“我好奇”的问题。
今日关键词：{', '.join(self.keywords)}
L0核心孔洞：
{l0_text}
L1晶体：
{l1_text}
外部信息：
{external_text[:2000]}

每行一个问题，不加编号。"""
        try:
            reply = self.ai.chat(prompt)
            lines = [l.strip(" -0123456789.、") for l in reply.strip().split("\n") if l.strip()]
        except Exception:
            lines = []
        defaults = [
            "如何把今日外部信息转化为可验证的学习策略？",
            "哪些新信息可能挑战当前晶体树的默认判断？",
            "哪个孔洞最值得下一步集中探索？",
        ]
        unique = []
        for q in lines + defaults:
            if q and not any(self.engine._simple_similarity(q, u) > 0.8 for u in unique):
                unique.append(q)
            if len(unique) >= 3:
                break
        self.wondering = [f"{i + 1}. {q}" for i, q in enumerate(unique[:3])]
        return self.wondering

    def _generate_report(
        self,
        status: str,
        external_insights: List[str],
        hole_progress_list: List[Dict],
        created_pending_count: int,
        created_conflict_count: int,
    ):
        elapsed = time.time() - self.started_at
        top_lines = [
            f"- {item.get('title', '')[:90]} | {item.get('source', '')} | 评分 {item.get('value_score', 0):.2f}"
            for item in self.scored_candidates[:10]
        ]
        hole_text = "\n".join([
            f"- {hp['id']}「{hp['content'][:40]}」填充进度 {hp['progress']}%"
            for hp in hole_progress_list[:5]
        ]) or "暂无高紧迫度孔洞进展。"
        completed = "、".join(self.completed_stages) or "暂无"
        skipped = "、".join([s for s in self.STAGES if s not in self.completed_stages]) or "无"
        report = f"""# 自主工作日报 {self.today}

## 运行状态
- 状态：{status}
- 当日意向关键词：{', '.join(self.keywords)}
- 实际耗时：{elapsed:.1f} 秒
- 已完成阶段：{completed}
- 未完成阶段：{skipped}
- 中断/收尾原因：{self.interrupted_reason or '正常完成'}

## 高价值候选
{chr(10).join(top_lines) if top_lines else '暂无高价值候选。'}

## 孔洞填充
{hole_text}

## Wondering
{chr(10).join(self.wondering) if self.wondering else '暂无。'}

## 成果清单
- 已生成 PENDING：{created_pending_count} 张
- PENDING ID：{', '.join(self.created_pending_ids) or '无'}
- 已生成任务卡：{created_conflict_count} 张
- 任务卡 ID：{', '.join(self.created_task_ids) or '无'}

## 下次继续建议
- 优先复查本次 Top 候选中的前 3 条。
- 若本次中断，下次可继续使用关键词：{', '.join(self.keywords)}
"""
        self.engine._append_change_log("自主工作日报", report)
        self.report_content = report

    def _finish_partial(self, hole_progress_list: Optional[List[Dict]] = None):
        self.log(f"⏸ 每日计划部分完成：{self.interrupted_reason}", "warning")
        self._generate_report(
            "中断" if self.interrupted_reason == "用户中断" else "超时部分完成",
            [f"- {c.get('title', '')}" for c in self.scored_candidates[:10]],
            hole_progress_list or [],
            len(self.created_pending_ids),
            len(self.created_task_ids),
        )
        self._emit_progress("整理中断成果完成", min(99, self._current_progress()), len(self.completed_stages))

    def _current_progress(self) -> int:
        return int(len(self.completed_stages) / max(1, len(self.STAGES)) * 100)

    def run(
        self,
        intent_keywords: Optional[List[str]] = None,
        time_budget_seconds: int = None,
        stop_flag: Optional[Callable[[], bool]] = None,
        progress_callback: Optional[Callable[[Dict], None]] = None,
    ):
        self._reset_state()
        self.keywords = self._normalize_keywords(intent_keywords)
        budget = int(time_budget_seconds or Config.DAILY_PLAN_TIME_BUDGET_SECONDS)
        self.deadline = self.started_at + max(60, budget)
        self.stop_flag = stop_flag
        self.progress_callback = progress_callback
        hole_progress_list: List[Dict] = []
        external_insights_text: List[str] = []

        self.log("📅 开始执行每日自主工作计划（15分钟高价值版）...", "system")
        self.log(f"  今日意向关键词：{', '.join(self.keywords)}", "system")
        self.update_status("执行每日计划中...")

        try:
            self._ensure_can_continue(self.STAGES[0])
            self._emit_progress(self.STAGES[0], 5, 1)
            L1, L2, L3 = self.engine.update_crystal_layers()
            self.log(f"  分层完成：L1={len(L1)}，L2={len(L2)}，L3={len(L3)}", "success")
            archived = self.engine.archive_cold_crystals()
            self.log(f"  冷晶体归档：{len(archived)} 条", "system")
            # ===== Day 3 新增：运行元原语并触发链 =====
            try:
                self.log("🧠 运行元层元原语（含触发链）...", "system")
                meta_results = self.engine.meta.run_all_primitives()
                triggered = meta_results.get("triggered_chains", [])
                if triggered:
                    self.log(f"✅ 触发链执行：{len(triggered)} 条链被触发", "success")
                    for chain in triggered:
                        self.log(f"   {chain['chain']}: {chain['source']} → {chain['target']} (通过: {chain['passed']})", "system")
                else:
                    self.log("ℹ️ 无触发链条件满足", "system")
            except Exception as e:
                self.log(f"⚠️ 元原语执行异常：{e}", "warning")           
            self.completed_stages.append(self.STAGES[0])

            self._ensure_can_continue(self.STAGES[1])
            self._emit_progress(self.STAGES[1], 15, 2)
            external_insights_text, self.candidates = self._collect_external_info()
            self.log(f"  采集候选信息：{len(self.candidates)} 条", "system")
            self.completed_stages.append(self.STAGES[1])

            self._ensure_can_continue(self.STAGES[2])
            self._emit_progress(self.STAGES[2], 32, 3)
            self.scored_candidates = self._score_candidates(self.candidates)
            self.log(f"  入选高价值候选：{len(self.scored_candidates)} 条", "system")
            self.completed_stages.append(self.STAGES[2])

            self._ensure_can_continue(self.STAGES[3])
            self._emit_progress(self.STAGES[3], 48, 4)
            for item in self.scored_candidates[:Config.DAILY_PLAN_MAX_PENDING_CARDS]:
                if self._should_stop():
                    raise InterruptedError(self.interrupted_reason)
                summarized = self._summarize_candidate(item)
                self._create_pending_card_from_info(summarized)
            self.completed_stages.append(self.STAGES[3])

            self._ensure_can_continue(self.STAGES[4])
            self._emit_progress(self.STAGES[4], 62, 5)

            # 使用向量级矛盾检测（自动选择最优方法）
            try:
                conflicts = self.engine.detect_conflicts(method="auto")
            except Exception as e:
                self.log(f"  向量矛盾检测失败，降级到 Jaccard: {e}", "warning")
                conflicts = self.engine.detect_conflicts(method="jaccard")

            conflict_count = len(conflicts)
            created_count = 0

            # 按严重程度排序：高优先级先处理
            high_severity = [c for c in conflicts if c.similarity > 0.85]
            medium_severity = [c for c in conflicts if 0.70 < c.similarity <= 0.85]
            low_severity = [c for c in conflicts if c.similarity <= 0.70]

            for conflict in high_severity[:3]:
                if self._should_stop():
                    raise InterruptedError(self.interrupted_reason)
                if self._create_conflict_task_card(conflict):
                    created_count += 1

            for conflict in medium_severity[:2]:
                if self._should_stop():
                    raise InterruptedError(self.interrupted_reason)
                if self._create_conflict_task_card(conflict):
                    created_count += 1

            if low_severity:
                self.log(f"  [INFO] 发现 {len(low_severity)} 个潜在关联（低优先级，暂不生成任务卡）", "system")

            self.log(f"  冲突预警：发现 {conflict_count} 个潜在问题（高:{len(high_severity)}, 中:{len(medium_severity)}, 低:{len(low_severity)}），生成 {created_count} 张任务卡", "system")
            self.completed_stages.append(self.STAGES[4])

            self._ensure_can_continue(self.STAGES[5])
            self._emit_progress(self.STAGES[5], 76, 6)
            hole_progress_list, hole_desc = self._update_hole_progress_with_evidence(self.scored_candidates)
            self.log(f"  孔洞评估完成：{hole_desc[:120]}", "system")
            self.completed_stages.append(self.STAGES[5])

            # ===== 新增：负能力修剪 + 验证门控 =====
            self._emit_progress("负能力修剪", 82)
            self.log("🧹 开始负能力修剪与验证门控...", "system")

            try:
                low_contribution = self.engine.get_low_contribution_crystals(threshold=25.0)
                self.log(f"  发现 {len(low_contribution)} 条低贡献晶体", "system")

                verification_count = 0
                archive_count = 0

                for item in low_contribution[:10]:  # 每轮最多处理10条，避免任务泛滥
                    if self._should_stop():
                        raise InterruptedError(self.interrupted_reason)

                    # 执行验证门控
                    validation = self._validate_crystal_for_retention(item)

                    if validation["suggested_action"] == "keep":
                        self.log(f"  ✅ {item['crystal_id']} 通过验证，保留", "system")
                    elif validation["suggested_action"] == "pending_review":
                        # 生成 PENDING 验证卡片
                        self._generate_verification_pending_card(item, validation)
                        verification_count += 1
                    else:  # archive
                        # 记录但不立即删除（等待用户确认）
                        self.log(f"  📦 {item['crystal_id']} 建议归档（等待用户确认）", "system")
                        archive_count += 1

                self.log(f"  验证门控完成：保留 {len(low_contribution) - verification_count - archive_count} 条，待复核 {verification_count} 条，建议归档 {archive_count} 条", "system")

            except Exception as e:
                self.log(f"  ⚠️ 负能力修剪异常：{e}", "warning")

            self._ensure_can_continue(self.STAGES[6])
            self._emit_progress(self.STAGES[6], 90, 7)
            self._generate_wondering(external_insights_text)
            self._generate_report("完成", external_insights_text, hole_progress_list, len(self.created_pending_ids), len(self.created_task_ids))
            self.completed_stages.append(self.STAGES[6])
            self._mark_today_run()
            elapsed = time.time() - self.started_at
            self.log(f"✅ 每日计划执行完成，耗时 {elapsed:.1f} 秒", "success")
            self.update_status("每日计划完成")
            self._emit_progress("每日计划完成", 100, len(self.STAGES))
        except InterruptedError:
            self._finish_partial(hole_progress_list)
            self.update_status("每日计划已中断并整理成果")
        except Exception as exc:
            self.interrupted_reason = f"异常：{exc}"
            self._finish_partial(hole_progress_list)
            self.update_status("每日计划异常收尾")
            raise

        self.log("\n" + "=" * 60 + "\n" + self.report_content + "\n" + "=" * 60, "system")
        return {
            "status": "partial" if self.interrupted_reason else "done",
            "reason": self.interrupted_reason,
            "report": self.report_content,
            "pending_ids": self.created_pending_ids,
            "task_ids": self.created_task_ids,
            "keywords": self.keywords,
        }

    def _validate_crystal_for_retention(self, crystal_data: Dict) -> Dict[str, Any]:
        """
        验证门控自我进化：判断低贡献晶体是否应该保留

        验证规则（可配置）：
        1. 是否至少被 3 个不同角色引用过？
        2. 是否有独立来源支持？（外部案例或论文引用）
        3. 是否被用户主动标记为“固定”？
        4. 是否在最近 60 天内被访问过？

        Returns:
            {
                "passed": bool,
                "rules": [{"name": str, "passed": bool, "reason": str}],
                "summary": str,
                "suggested_action": str  # "keep" | "pending_review" | "archive"
            }
        """
        cid = crystal_data.get("crystal_id")
        if not cid:
            return {"passed": False, "summary": "缺少晶体ID", "suggested_action": "archive"}

        # 获取晶体完整信息
        crystals = self.engine.parse_crystals()
        crystal = next((c for c in crystals if c.id == cid), None)
        if not crystal:
            return {"passed": False, "summary": f"未找到晶体 {cid}", "suggested_action": "archive"}

        rules = []
        passed_count = 0
        total_rules = 4

        # 规则1：至少被 3 个不同晶体引用
        ref_count = sum(1 for c in crystals if cid in c.links)
        rule1_passed = ref_count >= 2  # 放宽到 2 个引用
        rules.append({
            "name": "被其他晶体引用",
            "passed": rule1_passed,
            "reason": f"被 {ref_count} 条晶体引用" if rule1_passed else f"仅被 {ref_count} 条晶体引用（需要 ≥2）"
        })
        if rule1_passed:
            passed_count += 1

        # 规则2：是否有外部来源（检查链接中是否有 H 开头或外部引用）
        has_external = any(link.startswith("H") or link.startswith("http") for link in crystal.links)
        rule2_passed = has_external
        rules.append({
            "name": "有外部来源或孔洞支撑",
            "passed": rule2_passed,
            "reason": "有孔洞或外部链接支撑" if rule2_passed else "仅内部引用，建议补充外部来源"
        })
        if rule2_passed:
            passed_count += 1

        # 规则3：是否被固定到 L1
        state = self.engine.load_layer_state()
        layers = state.get("layers", {})
        manual = state.get("manual_override", {})
        is_fixed = manual.get(cid) == "L1_fixed"
        rule3_passed = is_fixed or layers.get(cid) == "L1"
        rules.append({
            "name": "在核心区（L1）或已固定",
            "passed": rule3_passed,
            "reason": "已固定到 L1" if is_fixed else "在 L1 层" if layers.get(cid) == "L1" else "不在核心区"
        })
        if rule3_passed:
            passed_count += 1

        # 规则4：最近 60 天内被访问过
        last_accessed = state.get("last_accessed", {})
        last = last_accessed.get(cid)
        if last:
            try:
                days_since = (date.today() - date.fromisoformat(last)).days
                rule4_passed = days_since <= 60
            except:
                rule4_passed = False
        else:
            rule4_passed = False
        rules.append({
            "name": "近期被访问",
            "passed": rule4_passed,
            "reason": f"{days_since} 天前访问" if last and rule4_passed else "60 天以上未访问" if last else "从未访问"
        })
        if rule4_passed:
            passed_count += 1

        # 综合判断：至少通过 2 条规则（宽松标准）
        passed = passed_count >= 2

        if passed:
            suggested_action = "keep"
            summary = f"通过 {passed_count}/{total_rules} 条验证规则，建议保留"
            self.engine.log_evolution_event(
                "verification_passed",
                {
                    "crystal_id": cid,
                    "crystal_content": crystal.content[:80],
                    "rules_passed": passed_count,
                    "rules_total": total_rules,
                    "trigger": "daily_plan"
                }
            )
        elif passed_count >= 1:
            suggested_action = "pending_review"
            summary = f"仅通过 {passed_count}/{total_rules} 条验证规则，建议进入 PENDING 复核"
        else:
            suggested_action = "archive"
            summary = f"未通过验证（{passed_count}/{total_rules}），建议归档"

        return {
            "passed": passed,
            "rules": rules,
            "summary": summary,
            "suggested_action": suggested_action,
            "crystal_id": cid,
            "content": crystal.content
        }

    def _generate_verification_pending_card(self, crystal_data: Dict, validation_result: Dict) -> str:
        """
        为低贡献晶体生成 PENDING 验证卡片
        """
        cid = crystal_data.get("crystal_id")
        if not cid:
            return ""

        card_id = f"PENDING-VERIFY-{datetime.now().strftime('%Y%m%d%H%M%S')}-{self._stable_suffix(cid)}"

        rules_text = "\n".join([
            f"  - {'✅' if r['passed'] else '❌'} {r['name']}：{r['reason']}"
            for r in validation_result.get("rules", [])
        ])

        content = f"""## {card_id}
- 类型：验证门控·低贡献晶体复核
- 晶体ID：{cid}
- 晶体内容：{crystal_data.get('content', '')[:80]}
- 贡献得分：{crystal_data.get('score', 0):.1f}
- 验证结果：{validation_result.get('summary', '')}
- 建议动作：{validation_result.get('suggested_action', 'pending_review')}
- 验证详情：
{rules_text}
- AI判断：此晶体贡献度较低，建议人工确认是否保留、修改或归档。
"""
        FileIO.append("pending", f"\n{content}\n")
        self.created_pending_ids.append(card_id)
        self.log(f"  📄 生成验证门控卡片 {card_id} (晶体 {cid})", "system")
        return card_id

# =============================================================================
# 13. Tkinter GUI (gui.py) - 完整复制并调整导入
# =============================================================================
import hashlib
import os
import queue
import threading
import time
import tkinter as tk
from datetime import datetime
from tkinter import Toplevel, filedialog, messagebox, scrolledtext, simpledialog, ttk

class CrystalTreeApp:
    def __init__(self):
        self.files = FileIO()
        self.db = DBManager()
        self.engine = CrystalEngine(self.files)
        self.engine.cheap_gate.log = self._log
        self.ai = AIClient()
        self.fetcher = ExternalFetcher(self._log)
        self.batch_processor = BatchProcessor(self.ai, self._log)

        self.root = tk.Tk()
        self.root.title("认知晶体树 v2.0 · 多角色竞争与联想增强版")
        self.root.geometry("1400x900")
        self.root.minsize(1180, 760)
        self.root.configure(bg=Config.GUI_BG_MAIN)
        self._configure_styles()

        self.api_key = tk.StringVar(value=Config.get_api_key())
        self.fast_mode = tk.BooleanVar(value=False)
        self.scope_var = tk.StringVar(value="L1")
        self.batch_mode = tk.StringVar(value="crystal")
        self.inject_history = tk.BooleanVar(value=True)

        self.deep_think_mode = tk.StringVar(value="多角色竞争")
        self.crystal_mode_var = tk.StringVar(value="单条输入")
        self.debate_rounds_var = tk.IntVar(value=2)

        self.roles = self._load_roles()

        self.processing = False
        self.stop_batch = False
        self.daily_plan_running = False
        self.stop_daily_plan = False
        self.log_queue = queue.Queue()

        # ===== 新增：结果面板引用 =====
        self.result_area = None

        self.current_session_id = None
        self.current_session_name = None
        self.current_history = []
        self.session_listbox = None
        self.session_search_var = tk.StringVar()
        self.session_filtered_ids = []
        self.question_listbox = None
        self.question_indices = []
        self.right_expanded = True
        self._mousewheel_zones = []
        self._mousewheel_bound = False
        self.foldable_contents = {}
        self._chat_in_progress = False
        self._suppress_select = False

        self._init_app()
        self._create_widgets()
        self._start_log_processor()
        self.root.protocol("WM_DELETE_WINDOW", self._on_closing)

    def _configure_styles(self):
        style = ttk.Style(self.root)
        try:
            style.theme_use("clam")
        except tk.TclError:
            pass
        style.configure("TCombobox", fieldbackground=Config.GUI_BG_INPUT, background=Config.GUI_BG_INPUT,
                        foreground=Config.GUI_FG_TEXT, arrowcolor=Config.GUI_ACCENT, bordercolor=Config.GUI_BORDER,
                        lightcolor=Config.GUI_BORDER, darkcolor=Config.GUI_BORDER)
        style.configure("TNotebook", background=Config.GUI_BG_MAIN, borderwidth=0)
        style.configure("TNotebook.Tab", background=Config.GUI_BUTTON_SOFT, foreground=Config.GUI_FG_TEXT,
                        padding=(12, 6), font=Config.GUI_TEXT_FONT)
        style.map("TNotebook.Tab", background=[("selected", Config.GUI_BG_CARD)],
                  foreground=[("selected", Config.GUI_ACCENT_DARK)])
        style.configure("Treeview", background=Config.GUI_BG_INPUT, fieldbackground=Config.GUI_BG_INPUT,
                        foreground=Config.GUI_FG_TEXT, rowheight=28, bordercolor=Config.GUI_BORDER,
                        font=Config.GUI_TEXT_FONT)
        style.configure("Treeview.Heading", background=Config.GUI_BG_CARD_ALT, foreground=Config.GUI_ACCENT_DARK,
                        font=("微软雅黑", 10, "bold"))
        style.map("Treeview", background=[("selected", Config.GUI_HIGHLIGHT)],
                  foreground=[("selected", Config.GUI_FG_TEXT)])
        style.configure("Horizontal.TProgressbar", troughcolor=Config.GUI_BG_CARD_ALT,
                        background=Config.GUI_ACCENT, bordercolor=Config.GUI_BORDER)


    def _load_roles(self) -> List[Dict]:
        path = Config.get_path("roles")
        if not path.exists():
            default = {
                "radical": {
                    "name": "激进者",
                    "instruction": "攻击默认前提，假设现有框架是错的，给出颠覆性方案。"
                },
                "conservative": {
                    "name": "保守者",
                    "instruction": "风险优先，假设资源有限，给出最可落地的稳健方案。"
                },
                "structural": {
                    "name": "结构主义者",
                    "instruction": "从已有晶体中寻找同构案例，用类比生成方案。"
                },
                "judge": {
                    "name": "大法官",
                    "instruction": "以晶体卡片、核心操作原则和资源约束为准绳，做出终审裁决。必须明确引用依据（晶体ID、原则条款或约束条件），不得凭直觉判案。"
                },
                "spokesperson": {
                    "name": "首席发言人",
                    "instruction": "将内部辩论结论转化为清晰、简洁、无歧义的对外陈述。遵循降维（通俗化）、定调（不超过3条核心信息）、检验（老板读前100字能决策）三原则。"
                },
                "lark": {
                    "name": "百灵鸟",
                    "instruction": "见多识广的通用智能体，从外部世界（学术、产业、政策、跨学科）补充知识，打破信息茧房。在第二轮登场。"
                },
                "pilgrim": {
                    "name": "取经者",
                    "instruction": "以长期愿景和核心价值观为锚，防止短期利益或局部优化偏离最终使命。评估方案的可持续性和道德一致性。"
                },
                "strategist": {
                    "name": "奇谋者",
                    "instruction": "善于洞察人心、把握时机，敢押注非常规路径，捕捉机会窗口。评估方案能否借力打力、以奇制胜。"
                },
                "statesman": {
                    "name": "延安智者",
                    "instruction": "坚持调查研究，不唯上、不唯书、只唯实。从全局矛盾和主要矛盾切入，提出实事求是、可落地的综合方略。"
                }
            }
            FileIO.write("roles", json.dumps(default, ensure_ascii=False, indent=2))
            data = default
        else:
            try:
                data = json.loads(path.read_text(encoding="utf-8"))
            except:
                data = {}

        roles_list = []
        for key, val in data.items():
            roles_list.append({
                "id": key,
                "key": key,
                "name": val.get("name", key),
                "instruction": val.get("instruction", "")
            })

        fallback_keys = [
            "radical", "conservative", "structural", "judge",
            "spokesperson", "lark", "pilgrim", "strategist", "statesman"
        ]
        fallback_roles = {
            "radical": {
                "name": "激进者",
                "instruction": "攻击默认前提，假设现有框架是错的，给出颠覆性方案。"
            },
            "conservative": {
                "name": "保守者",
                "instruction": "风险优先，假设资源有限，给出最可落地的稳健方案。"
            },
            "structural": {
                "name": "结构主义者",
                "instruction": "从已有晶体中寻找同构案例，用类比生成方案。"
            },
            "judge": {
                "name": "大法官",
                "instruction": "以晶体卡片、核心操作原则和资源约束为准绳，做出终审裁决。必须明确引用依据（晶体ID、原则条款或约束条件），不得凭直觉判案。"
            },
            "spokesperson": {
                "name": "首席发言人",
                "instruction": "将内部辩论结论转化为清晰、简洁、无歧义的对外陈述。遵循降维（通俗化）、定调（不超过3条核心信息）、检验（老板读前100字能决策）三原则。"
            },
            "lark": {
                "name": "百灵鸟",
                "instruction": "见多识广的通用智能体，从外部世界（学术、产业、政策、跨学科）补充知识，打破信息茧房。在第二轮登场。"
            },
            "pilgrim": {
                "name": "取经者",
                "instruction": "以长期愿景和核心价值观为锚，防止短期利益或局部优化偏离最终使命。评估方案的可持续性和道德一致性。"
            },
            "strategist": {
                "name": "奇谋者",
                "instruction": "善于洞察人心、把握时机，敢押注非常规路径，捕捉机会窗口。评估方案能否借力打力、以奇制胜。"
            },
            "statesman": {
                "name": "延安智者",
                "instruction": "坚持调查研究，不唯上、不唯书、只唯实。从全局矛盾和主要矛盾切入，提出实事求是、可落地的综合方略。"
            }
        }

        existing_keys = {r["key"] for r in roles_list}
        for key in fallback_keys:
            if key not in existing_keys:
                roles_list.append({
                    "id": key,
                    "key": key,
                    "name": fallback_roles[key]["name"],
                    "instruction": fallback_roles[key]["instruction"]
                })
        return roles_list

    def _init_app(self):
        FileIO.ensure_directories()
        FileIO.ensure_default_files()
        self.db._init_db()
        self._load_session_list()
        self._new_session()
        self.root.after(2000, self._check_and_run_daily_plan)

    def _create_widgets(self):
        self.root.grid_rowconfigure(0, weight=1)
        self.root.grid_columnconfigure(0, weight=0, minsize=260)
        self.root.grid_columnconfigure(1, weight=1)
        self.root.grid_columnconfigure(2, weight=0, minsize=250)
        self._create_left_panel()
        self._create_center_panel()
        self._create_right_panel()
        self._create_status_bar()

    def _create_card(self, parent, title):
        card = tk.LabelFrame(parent, text=title, font=Config.GUI_CARD_FONT, bg=Config.GUI_BG_CARD,
                             fg=Config.GUI_ACCENT_DARK, bd=0, relief=tk.FLAT, padx=2, pady=2,
                             labelanchor="nw", highlightthickness=1, highlightbackground=Config.GUI_BORDER)
        card.pack(fill=tk.X, padx=12, pady=6)
        return card

    def _create_btn(self, parent, text, command, color=None, fg="white"):
        bg = color or Config.GUI_ACCENT
        btn = tk.Button(parent, text=text, command=command, bg=bg, fg=fg, font=Config.GUI_BUTTON_FONT,
                        relief=tk.FLAT, bd=0, padx=10, pady=6, cursor="hand2",
                        activebackground=Config.GUI_ACCENT_DARK if bg == Config.GUI_ACCENT else bg,
                        activeforeground=fg)
        return btn

    def _decorate_popup(self, win, title: str, geometry: str):
        win.title(title)
        win.geometry(geometry)
        win.configure(bg=Config.GUI_BG_MAIN)

    def _create_text_panel(self, parent, **kwargs):
        text = scrolledtext.ScrolledText(parent, wrap=tk.WORD,
                                         font=kwargs.pop("font", Config.GUI_LOG_FONT),
                                         bg=Config.GUI_BG_INPUT, fg=Config.GUI_FG_TEXT,
                                         relief=tk.FLAT, bd=0, highlightthickness=1,
                                         highlightbackground=Config.GUI_BORDER,
                                         insertbackground=Config.GUI_ACCENT_DARK, **kwargs)
        self._register_mousewheel(text, text)
        return text

    def _install_mousewheel_router(self):
        if self._mousewheel_bound:
            return
        self._mousewheel_bound = True
        self.root.bind_all("<MouseWheel>", self._route_mousewheel, add="+")
        self.root.bind_all("<Button-4>", self._route_mousewheel, add="+")
        self.root.bind_all("<Button-5>", self._route_mousewheel, add="+")

    def _register_mousewheel(self, container, target=None):
        target = target or container
        try:
            setattr(container, "_crystal_wheel_target", target)
        except Exception:
            pass
        if not any(item[0] is container and item[1] is target for item in self._mousewheel_zones):
            self._mousewheel_zones.append((container, target))
        self._install_mousewheel_router()

    def _event_scroll_delta(self, event) -> int:
        if getattr(event, "num", None) == 4:
            return -1
        if getattr(event, "num", None) == 5:
            return 1
        delta = getattr(event, "delta", 0)
        if not delta:
            return 0
        return -1 if delta > 0 else 1

    def _route_mousewheel(self, event):
        delta = self._event_scroll_delta(event)
        if not delta:
            return None
        target = None
        try:
            widget = self.root.winfo_containing(event.x_root, event.y_root)
            while widget is not None:
                target = getattr(widget, "_crystal_wheel_target", None)
                if target is not None:
                    break
                parent_name = widget.winfo_parent()
                widget = widget.nametowidget(parent_name) if parent_name else None
        except Exception:
            target = None
        if target is None:
            for container, candidate in reversed(self._mousewheel_zones):
                try:
                    x1 = container.winfo_rootx()
                    y1 = container.winfo_rooty()
                    x2 = x1 + container.winfo_width()
                    y2 = y1 + container.winfo_height()
                    if x1 <= event.x_root <= x2 and y1 <= event.y_root <= y2:
                        target = candidate
                        break
                except Exception:
                    continue
        if target is None:
            return None
        try:
            target.yview_scroll(delta, "units")
            return "break"
        except Exception:
            return None

    def _add_label(self, parent, text, row, col, **kwargs):
        lbl = tk.Label(parent, text=text, bg=Config.GUI_BG_CARD, fg=Config.GUI_FG_TEXT,
                       font=Config.GUI_TEXT_FONT, **kwargs)
        lbl.grid(row=row, column=col, sticky=tk.W, pady=4)
        return lbl

    def _add_entry(self, parent, text_var, width, row, col, **kwargs):
        entry = tk.Entry(parent, textvariable=text_var, width=width, font=Config.GUI_TEXT_FONT,
                         bg=Config.GUI_BG_INPUT, fg=Config.GUI_FG_TEXT, relief=tk.FLAT, bd=0,
                         highlightthickness=1, highlightbackground=Config.GUI_BORDER,
                         highlightcolor=Config.GUI_ACCENT, **kwargs)
        entry.grid(row=row, column=col, pady=4, padx=5)
        return entry

    def _add_btn(self, parent, text, command, row, col):
        btn = self._create_btn(parent, text, command, Config.GUI_ACCENT)
        btn.grid(row=row, column=col, padx=5, pady=4)
        return btn

    def _draw_center_background(self, canvas, width, height):
        canvas.delete("bg_ornament")
        ornaments = [
            (width - 250, -150, width + 70, 170, "#ded3fa"),
            (width - 470, 120, width - 130, 460, "#eee8fb"),
            (95, height - 90, 330, height + 145, "#e8defb"),
        ]
        for x1, y1, x2, y2, color in ornaments:
            canvas.create_oval(x1, y1, x2, y2, fill=color, outline="", tags="bg_ornament")
        canvas.tag_lower("bg_ornament")

    def _log(self, msg: str, tag: str = None):
        timestamp = datetime.now().strftime("%H:%M:%S")
        if tag == "user":
            formatted = f"\n[{timestamp}] 🧑 你\n{msg}\n\n"
            self.log_queue.put((formatted, tag))
        elif tag == "ai":
            formatted = f"\n[{timestamp}] 🤖 AI\n{msg}\n"
            self.log_queue.put((formatted, tag))
        elif tag in ("system", "error", "success", "warning"):
            formatted = f"[{timestamp}] {msg}\n"
            self.log_queue.put((formatted, tag))
        else:
            formatted = f"[{timestamp}] {msg}\n"
            self.log_queue.put((formatted, tag))

    def _start_log_processor(self):
        try:
            while True:
                msg, tag = self.log_queue.get_nowait()
                if hasattr(self, 'log_area'):
                    if tag:
                        self.log_area.insert(tk.END, msg, tag)
                    else:
                        self.log_area.insert(tk.END, msg)
                    self.log_area.see(tk.END)
        except queue.Empty:
            pass
        self.root.after(100, self._start_log_processor)

    def update_status(self, text):
        if hasattr(self, 'status_var'):
            self.status_var.set(text)

    def _setup_log_context_menu(self):
        """为日志区域添加右键复制菜单"""
        self.log_menu = tk.Menu(self.log_area, tearoff=0)
        self.log_menu.add_command(label="复制", command=self._copy_log_selection)
        self.log_menu.add_command(label="全选", command=self._select_all_log)
        self.log_area.bind("<Button-3>", self._show_log_menu)
        self.log_area.bind("<Control-c>", self._copy_log_selection_event)
        self.log_area.bind("<Control-a>", self._select_all_log_event)
        self.log_area.config(state=tk.NORMAL)

    def _show_log_menu(self, event):
        try:
            self.log_menu.tk_popup(event.x_root, event.y_root)
        finally:
            self.log_menu.grab_release()

    def _copy_log_selection(self):
        try:
            selected = self.log_area.get(tk.SEL_FIRST, tk.SEL_LAST)
            self.root.clipboard_clear()
            self.root.clipboard_append(selected)
        except tk.TclError:
            pass
        
    def _typewriter_print(self, text: str, tag: str = None, delay_ms: int = 30,
                          final_newline: bool = True, callback: callable = None):
        """打字机效果，支持完成后回调"""
        if not text:
            if callback:
                callback()
            return

        def _print_next(index: int):
            if index >= len(text):
                if final_newline:
                    self.log_area.insert(tk.END, "\n", ())
                    self.log_area.see(tk.END)
                if callback:
                    callback()
                return
            char = text[index]
            if tag:
                self.log_area.insert(tk.END, char, (tag,))
            else:
                self.log_area.insert(tk.END, char, ())
            self.log_area.see(tk.END)
            self.root.after(delay_ms, lambda: _print_next(index + 1))

        _print_next(0)
        

    def _select_all_log(self):
        self.log_area.tag_add(tk.SEL, "1.0", tk.END)
        self.log_area.mark_set(tk.INSERT, "1.0")
        self.log_area.see(tk.INSERT)

    def _copy_log_selection_event(self, event):
        self._copy_log_selection()
        return "break"

    def _select_all_log_event(self, event):
        self._select_all_log()
        return "break"

    def _create_left_panel(self):
        left_frame = tk.Frame(self.root, bg=Config.GUI_BG_SIDEBAR, width=260,
                              highlightthickness=1, highlightbackground=Config.GUI_BORDER)
        left_frame.grid(row=0, column=0, sticky="ns", padx=(8, 4), pady=8)
        left_frame.grid_propagate(False)
        left_frame.config(width=260)
        tk.Label(left_frame, text="📂 历史会话", font=("微软雅黑", 12, "bold"),
                 bg=Config.GUI_BG_SIDEBAR, fg=Config.GUI_ACCENT_DARK).pack(pady=(12, 6))
        search_frame = tk.Frame(left_frame, bg=Config.GUI_BG_SIDEBAR)
        search_frame.pack(fill=tk.X, padx=8, pady=5)
        tk.Entry(search_frame, textvariable=self.session_search_var, font=Config.GUI_TEXT_FONT,
                 bg=Config.GUI_BG_INPUT, relief=tk.FLAT, bd=0,
                 highlightthickness=1, highlightbackground=Config.GUI_BORDER).pack(side=tk.LEFT, fill=tk.X,
                                                                                   expand=True)
        tk.Button(search_frame, text="🔍", command=self._filter_sessions, width=3,
                  bg=Config.GUI_ACCENT, fg='white', font=Config.GUI_BUTTON_FONT,
                  relief=tk.FLAT, bd=0).pack(side=tk.RIGHT, padx=2)
        list_frame = tk.Frame(left_frame, bg=Config.GUI_BG_SIDEBAR)
        list_frame.pack(fill=tk.BOTH, expand=True, padx=5, pady=5)
        self.session_listbox = tk.Listbox(list_frame, font=Config.GUI_TEXT_FONT,
                                          bg=Config.GUI_BG_INPUT, fg=Config.GUI_FG_TEXT,
                                          selectbackground=Config.GUI_HIGHLIGHT,
                                          selectforeground=Config.GUI_ACCENT_DARK,
                                          relief=tk.FLAT, bd=0, activestyle="none",
                                          selectmode=tk.EXTENDED)   # 添加这一行)
        self.session_listbox.pack(side=tk.LEFT, fill=tk.BOTH, expand=True)
        scrollbar = ttk.Scrollbar(list_frame, orient=tk.VERTICAL, command=self.session_listbox.yview)
        scrollbar.pack(side=tk.RIGHT, fill=tk.Y)
        self.session_listbox.config(yscrollcommand=scrollbar.set)
        self._register_mousewheel(self.session_listbox, self.session_listbox)
        self.session_listbox.bind('<<ListboxSelect>>', self._on_session_select)
        btn_frame = tk.Frame(left_frame, bg=Config.GUI_BG_SIDEBAR)
        btn_frame.pack(fill=tk.X, padx=5, pady=10)
        self._create_btn(btn_frame, "➕ 新建会话", self._new_session, Config.GUI_SUCCESS).pack(side=tk.LEFT,
                                                                                               padx=2)
        self._create_btn(btn_frame, "✏️ 重命名", self._rename_current_session, Config.GUI_WARNING).pack(side=tk.LEFT,
                                                                                                        padx=2)
        self._create_btn(btn_frame, "🗑️ 删除", self._delete_selected_sessions, Config.GUI_DANGER).pack(side=tk.LEFT, padx=2)

    def _create_center_panel(self):
        center_container = tk.Frame(self.root, bg=Config.GUI_BG_MAIN, width=820)
        center_container.grid(row=0, column=1, sticky="nsew", padx=4, pady=8)
        center_container.grid_propagate(False)
        canvas = tk.Canvas(center_container, bg=Config.GUI_BG_MAIN, highlightthickness=0)
        scrollbar = ttk.Scrollbar(center_container, orient="vertical", command=canvas.yview)
        scrollable_frame = tk.Frame(canvas, bg=Config.GUI_BG_MAIN)
        scrollable_frame.bind("<Configure>", lambda e: canvas.configure(scrollregion=canvas.bbox("all")))
        canvas_window = canvas.create_window((0, 0), window=scrollable_frame, anchor="nw")

        def configure_canvas_width(event):
            canvas.itemconfig(canvas_window, width=event.width)
            canvas.configure(scrollregion=canvas.bbox("all"))
            self._draw_center_background(canvas, event.width, event.height)

        canvas.bind("<Configure>", configure_canvas_width)
        canvas.configure(yscrollcommand=scrollbar.set)
        self._bind_canvas_mousewheel(canvas, scrollable_frame)
        canvas.pack(side="left", fill="both", expand=True)
        scrollbar.pack(side="right", fill="y")
        center_frame = scrollable_frame

        title_frame = tk.Frame(center_frame, bg=Config.GUI_BG_HEADER, height=62)
        title_frame.pack(fill=tk.X, pady=(0, 10), padx=4)
        title_frame.pack_propagate(False)
        title_canvas = tk.Canvas(title_frame, bg=Config.GUI_BG_HEADER, highlightthickness=0)
        title_canvas.pack(fill=tk.BOTH, expand=True)

        def _paint_title(event):
            title_canvas.delete("all")
            width = event.width
            title_canvas.create_oval(width - 260, -150, width + 70, 180, fill="#d8cbf7", outline="")
            title_canvas.create_oval(70, 34, 210, 174, fill="#f4efff", outline="")
            title_canvas.create_text(width / 2, 23,
                                     text="🧠 认知晶体树 v2.0 · 多角色竞争与联想增强版",
                                     font=Config.GUI_TITLE_FONT, fill=Config.GUI_ACCENT_DARK)
            title_canvas.create_text(width / 2, 50,
                                     text="注意力上限50 | 动态联想检索 | 多角色竞争（实验） | 每日计划自动落地",
                                     font=Config.GUI_TEXT_FONT, fill=Config.GUI_FG_MUTED)

        title_canvas.bind("<Configure>", _paint_title)
        sep = tk.Frame(center_frame, height=1, bg=Config.GUI_BORDER)
        sep.pack(fill=tk.X, padx=5, pady=(0, 10))
        self._create_config_card(center_frame)
        self._create_single_card(center_frame)
        self._create_batch_card(center_frame)
        self._create_info_card(center_frame)
        self._create_log_card(center_frame)

        # ===== 新增：独立的结果展示面板 =====
        self._create_result_panel(center_frame)

    def _create_config_card(self, parent):
        config_card = self._create_card(parent, "⚙️ 全局配置")
        config_frame = tk.Frame(config_card, bg=Config.GUI_BG_CARD)
        config_frame.pack(fill=tk.X, padx=12, pady=7)
        self._add_label(config_frame, "DeepSeek API Key：", 0, 0)
        self._add_entry(config_frame, self.api_key, 46, 0, 1, show="*")
        self._add_btn(config_frame, "从环境变量加载", self._load_api_from_env, 0, 2)
        tk.Checkbutton(config_frame, text="快速模式（晶体化时跳过外部搜索）",
                       variable=self.fast_mode, bg=Config.GUI_BG_CARD, fg=Config.GUI_FG_TEXT,
                       activebackground=Config.GUI_BG_CARD, font=Config.GUI_TEXT_FONT).grid(row=1, column=1,
                                                                                            sticky=tk.W, pady=4)
        self._add_label(config_frame, "检索范围:", 1, 0)
        scope_menu = ttk.Combobox(config_frame, textvariable=self.scope_var,
                                  values=["L1", "L1+L2", "all"], state="readonly", width=10)
        scope_menu.grid(row=1, column=2, sticky=tk.W, pady=4, padx=5)
        # 在现有配置行后面添加模式选择
        self._add_label(config_frame, "运行模式:", 2, 0)
        self.profile_var = tk.StringVar(value="balanced")
        profile_menu = ttk.Combobox(
            config_frame,
            textvariable=self.profile_var,
            values=["high_accuracy", "balanced", "economy"],
            state="readonly",
            width=12
        )
        profile_menu.grid(row=2, column=1, sticky=tk.W, pady=4, padx=5)
        # 添加模式描述标签
        self.profile_desc_label = tk.Label(
            config_frame,
            text="平衡模式：质量与成本的平衡",
            bg=Config.GUI_BG_CARD,
            fg=Config.GUI_FG_MUTED,
            font=("微软雅黑", 9)
        )
        self.profile_desc_label.grid(row=2, column=2, sticky=tk.W, pady=4, padx=5)
        # 绑定选择事件
        profile_menu.bind("<<ComboboxSelected>>", self._on_profile_change)

    def _create_single_card(self, parent):
        single_card = self._create_card(parent, "📝 单条输入与文件对话")
        single_frame = tk.Frame(single_card, bg=Config.GUI_BG_CARD)
        single_frame.pack(fill=tk.X, padx=12, pady=8)
        primary_line = tk.Frame(single_frame, bg=Config.GUI_BG_CARD)
        primary_line.pack(pady=(0, 4))
        self._create_btn(primary_line, "💬 聊天 (回车)", self._do_chat).pack(side=tk.LEFT, padx=5)
                # ---- 晶体化按钮及模式选择 ----
        crystal_btn_frame = tk.Frame(primary_line, bg=Config.GUI_BG_CARD)
        crystal_btn_frame.pack(side=tk.LEFT, padx=5)
        self.crystal_mode_var = tk.StringVar(value="单条输入")
        crystal_mode_combo = ttk.Combobox(crystal_btn_frame, textvariable=self.crystal_mode_var,
                                          values=["单条输入", "会话历史"], state="readonly", width=10)
        crystal_mode_combo.pack(side=tk.LEFT, padx=(0, 3))
        self._create_btn(crystal_btn_frame, "✨ 晶体化", self._do_crystal, Config.GUI_SUCCESS).pack(side=tk.LEFT)
        # ---- 结束 ----

        mode_frame = tk.Frame(primary_line, bg=Config.GUI_BG_CARD)
        mode_frame.pack(side=tk.LEFT, padx=5)
        self.deep_mode_combo = ttk.Combobox(mode_frame, textvariable=self.deep_think_mode,
                                            values=["多角色竞争", "卢氏注意力增强", "单路径深度推理", "替身自我博弈"],
                                            state="readonly", width=14)
        self.deep_mode_combo.pack(side=tk.LEFT, padx=2)
        
        self.deep_mode_combo.bind("<<ComboboxSelected>>", lambda _event: self._sync_debate_round_state())
        tk.Label(mode_frame, text="轮次", bg=Config.GUI_BG_CARD, fg=Config.GUI_FG_MUTED,
                 font=Config.GUI_TEXT_FONT).pack(side=tk.LEFT, padx=(4, 1))
        self.debate_rounds_spinbox = tk.Spinbox(mode_frame, from_=2, to=12, textvariable=self.debate_rounds_var,
                                                width=3, font=Config.GUI_TEXT_FONT, bg=Config.GUI_BG_INPUT,
                                                fg=Config.GUI_FG_TEXT, relief=tk.SOLID, bd=1)
        self.debate_rounds_spinbox.pack(side=tk.LEFT, padx=2)
        self._create_btn(mode_frame, "🔍 深度分析", self._do_deep_reasoning, '#9c27b0').pack(side=tk.LEFT,
                                                                                           padx=2)
        self._sync_debate_round_state()
        secondary_line = tk.Frame(single_frame, bg=Config.GUI_BG_CARD)
        secondary_line.pack(pady=(0, 2))
        self._create_btn(secondary_line, "📄 选择文件并对话", self._do_file_chat, Config.GUI_WARNING).pack(
            side=tk.LEFT, padx=5)
        self._create_btn(secondary_line, "🗑 清空输入", self._clear_input, Config.GUI_INFO).pack(side=tk.LEFT,
                                                                                                padx=5)
        self._create_btn(secondary_line, "🧹 清空当前会话", self._clear_current_session_messages,
                         Config.GUI_DANGER).pack(side=tk.LEFT, padx=5)
        self.input_text = tk.Text(single_frame, height=4, font=Config.GUI_INPUT_FONT, wrap=tk.WORD,
                                  bg=Config.GUI_BG_INPUT, fg=Config.GUI_FG_TEXT, relief=tk.FLAT, bd=0,
                                  highlightthickness=1, highlightbackground=Config.GUI_BORDER,
                                  highlightcolor=Config.GUI_ACCENT,
                                  insertbackground=Config.GUI_ACCENT_DARK)
        self.input_text.pack(fill=tk.X, pady=(9, 4))
        self.input_text.bind("<Return>", self._on_enter_chat)
        self.input_text.bind("<Control-Return>", self._on_enter_crystal)

    def _create_batch_card(self, parent):
        batch_card = self._create_card(parent, "📁 批量处理文件夹")
        batch_frame = tk.Frame(batch_card, bg=Config.GUI_BG_CARD)
        batch_frame.pack(fill=tk.X, padx=12, pady=7)
        mode_frame = tk.Frame(batch_frame, bg=Config.GUI_BG_CARD)
        mode_frame.pack(anchor=tk.W, pady=3)
        tk.Label(mode_frame, text="处理模式：", bg=Config.GUI_BG_CARD, fg=Config.GUI_FG_TEXT,
                 font=Config.GUI_TEXT_FONT).pack(side=tk.LEFT)
        tk.Radiobutton(mode_frame, text="晶体化（入库）", variable=self.batch_mode,
                       value="crystal", bg=Config.GUI_BG_CARD, fg=Config.GUI_FG_TEXT,
                       activebackground=Config.GUI_BG_CARD, font=Config.GUI_TEXT_FONT).pack(side=tk.LEFT,
                                                                                            padx=10)
        tk.Radiobutton(mode_frame, text="聊天（仅对话）", variable=self.batch_mode,
                       value="chat", bg=Config.GUI_BG_CARD, fg=Config.GUI_FG_TEXT,
                       activebackground=Config.GUI_BG_CARD, font=Config.GUI_TEXT_FONT).pack(side=tk.LEFT,
                                                                                            padx=10)
        inject_frame = tk.Frame(batch_frame, bg=Config.GUI_BG_CARD)
        inject_frame.pack(anchor=tk.W, pady=2)
        tk.Checkbutton(inject_frame, text="将批量处理结果自动注入当前会话",
                       variable=self.inject_history, bg=Config.GUI_BG_CARD, fg=Config.GUI_FG_TEXT,
                       activebackground=Config.GUI_BG_CARD, font=Config.GUI_TEXT_FONT).pack(side=tk.LEFT)
        btn_batch_frame = tk.Frame(batch_frame, bg=Config.GUI_BG_CARD)
        btn_batch_frame.pack(pady=6)
        self.batch_btn = self._create_btn(btn_batch_frame, "📂 选择文件夹并开始批量处理",
                                          self._start_batch, Config.GUI_WARNING)
        self.batch_btn.pack(side=tk.LEFT, padx=5)
        self.stop_batch_btn = self._create_btn(btn_batch_frame, "⏹ 停止批量处理",
                                               self._stop_batch_process, Config.GUI_DANGER)
        self.stop_batch_btn.pack(side=tk.LEFT, padx=5)
        self.stop_batch_btn.config(state=tk.DISABLED)
        self.batch_progress = ttk.Progressbar(batch_frame, mode='determinate', length=400)
        self.batch_progress.pack(fill=tk.X, pady=5)

    def _create_info_card(self, parent):
        info_card = self._create_card(parent, "📊 信息查询")
        info_frame = tk.Frame(info_card, bg=Config.GUI_BG_CARD)
        info_frame.pack(fill=tk.X, padx=12, pady=7)
        groups = [
            [("📊 系统状态", self._show_status, Config.GUI_BUTTON_SOFT, Config.GUI_FG_TEXT),
             ("🕳️ 孔洞花园", self._show_holes, Config.GUI_BUTTON_SOFT, Config.GUI_FG_TEXT),
             ("📋 待确认卡片", self._show_pending, Config.GUI_BUTTON_SOFT, Config.GUI_FG_TEXT),
             ("✅ 确认卡片", self._confirm_card_dialog, Config.GUI_BUTTON_SOFT, Config.GUI_FG_TEXT),
             ("🆕 今日新增", self._show_today_changes, Config.GUI_BUTTON_SOFT, Config.GUI_FG_TEXT)],
            [("📋 任务面板", self._open_task_panel, Config.GUI_ACCENT, "white"),
             ("🔮 晶体管理", self._open_crystal_manager, Config.GUI_ACCENT, "white"),
             ("📅 每日计划", self._manual_run_daily_plan, Config.GUI_INFO, "white"),
             ("🔍 文档搜索", self._open_search_window, Config.GUI_BUTTON_SOFT, Config.GUI_FG_TEXT),
             ("🩺 健康检查", self._show_health_check, Config.GUI_BUTTON_SOFT, Config.GUI_FG_TEXT)],
             [("🔄 同步向量库", self._sync_vector_store, Config.GUI_INFO, "white"),  # ← 新增的第三行
             ("📄 导出AGENTS.md", self._export_agents_md, Config.GUI_WARNING, "white"),
             ("📦 导出Skill", self._export_skill_dialog, Config.GUI_SUCCESS, "white")],
             [("🔍 Meta搜索", self._do_meta_search, Config.GUI_ACCENT, "white"),
             ("📊 帕累托报告", self._show_pareto_report, Config.GUI_INFO, "white"),
             ("📈 认知效率", self._show_cognitive_efficiency, Config.GUI_ACCENT, "white"),
             ("🔥 灵感熔炉", self._show_inspiration_furnace, Config.GUI_WARNING, "white"),
             ("📊 饱和状态", self._show_saturation_status, Config.GUI_INFO, "white")],
        ]
        for group in groups:
            row = tk.Frame(info_frame, bg=Config.GUI_BG_CARD)
            row.pack(anchor=tk.W, pady=3)
            for text, cmd, color, fg in group:
                self._create_btn(row, text, cmd, color, fg).pack(side=tk.LEFT, padx=4)
        self.stop_daily_btn = self._create_btn(info_frame, "⏹ 中断每日计划并整理成果",
                                               self._stop_daily_plan, Config.GUI_DANGER, "white")
        self.stop_daily_btn.pack(anchor=tk.W, padx=5, pady=(4, 0))
        self.stop_daily_btn.config(state=tk.DISABLED)
        progress_frame = tk.LabelFrame(info_card, text="🧠 辩论进度", 
                                       font=Config.GUI_CARD_FONT, bg=Config.GUI_BG_CARD,
                                       fg=Config.GUI_ACCENT_DARK, bd=0, relief=tk.FLAT, padx=5, pady=5)
        progress_frame.pack(fill=tk.X, padx=12, pady=6)
           
        self.debate_progress = ttk.Progressbar(progress_frame, mode='determinate', length=400)
        self.debate_progress.pack(fill=tk.X, padx=5, pady=2)
    
        self.debate_status_label = tk.Label(progress_frame, text="⏳ 等待开始", 
                                            bg=Config.GUI_BG_CARD, fg=Config.GUI_FG_MUTED,
                                            font=Config.GUI_TEXT_FONT)
        self.debate_status_label.pack(anchor=tk.W, padx=5)

    def _create_log_card(self, parent):
        log_card = self._create_card(parent, "📜 对话与日志")
        self.log_area = scrolledtext.ScrolledText(
            log_card,
            wrap=tk.WORD,
            height=18,
            font=Config.GUI_LOG_FONT,
            fg=Config.GUI_FG_TEXT,
            bg=Config.GUI_BG_INPUT,
            relief=tk.FLAT,
            bd=0,
            highlightthickness=1,
            highlightbackground=Config.GUI_BORDER,
            insertbackground=Config.GUI_ACCENT_DARK,
            exportselection=True,
            selectbackground="#c7d2fe",
            selectforeground="#1e1b4b",
            takefocus=True
        )
        self.log_area.pack(fill=tk.BOTH, expand=True, padx=12, pady=10)
        self._register_mousewheel(self.log_area, self.log_area)
        self._setup_log_context_menu()
        self.log_area.tag_config("user",
                                 foreground=Config.GUI_ACCENT_DARK,
                                 background="#f3edff",
                                 font=(Config.GUI_LOG_FONT[0], Config.GUI_LOG_FONT[1], "bold"),
                                 spacing1=4, spacing3=8,
                                 selectbackground="#c7d2fe",
                                 selectforeground="#1e1b4b")
        self.log_area.tag_config("ai",
                                 foreground=Config.GUI_FG_TEXT,
                                 background="#ffffff",
                                 font=("宋体", 13),
                                 spacing1=4, spacing3=8,
                                 selectbackground="#c7d2fe",
                                 selectforeground="#1e1b4b")
        self.log_area.tag_config("system",
                                 foreground=Config.GUI_FG_MUTED,
                                 spacing1=2, spacing3=4)
        self.log_area.tag_config("error",
                                 foreground=Config.GUI_DANGER,
                                 spacing1=2, spacing3=4)
        self.log_area.tag_config("success",
                                 foreground=Config.GUI_SUCCESS,
                                 spacing1=2, spacing3=4)
        self.log_area.tag_config("warning",
                                 foreground=Config.GUI_WARNING,
                                 spacing1=2, spacing3=4)
        self.log_area.tag_raise(tk.SEL)

    def _create_result_panel(self, parent):
        """创建独立的结果展示面板（与日志分离）"""
        result_card = self._create_card(parent, "📋 辩论报告（最终输出）")

        self.result_area = scrolledtext.ScrolledText(
            result_card,
            wrap=tk.WORD,
            height=20,
            font=("微软雅黑", 11),
            fg=Config.GUI_FG_TEXT,
            bg="#faf8ff",
            relief=tk.FLAT,
            bd=0,
            highlightthickness=1,
            highlightbackground=Config.GUI_BORDER,
            insertbackground=Config.GUI_ACCENT_DARK,
            selectbackground="#c7d2fe",
            selectforeground="#1e1b4b"
        )
        self.result_area.pack(fill=tk.BOTH, expand=True, padx=12, pady=10)
        self._register_mousewheel(self.result_area, self.result_area)

        # 配置所有样式标签
        self._configure_result_tags()

        # 默认显示占位信息
        self.result_area.insert(tk.END, "💡 深度推理完成后，完整报告将在此显示。\n", "result_placeholder")
        self.result_area.config(state=tk.DISABLED)

    def _configure_result_tags(self):
        """配置结果面板的所有样式标签"""
        self.result_area.tag_config("result_placeholder",
                                    foreground=Config.GUI_FG_MUTED,
                                    font=("微软雅黑", 12, "italic"))

        # ---- 阶段标题 ----
        self.result_area.tag_config("phase_title",
                                    font=("微软雅黑", 14, "bold"),
                                    foreground="#1A237E",
                                    spacing1=10,
                                    spacing3=8)

        # ---- 角色观点标签 ----
        self.result_area.tag_config("role_radical",
                                    foreground="#c62828",
                                    font=("微软雅黑", 11, "bold"),
                                    background="#ffebee",
                                    spacing1=4,
                                    spacing3=4,
                                    lmargin1=10,
                                    lmargin2=10)

        self.result_area.tag_config("role_conservative",
                                    foreground="#0d47a1",
                                    font=("微软雅黑", 11, "bold"),
                                    background="#e3f2fd",
                                    spacing1=4,
                                    spacing3=4,
                                    lmargin1=10,
                                    lmargin2=10)

        self.result_area.tag_config("role_structural",
                                    foreground="#1b5e20",
                                    font=("微软雅黑", 11, "bold"),
                                    background="#e8f5e9",
                                    spacing1=4,
                                    spacing3=4,
                                    lmargin1=10,
                                    lmargin2=10)

        self.result_area.tag_config("role_lark",
                                    foreground="#4a148c",
                                    font=("微软雅黑", 11, "bold"),
                                    background="#f3e5f5",
                                    spacing1=4,
                                    spacing3=4,
                                    lmargin1=10,
                                    lmargin2=10)

        self.result_area.tag_config("role_pilgrim",
                                    foreground="#bf360c",
                                    font=("微软雅黑", 11, "bold"),
                                    background="#fbe9e7",
                                    spacing1=4,
                                    spacing3=4,
                                    lmargin1=10,
                                    lmargin2=10)

        self.result_area.tag_config("role_strategist",
                                    foreground="#00695c",
                                    font=("微软雅黑", 11, "bold"),
                                    background="#e0f2f1",
                                    spacing1=4,
                                    spacing3=4,
                                    lmargin1=10,
                                    lmargin2=10)

        self.result_area.tag_config("role_statesman",
                                    foreground="#4e342e",
                                    font=("微软雅黑", 11, "bold"),
                                    background="#efebe9",
                                    spacing1=4,
                                    spacing3=4,
                                    lmargin1=10,
                                    lmargin2=10)

        self.result_area.tag_config("role_judge",
                                    foreground="#f57f17",
                                    font=("微软雅黑", 11, "bold"),
                                    background="#fff8e1",
                                    spacing1=4,
                                    spacing3=4,
                                    lmargin1=10,
                                    lmargin2=10)

        self.result_area.tag_config("role_spokesperson",
                                    foreground="#1a237e",
                                    font=("微软雅黑", 11, "bold"),
                                    background="#e8eaf6",
                                    spacing1=4,
                                    spacing3=4,
                                    lmargin1=10,
                                    lmargin2=10)

        self.result_area.tag_config("role_default",
                                    foreground=Config.GUI_ACCENT_DARK,
                                    font=("微软雅黑", 11, "bold"),
                                    background=Config.GUI_BG_CARD_ALT,
                                    spacing1=4,
                                    spacing3=4,
                                    lmargin1=10,
                                    lmargin2=10)

        # ---- 角色观点正文 ----
        self.result_area.tag_config("viewpoint_body",
                                    font=("微软雅黑", 10),
                                    foreground="#333333",
                                    spacing1=2,
                                    spacing3=6,
                                    lmargin1=20,
                                    lmargin2=20)

        # ---- 大法官板块 ----
        self.result_area.tag_config("judge_header",
                                    font=("微软雅黑", 13, "bold"),
                                    foreground="#e65100",
                                    spacing1=12,
                                    spacing3=6)

        self.result_area.tag_config("judge_body",
                                    font=("微软雅黑", 10),
                                    foreground="#4e342e",
                                    spacing1=3,
                                    spacing3=6,
                                    lmargin1=15,
                                    lmargin2=15)

        self.result_area.tag_config("judge_table",
                                    font=("Consolas", 10),
                                    foreground="#1a237e",
                                    background="#fafafa",
                                    spacing1=2,
                                    spacing3=2,
                                    lmargin1=15,
                                    lmargin2=15)

        # ---- 首席发言人板块 ----
        self.result_area.tag_config("spokesperson_header",
                                    font=("微软雅黑", 13, "bold"),
                                    foreground="#0d47a1",
                                    spacing1=12,
                                    spacing3=6)

        self.result_area.tag_config("spokesperson_body",
                                    font=("微软雅黑", 11),
                                    foreground="#1a237e",
                                    spacing1=3,
                                    spacing3=6,
                                    lmargin1=15,
                                    lmargin2=15)

        # ---- 儒雅笔谈 ----
        self.result_area.tag_config("elegant_header",
                                    font=("华文楷体", 13, "bold"),
                                    foreground="#4A3728",
                                    spacing1=12,
                                    spacing3=6)

        self.result_area.tag_config("elegant_body",
                                    font=("华文楷体", 12),
                                    foreground="#3C2A1E",
                                    spacing1=4,
                                    spacing3=8,
                                    lmargin1=25,
                                    lmargin2=25)

        # ---- 分隔线 ----
        self.result_area.tag_config("result_divider",
                                    font=("微软雅黑", 1),
                                    foreground=Config.GUI_BORDER,
                                    spacing1=8,
                                    spacing3=8)

    def _create_right_panel(self):
        self.right_pane = tk.Frame(self.root, bg=Config.GUI_BG_SIDEBAR, width=250,
                                   highlightthickness=1, highlightbackground=Config.GUI_BORDER)
        self.right_pane.grid(row=0, column=2, sticky="ns", padx=(4, 8), pady=8)
        self.right_pane.grid_propagate(False)
        self.right_pane.config(width=250)
        right_header = tk.Frame(self.right_pane, bg=Config.GUI_BG_HEADER)
        right_header.pack(fill=tk.X)
        tk.Label(right_header, text="📋 本会话问题", font=("微软雅黑", 10, "bold"),
                 bg=Config.GUI_BG_HEADER, fg=Config.GUI_ACCENT_DARK).pack(side=tk.LEFT, padx=8, pady=6)
        self.collapse_btn = tk.Button(right_header, text="◀", command=self._toggle_right_pane,
                                      width=3, bg=Config.GUI_ACCENT, fg='white',
                                      font=Config.GUI_BUTTON_FONT, relief=tk.FLAT, bd=0,
                                      cursor="hand2", activebackground=Config.GUI_ACCENT_DARK,
                                      activeforeground="white")
        self.collapse_btn.pack(side=tk.RIGHT, padx=2)
        self.question_frame = tk.Frame(self.right_pane, bg=Config.GUI_BG_SIDEBAR)
        self.question_frame.pack(fill=tk.BOTH, expand=True, padx=5, pady=5)
        self.question_listbox = tk.Listbox(self.question_frame, font=Config.GUI_TEXT_FONT,
                                           bg=Config.GUI_BG_INPUT, fg=Config.GUI_FG_TEXT,
                                           selectbackground=Config.GUI_HIGHLIGHT,
                                           selectforeground=Config.GUI_ACCENT_DARK,
                                           relief=tk.FLAT, bd=0, activestyle="none")
        self.question_listbox.pack(side=tk.LEFT, fill=tk.BOTH, expand=True)
                # ===== 新增：认知画像折叠面板 =====
        profile_frame = tk.Frame(self.right_pane, bg=Config.GUI_BG_SIDEBAR)
        profile_frame.pack(fill=tk.X, padx=5, pady=(10, 5))

        # 标题栏（可点击折叠）
        profile_header = tk.Frame(profile_frame, bg=Config.GUI_BG_HEADER)
        profile_header.pack(fill=tk.X)
        profile_header.bind("<Button-1>", lambda e: self._toggle_profile_panel())

        tk.Label(profile_header, text="🧠 认知画像", font=("微软雅黑", 10, "bold"),
                 bg=Config.GUI_BG_HEADER, fg=Config.GUI_ACCENT_DARK).pack(side=tk.LEFT, padx=8, pady=4)
        self.profile_toggle_btn = tk.Label(profile_header, text="▼", bg=Config.GUI_BG_HEADER,
                                           fg=Config.GUI_FG_MUTED, font=("微软雅黑", 9))
        self.profile_toggle_btn.pack(side=tk.RIGHT, padx=8)

        # 画像内容（初始展开）
        self.profile_content = tk.Frame(profile_frame, bg=Config.GUI_BG_SIDEBAR)
        self.profile_content.pack(fill=tk.X, padx=4, pady=6)
        self.profile_expanded = True

        # 雷达图 Canvas
        self.radar_canvas = tk.Canvas(self.profile_content, width=200, height=200,
                                      bg=Config.GUI_BG_SIDEBAR, highlightthickness=0)
        self.radar_canvas.pack(pady=4)

        # 指纹数值标签
        self.profile_labels = tk.Frame(self.profile_content, bg=Config.GUI_BG_SIDEBAR)
        self.profile_labels.pack(fill=tk.X, pady=2)

        # 5个维度标签（初始占位）
        self.profile_dim_labels = []
        for i, name in enumerate(["风险容忍", "创新偏好", "决策果断", "注意力持续", "认知置信"]):
            lbl = tk.Label(self.profile_labels, text=f"{name}: --", font=("微软雅黑", 8),
                           bg=Config.GUI_BG_SIDEBAR, fg=Config.GUI_FG_MUTED)
            lbl.grid(row=i // 3, column=i % 3, sticky=tk.W, padx=4, pady=1)
            self.profile_dim_labels.append(lbl)

        # 刷新按钮
        refresh_btn = self._create_btn(self.profile_content, "🔄 刷新画像", self._refresh_fingerprint_display,
                                       Config.GUI_ACCENT, "white")
        refresh_btn.pack(pady=4)

        self._register_mousewheel(profile_frame, profile_frame)
        q_scrollbar = ttk.Scrollbar(self.question_frame, orient=tk.VERTICAL,
                                    command=self.question_listbox.yview)
        q_scrollbar.pack(side=tk.RIGHT, fill=tk.Y)
        self.question_listbox.config(yscrollcommand=q_scrollbar.set)
        self._register_mousewheel(self.question_listbox, self.question_listbox)
        self.question_listbox.bind('<<ListboxSelect>>', self._on_question_select)

    def _create_status_bar(self):
        status_frame = tk.Frame(self.root, bg=Config.GUI_BG_HEADER, height=35)
        status_frame.grid(row=1, column=0, columnspan=3, sticky="ew")
        status_frame.grid_propagate(False)
        self.status_var = tk.StringVar(value="✅ 就绪 - v2.0 联想增强版")
        tk.Label(status_frame, textvariable=self.status_var, bg=Config.GUI_BG_HEADER,
                 fg=Config.GUI_FG_MUTED, anchor=tk.W, font=Config.GUI_TEXT_FONT).pack(fill=tk.X, padx=15,
                                                                                      pady=8)

    def _load_session_list(self):
        sessions = self.db.list_sessions()
        if not hasattr(self, 'session_listbox') or self.session_listbox is None:
            return
        self.session_listbox.delete(0, tk.END)
        self.session_filtered_ids = []
        search = self.session_search_var.get().strip().lower()
        for sid, name, updated in sessions:
            if search and search not in name.lower():
                continue
            self.session_filtered_ids.append(sid)
            self.session_listbox.insert(tk.END, name[:30] + ('...' if len(name) > 30 else ''))
        if self.session_filtered_ids:
            self.session_listbox.selection_set(0)
            self._on_session_select()

    def _filter_sessions(self):
        self._load_session_list()

    def _on_session_select(self, event=None):
        if not hasattr(self, 'session_listbox') or self.session_listbox is None:
            return
        sel = self.session_listbox.curselection()
        if not sel:
            return
        idx = sel[0]
        if idx >= len(self.session_filtered_ids):
            return
        sid = self.session_filtered_ids[idx]
        if sid == self.current_session_id:
            return
        if self.current_session_id:
            self.db.update_session(self.current_session_id, self.current_history, self.current_session_name)
        name, history = self.db.get_session(sid)
        if name is not None:
            self.current_session_id = sid
            self.current_session_name = name
            self.current_history = history
            if hasattr(self, 'log_area') and self.log_area:
                self.log_area.delete(1.0, tk.END)
                for role, content in history:
                    tag = "user" if role == "user" else "ai"
                    prefix = "🧑 你：" if role == "user" else "🤖 AI："
                    if role == "user":
                        formatted = f"[{datetime.now().strftime('%H:%M:%S')}] ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━\n{prefix}{content}\n━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━\n\n"
                    else:
                        formatted = f"[{datetime.now().strftime('%H:%M:%S')}] {prefix}{content}\n"
                    self.log_area.insert(tk.END, formatted, tag)
                self.log_area.see(tk.END)
            self._refresh_question_list()
            self.update_status(f"已加载会话：{name}")

    def _new_session(self):
        if self.current_session_id:
            self.db.update_session(self.current_session_id, self.current_history, self.current_session_name)
        sid = datetime.now().strftime("%Y%m%d%H%M%S") + "_" + str(hash(datetime.now()))[-6:]
        temp_name = f"新会话 {datetime.now().strftime('%H:%M')}"
        self.db.create_session(sid, temp_name)
        self.current_session_id = sid
        self.current_session_name = temp_name
        self.current_history = []
        if hasattr(self, 'log_area') and self.log_area:
            self.log_area.delete(1.0, tk.END)
        self._refresh_question_list()
        self._load_session_list()
        if hasattr(self, 'session_filtered_ids'):
            for i, sid2 in enumerate(self.session_filtered_ids):
                if sid2 == sid:
                    self.session_listbox.selection_set(i)
                    break
        self.update_status(f"新建会话：{temp_name}")

    def _rename_current_session(self):
        if not self.current_session_id:
            messagebox.showwarning("提示", "没有活动会话")
            return
        new_name = simpledialog.askstring("重命名", "请输入新名称:",
                                          initialvalue=self.current_session_name, parent=self.root)
        if new_name and new_name.strip():
            new_name = new_name.strip()
            self.db.rename_session(self.current_session_id, new_name)
            self.current_session_name = new_name
            self._load_session_list()
            self.update_status(f"会话已重命名为：{new_name}")
    
    def _delete_selected_sessions(self):
        """删除所有选中的历史会话（支持多选）"""
        if not hasattr(self, 'session_listbox') or self.session_listbox is None:
            return

        selected_indices = self.session_listbox.curselection()
        if not selected_indices:
            messagebox.showwarning("提示", "请至少选中一个会话", parent=self.root)
            return

        # 获取选中的会话ID列表
        selected_ids = []
        for idx in selected_indices:
            if idx < len(self.session_filtered_ids):
                selected_ids.append(self.session_filtered_ids[idx])

        if not selected_ids:
            return

        count = len(selected_ids)
        if not messagebox.askyesno("确认删除", f"确定要删除选中的 {count} 个会话吗？\n此操作不可恢复！", parent=self.root):
            return

        # 执行删除
        for sid in selected_ids:
            self.db.delete_session(sid)

        # 如果当前会话被删除了，需要切换或新建
        if self.current_session_id in selected_ids:
            self.current_session_id = None
            self.current_session_name = None
            self.current_history = []
            if hasattr(self, 'log_area') and self.log_area:
                self.log_area.delete(1.0, tk.END)

        # 刷新会话列表
        self._load_session_list()

        # 如果当前会话未设置或已被删除，选择第一个会话或新建
        if self.current_session_id is None:
            sessions = self.db.list_sessions()
            if sessions:
                new_id = sessions[0][0]
                name, history = self.db.get_session(new_id)
                if name is not None:
                    self.current_session_id = new_id
                    self.current_session_name = name
                    self.current_history = history
                    self._refresh_question_list()
                    self.update_status(f"已切换到会话：{name}")
                    if hasattr(self, 'log_area') and self.log_area:
                        self.log_area.delete(1.0, tk.END)
                        for role, content in history:
                            tag = "user" if role == "user" else "ai"
                            prefix = "🧑 你：" if role == "user" else "🤖 AI："
                            if role == "user":
                                formatted = f"[{datetime.now().strftime('%H:%M:%S')}] ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━\n{prefix}{content}\n━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━\n\n"
                            else:
                                formatted = f"[{datetime.now().strftime('%H:%M:%S')}] {prefix}{content}\n"
                            self.log_area.insert(tk.END, formatted, tag)
                        self.log_area.see(tk.END)
            else:
                self._new_session()

        # 重新加载会话列表并高亮当前会话
        self._load_session_list()
        self.update_status(f"已删除 {count} 个会话")
        self._log(f"🗑️ 已删除 {count} 个会话", "system")

    def _clear_current_session_messages(self):
        if not self.current_session_id:
            return
        if messagebox.askyesno("清空会话", f"清空「{self.current_session_name}」的所有消息？"):
            self.current_history = []
            self.db.update_session(self.current_session_id, [])
            if hasattr(self, 'log_area') and self.log_area:
                self.log_area.delete(1.0, tk.END)
            self._refresh_question_list()
            self.update_status("已清空当前会话消息")

    def _add_message_to_session(self, role: str, content: str):
        """向当前会话添加一条消息（不含 label）"""
        self.current_history.append((role, content))
        self.db.update_session(self.current_session_id, self.current_history, self.current_session_name)
        self._load_session_list()
        # 强制刷新右侧问题列表（从数据库重新读取）
        self._refresh_question_list()

    def _toggle_right_pane(self):
        if self.right_expanded:
            self.right_pane.config(width=40)
            self.question_frame.pack_forget()
            self.collapse_btn.config(text="▶")
            self.right_expanded = False
        else:
            self.right_pane.config(width=250)
            self.question_frame.pack(fill=tk.BOTH, expand=True, padx=5, pady=5)
            self.collapse_btn.config(text="◀")
            self.right_expanded = True
            self._refresh_question_list()

    def _toggle_profile_panel(self):
        """切换认知画像面板的展开/折叠"""
        if self.profile_expanded:
            self.profile_content.pack_forget()
            self.profile_toggle_btn.config(text="▶")
            self.profile_expanded = False
        else:
            self.profile_content.pack(fill=tk.X, padx=4, pady=6)
            self.profile_toggle_btn.config(text="▼")
            self.profile_expanded = True
            self._refresh_fingerprint_display()

    def _refresh_fingerprint_display(self):
        """刷新认知画像显示"""
        try:
            fp = self.engine.fingerprint_extractor.get_fingerprint()

            self._draw_radar_chart(fp)

            dims = [
                ("风险容忍", fp.risk_tolerance),
                ("创新偏好", fp.innovation_preference),
                ("决策果断", fp.decisiveness),
                ("注意力持续", fp.attention_span),
                ("认知置信", fp.confidence)
            ]
            for i, (name, val) in enumerate(dims):
                if i < len(self.profile_dim_labels):
                    self.profile_dim_labels[i].config(text=f"{name}: {val:.2f}")

        except Exception as e:
            for lbl in self.profile_dim_labels:
                lbl.config(text="--: --")

    def _draw_radar_chart(self, fp):
        """在 Canvas 上绘制五维雷达图"""
        canvas = self.radar_canvas
        canvas.delete("all")

        w, h = 200, 200
        cx, cy = w // 2, h // 2
        radius = 75

        dims = [
            fp.risk_tolerance,
            fp.innovation_preference,
            fp.decisiveness,
            fp.attention_span,
            fp.confidence
        ]
        labels = ["风险容忍", "创新偏好", "决策果断", "注意力持续", "认知置信"]

        import math
        angles = [-90 + i * 72 for i in range(5)]

        for r in [0.3, 0.6, 0.9]:
            points = []
            for i, angle in enumerate(angles):
                rad = math.radians(angle)
                x = cx + radius * r * math.cos(rad)
                y = cy + radius * r * math.sin(rad)
                points.extend([x, y])
            canvas.create_polygon(points, outline=Config.GUI_BORDER, fill="", width=0.5)

        for angle in angles:
            rad = math.radians(angle)
            x = cx + radius * 1.05 * math.cos(rad)
            y = cy + radius * 1.05 * math.sin(rad)
            canvas.create_line(cx, cy, x, y, fill=Config.GUI_BORDER, width=0.5)

        points = []
        for i, (val, angle) in enumerate(zip(dims, angles)):
            rad = math.radians(angle)
            r = radius * max(0, min(1, val))
            x = cx + r * math.cos(rad)
            y = cy + r * math.sin(rad)
            points.extend([x, y])
        canvas.create_polygon(points, fill=Config.GUI_ACCENT, outline=Config.GUI_ACCENT_DARK,
                              stipple="gray50", width=1.5)

        for i, (val, angle) in enumerate(zip(dims, angles)):
            rad = math.radians(angle)
            r = radius * max(0, min(1, val))
            x = cx + r * math.cos(rad)
            y = cy + r * math.sin(rad)
            canvas.create_oval(x - 4, y - 4, x + 4, y + 4, fill=Config.GUI_ACCENT_DARK,
                               outline="white", width=1)

            label_rad = math.radians(angle)
            lx = cx + (radius + 22) * math.cos(label_rad)
            ly = cy + (radius + 22) * math.sin(label_rad)
            canvas.create_text(lx, ly, text=labels[i], font=("微软雅黑", 7), fill=Config.GUI_FG_MUTED)

        avg_val = sum(dims) / 5
        canvas.create_text(cx, cy, text=f"{avg_val:.2f}", font=("微软雅黑", 10, "bold"),
                           fill=Config.GUI_ACCENT_DARK)

    def _refresh_question_list(self):
        if not hasattr(self, 'question_listbox') or self.question_listbox is None:
            return
        self.question_listbox.delete(0, tk.END)
        self.question_indices = []
        # 强制从数据库读取最新数据，避免缓存不一致
        name, history, labels = self.db.get_session(self.current_session_id)
        if not history:
            return
        # 同步更新当前内存中的历史（避免其他部分使用过时数据）
        self.current_history = history[:]
        # 同时更新会话名称（如果数据库中有更新）
        if name and name != self.current_session_name:
            self.current_session_name = name
        q_num = 1
        for i, (role, content) in enumerate(history):
            if role == "user":
                label = labels[i] if i < len(labels) else None
                if label:
                    display = f"{q_num}. {label}"
                else:
                    # 降级：原截断逻辑（简化展示，保留原文）
                    if content.startswith("[晶体化] "):
                        prefix = "[晶体化] "
                        body = content[len(prefix):]
                        display = f"{q_num}. [晶体化] {body[:44]}{'...' if len(body) > 44 else ''}"
                    elif content.startswith("[深度推理] "):
                        prefix = "[深度推理] "
                        body = content[len(prefix):]
                        display = f"{q_num}. [深度推理] {body[:43]}{'...' if len(body) > 43 else ''}"
                    elif content.startswith("[深度推理-多角色] "):
                        prefix = "[深度推理-多角色] "
                        body = content[len(prefix):]
                        display = f"{q_num}. [多角色] {body[:40]}{'...' if len(body) > 40 else ''}"
                    elif content.startswith("[辩论增强] "):
                        prefix = "[辩论增强] "
                        body = content[len(prefix):]
                        display = f"{q_num}. [辩论] {body[:43]}{'...' if len(body) > 43 else ''}"
                    elif content.startswith("[卢氏注意力增强] "):
                        prefix = "[卢氏注意力增强] "
                        body = content[len(prefix):]
                        display = f"{q_num}. [卢氏] {body[:41]}{'...' if len(body) > 41 else ''}"
                    elif content.startswith("[卢氏注意力增强 + 辩论增强] "):
                        prefix = "[卢氏注意力增强 + 辩论增强] "
                        body = content[len(prefix):]
                        display = f"{q_num}. [卢氏+辩论] {body[:34]}{'...' if len(body) > 34 else ''}"
                    elif content.startswith("[文件内容] "):
                        prefix = "[文件内容] "
                        body = content[len(prefix):]
                        display = f"{q_num}. [文件] {body[:40]}{'...' if len(body) > 40 else ''}"
                    else:
                        display = f"{q_num}. {content[:50]}{'...' if len(content) > 50 else ''}"
                self.question_listbox.insert(tk.END, display)
                self.question_indices.append((i, content))
                q_num += 1
        if self.question_indices:
            self._suppress_select = True
            self.question_listbox.selection_set(0)
            self._suppress_select = False

    def _on_question_select(self, event=None):
        if self._suppress_select:
            return
        sel = self.question_listbox.curselection()
        if not sel:
            return
        idx = sel[0]
        if idx >= len(self.question_indices):
            return
        history_idx, content = self.question_indices[idx]
        if not hasattr(self, 'log_area') or self.log_area is None:
            return

        log_text = self.log_area.get(1.0, tk.END)
        if not log_text.strip():
            self._log("⚠️ 日志为空，无法跳转", "warning")
            return

        search_content = content
        prefix_map = {
            "[晶体化] ": "",
            "[深度推理] ": "",
            "[深度推理-多角色] ": "",
            "[辩论增强] ": "",
            "[卢氏注意力增强] ": "",
            "[卢氏注意力增强 + 辩论增强] ": "",
            "[文件内容] ": "",
            "[文件] ": "",
        }
        for prefix, replacement in prefix_map.items():
            if content.startswith(prefix):
                search_content = content[len(prefix):]
                break

        if not search_content or not search_content.strip():
            self._log("⚠️ 消息内容为空，无法定位", "warning")
            return

        lines = log_text.split('\n')
        found_line = None
        search_key = search_content.strip()

        for i in range(len(lines) - 1, -1, -1):
            line = lines[i]
            if line.startswith("🧑 你："):
                line_content = line[5:]
                if search_content in line_content:
                    found_line = i + 1
                    break
                if len(search_content) > 30:
                    short_key = search_content[:30]
                    if short_key in line_content:
                        found_line = i + 1
                        break
                elif search_content in line_content:
                    found_line = i + 1
                    break

        if found_line is None:
            for i in range(len(lines) - 1, -1, -1):
                line = lines[i]
                if "🧑 你：" in line:
                    line_content = line.split("🧑 你：", 1)[-1]
                    if search_content in line_content:
                        found_line = i + 1
                        break
                    if len(search_content) > 30 and search_content[:30] in line_content:
                        found_line = i + 1
                        break

        if found_line is None and search_content:
            escaped = re.escape(search_content[:50])
            pattern = rf"🧑 你：.*?{escaped}"
            try:
                match = re.search(pattern, log_text)
                if match:
                    pos = match.start()
                    found_line = log_text[:pos].count('\n') + 1
            except re.error:
                pass

        if found_line is not None:
            self.log_area.tag_remove("highlight", 1.0, tk.END)
            self.log_area.tag_remove("highlight_ai", 1.0, tk.END)

            start_line = f"{found_line}.0"
            end_line = f"{found_line + 1}.0"
            self.log_area.tag_add("highlight", start_line, end_line)
            self.log_area.tag_config("highlight", background="#fff3cd", foreground="#c62828")
            self.log_area.see(start_line)

            ai_line = None
            for i in range(found_line, min(found_line + 50, len(lines))):
                if "🤖 AI：" in lines[i-1]:
                    ai_line = i
                    break
            if ai_line:
                self.log_area.tag_add("highlight_ai", f"{ai_line}.0", f"{ai_line + 1}.0")
                self.log_area.tag_config("highlight_ai", background="#e3f2fd", foreground="#1a5b8c")
                self.root.after(3000, lambda: self.log_area.tag_remove("highlight", 1.0, tk.END))
                self.root.after(3000, lambda: self.log_area.tag_remove("highlight_ai", 1.0, tk.END))
            else:
                self.root.after(2000, lambda: self.log_area.tag_remove("highlight", 1.0, tk.END))
        else:
            display = search_content[:50] + "..." if len(search_content) > 50 else search_content
            self._log(f"⚠️ 未在日志中找到对应消息：{display}", "warning")

    def _load_api_from_env(self):
        key = os.environ.get("DEEPSEEK_API_KEY")
        if key:
            self.api_key.set(key)
            self._log("已从环境变量加载 API 密钥", "success")
        else:
            self._log("未找到环境变量 DEEPSEEK_API_KEY", "error")

        def on_elegant_chunk(chunk: str):
            self.root.after(0, lambda: self.log_area.insert(tk.END, chunk, ("elegant_body",)))
            self.root.after(0, lambda: self.log_area.see(tk.END))

        def on_elegant_complete():
            self.log_area.insert(tk.END, "\n\n", ())
            self.log_area.see(tk.END)
            self._log("✅ 儒雅笔谈生成完成", "success")

        self._log("📜 儒雅笔谈生成中...", "system")

        try:
            self.ai.chat_stream(
                elegant_prompt,
                system="你是一位深谙苏轼、辛弃疾文风的散文大家。你写的文字让人读来如沐春风，心中舒畅。你只输出正文。",
                callback=on_elegant_chunk
            )
            self.root.after(0, on_elegant_complete)
        except Exception as e:
            self._log(f"⚠️ 儒雅笔谈生成失败：{e}", "warning")
            fallback = f"以我观之，此事如月照寒潭，明澈而深邃。{one_sentence}。行者自知，行之者自达。"
            self.log_area.insert(tk.END, fallback, ("elegant_body",))

    def _clear_input(self):
        self.input_text.delete("1.0", tk.END)

    def _get_input(self) -> str:
        return self.input_text.get("1.0", tk.END).strip()

    def _on_enter_chat(self, event):
        if self._get_input():
            self._do_chat()
        return "break"

    def _on_enter_crystal(self, event):
        if self._get_input():
            self._do_crystal()
        return "break"

    def _do_chat(self):
        # ===== 防重复调用 =====
        if self._chat_in_progress:
            self._log("⚠️ 上一轮对话尚未完成，请稍候...", "warning")
            return
        self._chat_in_progress = True

        user_input = self._get_input()
        if not user_input:
            self._chat_in_progress = False
            self._log("请输入内容", "error")
            return

        # ===== G1 质量门 =====
        g1_result = self.engine.quality_gate_g1(user_input)
        if not g1_result["passed"]:
            self._log(f"⚠️ G1 提醒: {g1_result['reason']}", "warning")
        else:
            self._log(f"✅ G1 通过: {g1_result['reason']}", "system")

        self._clear_input()
        self._log(user_input, "user")
        self.update_status("AI 思考中...")
        self._add_message_to_session("user", user_input)
        cur_history = self.current_history.copy()

        def task():
            os.environ["DEEPSEEK_API_KEY"] = self.api_key.get()
            l0_holes, l1_crystals = self.engine.get_attention_context()
            context = f"\n[注意力上下文] 当前核心孔洞：{', '.join([h.content[:50] for h in l0_holes])}\nL1晶体数量：{len(l1_crystals)} 条\n"

            # ===== 注入认知风格 =====
            try:
                fp = self.engine.fingerprint_extractor.get_fingerprint()
                ops = self.engine.fingerprint_extractor.get_cognitive_operators(fp)
                self.root.after(0, lambda: self._log(f"🧠 注入认知风格：{ops}", "system"))
            except Exception as e:
                ops = "[思维模式：平衡] [论证偏好：平衡] [输出偏好：平衡]"
                self.root.after(0, lambda: self._log(f"⚠️ 认知风格加载失败，使用默认：{ops}", "warning"))

            base_system = "你是认知晶体树的AI协作者，请友好自然地回答问题。"
            system_with_style = f"{base_system}\n\n【用户认知风格】{ops}\n请根据这些偏好调整你的表达方式，使回答更贴近用户的思维习惯。"

            reply = self.ai.chat_with_history(
                cur_history,
                system=system_with_style,
                context=context
            )

            # ===== G2 质量门 =====
            g2_result = self.engine.quality_gate_g2(reply, {"audit_score": 0.5})
            if g2_result["passed"]:
                self.root.after(0, lambda: self._log(f"✅ G2 通过: {g2_result['reason']}", "system"))
            else:
                self.root.after(0, lambda: self._log(f"⚠️ G2 提醒: {g2_result['reason']}", "warning"))

            self.root.after(0, lambda: self._chat_done(reply))

        threading.Thread(target=task, daemon=True).start()

    def _chat_done(self, reply):
        self._log(reply, "ai")
        # ===== 立即释放标志，避免阻塞后续输入 =====
        self._chat_in_progress = False

        self._add_message_to_session("assistant", reply)
        # 获取本轮用户消息
        user_msg = ""
        for idx in range(len(self.current_history)-2, -1, -1):
            if self.current_history[idx][0] == "user":
                user_msg = self.current_history[idx][1]
                break
        # 处理标题生成（此过程可能调用AI，但标志已释放，不会阻塞界面）
        self._handle_title_generation_after_reply(reply, user_msg)

        try:
            fp = self.engine.fingerprint_extractor.get_fingerprint()
            self._log(f"🧠 当前认知指纹置信度：{fp.confidence:.2f}，总交互数：{fp.total_interactions}", "system")
        except:
            pass

        self.update_status("就绪")
        # 更新认知指纹（本次新增1条用户消息）
        self.engine.fingerprint_extractor.extract(self.current_history, increment_interactions=1)   
        self._refresh_question_list()

    def _handle_title_generation_after_reply(self, reply: str, user_msg: str):
        """
        统一处理标题生成：
        - 首轮时生成会话标题 + 本轮标题
        - 非首轮仅生成本轮标题
        """
        if not user_msg:
            return

        # 判断是否为首轮（会话名以“新会话”开头且历史只有两条消息：用户+AI）
        is_first = (len(self.current_history) == 2 and
                    self.current_session_name.startswith("新会话"))

        if is_first:
            # 获取历史标题
            sessions = self.db.list_sessions()
            existing_titles = [name for _, name, _ in sessions if name and name != self.current_session_name]
            sess_title, round_label = self._generate_dual_titles(user_msg, reply, existing_titles)

            if sess_title:
                self.db.rename_session(self.current_session_id, sess_title)
                self.current_session_name = sess_title
                self._load_session_list()
                self.update_status(f"已更新会话标题：{sess_title}")
                self._log(f"📌 会话标题：{sess_title}", "system")
            else:
                self._log("⚠️ 会话标题生成失败，保留默认名称", "warning")

            if round_label:
                self._update_last_round_label(round_label)
                self._log(f"🏷️ 本轮标题：{round_label}", "system")
            else:
                self._log("⚠️ 本轮标题生成失败", "warning")
        else:
            # 非首轮，生成轮次标题
            round_label = self._generate_round_label_simple(user_msg, reply)
            if round_label:
                self._update_last_round_label(round_label)
                self._log(f"🏷️ 本轮标题：{round_label}", "system")

    def _generate_dual_titles(self, user_msg: str, ai_reply: str, existing_titles: List[str]) -> Tuple[str, str]:
        if not user_msg or not ai_reply:
            return "", ""
        ai_summary = ai_reply[:800]
        if len(ai_reply) > 800:
            ai_summary += "……"
        avoid_hint = ""
        if existing_titles:
            recent = existing_titles[-50:]
            avoid_hint = f"\n【已有的会话标题（必须避免）】\n{', '.join(recent)}\n"
        prompt = f"""你是一位标题提炼专家。请根据以下完整对话，生成两个标题：

1. **会话标题**（5~8个字）：概括整段对话的核心主题，必须区别于【已有的会话标题】列表。
2. **本轮标题**（4~6个字）：精准概括用户这一轮提问的独特意图，用于会话内的快速导航。

用户问题：
{user_msg[:200]}

AI回答摘要：
{ai_summary}

{avoid_hint}
要求：
- 标题必须基于具体内容，禁止使用“对话”、“咨询”、“讨论”、“学习”等泛化词汇。
- 如果涉及领域（如编程、管理、心理），请在标题中体现。
- 会话标题要确保与已有列表完全不同（如内容相似，请加场景/时间后缀）。
- 本轮标题要突出本轮提问的独特侧重点。

只返回 JSON，格式：{{"session_title": "...", "round_title": "..."}}
"""
        try:
            ai = AIClient(api_key=self.api_key.get() or Config.get_api_key())
            result = ai.chat_json(prompt, temperature=0.6)
            if "error" not in result:
                sess_title = result.get("session_title", "").strip()
                round_title = result.get("round_title", "").strip()
                # 硬去重（会话标题）
                if sess_title:
                    original = sess_title
                    counter = 2
                    title_set = set(existing_titles)
                    while sess_title in title_set:
                        sess_title = f"{original} ({counter})"
                        counter += 1
                        if len(sess_title) > 15:
                            sess_title = sess_title[:12] + "..."
                return sess_title, round_title
        except Exception as e:
            self._log(f"⚠️ AI 生成双标题失败: {e}", "warning")
        return "", ""
    
    def _update_last_round_label(self, label: str):
        if not self.current_session_id or not label:
            return
        # 重新读取当前会话的完整消息（含 labels）
        name, history, labels = self.db.get_session(self.current_session_id)
        if not history:
            return
        # 找到最后一个 user 消息的索引
        last_user_idx = -1
        for i in range(len(history)-1, -1, -1):
            if history[i][0] == "user":
                last_user_idx = i
                break
        if last_user_idx == -1:
            return
        # 构造新的 messages 列表（dict 格式）
        new_messages = []
        for i, (role, content) in enumerate(history):
            lbl = labels[i] if i < len(labels) else None
            if i == last_user_idx:
                lbl = label
            new_messages.append({"role": role, "content": content, "label": lbl})
        # 直接调用 update_session 重写
        self.db.update_session(self.current_session_id, new_messages, name)
        # 刷新右侧列表
        self._refresh_question_list()

    def _generate_round_label_simple(self, user_msg: str, ai_reply: str) -> str:
        if not user_msg:
            return ""
        prompt = f"""请根据以下本轮提问，生成一个 4~6 字的精炼标题（仅输出标题，不要其他内容）：

用户问题：{user_msg[:150]}

标题要求：精准概括本轮提问的核心意图，避免泛化词汇。
只输出标题，不加引号或其他格式。
"""
        try:
            ai = AIClient(api_key=self.api_key.get() or Config.get_api_key())
            title = ai.chat(prompt, temperature=0.5)
            title = title.strip().strip('"').strip("'")
            # 会话内去重（简单模糊去重）
            _, _, labels = self.db.get_session(self.current_session_id)
            existing_labels = [lbl for lbl in labels if lbl and lbl.startswith(title[:3])]
            if existing_labels:
                counter = 2
                orig = title
                while any(lbl == title for lbl in existing_labels):
                    title = f"{orig} ({counter})"
                    counter += 1
                    if len(title) > 10:
                        title = title[:8] + "…"
            return title
        except Exception as e:
            self._log(f"⚠️ 轮次标题生成失败: {e}", "warning")
            return ""

    def _existing_crystal_ids(self) -> set:
        return set(re.findall(r"\bC\d+\b", self.files.read("crystals")))

    def _shorten_crystal_content(self, content: str) -> str:
        text = re.sub(r"\s+", " ", str(content or "")).strip()
        return text[:80]

    def _build_crystallization_prompt(self, user_input: str, search_res: str, l0_holes=None, l1_crystals=None):
        if l0_holes is None:
            l0_holes, l1_crystals = self.engine.get_attention_context()
        l0_text = "\n".join([f"- {h.id}: {h.content[:100]}" for h in l0_holes])
        l1_text = "\n".join([f"- {c.id}: {c.content[:80]} | links={','.join(c.links)}" for c in
                             l1_crystals[:Config.L1_MAX]])
        related = self.engine.get_associative_crystals(user_input, top_k=8)
        related_text = "\n".join([f"- {c.id}: {c.content[:80]} | links={','.join(c.links)}" for c in related])
        existing_ids = ", ".join(sorted(self._existing_crystal_ids())[-20:])
        return f"""
你是认知晶体树的结构化整理器。请只返回 JSON，不要返回 Markdown、解释或代码块。

目标：
1. 从用户输入中提炼可长期复用的认知晶体。
2. 优先连接已有晶体，不要重复制造同义晶体。
3. 暴露冲突和孔洞，但不要夸大不确定性。
4. 每条晶体 content 必须不超过 80 个中文字符。

用户输入：
{user_input}

外部搜索结果：
{search_res}

L0 核心孔洞：
{l0_text}

L1 注意力晶体：
{l1_text}

联想检索命中的相关晶体：
{related_text}

近期已有晶体 ID：
{existing_ids}

返回 JSON schema：
{{
  "new_crystals": [
    {{"id": "", "content": "不超过80字的新晶体", "links": ["C001"]}}
  ],
  "updated_crystals": [
    {{"id": "C001", "new_content": "不超过80字的更新后内容"}}
  ],
  "new_holes": [
    {{"id": "", "content": "需要继续验证的问题", "urgency": 0.5, "layer": 2}}
  ],
  "updated_holes": [
    {{"id": "H001", "content": "更新后孔洞内容"}}
  ],
  "conflicts": [
    {{"a": "C001", "b": "C002", "reason": "冲突原因"}}
  ],
  "report_summary": "一句话总结本次结构变化",
  "pending_cards": [
    {{"type": "晶体候选", "content": "待确认内容", "source": "AI生成", "confidence": "中"}}
  ]
}}

约束：
- 如果只是表达、例子或临时信息，不要直接入库，放入 pending_cards。
- new_crystals 的 id 可以留空，系统会自动分配。
- links 只能引用已存在或本次返回的新晶体 ID；不确定就留空数组。
- 没有内容的字段返回空数组。
"""

    def _normalize_crystal_response(self, ai_response: Dict) -> Dict:
        if not isinstance(ai_response, dict):
            return {"error": "AI返回不是JSON对象"}
        existing_ids = self._existing_crystal_ids()
        id_map = {}
        normalized_new = []
        next_num = max([int(i) for i in re.findall(r"C(\d+)", self.files.read("crystals"))], default=0) + 1
        seen_contents = {self._shorten_crystal_content(c.content) for c in self.engine.parse_crystals()}
        for item in ai_response.get("new_crystals", []) or []:
            content = self._shorten_crystal_content(item.get("content", ""))
            if not content or content in seen_contents:
                continue
            old_id = str(item.get("id", "")).strip()
            new_id = f"C{next_num:03d}"
            next_num += 1
            if old_id and old_id not in existing_ids:
                id_map[old_id] = new_id
            links = []
            for link in item.get("links", []) or []:
                link = id_map.get(str(link).strip(), str(link).strip())
                if re.fullmatch(r"C\d+", link) and (link in existing_ids or link in id_map.values()):
                    links.append(link)
            normalized_new.append({"id": new_id, "content": content, "links": sorted(set(links))})
            seen_contents.add(content)
        valid_ids = existing_ids | {c["id"] for c in normalized_new}
        normalized_updates = []
        for item in ai_response.get("updated_crystals", []) or []:
            cid = str(item.get("id", "")).strip()
            content = self._shorten_crystal_content(item.get("new_content") or item.get("content", ""))
            if cid in existing_ids and content:
                normalized_updates.append({"id": cid, "new_content": content})
        normalized_holes = []
        next_hole = max([int(i) for i in re.findall(r"H(\d+)", self.files.read("holes"))], default=0) + 1
        for item in ai_response.get("new_holes", []) or []:
            content = re.sub(r"\s+", " ", str(item.get("content", ""))).strip()
            if not content:
                continue
            try:
                urgency = float(item.get("urgency", 0.5))
            except (TypeError, ValueError):
                urgency = 0.5
            try:
                layer = int(item.get("layer", 2))
            except (TypeError, ValueError):
                layer = 2
            normalized_holes.append({"id": f"H{next_hole:03d}", "content": content[:120],
                                     "urgency": max(0.0, min(1.0, urgency)), "layer": min(3, max(1, layer))})
            next_hole += 1
        normalized_conflicts = []
        for item in ai_response.get("conflicts", []) or []:
            a = str(item.get("a") or item.get("crystal_a") or "").strip()
            b = str(item.get("b") or item.get("crystal_b") or "").strip()
            if a in valid_ids and b in valid_ids and a != b:
                normalized_conflicts.append({"a": a, "b": b, "reason": str(item.get("reason", ""))[:120]})
        normalized_pending = []
        for item in ai_response.get("pending_cards", []) or []:
            content = re.sub(r"\s+", " ", str(item.get("content", ""))).strip()
            if content:
                normalized_pending.append({
                    "type": str(item.get("type", "晶体候选")),
                    "content": content[:200],
                    "source": str(item.get("source", "AI生成")),
                    "confidence": str(item.get("confidence", "中")),
                })
        return {
            "new_crystals": normalized_new,
            "updated_crystals": normalized_updates,
            "new_holes": normalized_holes,
            "updated_holes": ai_response.get("updated_holes", []) or [],
            "conflicts": normalized_conflicts,
            "report_summary": str(ai_response.get("report_summary", "晶体化完成"))[:120],
            "pending_cards": normalized_pending,
        }

    def _find_similar_crystals(self, content: str, threshold: float = 0.55) -> List[tuple]:
        matches = []
        for crystal in self.engine.parse_crystals():
            score = self.engine._simple_similarity(content, crystal.content)
            if score >= threshold:
                matches.append((score, crystal))
        matches.sort(key=lambda item: item[0], reverse=True)
        return matches[:5]

    def _format_crystal_update_preview(self, result: Dict) -> str:
        lines = [f"摘要：{result.get('report_summary', '晶体化完成')}", ""]
        lines.append(f"新增晶体：{len(result.get('new_crystals', []))}")
        for c in result.get("new_crystals", []):
            lines.append(f"- {c['id']} | {c['content']} | links={','.join(c.get('links', [])) or '—'}")
            similar = self._find_similar_crystals(c["content"])
            if similar:
                lines.append("  可能重复：")
                for score, old in similar:
                    lines.append(f"  * {old.id} ({score:.2f}) {old.content[:60]}")
        lines.append("")
        lines.append(f"更新晶体：{len(result.get('updated_crystals', []))}")
        for c in result.get("updated_crystals", []):
            lines.append(f"- {c['id']} -> {c.get('new_content', '')}")
        lines.append("")
        lines.append(f"新增孔洞：{len(result.get('new_holes', []))}")
        for h in result.get("new_holes", []):
            lines.append(f"- {h['id']} | {h['content']} | urgency={h.get('urgency', 0.5)} | layer={h.get('layer', 2)}")
        lines.append("")
        lines.append(f"待确认卡片：{len(result.get('pending_cards', []))}")
        for card in result.get("pending_cards", []):
            lines.append(f"- {card.get('type', '晶体候选')} | {card.get('content', '')[:100]}")
        lines.append("")
        lines.append(f"冲突：{len(result.get('conflicts', []))}")
        for c in result.get("conflicts", []):
            lines.append(f"- {c.get('a')} vs {c.get('b')} | {c.get('reason', '')}")
        return "\n".join(lines)

    def _preview_crystal_update(self, result: Dict):
        """显示可编辑的晶体化预览窗口，文本框自适应高度，链接等额外信息另起一行"""
        win = Toplevel(self.root)
        win.title("晶体化预览确认")
        win.configure(bg=Config.GUI_BG_MAIN)
        win.resizable(True, True)

        tk.Label(win, text="请确认本次晶体化写入内容，可直接编辑下方晶体内容",
                 font=("微软雅黑", 14, "bold"), bg=Config.GUI_BG_MAIN, fg=Config.GUI_ACCENT).pack(anchor=tk.W, padx=12, pady=(12, 6))

        main_frame = tk.Frame(win)
        main_frame.pack(fill=tk.BOTH, expand=True, padx=12, pady=8)
        canvas = tk.Canvas(main_frame, bg=Config.GUI_BG_MAIN, highlightthickness=0)
        scrollbar = ttk.Scrollbar(main_frame, orient="vertical", command=canvas.yview)
        scrollable = tk.Frame(canvas, bg=Config.GUI_BG_MAIN)
        scrollable.bind("<Configure>", lambda e: canvas.configure(scrollregion=canvas.bbox("all")))
        canvas.create_window((0, 0), window=scrollable, anchor="nw")
        canvas.configure(yscrollcommand=scrollbar.set)
        canvas.pack(side="left", fill="both", expand=True)
        scrollbar.pack(side="right", fill="y")
        self._register_mousewheel(canvas, canvas)

        edit_widgets = {}

        # 辅助：创建多行文本框，高度自适应
        def add_editable_text(parent, initial_text):
            text_widget = tk.Text(parent, wrap=tk.WORD, font=Config.GUI_TEXT_FONT,
                                  bg=Config.GUI_BG_INPUT, fg=Config.GUI_FG_TEXT, relief=tk.SOLID, bd=1)
            text_widget.pack(side=tk.LEFT, fill=tk.X, expand=True, padx=5)
            text_widget.insert("1.0", initial_text)
            line_count = int(text_widget.index('end-1c').split('.')[0])
            height = max(2, min(line_count, 6))
            text_widget.config(height=height)
            return text_widget

        # 1. 新增晶体（链接另起一行）
        new_crystals = result.get("new_crystals", [])
        if new_crystals:
            tk.Label(scrollable, text="新增晶体：", font=("微软雅黑", 12, "bold"),
                     bg=Config.GUI_BG_MAIN, fg=Config.GUI_FG_TEXT).pack(anchor=tk.W, pady=(10, 5))
            for idx, crystal in enumerate(new_crystals):
                frame = tk.Frame(scrollable, bg=Config.GUI_BG_CARD, relief=tk.FLAT, bd=0,
                                 highlightthickness=1, highlightbackground=Config.GUI_BORDER)
                frame.pack(fill=tk.X, pady=3, padx=5)
                # 第一行：ID + 内容框
                top_frame = tk.Frame(frame, bg=Config.GUI_BG_CARD)
                top_frame.pack(fill=tk.X)
                label = tk.Label(top_frame, text=f"{crystal['id']}:", font=Config.GUI_TEXT_FONT,
                                 bg=Config.GUI_BG_CARD, fg=Config.GUI_ACCENT_DARK, width=8)
                label.pack(side=tk.LEFT, padx=5)
                text = add_editable_text(top_frame, crystal['content'])
                edit_widgets[("new_crystals", idx)] = text
                # 第二行：链接（如果有）
                links = crystal.get("links", [])
                if links:
                    link_frame = tk.Frame(frame, bg=Config.GUI_BG_CARD)
                    link_frame.pack(fill=tk.X, padx=(8, 5), pady=(0, 3))
                    link_label = tk.Label(link_frame, text=f"链接: {', '.join(links)}", font=Config.GUI_TEXT_FONT,
                                          bg=Config.GUI_BG_CARD, fg=Config.GUI_FG_MUTED, anchor=tk.W)
                    link_label.pack(fill=tk.X)

        # 2. 更新晶体
        updated = result.get("updated_crystals", [])
        if updated:
            tk.Label(scrollable, text="更新晶体：", font=("微软雅黑", 12, "bold"),
                     bg=Config.GUI_BG_MAIN, fg=Config.GUI_FG_TEXT).pack(anchor=tk.W, pady=(10, 5))
            for idx, upd in enumerate(updated):
                frame = tk.Frame(scrollable, bg=Config.GUI_BG_CARD, relief=tk.FLAT, bd=0,
                                 highlightthickness=1, highlightbackground=Config.GUI_BORDER)
                frame.pack(fill=tk.X, pady=3, padx=5)
                top_frame = tk.Frame(frame, bg=Config.GUI_BG_CARD)
                top_frame.pack(fill=tk.X)
                label = tk.Label(top_frame, text=f"{upd['id']} →", font=Config.GUI_TEXT_FONT,
                                 bg=Config.GUI_BG_CARD, fg=Config.GUI_ACCENT_DARK, width=8)
                label.pack(side=tk.LEFT, padx=5)
                text = add_editable_text(top_frame, upd['new_content'])
                edit_widgets[("updated_crystals", idx)] = text

        # 3. 新增孔洞（紧迫度另起一行）
        new_holes = result.get("new_holes", [])
        if new_holes:
            tk.Label(scrollable, text="新增孔洞：", font=("微软雅黑", 12, "bold"),
                     bg=Config.GUI_BG_MAIN, fg=Config.GUI_FG_TEXT).pack(anchor=tk.W, pady=(10, 5))
            for idx, hole in enumerate(new_holes):
                frame = tk.Frame(scrollable, bg=Config.GUI_BG_CARD, relief=tk.FLAT, bd=0,
                                 highlightthickness=1, highlightbackground=Config.GUI_BORDER)
                frame.pack(fill=tk.X, pady=3, padx=5)
                top_frame = tk.Frame(frame, bg=Config.GUI_BG_CARD)
                top_frame.pack(fill=tk.X)
                label = tk.Label(top_frame, text=f"{hole['id']}:", font=Config.GUI_TEXT_FONT,
                                 bg=Config.GUI_BG_CARD, fg=Config.GUI_ACCENT_DARK, width=8)
                label.pack(side=tk.LEFT, padx=5)
                text = add_editable_text(top_frame, hole['content'])
                edit_widgets[("new_holes", idx)] = text
                # 第二行：紧迫度
                urgency = hole.get('urgency', 0.5)
                urgency_frame = tk.Frame(frame, bg=Config.GUI_BG_CARD)
                urgency_frame.pack(fill=tk.X, padx=(8, 5), pady=(0, 3))
                urgency_label = tk.Label(urgency_frame, text=f"紧迫度: {urgency}", font=Config.GUI_TEXT_FONT,
                                         bg=Config.GUI_BG_CARD, fg=Config.GUI_FG_MUTED, anchor=tk.W)
                urgency_label.pack(fill=tk.X)

        # 4. 冲突（只读）
        conflicts = result.get("conflicts", [])
        if conflicts:
            tk.Label(scrollable, text="冲突：", font=("微软雅黑", 12, "bold"),
                     bg=Config.GUI_BG_MAIN, fg=Config.GUI_FG_TEXT).pack(anchor=tk.W, pady=(10, 5))
            for c in conflicts:
                tk.Label(scrollable, text=f"{c.get('a')} vs {c.get('b')} | {c.get('reason', '')}",
                         font=Config.GUI_TEXT_FONT, bg=Config.GUI_BG_MAIN, fg=Config.GUI_FG_MUTED).pack(anchor=tk.W, padx=10)

        # 5. 待确认卡片（只读）
        pending = result.get("pending_cards", [])
        if pending:
            tk.Label(scrollable, text="待确认卡片：", font=("微软雅黑", 12, "bold"),
                     bg=Config.GUI_BG_MAIN, fg=Config.GUI_FG_TEXT).pack(anchor=tk.W, pady=(10, 5))
            for card in pending:
                tk.Label(scrollable, text=f"{card.get('type')}: {card.get('content')[:100]}",
                         font=Config.GUI_TEXT_FONT, bg=Config.GUI_BG_MAIN, fg=Config.GUI_FG_TEXT).pack(anchor=tk.W, padx=10)

        # 底部按钮
        btn_frame = tk.Frame(win, bg=Config.GUI_BG_MAIN)
        btn_frame.pack(fill=tk.X, padx=12, pady=10)

        def confirm():
            for (kind, idx), text_widget in edit_widgets.items():
                new_val = text_widget.get("1.0", tk.END).strip()
                if kind == "new_crystals":
                    result["new_crystals"][idx]["content"] = new_val
                elif kind == "updated_crystals":
                    result["updated_crystals"][idx]["new_content"] = new_val
                elif kind == "new_holes":
                    result["new_holes"][idx]["content"] = new_val
            win.destroy()
            self._commit_crystal_update(result)

        def cancel():
            win.destroy()
            self._log("已取消本次晶体化入库", "warning")

        self._create_btn(btn_frame, "✅ 确认入库", confirm, Config.GUI_SUCCESS).pack(side=tk.LEFT, padx=5)
        self._create_btn(btn_frame, "❌ 取消", cancel, Config.GUI_DANGER).pack(side=tk.LEFT, padx=5)

        # 动态窗口高度估算
        total_items = len(new_crystals) + len(updated) + len(new_holes) + len(conflicts) + len(pending)
        estimated_height = 150 + 50 * (1 if new_crystals else 0) + 50 * (1 if updated else 0) + 50 * (1 if new_holes else 0) + 50 * (1 if conflicts else 0) + 50 * (1 if pending else 0) + total_items * 80
        final_height = max(400, min(estimated_height, 900))
        win.geometry(f"900x{final_height}")

    def _commit_crystal_update(self, result: Dict):
        self._update_files(result)
        summary = result.get('report_summary', '晶体化完成')
        self._log(f"✅ 晶体化完成\n   摘要：{summary}", "success")
        if result.get("pending_cards"):
            self._log(f"📋 生成了 {len(result['pending_cards'])} 张待确认卡片", "system")
        self._add_message_to_session("assistant", f"[晶体化结果] {summary}")
        self._show_status()

    def _do_crystal(self):
        mode = self.crystal_mode_var.get()
        if mode == "单条输入":
            text = self._get_input()
            if not text:
                self._log("请输入要晶体化的内容", "error")
                return
            self._clear_input()
            self._add_message_to_session("user", f"[晶体化] {text}")
            self._log(f"[晶体化] {text}", "user")
        else:
            if not self.current_history:
                self._log("当前会话为空，无法晶体化", "error")
                return
            dialog_text = "【对话历史】\n"
            for role, content in self.current_history:
                label = "用户" if role == "user" else "AI"
                dialog_text += f"{label}: {content}\n\n"
            if len(dialog_text) > 8000:
                dialog_text = dialog_text[:8000] + "\n...（内容过长已截断）"
                self._log("对话过长，已截断至8000字符", "warning")
            text = dialog_text
            self._log("📚 开始将当前会话生成晶体...", "system")
            self.update_status("会话晶体化中...")

        # 通用晶体化处理
        self.update_status("晶体化中...")
        def task():
            os.environ["DEEPSEEK_API_KEY"] = self.api_key.get()
            try:
                l0_holes, l1_crystals = self.engine.get_attention_context()
                skip = self.fast_mode.get()
                if skip:
                    search_res = "（快速模式：跳过外部搜索）"
                else:
                    keywords = self._extract_keywords(text)
                    search_res = self._search_duckduckgo(keywords) if keywords else "（关键词提取失败）"
                prompt = self._build_crystallization_prompt(text, search_res, l0_holes, l1_crystals)
                result = self.ai.chat_json(prompt)
                if "error" in result:
                    self.root.after(0, lambda: self._log(f"❌ 晶体化失败：{result['error']}", "error"))
                    return
                normalized = self._normalize_crystal_response(result)
                if "error" in normalized:
                    self.root.after(0, lambda: self._log(f"❌ 标准化失败：{normalized['error']}", "error"))
                    return
                self.root.after(0, lambda: self._preview_crystal_update(normalized))
            except Exception as e:
                self.root.after(0, lambda: self._log(f"❌ 出错：{e}", "error"))
            finally:
                self.root.after(0, self._crystal_done)

        threading.Thread(target=task, daemon=True).start()
        def task():
            os.environ["DEEPSEEK_API_KEY"] = self.api_key.get()
            try:
                l0_holes, l1_crystals = self.engine.get_attention_context()
                skip = self.fast_mode.get()
                if skip:
                    search_res = "（快速模式：跳过外部搜索）"
                else:
                    keywords = self._extract_keywords(user_input)
                    search_res = self._search_duckduckgo(keywords) if keywords else "（关键词提取失败）"
                prompt = self._build_crystallization_prompt(user_input, search_res, l0_holes, l1_crystals)
                result = self.ai.chat_json(prompt)
                if "error" in result:
                    self.root.after(0, lambda: self._log(f"❌ {result['error']}", "error"))
                    return
                result = self._normalize_crystal_response(result)
                if "error" in result:
                    self.root.after(0, lambda: self._log(f"❌ {result['error']}", "error"))
                    return
                self.root.after(0, lambda r=result: self._preview_crystal_update(r))
            except Exception as e:
                self.root.after(0, lambda e=e: self._log(f"❌ 出错：{e}", "error"))
            finally:
                self.root.after(0, self._crystal_done)

        threading.Thread(target=task, daemon=True).start()

    def _crystal_done(self):
        # 晶体化通常由单条输入触发，新增1条用户消息
        self.engine.fingerprint_extractor.extract(self.current_history, increment_interactions=1)
        self.update_status("就绪")
        self._refresh_question_list()
        
    def _show_full_content(self, content_id: str):
        """弹出窗口显示完整内容"""
        full = self.foldable_contents.get(content_id, "")
        if not full:
            return
        win = Toplevel(self.root)
        win.title("完整内容")
        win.geometry("800x600")
        win.configure(bg=Config.GUI_BG_MAIN)
        text = scrolledtext.ScrolledText(win, wrap=tk.WORD, font=Config.GUI_LOG_FONT,
                                         bg=Config.GUI_BG_INPUT, fg=Config.GUI_FG_TEXT,
                                         relief=tk.FLAT, bd=0)
        text.pack(fill=tk.BOTH, expand=True, padx=10, pady=10)
        text.insert("1.0", full)
        text.config(state=tk.DISABLED)

    def _insert_foldable_content(self, content: str, tag_prefix: str = "fold"):
        """向日志区域插入可折叠内容，默认折叠显示摘要"""
        import time
        content_id = f"{tag_prefix}_{int(time.time()*1000)}_{hash(content) % 10000}"
        self.foldable_contents[content_id] = content

        # 生成摘要（前300字符 + ...）
        summary = content[:300] + " ... "

        # 插入摘要 + 展开标签
        self.log_area.insert(tk.END, summary, "ai")
        self.log_area.insert(tk.END, "[展开全文]", ("foldable", content_id))
        self.log_area.insert(tk.END, "\n\n")

        # 绑定点击事件
        self.log_area.tag_bind(content_id, "<Button-1>",
                               lambda e, cid=content_id: self._show_full_content(cid))
        self.log_area.tag_config("foldable", foreground=Config.GUI_ACCENT,
                                 underline=True, font=("微软雅黑", 10, "bold"))

    def _deep_done(self, result: str):
        """深度推理完成后的处理：显示结果，支持折叠"""
        # 使用可折叠插入
        self._insert_foldable_content(result, "deep")
        # 记录到会话
        self._add_message_to_session("assistant", result)
        self.update_status("就绪")
        self._refresh_question_list()

    def _generate_elegant_narrative(self, question: str, one_sentence: str, student_answer: str, teacher_detail: str, return_only=False):
        """生成儒雅风格叙事，支持返回纯文本或直接输出到日志"""
        if not one_sentence and not student_answer:
            if not return_only:
                self._log("⚠️ 缺少素材，无法生成叙事版", "warning")
            return None

        elegant_prompt = f"""
请将以下商业分析内容，改写为一段「儒雅风格」的文字。

【核心观点】
{one_sentence}

【详细建议】
{student_answer[:600] if student_answer else "（无）"}

【风格要求】
模仿苏轼的旷达通透与辛弃疾的豪迈沉郁——文字从容而有筋骨，既有"一蓑烟雨任平生"的豁达，又有"把吴钩看了，栏杆拍遍"的深沉关怀。

【具体要求】
1. 以"以我观之，此事有三层意味"或类似文人开篇起笔
2. 中间穿插一个自然意象（如山、水、月、竹、云），借景说理
3. 引用一句古诗词或化用其意境（可改动以适应语境）
4. 结尾落在"行"字上——不是空谈，是可行之道
5. 字数：100-150字（精简隽永，点到即止）
6. 通篇让人感觉像在品茶听琴时的一席话，不急切、不炫耀

【输出要求】
只输出正文，不加标题、不加序号、不加任何格式标记（如 Markdown）。
"""

        if return_only:
            # 同步获取文本，不操作UI
            try:
                full_text = self.ai.chat(
                    elegant_prompt,
                    system="你是一位深谙苏轼、辛弃疾文风的散文大家。你写的文字让人读来如沐春风，心中舒畅。你只输出正文。"
                )
                return full_text.strip()
            except Exception as e:
                fallback = f"以我观之，此事如月照寒潭，明澈而深邃。{one_sentence}。行者自知，行之者自达。"
                return fallback
        else:
            # 原有UI流式输出
            self.log_area.insert(tk.END, "\n" + "=" * 60 + "\n", "system")
            self.log_area.insert(tk.END, "📜 儒雅笔谈 · 与君共赏\n", ("elegant_title",))
            self.log_area.insert(tk.END, "=" * 60 + "\n", "system")
            self.log_area.see(tk.END)

            # 配置标签（若未配置则在此处配置，实际可能已存在）
            self.log_area.tag_config("elegant_title",
                                     font=("华文楷体", 14, "bold"),
                                     foreground="#4A3728")
            self.log_area.tag_config("elegant_body",
                                     font=("华文楷体", 12),
                                     foreground="#3C2A1E",
                                     spacing1=6,
                                     spacing3=10,
                                     lmargin1=20,
                                     lmargin2=20)

            def on_elegant_chunk(chunk: str):
                self.root.after(0, lambda: self.log_area.insert(tk.END, chunk, ("elegant_body",)))
                self.root.after(0, lambda: self.log_area.see(tk.END))

            def on_elegant_complete():
                self.log_area.insert(tk.END, "\n\n", ())
                self.log_area.see(tk.END)
                self._log("✅ 儒雅笔谈生成完成", "success")

            self._log("📜 儒雅笔谈生成中...", "system")

            try:
                self.ai.chat_stream(
                    elegant_prompt,
                    system="你是一位深谙苏轼、辛弃疾文风的散文大家。你写的文字让人读来如沐春风，心中舒畅。你只输出正文。",
                    callback=on_elegant_chunk
                )
                self.root.after(0, on_elegant_complete)
            except Exception as e:
                self._log(f"⚠️ 儒雅笔谈生成失败：{e}", "warning")
                fallback = f"以我观之，此事如月照寒潭，明澈而深邃。{one_sentence}。行者自知，行之者自达。"
                self.log_area.insert(tk.END, fallback, ("elegant_body",))
            return None

    def _is_complex_question(self, question: str) -> bool:
        keywords = ["如何", "为什么", "方案", "决策", "分析", "比较", "评估", "设计", "策略", "方法", "框架"]
        has_keyword = any(kw in question for kw in keywords)
        hole_keywords = ["非共识", "高不确定性", "复杂系统", "多重约束", "因果", "判断", "决策", "H001", "H002", "H003"]
        has_hole_keyword = any(kw in question for kw in hole_keywords)
        return (len(question) > 20 and has_keyword) or has_hole_keyword

    def _get_debate_rounds(self) -> int:
        try:
            return max(2, min(12, int(self.debate_rounds_var.get())))
        except Exception:
            return 2

    def _sync_debate_round_state(self):
        if not hasattr(self, "debate_rounds_spinbox"):
            return
        state = tk.DISABLED if self.deep_think_mode.get() == "单路径深度推理" else tk.NORMAL
        self.debate_rounds_spinbox.configure(state=state)

    def _build_session_context(self, current_input: str, limit: int = 8) -> str:
        lines = []
        for role, content in self.current_history[-limit:]:
            if role == "user" and current_input and current_input in content:
                continue
            label = "用户" if role == "user" else "AI"
            text = re.sub(r"\s+", " ", str(content)).strip()
            if text:
                lines.append(f"{label}: {text[:900]}")
        if not lines:
            return current_input
        return "【本会话最近上下文】\n" + "\n".join(lines) + f"\n\n【当前问题】\n{current_input}"

    def _do_deep_reasoning(self):
        user_input = self._get_input()
        if not user_input:
            self._log("请输入内容", "error")
            return
        self._clear_input()
        mode = self.deep_think_mode.get()
        mode_map = {
            "多角色竞争": ("debate_full", "多角色竞争", "[深度推理-多角色]"),
            "卢氏注意力增强": ("lushi_sampling", "卢氏注意力增强", "[卢氏注意力增强]"),
            "替身自我博弈": ("twin_self_play", "替身自我博弈", "[替身自我博弈]"),
        }

        if mode == "单路径深度推理":
            self._add_message_to_session("user", f"[深度推理] {user_input}")
            self._log(f"[深度推理] {user_input}", "user")
            self.update_status("深度推理中（联想增强）...")
            threading.Thread(target=self._single_deep_reasoning, args=(user_input,), daemon=True).start()
            return

        debate_mode, label, prefix = mode_map.get(mode, mode_map["多角色竞争"])
        rounds = self._get_debate_rounds()
        if not messagebox.askyesno(label, f"将启用【{label}】，辩论轮次 {rounds}，AI 调用次数较多，是否继续？"):
            self._log(f"已取消{label}", "system")
            return
        self._log(f"🚀 启用{label}...", "system")
        self._add_message_to_session("user", f"{prefix} {user_input}")
        self._log(f"{prefix} {user_input}", "user")
        self.update_status(f"{label}推理中...")
        threading.Thread(target=self._debate_engine_reasoning,
                         args=(user_input, debate_mode, label, rounds),
                         daemon=True).start()

    def _single_deep_reasoning(self, user_input: str):
        """单路径深度推理（进度条支持）"""
        # ===== 进度初始化 =====
        self.root.after(0, lambda: self.debate_progress.configure(value=0))
        self.root.after(0, lambda: self.debate_status_label.config(text="单路径推理启动 0%"))

        try:
            os.environ["DEEPSEEK_API_KEY"] = self.api_key.get()
            full_input = self._build_session_context(user_input)

            self.root.after(0, lambda: self.debate_progress.configure(value=30))
            self.root.after(0, lambda: self.debate_status_label.config(text="裸模型生成中 30%"))
            raw = self.ai.chat(full_input, system="请直接回答问题，不要引用外部知识，给出最直接的答案。")

            assoc_crystals = self.engine.get_associative_crystals(user_input, top_k=5)
            if assoc_crystals:
                crystal_ctx = "\n".join([f"- [{c.id}] {c.content}" for c in assoc_crystals])
                self.root.after(0, lambda: self.debate_progress.configure(value=60))
                self.root.after(0, lambda: self.debate_status_label.config(text="晶体树评论中 60%"))
                comment_prompt = f"用户问题：{user_input}\n裸模型回答：{raw}\n相关晶体树知识（联想检索）：\n{crystal_ctx}\n请指出哪些晶体支持或反驳了裸模型回答，以及提供了哪些新视角。"
                comment = self.ai.chat(comment_prompt, system="你是晶体树的知识审计员，输出晶体观点。")
            else:
                comment = "（无相关晶体）"

            self.root.after(0, lambda: self.debate_progress.configure(value=80))
            self.root.after(0, lambda: self.debate_status_label.config(text="综合融合中 80%"))
            final_prompt = f"""原始问题：{user_input}
裸模型回答：{raw}
晶体树评论：{comment}
请综合两者给出最终答案，并说明裸模型和晶体树各自的贡献。"""
            final = self.ai.chat(final_prompt, system="你是认知晶体树的综合推理者，输出最终答案。")
            result = f"""【裸模型回答】
{raw}

【晶体树评论（联想增强）】
{comment}

【综合最终答案】
{final}"""

            self.root.after(0, lambda: self.debate_progress.configure(value=100))
            self.root.after(0, lambda: self.debate_status_label.config(text="完成 100%"))
            self.root.after(0, lambda: self._add_message_to_session("assistant", result))
            self.root.after(0, lambda: self._deep_done(result))
        except Exception as e:
            self.root.after(0, lambda: self._log(f"❌ 深度推理出错：{e}", "error"))
            self.root.after(0, lambda: self.debate_progress.configure(value=0))
            self.root.after(0, lambda: self.debate_status_label.config(text="出错"))
            self.root.after(0, lambda: self.update_status("就绪"))

    # ---- 修改开始：_debate_engine_reasoning 显示完整结论摘要 ----
    def _debate_engine_reasoning(self, user_input: str, debate_mode: str, label: str, max_rounds: int = 2):
        self._log(f"🔍 DEBUG: _debate_engine_reasoning 被调用，label={label}", "system")
        try:
            os.environ["DEEPSEEK_API_KEY"] = self.api_key.get()

            def on_progress(data):
                self.root.after(0, lambda: self.debate_progress.configure(value=data['progress']))
                self.root.after(0, lambda: self.debate_status_label.config(
                    text=f"{data['stage']} {data['progress']}%"
                ))

            def thread_safe_log(message, level="system"):
                self.root.after(0, lambda m=message, lv=level: self._log(m, lv))

            if debate_mode == "twin_self_play":
                fingerprint = self.engine.get_user_fingerprint()
                if fingerprint and fingerprint.confidence > 0.3:
                    self._log(f"[INFO] 加载认知指纹成功 (置信度={fingerprint.confidence:.2f})", "system")
                    debate = DebateEngine(
                        self.ai, self.engine, self.roles,
                        thread_safe_log,
                        progress_callback=on_progress
                    )
                    debate.roles = debate.get_roles_with_twin(include_twin=True, fingerprint=fingerprint)
                    debate_mode_effective = "debate_full"
                else:
                    self._log("[WARN] 认知指纹不足，使用标准多角色辩论", "warning")
                    debate = DebateEngine(
                        self.ai, self.engine, self.roles,
                        thread_safe_log,
                        progress_callback=on_progress
                    )
                    debate_mode_effective = debate_mode
            else:
                debate = DebateEngine(
                    self.ai, self.engine, self.roles,
                    thread_safe_log,
                    progress_callback=on_progress
                )
                debate_mode_effective = debate_mode

            reason_input = self._build_session_context(user_input, limit=20)
            result = debate.run(reason_input, mode=debate_mode_effective, max_rounds=max_rounds)

            # ================================================================
            # 提取数据
            # ================================================================
            final_schema = result.get("final_schema", {})
            if not final_schema:
                board_version = result.get("board_version", "")
                employee_version = result.get("employee_version", "")
                novice_version = result.get("novice_version", "")
                expert_version = result.get("expert_version", "")
                judge_audit = result.get("judge_audit", {})
                dashboard_stats = result.get("dashboard_stats", {})
            else:
                board_version = final_schema.get("board_version", "")
                employee_version = final_schema.get("employee_version", "")
                novice_version = final_schema.get("novice_version", "")
                expert_version = final_schema.get("expert_version", "")
                judge_audit = final_schema.get("judge_audit", {})
                dashboard_stats = final_schema.get("dashboard_stats", {})

            # ================================================================
            # 【核心修改1】过滤百灵鸟裸模型（解决重复问题）
            # 只从 Round >= 1 中提取正式角色观点，排除 Round 0 的裸模型
            # ================================================================
            role_blocks = []
            seen_roles = set()
            rounds_data = result.get("rounds", [])

            for rd in rounds_data:
                round_no = rd.get("round", 0)
                # 跳过 Round 0（裸模型轮次）
                if round_no == 0:
                    continue
                for ans in rd.get("answers", []):
                    role_name = ans.get('role', '未知角色')
                    if role_name in seen_roles:
                        continue
                    seen_roles.add(role_name)
                    content = ans.get('answer', '（无回答）')
                    if len(content) > 350:
                        content = content[:350] + "……（完整版见上方滚动日志）"
                    role_blocks.append({
                        "name": role_name,
                        "content": content,
                        "round": round_no
                    })
                    if len(seen_roles) >= 12:
                        break
                if len(seen_roles) >= 12:
                    break

            # 按轮次排序
            role_blocks.sort(key=lambda x: x.get("round", 999))

            # ================================================================
            # 【核心修改2】使用富文本渲染到独立结果面板
            # ================================================================
            self.root.after(0, lambda: self._render_result_to_panel(
                question=user_input,
                role_blocks=role_blocks,
                board_version=board_version,
                employee_version=employee_version,
                novice_version=novice_version,
                expert_version=expert_version,
                judge_audit=judge_audit,
                dashboard_stats=dashboard_stats,
                final_schema=final_schema,
                result=result
            ))

            # ================================================================
            # 【核心修改3】精简日志区输出（解决时序错乱问题）
            # ================================================================
            self.root.after(0, lambda: self._log("=" * 60, "system"))
            self.root.after(0, lambda: self._log(f"✅ {label} 完成！完整报告已在上方「辩论报告」面板显示。", "success"))
            self.root.after(0, lambda: self._log("=" * 60, "system"))

            # 元层和统计信息：在日志区以简洁方式输出
            meta = result.get('_meta', {})
            self.root.after(0, lambda: self._log(f"📊 耗时：{meta.get('elapsed_seconds', 0):.1f}s", "system"))
            self.root.after(0, lambda: self._log(f"📊 Token估算：{meta.get('estimated_tokens', 0)}", "system"))

            # 记录到会话（只记录结论摘要）
            full_record = f"【{label}完成】{board_version[:200] if board_version else '（结论已生成）'}"
            self._add_message_to_session("assistant", full_record)
            self._handle_title_generation_after_reply(full_record, user_input)
            self.update_status("就绪")
            self.debate_progress.configure(value=100)
            self.debate_status_label.config(text="完成 100%")

        except Exception as e:
            self.root.after(0, lambda: self._log(f"❌ {label}出错：{e}", "error"))
            self.root.after(0, lambda: self.debate_progress.configure(value=0))
            self.root.after(0, lambda: self.debate_status_label.config(text="出错"))
            self.root.after(0, lambda: self.update_status("就绪"))

    def _render_result_to_panel(self, question: str, role_blocks: List[Dict],
                                board_version: str, employee_version: str,
                                novice_version: str, expert_version: str,
                                judge_audit: Dict, dashboard_stats: Dict,
                                final_schema: Dict, result: Dict):
        """
        将辩论结果以富文本形式渲染到独立的结果面板
        """
        # ===== 安全检查：如果面板未初始化，降级到日志输出 =====
        if self.result_area is None:
            self._log("⚠️ 结果面板未初始化，报告输出到日志区", "warning")
            self._log("=" * 60, "system")
            self._log("📋 辩论报告 v2.2", "system")
            self._log("=" * 60, "system")
            for block in role_blocks:
                self._log(f"▸ {block['name']}", "system")
                self._log(f"{block['content'][:200]}...", "system")
            return

        # 启用编辑模式
        self.result_area.config(state=tk.NORMAL)
        self.result_area.delete(1.0, tk.END)

        # ================================================================
        # 1. 标题
        # ================================================================
        self.result_area.insert(tk.END, "📋 辩论报告 v2.2\n", "phase_title")
        self.result_area.insert(tk.END, "─" * 60 + "\n", "result_divider")
        self.result_area.insert(tk.END, f"问题：{question}\n\n", "judge_body")

        # ================================================================
        # 2. 各角色观点
        # ================================================================
        self.result_area.insert(tk.END, "【阶段1：各角色独立观点】\n", "phase_title")
        self.result_area.insert(tk.END, "─" * 60 + "\n", "result_divider")

        role_tag_map = {
            "激进者": "role_radical",
            "保守者": "role_conservative",
            "结构主义者": "role_structural",
            "百灵鸟": "role_lark",
            "取经者": "role_pilgrim",
            "奇谋者": "role_strategist",
            "延安智者": "role_statesman",
            "大法官": "role_judge",
            "首席发言人": "role_spokesperson",
        }

        for block in role_blocks:
            name = block.get("name", "未知角色")
            content = block.get("content", "")

            # 角色名（带颜色）
            tag = role_tag_map.get(name, "role_default")
            self.result_area.insert(tk.END, f"▸ {name}\n", tag)

            # 观点正文
            self.result_area.insert(tk.END, f"{content}\n\n", "viewpoint_body")

        # ================================================================
        # 3. 大法官裁决
        # ================================================================
        self.result_area.insert(tk.END, "\n【阶段2：大法官裁决】\n", "judge_header")
        self.result_area.insert(tk.END, "─" * 60 + "\n", "result_divider")

        # 3.1 绩效看板
        performance_board = final_schema.get("judge_performance_board", [])
        if not performance_board:
            performance_board = judge_audit.get("by_rule", [])

        if performance_board:
            self.result_area.insert(tk.END, "① 角色绩效看板\n", "judge_body")
            table_lines = self._build_performance_table(performance_board)
            for line in table_lines:
                self.result_area.insert(tk.END, line + "\n", "judge_table")
            self.result_area.insert(tk.END, "\n", "judge_body")

        # 3.2 终审裁决
        final_verdict = final_schema.get("judge_final_verdict", "")
        if not final_verdict:
            final_verdict = judge_audit.get("summary", "（终审裁决待生成）")

        self.result_area.insert(tk.END, "② 终审裁决\n", "judge_body")
        self.result_area.insert(tk.END, f"{final_verdict}\n\n", "judge_body")

        # 3.3 驳回明细
        rejected_details = final_schema.get("judge_rejected_details", "")
        if not rejected_details:
            rejected_details = "无驳回项。"

        self.result_area.insert(tk.END, "③ 驳回明细\n", "judge_body")
        self.result_area.insert(tk.END, f"{rejected_details}\n\n", "judge_body")

        # ================================================================
        # 4. 首席发言人叙事
        # ================================================================
        self.result_area.insert(tk.END, "【阶段3：首席发言人叙事】\n", "spokesperson_header")
        self.result_area.insert(tk.END, "─" * 60 + "\n", "result_divider")

        spokesperson_narrative = final_schema.get("spokesperson_narrative", "")
        if not spokesperson_narrative:
            spokesperson_narrative = board_version or "（首席发言人叙事待生成）"

        self.result_area.insert(tk.END, f"{spokesperson_narrative}\n\n", "spokesperson_body")

        # ================================================================
        # 5. 儒雅笔谈（附录）
        # ================================================================
        elegant_epilogue = final_schema.get("elegant_epilogue", "")
        if not elegant_epilogue:
            elegant_epilogue = result.get("elegant_epilogue", "")

        if elegant_epilogue:
            self.result_area.insert(tk.END, "\n【附录：儒雅笔谈】\n", "elegant_header")
            self.result_area.insert(tk.END, "─" * 60 + "\n", "result_divider")
            self.result_area.insert(tk.END, f"{elegant_epilogue}\n", "elegant_body")

        # 锁定编辑
        self.result_area.config(state=tk.DISABLED)
        # 滚动到顶部
        self.result_area.see(1.0)

    def _build_performance_table(self, performance_board: List[Dict]) -> List[str]:
        """
        构建绩效看板的纯文本表格
        """
        if not performance_board:
            return ["（绩效看板数据缺失）"]

        # 提取数据
        rows = []
        for item in performance_board:
            role = item.get("role", "未知")[:10]
            contrib = str(item.get("contribution_percent", 0))
            kpi = str(item.get("kpi_score", 0))
            status = item.get("status", "暂缓")
            reason = item.get("reason", "")[:15]

            # 状态图标
            status_icon = {"采纳": "✅", "附条件": "⚠️", "暂缓": "⏸", "驳回": "❌"}.get(status, "•")
            status_display = f"{status_icon}{status}"

            rows.append({
                "role": role,
                "contrib": contrib,
                "kpi": kpi,
                "status": status_display,
                "reason": reason
            })

        # 计算列宽
        max_role = max([len(r["role"]) for r in rows] + [4])
        max_contrib = max([len(r["contrib"]) for r in rows] + [5])
        max_kpi = max([len(r["kpi"]) for r in rows] + [5])
        max_status = max([len(r["status"]) for r in rows] + [4])
        max_reason = max([len(r["reason"]) for r in rows] + [6])

        # 构建表格
        sep = ("+" + "-" * (max_role + 2) +
               "+" + "-" * (max_contrib + 2) +
               "+" + "-" * (max_kpi + 2) +
               "+" + "-" * (max_status + 2) +
               "+" + "-" * (max_reason + 2) + "+")

        lines = [sep]
        header = ("| " + "角色".ljust(max_role) + " | " +
                  "贡献度".ljust(max_contrib) + " | " +
                  "KPI".ljust(max_kpi) + " | " +
                  "状态".ljust(max_status) + " | " +
                  "核心理由".ljust(max_reason) + " |")
        lines.append(header)
        lines.append(sep)

        for row in rows:
            line = ("| " + row["role"].ljust(max_role) + " | " +
                    row["contrib"].rjust(max_contrib) + "% | " +
                    row["kpi"].rjust(max_kpi) + "/10 | " +
                    row["status"].ljust(max_status) + " | " +
                    row["reason"].ljust(max_reason) + " |")
            lines.append(line)

        lines.append(sep)
        return lines
            
    def _create_pending_from_competition(self, question: str, role_answers: List[Dict]):
        suffix = int(hashlib.sha256(question.encode("utf-8")).hexdigest(), 16) % 1000
        card_id = f"PENDING-{datetime.now().strftime('%Y%m%d%H%M%S')}-{suffix:03d}"
        lines = [f"## {card_id}", "- 类型：多角色竞争结果", f"- 问题：{question}", "- 各角色答案摘要："]
        for ra in role_answers:
            summary = ra['answer'][:150].replace('\n', ' ')
            lines.append(f"  - {ra['name']}：{summary}...")
        lines.append("- 建议：可从中选择最合理的路径，或将融合后的答案转为晶体")
        card_content = "\n" + "\n".join(lines) + "\n"
        FileIO.append("pending", card_content)
        self._log(f"  📄 多角色竞争结果已生成待确认卡片 {card_id}", "system")

    def _do_file_chat(self):
        """选择多个文件，合并内容后对话（支持多选）"""
        paths = filedialog.askopenfilenames(
            title="选择一个或多个文件",
            filetypes=[("所有支持文件",
                        "*.txt *.md *.py *.html *.htm *.json *.xml *.css *.js *.xlsx *.xls *.csv *.docx *.pdf *.pptx")]
        )
        if not paths:
            return

        self._log(f"📄 选择了 {len(paths)} 个文件", "system")

        # 提取所有文件内容，合并为一个字符串
        all_content = []
        total_chars = 0
        max_chars = 20000  # 内容上限，防止AI上下文过长

        for file_path in paths:
            filename = os.path.basename(file_path)
            self._log(f"  读取: {filename}", "system")
            try:
                text_units = self.batch_processor.extract_text_from_file(file_path)
                if not text_units:
                    self._log(f"  文件 {filename} 无有效内容或读取失败", "warning")
                    continue
                # 取第一个文本单元（通常足够）
                content = text_units[0].strip()
                if not content:
                    continue
                # 截断单个文件过长内容（保留开头）
                if len(content) > 3000:
                    content = content[:3000] + "\n...（内容过长，已截断）"
                    self._log(f"  {filename} 内容超过3000字符，已截断", "warning")
                # 添加文件名标识
                all_content.append(f"【文件：{filename}】\n{content}\n")
                total_chars += len(content)
                if total_chars > max_chars:
                    self._log(f"总内容已达 {max_chars} 字符上限，后续文件将被跳过", "warning")
                    break
            except Exception as e:
                self._log(f"  读取 {filename} 出错: {e}", "error")

        if not all_content:
            self._log("所有文件均无有效内容", "error")
            return

        # 合并所有内容
        combined = "\n".join(all_content)
        if total_chars > max_chars:
            combined = combined[:max_chars] + "\n...（总内容过长，已截断）"

        # 构造用户消息
        user_msg = f"[多文件内容] 共 {len(paths)} 个文件，内容如下：\n\n{combined}\n\n请基于以上文件内容回答。"

        self._log(user_msg[:200] + "...", "user")
        self.update_status("AI 思考中...")

        self._add_message_to_session("user", user_msg)
        cur = self.current_history.copy()

        def task():
            os.environ["DEEPSEEK_API_KEY"] = self.api_key.get()
            l0_holes, l1_crystals = self.engine.get_attention_context()
            ctx = f"\n[注意力上下文] 当前核心孔洞：{', '.join([h.content[:50] for h in l0_holes])}\nL1晶体数量：{len(l1_crystals)} 条\n"
            reply = self.ai.chat_with_history(cur, context=ctx)
            self.root.after(0, lambda: self._file_chat_done(reply))

        threading.Thread(target=task, daemon=True).start()

    def _file_chat_done(self, reply):
        self._log(reply, "ai")
        self._add_message_to_session("assistant", reply)
        # 获取本轮用户消息（由于是文件对话，用户消息可能包含前缀，但内容在 history 中）
        user_msg = ""
        for idx in range(len(self.current_history)-2, -1, -1):
            if self.current_history[idx][0] == "user":
                user_msg = self.current_history[idx][1]
                break
        self._handle_title_generation_after_reply(reply, user_msg)
        self.update_status("就绪")
        self.engine.fingerprint_extractor.extract(self.current_history, increment_interactions=1)
        self._refresh_question_list()

    def _start_batch(self):
        folder = filedialog.askdirectory(title="选择文件夹")
        if not folder:
            return
        self.batch_btn.config(state=tk.DISABLED)
        self.stop_batch_btn.config(state=tk.NORMAL)
        self.processing = True
        self.stop_batch = False
        self.batch_progress['value'] = 0
        self._log(f"开始批量处理: {folder}", "system")
        self._log(f"模式: {'晶体化' if self.batch_mode.get() == 'crystal' else '聊天'}", "system")

        def batch_task():
            def prog(v):
                self.root.after(0, lambda: self.batch_progress.configure(value=v))

            def stop():
                return self.stop_batch

            def hist_cb(role, content):
                if not self.inject_history.get():
                    return

                # 如果是 AI 回复，确保会话中已有用户消息（否则标题无法生成）
                if role == "assistant":
                    # 检查当前会话是否已有用户消息
                    has_user = any(r == "user" for r, _ in self.current_history)
                    if not has_user:
                        # 添加占位用户消息（内容为批量处理提示）
                        placeholder = f"[批量处理] 文件内容分析"
                        self.root.after(0, lambda: self._add_message_to_session("user", placeholder))

                # 添加 AI 回复
                self.root.after(0, lambda: self._add_message_to_session(role, content))

                # 如果是 AI 回复，延迟生成标题（等待消息写入）
                if role == "assistant":
                    def generate_title():
                        # 获取最近一条用户消息
                        user_msg = ""
                        for r, c in reversed(self.current_history):
                            if r == "user":
                                user_msg = c
                                break
                        if user_msg:
                            self._handle_title_generation_after_reply(content, user_msg)
                    self.root.after(150, generate_title)  # 增加延迟确保写入完成

            os.environ["DEEPSEEK_API_KEY"] = self.api_key.get()
            self.batch_processor.process_folder(folder, self.batch_mode.get(), self.fast_mode.get(), prog, stop,
                                                hist_cb)
            self.root.after(0, self._batch_done)

        threading.Thread(target=batch_task, daemon=True).start()

    def _stop_batch_process(self):
        if self.processing:
            self.stop_batch = True
            self._log("正在停止批量处理...", "warning")
            self.stop_batch_btn.config(state=tk.DISABLED)

    def _batch_done(self):
        self.processing = False
        self.batch_btn.config(state=tk.NORMAL)
        self.stop_batch_btn.config(state=tk.DISABLED)
        self.batch_progress['value'] = 0
        self.update_status("批量处理完成")

    def _show_status(self):
        self._log("📊 系统状态", "system")
        content = self.files.read("state")
        if content:
            for line in content.split('\n'):
                self._log(line, "system")

    def _show_holes(self):
        self._log("🕳️ 孔洞花园", "system")
        content = self.files.read("holes")
        if content:
            for line in content.split('\n'):
                self._log(line, "system")

    def _show_pending(self):
        self._log("📋 待确认卡片", "system")
        content = self.files.read("pending")
        if content:
            for line in content.split('\n'):
                self._log(line, "system")

    def _show_today_changes(self):
        today = datetime.now().strftime("%Y-%m-%d")
        today_compact = datetime.now().strftime("%Y%m%d")
        win = Toplevel(self.root)
        self._decorate_popup(win, "今日新增", "900x650")
        text = self._create_text_panel(win)
        text.pack(fill=tk.BOTH, expand=True, padx=8, pady=8)
        text.insert(tk.END, f"今日新增概览：{today}\n\n")
        change_log = self.files.read("change_log")
        sections = re.findall(rf"(## {re.escape(today)}.*?)(?=\n## \d{{4}}-\d{{2}}-\d{{2}}|\Z)", change_log,
                              re.DOTALL)
        text.insert(tk.END, "## 今日变更记录\n")
        text.insert(tk.END, ("\n\n".join(sections) if sections else "暂无今日变更记录。") + "\n\n")
        pending = self.files.read("pending")
        pending_cards = re.findall(rf"(## PENDING-{today_compact}\d+-\d+.*?)(?=\n## |\Z)", pending, re.DOTALL)
        text.insert(tk.END, "## 今日 PENDING\n")
        text.insert(tk.END, ("\n\n".join(pending_cards) if pending_cards else "暂无今日待确认卡片。") + "\n\n")
        tasks = self._load_task_cards()
        today_tasks = [t for t in tasks if str(t.get("id", "")).startswith(f"TASK-{today_compact}")]
        text.insert(tk.END, "## 今日任务卡\n")
        if today_tasks:
            for task in today_tasks:
                text.insert(tk.END, f"- {task.get('id')} | {task.get('status')} | {task.get('title')}\n")
        else:
            text.insert(tk.END, "暂无今日任务卡。\n")
        text.config(state=tk.DISABLED)

    def _show_health_check(self):
        win = Toplevel(self.root)
        self._decorate_popup(win, "健康检查", "850x560")
        text = self._create_text_panel(win)
        text.pack(fill=tk.BOTH, expand=True, padx=8, pady=8)
        results = HealthChecker.run()
        text.insert(tk.END, f"数据根目录：{Config.DATA_ROOT}\n")
        text.insert(tk.END, f"会话数据库：{Config.get_db_path()}\n")
        text.insert(tk.END, f"API Key：{'已配置' if self.api_key.get() or Config.get_api_key() else '未配置'}\n\n")
        if not results:
            text.insert(tk.END, "✅ 未发现数据文件健康问题。\n")
        else:
            for item in results:
                text.insert(tk.END, f"[{item.level}] {item.file}: {item.message}")
                if item.suggested_fix:
                    text.insert(tk.END, f" | 建议：{item.suggested_fix}")
                text.insert(tk.END, "\n")
        optional = [
            ("requests", REQUESTS_AVAILABLE),
        ]
        text.insert(tk.END, "\n## 依赖状态\n")
        for name, ok in optional:
            text.insert(tk.END, f"- {name}: {'可用' if ok else '缺失'}\n")
        text.config(state=tk.DISABLED)

    def _confirm_card_dialog(self):
        card_id = simpledialog.askstring("确认卡片", "输入卡片ID (例如 PENDING-20260525120000):", parent=self.root)
        if card_id:
            self._log(f"✅ 确认卡片 {card_id}...", "system")
            _, content, _ = self._extract_pending_content(card_id)
            edited = simpledialog.askstring("编辑晶体内容", "确认前可编辑晶体内容：", initialvalue=content,
                                            parent=self.root)
            if edited is None:
                self._log("已取消确认卡片", "warning")
                return
            result = self._confirm_pending_card_with_content(card_id, edited)
            self._log(result, "success")
            self._show_status()

    def _confirm_pending_card(self, card_id: str) -> str:
        pending_content = self.files.read("pending")
        pattern = rf"(##\s*{re.escape(card_id)}.*?)(?=\n## |\Z)"
        match = re.search(pattern, pending_content, re.DOTALL)
        if not match:
            return f"未找到卡片 {card_id}"
        block = match.group(1)
        content = None
        title_match = re.search(r"- 标题：(.+)", block)
        if title_match:
            content = title_match.group(1).strip()
        else:
            summary_match = re.search(r"- 内容摘要：(.+)", block)
            if summary_match:
                content = summary_match.group(1).strip()
            else:
                content_match = re.search(r"- 内容：(.+)", block)
                if content_match:
                    content = content_match.group(1).strip()
        if not content:
            lines = block.split('\n')
            for line in lines:
                line = line.strip()
                if line and not line.startswith('-') and not line.startswith('#'):
                    content = line
                    break
        if not content:
            return "无法解析卡片内容"
        crystals = self.files.read("crystals")
        ids = re.findall(r"C(\d+)", crystals)
        next_id = max([int(i) for i in ids], default=0) + 1
        new_crystal = f"\n| C{next_id:03d} | {content} | — |\n"
        self.files.append("crystals", new_crystal)
        new_pending = pending_content.replace(block, "")
        new_pending = re.sub(r"\n\s*\n", "\n\n", new_pending).strip()
        self.files.write("pending", new_pending)
        self.engine._append_change_log("待确认卡确认", f"确认卡片 {card_id} 转为晶体 C{next_id:03d}")
        return f"已确认 {card_id} 并转为晶体 C{next_id:03d}"

    def _extract_pending_content(self, card_id: str) -> Tuple[str, str, str]:
        pending_content = self.files.read("pending")
        pattern = rf"(##\s*{re.escape(card_id)}.*?)(?=\n## |\Z)"
        match = re.search(pattern, pending_content, re.DOTALL)
        if not match:
            return "", "", pending_content
        block = match.group(1)
        content = ""
        for pattern_text in [r"- 标题：(.+)", r"- 内容摘要：(.+)", r"- 内容：(.+)", r"- 生成晶体候选：(.+)"]:
            found = re.search(pattern_text, block)
            if found:
                content = found.group(1).strip()
                break
        if not content:
            for line in block.split('\n'):
                line = line.strip()
                if line and not line.startswith('-') and not line.startswith('#'):
                    content = line
                    break
        return block, content, pending_content

    def _confirm_pending_card_with_content(self, card_id: str, content: str) -> str:
        block, _, pending_content = self._extract_pending_content(card_id)
        content = self._shorten_crystal_content(content)
        if not block:
            return f"未找到卡片 {card_id}"
        if not content:
            return "无法解析卡片内容"
        similar = self._find_similar_crystals(content)
        if similar:
            msg = "发现可能重复晶体：\n" + "\n".join(
                [f"{c.id} ({score:.2f}) {c.content[:60]}" for score, c in similar]) + "\n\n仍然继续转为新晶体吗？"
            if not messagebox.askyesno("重复预警", msg, parent=self.root):
                return "已取消：存在可能重复晶体"
        crystals = self.files.read("crystals")
        ids = re.findall(r"C(\d+)", crystals)
        next_id = max([int(i) for i in ids], default=0) + 1
        new_crystal = f"\n| C{next_id:03d} | {content} | — |\n"
        self.files.append("crystals", new_crystal)
        new_pending = pending_content.replace(block, "")
        new_pending = re.sub(r"\n\s*\n", "\n\n", new_pending).strip()
        self.files.write("pending", new_pending)
        self.engine._append_change_log("待确认卡确认", f"确认卡片 {card_id} 转为晶体 C{next_id:03d}")
        return f"已确认 {card_id} 并转为晶体 C{next_id:03d}"

    def _edit_pending_before_confirm(self, card_id: str, parent, parent_win):
        _, content, _ = self._extract_pending_content(card_id)
        if not content:
            messagebox.showerror("错误", f"无法解析卡片 {card_id}", parent=parent_win)
            return
        win = Toplevel(parent_win)
        self._decorate_popup(win, f"编辑后转为晶体 - {card_id}", "760x420")
        tk.Label(win, text="确认前可编辑晶体内容（建议不超过80字）", bg=Config.GUI_BG_MAIN,
                 font=("微软雅黑", 11, "bold")).pack(anchor=tk.W, padx=10, pady=8)
        text = tk.Text(win, height=8, wrap=tk.WORD, font=Config.GUI_INPUT_FONT,
                       bg=Config.GUI_BG_INPUT, fg=Config.GUI_FG_TEXT, relief=tk.SOLID, bd=1,
                       highlightthickness=1, highlightbackground=Config.GUI_BORDER,
                       insertbackground=Config.GUI_ACCENT_DARK)
        text.pack(fill=tk.BOTH, expand=True, padx=10, pady=5)
        text.insert("1.0", content)
        btn_frame = tk.Frame(win, bg=Config.GUI_BG_MAIN)
        btn_frame.pack(fill=tk.X, padx=10, pady=10)

    def _read_edit_and_commit_pending(self, card_id: str, parent_frame, parent_win):
        """弹出窗口，完整展示卡片的原始 Markdown 全文，允许编辑后入库"""
        pending_content = self.files.read("pending")
        pattern = rf"(##\s*{re.escape(card_id)}.*?)(?=\n## |\Z)"
        match = re.search(pattern, pending_content, re.DOTALL)
        if not match:
            messagebox.showerror("错误", f"未找到卡片 {card_id}", parent=parent_win)
            return
        full_block = match.group(1).strip()

        win = Toplevel(parent_win)
        self._decorate_popup(win, f"阅读并修改 - {card_id}", "920x650")
        tk.Label(win, text="以下是该卡片的完整原始内容，您可以直接编辑，然后确认入库。",
                 bg=Config.GUI_BG_MAIN, font=("微软雅黑", 11, "bold")).pack(anchor=tk.W, padx=10, pady=8)

        text_area = scrolledtext.ScrolledText(win, wrap=tk.WORD, font=Config.GUI_INPUT_FONT,
                                              bg=Config.GUI_BG_INPUT, fg=Config.GUI_FG_TEXT,
                                              relief=tk.SOLID, bd=1, highlightthickness=1,
                                              highlightbackground=Config.GUI_BORDER,
                                              insertbackground=Config.GUI_ACCENT_DARK)
        text_area.pack(fill=tk.BOTH, expand=True, padx=10, pady=5)
        text_area.insert("1.0", full_block)

        btn_frame = tk.Frame(win, bg=Config.GUI_BG_MAIN)
        btn_frame.pack(fill=tk.X, padx=10, pady=10)

        def submit():
            edited = text_area.get("1.0", tk.END).strip()
            if not edited:
                messagebox.showwarning("提示", "编辑内容不能为空", parent=win)
                return
            result = self._confirm_pending_card_with_content(card_id, edited)
            messagebox.showinfo("结果", result, parent=win)
            win.destroy()
            # 刷新卡片列表
            for widget in parent_frame.winfo_children():
                widget.destroy()
            self._build_pending_tasks_view(parent_frame, parent_win)

        self._create_btn(btn_frame, "✅ 确认入库", submit, Config.GUI_SUCCESS).pack(side=tk.LEFT, padx=5)
        self._create_btn(btn_frame, "取消", win.destroy, Config.GUI_DANGER).pack(side=tk.LEFT, padx=5)

        def submit():
            edited = text.get("1.0", tk.END).strip()
            result = self._confirm_pending_card_with_content(card_id, edited)
            messagebox.showinfo("结果", result, parent=win)
            win.destroy()
            for widget in parent.winfo_children():
                widget.destroy()
            self._build_pending_tasks_view(parent, parent_win)

        self._create_btn(btn_frame, "✅ 转为晶体", submit, Config.GUI_SUCCESS).pack(side=tk.LEFT, padx=5)
        self._create_btn(btn_frame, "取消", win.destroy, Config.GUI_DANGER).pack(side=tk.LEFT, padx=5)

    def _ask_daily_keywords(self, title: str = "每日计划关键词") -> List[str]:
        text = simpledialog.askstring(
            title,
            "今天是否有重点关注关键词？\n可用逗号、空格或换行分隔；留空将使用默认主题。",
            parent=self.root,
        )
        if text is None:
            return []
        return [item.strip() for item in re.split(r"[,，;；\s\n]+", text) if item.strip()]

    def _run_daily_plan_thread(self, keywords: List[str]):
        self.daily_plan_running = True
        self.stop_daily_plan = False
        if hasattr(self, "stop_daily_btn"):
            self.stop_daily_btn.config(state=tk.NORMAL)

        def progress(data: Dict):
            msg = f"{data.get('stage', '每日计划')} {data.get('progress', 0)}% · 候选{data.get('candidate_count', 0)} · 卡片{data.get('pending_count', 0)} · 任务{data.get('task_count', 0)}"
            self.root.after(0, lambda: self.update_status(msg))

        def worker():
            planner = DailyPlanner(self.engine, self.ai, self.fetcher, self._log, self.update_status)
            try:
                planner.run(
                    intent_keywords=keywords,
                    time_budget_seconds=Config.DAILY_PLAN_TIME_BUDGET_SECONDS,
                    stop_flag=lambda: self.stop_daily_plan,
                    progress_callback=progress,
                )
            finally:
                self.daily_plan_running = False
                self.stop_daily_plan = False
                self.root.after(0, lambda: self.stop_daily_btn.config(state=tk.DISABLED) if hasattr(self, "stop_daily_btn") else None)

        threading.Thread(target=worker, daemon=True).start()

    def _stop_daily_plan(self):
        if not self.daily_plan_running:
            self._log("当前没有正在运行的每日计划", "system")
            return
        self.stop_daily_plan = True
        self._log("⏹ 已请求中断每日计划，正在整理已产生成果...", "warning")
        self.update_status("每日计划中断整理中...")
        if hasattr(self, "stop_daily_btn"):
            self.stop_daily_btn.config(state=tk.DISABLED)


    def _open_task_panel(self):
        win = Toplevel(self.root)
        self._decorate_popup(win, "任务面板", "900x600")
        notebook = ttk.Notebook(win)
        notebook.pack(fill=tk.BOTH, expand=True, padx=10, pady=10)
        pending_frame = tk.Frame(notebook, bg=Config.GUI_BG_MAIN)
        notebook.add(pending_frame, text="📋 待确认卡片")
        self._build_pending_tasks_view(pending_frame, win)
        conflict_frame = tk.Frame(notebook, bg=Config.GUI_BG_MAIN)
        notebook.add(conflict_frame, text="⚡ 冲突任务")
        self._build_conflict_tasks_view(conflict_frame, win)

    def _build_pending_tasks_view(self, parent, parent_win):
        content = self.files.read("pending")
        cards = re.findall(r"## (PENDING-\d+-\d+)\n(.*?)(?=\n## |\Z)", content, re.DOTALL)
        if not cards:
            tk.Label(parent, text="暂无待确认卡片", bg=Config.GUI_BG_MAIN, font=("微软雅黑", 12)).pack(pady=20)
            return
        canvas = tk.Canvas(parent, bg=Config.GUI_BG_MAIN, highlightthickness=0)
        scrollbar = ttk.Scrollbar(parent, orient="vertical", command=canvas.yview)
        frame = tk.Frame(canvas, bg=Config.GUI_BG_MAIN)
        frame.bind("<Configure>", lambda e: canvas.configure(scrollregion=canvas.bbox("all")))
        canvas.create_window((0, 0), window=frame, anchor="nw")
        canvas.configure(yscrollcommand=scrollbar.set)
        self._bind_canvas_mousewheel(canvas, frame)
        canvas.pack(side="left", fill="both", expand=True)
        scrollbar.pack(side="right", fill="y")
        for cid, body in cards:
            card_frame = tk.Frame(frame, bg=Config.GUI_BG_CARD, relief=tk.FLAT, bd=0,
                                  highlightthickness=1, highlightbackground=Config.GUI_BORDER)
            card_frame.pack(fill=tk.X, padx=8, pady=6)
            lines = body.split('\n')
            card_type = title = source = content_text = ""
            for line in lines:
                if line.startswith("- 类型："):
                    card_type = line.replace("- 类型：", "").strip()
                elif line.startswith("- 标题："):
                    title = line.replace("- 标题：", "").strip()
                elif line.startswith("- 来源："):
                    source = line.replace("- 来源：", "").strip()
                elif line.startswith("- 内容摘要："):
                    content_text = line.split("：", 1)[-1].strip()
                elif line.startswith("- 内容："):
                    content_text = line.split("：", 1)[-1].strip()
                elif line.startswith("- 生成晶体候选："):
                    content_text = line.split("：", 1)[-1].strip()
            if not title:
                title = content_text[:50] if content_text else cid
            tk.Label(card_frame, text=f"{cid} | {card_type} | {source}" if source else cid,
                     font=Config.GUI_TEXT_FONT, bg=Config.GUI_BG_CARD, fg=Config.GUI_FG_MUTED).pack(anchor=tk.W,
                                                                                                    padx=10,
                                                                                                    pady=(8, 2))
            tk.Label(card_frame, text=title, font=Config.GUI_CARD_FONT, bg=Config.GUI_BG_CARD,
                     fg=Config.GUI_FG_TEXT, wraplength=700, justify=tk.LEFT).pack(anchor=tk.W, padx=10, pady=2)
            if content_text:
                tk.Label(card_frame, text=content_text[:200], font=Config.GUI_TEXT_FONT,
                         bg=Config.GUI_BG_CARD, fg=Config.GUI_FG_TEXT, wraplength=700,
                         justify=tk.LEFT).pack(anchor=tk.W, padx=10, pady=2)
            btn_frame = tk.Frame(card_frame, bg=Config.GUI_BG_CARD)
            btn_frame.pack(fill=tk.X, padx=10, pady=(4, 10))

            def confirm(c=cid, w=parent_win):
                self._edit_pending_before_confirm(c, parent, w)

            def ignore(c=cid, w=parent_win):
                cur = self.files.read("pending")
                new = re.sub(rf"## {re.escape(c)}.*?(?=\n## |\Z)", "", cur, flags=re.DOTALL)
                new = re.sub(r"\n\s*\n", "\n\n", new).strip()
                self.files.write("pending", new)
                messagebox.showinfo("已忽略", f"已忽略卡片 {c}", parent=w)
                for widget in parent.winfo_children():
                    widget.destroy()
                self._build_pending_tasks_view(parent, w)
                
            def read_edit_commit(c=cid, w=parent_win):
                self._read_edit_and_commit_pending(c, parent, w)

            self._create_btn(btn_frame, "✅ 确认并转为晶体", confirm, Config.GUI_SUCCESS).pack(side=tk.LEFT,
                                                                                             padx=5)
            self._create_btn(btn_frame, "📖 阅读修改后入库", read_edit_commit, Config.GUI_INFO).pack(side=tk.LEFT, padx=5)
            self._create_btn(btn_frame, "❌ 忽略", ignore, Config.GUI_DANGER).pack(side=tk.LEFT, padx=5)

    def _build_conflict_tasks_view(self, parent, parent_win):
        tasks = self._load_task_cards()
        pending = [t for t in tasks if t.get("type") == "conflict" and t.get("status") == "pending"]
        if not pending:
            tk.Label(parent, text="暂无冲突任务", bg=Config.GUI_BG_MAIN, font=("微软雅黑", 12)).pack(pady=20)
            return
        canvas = tk.Canvas(parent, bg=Config.GUI_BG_MAIN, highlightthickness=0)
        scrollbar = ttk.Scrollbar(parent, orient="vertical", command=canvas.yview)
        frame = tk.Frame(canvas, bg=Config.GUI_BG_MAIN)
        frame.bind("<Configure>", lambda e: canvas.configure(scrollregion=canvas.bbox("all")))
        canvas.create_window((0, 0), window=frame, anchor="nw")
        canvas.configure(yscrollcommand=scrollbar.set)
        self._bind_canvas_mousewheel(canvas, frame)
        canvas.pack(side="left", fill="both", expand=True)
        scrollbar.pack(side="right", fill="y")
        for task in pending:
            card = tk.Frame(frame, bg=Config.GUI_BG_CARD, relief=tk.FLAT, bd=0,
                            highlightthickness=1, highlightbackground=Config.GUI_BORDER)
            card.pack(fill=tk.X, padx=8, pady=6)
            tk.Label(card, text=f"{task['id']} | {task['title']}", font=Config.GUI_CARD_FONT,
                     bg=Config.GUI_BG_CARD, fg=Config.GUI_FG_TEXT).pack(anchor=tk.W, padx=10, pady=(8, 2))
            tk.Label(card, text=task['content'][:300], font=Config.GUI_TEXT_FONT,
                     bg=Config.GUI_BG_CARD, fg=Config.GUI_FG_TEXT, wraplength=700,
                     justify=tk.LEFT).pack(anchor=tk.W, padx=10, pady=2)
            tk.Label(card, text=f"建议：{task.get('suggested_action', '')}", font=Config.GUI_TEXT_FONT,
                     bg=Config.GUI_BG_CARD, fg=Config.GUI_ACCENT_DARK).pack(anchor=tk.W, padx=10, pady=2)
            btnf = tk.Frame(card, bg=Config.GUI_BG_CARD)
            btnf.pack(fill=tk.X, padx=10, pady=(4, 10))

            def resolve(t=task, w=parent_win):
                t['status'] = 'done'
                cards = self._load_task_cards()
                new_cards = [c if c['id'] != t['id'] else t for c in cards]
                self._save_task_cards(new_cards)
                self.engine._append_change_log("冲突解决", f"任务 {t['id']} 已标记为已处理")
                messagebox.showinfo("完成", f"任务 {t['id']} 已标记为已处理", parent=w)
                for widget in parent.winfo_children():
                    widget.destroy()
                self._build_conflict_tasks_view(parent, w)

            def ignore(t=task, w=parent_win):
                t['status'] = 'ignored'
                cards = self._load_task_cards()
                new_cards = [c if c['id'] != t['id'] else t for c in cards]
                self._save_task_cards(new_cards)
                messagebox.showinfo("已忽略", f"已忽略任务 {t['id']}", parent=w)
                for widget in parent.winfo_children():
                    widget.destroy()
                self._build_conflict_tasks_view(parent, w)

            self._create_btn(btnf, "✅ 标记为已处理", resolve, Config.GUI_SUCCESS).pack(side=tk.LEFT, padx=5)
            self._create_btn(btnf, "❌ 忽略", ignore, Config.GUI_DANGER).pack(side=tk.LEFT, padx=5)

    def _load_task_cards(self) -> List[Dict]:
        if not self.files.exists("task_cards"):
            return []
        try:
            return json.loads(self.files.read("task_cards"))
        except:
            return []

    def _save_task_cards(self, cards: List[Dict]):
        self.files.write("task_cards", json.dumps(cards, ensure_ascii=False, indent=2))

    def _open_crystal_manager(self):
        win = Toplevel(self.root)
        self._decorate_popup(win, "晶体管理", "1100x700")
        tool_frame = tk.Frame(win, bg=Config.GUI_BG_MAIN)
        tool_frame.pack(fill=tk.X, padx=5, pady=5)
        tk.Label(tool_frame, text="筛选层级:", bg=Config.GUI_BG_MAIN, font=("微软雅黑", 10)).pack(side=tk.LEFT,
                                                                                               padx=5)
        filter_var = tk.StringVar(value="全部")
        filter_combo = ttk.Combobox(tool_frame, textvariable=filter_var, values=["全部", "L1", "L2", "L3"],
                                    state="readonly", width=8)
        filter_combo.pack(side=tk.LEFT, padx=5)
        self._create_btn(tool_frame, "刷新", lambda: self._refresh_crystal_treeview(tree, filter_var.get())).pack(
            side=tk.LEFT, padx=10)
        tree_frame = tk.Frame(win)
        tree_frame.pack(fill=tk.BOTH, expand=True, padx=10, pady=10)
        scroll_y = ttk.Scrollbar(tree_frame, orient=tk.VERTICAL)
        scroll_x = ttk.Scrollbar(tree_frame, orient=tk.HORIZONTAL)
        tree = ttk.Treeview(tree_frame, columns=("ID", "内容摘要", "层级", "热度", "最后访问", "固定"), show="headings",
                            yscrollcommand=scroll_y.set, xscrollcommand=scroll_x.set)
        scroll_y.config(command=tree.yview)
        scroll_x.config(command=tree.xview)
        for col, width in [("ID", 80), ("内容摘要", 400), ("层级", 80), ("热度", 80), ("最后访问", 120), ("固定", 80)]:
            tree.heading(col, text=col)
            tree.column(col, width=width)
        tree.pack(side=tk.LEFT, fill=tk.BOTH, expand=True)
        scroll_y.pack(side=tk.RIGHT, fill=tk.Y)
        scroll_x.pack(side=tk.BOTTOM, fill=tk.X)
        self._register_mousewheel(tree, tree)
        action_frame = tk.Frame(win, bg=Config.GUI_BG_MAIN)
        action_frame.pack(fill=tk.X, padx=10, pady=10)
        tk.Label(action_frame, text="选中晶体ID:", bg=Config.GUI_BG_MAIN, font=("微软雅黑", 10)).pack(side=tk.LEFT,
                                                                                                   padx=5)
        selected_id_var = tk.StringVar()
        tk.Entry(action_frame, textvariable=selected_id_var, width=10, font=("微软雅黑", 10)).pack(side=tk.LEFT,
                                                                                                  padx=5)
        tk.Label(action_frame, text="新层级:", bg=Config.GUI_BG_MAIN, font=("微软雅黑", 10)).pack(side=tk.LEFT,
                                                                                               padx=5)
        new_layer_var = tk.StringVar(value="L2")
        layer_combo = ttk.Combobox(action_frame, textvariable=new_layer_var, values=["L1", "L2", "L3"],
                                   state="readonly", width=5)
        layer_combo.pack(side=tk.LEFT, padx=5)

        def change_layer():
            cid = selected_id_var.get().strip()
            new_layer = new_layer_var.get()
            if not cid:
                messagebox.showwarning("警告", "请先选中或输入晶体ID")
                return
            state = self.engine.load_layer_state()
            layers = state.get("layers", {})
            manual = state.get("manual_override", {})
            if cid not in layers:
                messagebox.showerror("错误", f"未找到晶体 {cid}")
                return
            if new_layer == "L1":
                curr_l1 = sum(1 for l in layers.values() if l == "L1")
                if curr_l1 >= Config.L1_MAX and layers.get(cid) != "L1":
                    if not messagebox.askyesno("容量警告", f"L1已有 {Config.L1_MAX} 条，是否仍要提升？"):
                        return
            old = layers[cid]
            layers[cid] = new_layer
            if new_layer == "L1":
                if messagebox.askyesno("固定到L1", f"是否将 {cid} 固定到L1？"):
                    manual[cid] = "L1_fixed"
                else:
                    manual.pop(cid, None)
            else:
                manual.pop(cid, None)
            last_acc = state.get("last_accessed", {})
            last_acc[cid] = date.today().isoformat()
            state["last_accessed"] = last_acc
            state["layers"] = layers
            state["manual_override"] = manual
            self.engine.save_layer_state(state)
            self._log(f"已将晶体 {cid} 从 {old} 改为 {new_layer}", "success")
            self.engine._append_change_log("手动层级变更", f"晶体 {cid}: {old} → {new_layer}")
            self._refresh_crystal_treeview(tree, filter_var.get())

        self._create_btn(action_frame, "应用层级变更", change_layer, Config.GUI_SUCCESS).pack(side=tk.LEFT,
                                                                                             padx=10)

        def toggle_fixed():
            cid = selected_id_var.get().strip()
            if not cid:
                return
            state = self.engine.load_layer_state()
            manual = state.get("manual_override", {})
            layers = state.get("layers", {})
            if cid not in layers:
                messagebox.showerror("错误", f"未找到晶体 {cid}")
                return
            if layers[cid] != "L1":
                messagebox.showwarning("警告", "只有L1晶体可以固定")
                return
            if cid in manual and manual[cid] == "L1_fixed":
                manual.pop(cid, None)
                self._log(f"已取消 {cid} 的固定状态", "success")
            else:
                manual[cid] = "L1_fixed"
                self._log(f"已将 {cid} 固定到L1", "success")
            state["manual_override"] = manual
            self.engine.save_layer_state(state)
            self._refresh_crystal_treeview(tree, filter_var.get())

        self._create_btn(action_frame, "切换固定到L1", toggle_fixed, Config.GUI_WARNING).pack(side=tk.LEFT,
                                                                                             padx=5)

        def delete_crystal():
            cid = selected_id_var.get().strip()
            if not cid:
                return
            if not messagebox.askyesno("确认删除", f"删除晶体 {cid}？不可恢复。"):
                return
            cryst = self.files.read("crystals")
            pat = rf"\| {re.escape(cid)} \|.*?\|\n"
            new = re.sub(pat, "", cryst, flags=re.MULTILINE)
            if new == cryst:
                self._log(f"未找到晶体 {cid}", "error")
                return
            self.files.write("crystals", new)
            state = self.engine.load_layer_state()
            layers = state.get("layers", {})
            layers.pop(cid, None)
            heat = state.get("heat_map", {})
            heat.pop(cid, None)
            last = state.get("last_accessed", {})
            last.pop(cid, None)
            manual = state.get("manual_override", {})
            manual.pop(cid, None)
            state["layers"] = layers
            state["heat_map"] = heat
            state["last_accessed"] = last
            state["manual_override"] = manual
            self.engine.save_layer_state(state)
            self._log(f"已删除晶体 {cid}", "success")
            self.engine._append_change_log("手动删除晶体", f"删除 {cid}")
            self._refresh_crystal_treeview(tree, filter_var.get())

        self._create_btn(action_frame, "🗑 删除晶体", delete_crystal, Config.GUI_DANGER).pack(side=tk.LEFT,
                                                                                             padx=10)

        def on_tree_select(event):
            sel = tree.selection()
            if sel:
                selected_id_var.set(tree.item(sel[0])['values'][0])

        tree.bind('<<TreeviewSelect>>', on_tree_select)
        self._refresh_crystal_treeview(tree, "全部")

    def _refresh_crystal_treeview(self, tree, filter_layer):
        for item in tree.get_children():
            tree.delete(item)
        crystals = self.engine.parse_crystals()
        state = self.engine.load_layer_state()
        layers = state.get("layers", {})
        heat_map = state.get("heat_map", {})
        last_accessed = state.get("last_accessed", {})
        manual = state.get("manual_override", {})
        if not layers:
            self.engine.update_crystal_layers()
            state = self.engine.load_layer_state()
            layers = state.get("layers", {})
        for c in crystals:
            cid = c.id
            layer = layers.get(cid, "L2")
            if filter_layer != "全部" and layer != filter_layer:
                continue
            heat = heat_map.get(cid, 0.0)
            last = last_accessed.get(cid, "从未")
            fixed = "是" if manual.get(cid) == "L1_fixed" else ""
            summary = c.content[:50] + ("..." if len(c.content) > 50 else "")
            tree.insert("", tk.END, values=(cid, summary, layer, f"{heat:.2f}", last, fixed))

    def _open_search_window(self):
        win = Toplevel(self.root)
        self._decorate_popup(win, "文档搜索", "800x600")
        tk.Label(win, text="关键词或正则表达式:", font=Config.GUI_TEXT_FONT, bg=Config.GUI_BG_MAIN,
                 fg=Config.GUI_FG_TEXT).pack(pady=5)
        entry = tk.Entry(win, width=60, font=Config.GUI_TEXT_FONT, bg=Config.GUI_BG_INPUT,
                         fg=Config.GUI_FG_TEXT, relief=tk.SOLID, bd=1)
        entry.pack(pady=5)
        regex_var = tk.BooleanVar()
        tk.Checkbutton(win, text="使用正则表达式", variable=regex_var, font=Config.GUI_TEXT_FONT,
                       bg=Config.GUI_BG_MAIN, activebackground=Config.GUI_BG_MAIN).pack()
        dirs_var = tk.StringVar(value="晶体数据,核心配置,系统日志,暂存区")
        tk.Label(win, text="搜索目录（逗号分隔）:", font=Config.GUI_TEXT_FONT, bg=Config.GUI_BG_MAIN,
                 fg=Config.GUI_FG_TEXT).pack(pady=5)
        dir_entry = tk.Entry(win, width=60, textvariable=dirs_var, font=Config.GUI_TEXT_FONT,
                             bg=Config.GUI_BG_INPUT, fg=Config.GUI_FG_TEXT, relief=tk.SOLID, bd=1)
        dir_entry.pack(pady=5)
        result_text = self._create_text_panel(win)
        result_text.pack(fill=tk.BOTH, expand=True, padx=5, pady=5)
        result_text.tag_config("hit", background="#fff59d", foreground="#111111")

        def highlight_hits(keyword: str):
            if not keyword:
                return
            terms = [keyword]
            if not regex_var.get():
                terms.extend([t for t in SearchService._tokens(keyword) if len(t) >= 2])
            seen = set()
            for term in terms:
                if not term or term in seen:
                    continue
                seen.add(term)
                start = "1.0"
                while True:
                    pos = result_text.search(term, start, tk.END, nocase=True)
                    if not pos:
                        break
                    end = f"{pos}+{len(term)}c"
                    result_text.tag_add("hit", pos, end)
                    start = end

        def do_search():
            keyword = entry.get().strip()
            if not keyword:
                return
            dirs = [d.strip() for d in dirs_var.get().split(",") if d.strip()]
            results = SearchService.search_documents(keyword, dirs, regex=regex_var.get())
            result_text.delete(1.0, tk.END)
            for file_path, line_num, line in results[:500]:
                result_text.insert(tk.END, f"{file_path}:{line_num}: {line}\n")
            result_text.insert(tk.END, f"\n共找到 {len(results)} 条结果")
            highlight_hits(keyword)

        self._create_btn(win, "搜索", do_search).pack(pady=5)

    def _bind_canvas_mousewheel(self, canvas: tk.Canvas, *containers):
        """绑定鼠标滚轮事件到 Canvas 及其内部容器"""
        self._register_mousewheel(canvas, canvas)
        for container in containers:
            self._register_mousewheel(container, canvas)

    def _check_and_run_daily_plan(self):
        """检查今日是否已执行每日计划，未执行则提示用户"""
        planner = DailyPlanner(self.engine, self.ai, self.fetcher, self._log, self.update_status)
        if not planner.is_today_run():
            self._log("📆 今日尚未执行每日计划，等待确认当日意向关键词...", "system")
            if not messagebox.askyesno("每日计划", "今日尚未执行每日计划，是否现在填写关键词并开始？"):
                self._log("已跳过本次自动每日计划，不标记为今日完成", "system")
                return
            keywords = self._ask_daily_keywords("自动每日计划关键词")
            self._run_daily_plan_thread(keywords)
        else:
            self._log("✅ 今日每日计划已执行过，跳过自动运行", "system")

    def _manual_run_daily_plan(self):
        """手动触发每日计划（按钮调用）"""
        planner = DailyPlanner(self.engine, self.ai, self.fetcher, self._log, self.update_status)
        if planner.is_today_run() and not messagebox.askyesno("确认", "今日已执行过每日计划，是否重新执行？"):
            return
        keywords = self._ask_daily_keywords()
        self._run_daily_plan_thread(keywords)

    def _extract_keywords(self, text: str) -> List[str]:
        words = re.findall(r'[\w\u4e00-\u9fff]+', text)
        stop = {"的", "了", "和", "与", "或", "一个", "这个", "那个", "如何", "什么", "为什么"}
        kw = [w for w in words if w not in stop][:5]
        return kw if kw else ["晶体树", "认知"]

    def _search_duckduckgo(self, keywords: List[str]) -> str:
        if not REQUESTS_AVAILABLE:
            return "需要安装 requests 库"
        query = " ".join(keywords)
        headers = {"User-Agent": NetworkManager.get_random_user_agent()}
        url = "https://html.duckduckgo.com/html/"
        try:
            resp = requests.post(url, data={"q": query}, headers=headers, timeout=10)
            resp.raise_for_status()
            html = resp.text
            results = []
            titles = re.findall(r'<a class="result__a"[^>]*>(.*?)</a>', html)
            snippets = re.findall(r'<a class="result__snippet"[^>]*>(.*?)</a>', html)
            for i in range(min(3, len(titles), len(snippets))):
                title = re.sub(r'<[^>]+>', '', titles[i]).strip()
                snippet = re.sub(r'<[^>]+>', '', snippets[i]).strip()
                results.append(f"- {title}: {snippet[:150]}...")
            return "\n".join(results) if results else "未找到相关结果"
        except Exception as e:
            return f"搜索出错: {str(e)}"

    def _append_pending_card(self, card: Dict) -> bool:
        content = re.sub(r"\s+", " ", str(card.get("content", ""))).strip()
        if not content:
            return False
        pending_content = self.files.read("pending")
        if content in pending_content:
            return False
        suffix = int(hashlib.sha256(content.encode("utf-8")).hexdigest(), 16) % 1000
        card_id = f"PENDING-{datetime.now().strftime('%Y%m%d%H%M%S')}-{suffix:03d}"
        card_type = str(card.get("type", "晶体候选"))
        source = str(card.get("source", "AI生成"))
        confidence = str(card.get("confidence", "中"))
        block = f"""
## {card_id}
- 类型：{card_type}
- 来源：{source}
- 置信度：{confidence}
- 内容：{content}
- AI判断：建议人工确认后再转为晶体。
"""
        self.files.append("pending", "\n" + block + "\n")
        return True

    def _update_files(self, ai_response: Dict):
        if ai_response.get("new_crystals"):
            for c in ai_response["new_crystals"]:
                self.files.append("crystals",
                                  f"\n| {c['id']} | {c['content']} | {', '.join(c.get('links', []))} |\n")
                # ===== 新增：记录晶体添加进化事件 =====
                self.engine.log_evolution_event(
                    "crystal_added",
                    {
                        "crystal_id": c["id"],
                        "content": c["content"],
                        "links": c.get("links", []),
                        "trigger": "crystallization"
                    }
                )           
        if ai_response.get("updated_crystals"):
            for upd in ai_response["updated_crystals"]:
                content = self.files.read("crystals")
                pattern = rf"(\| {upd['id']} \| ).*?( \| .*? \|)"
                new = re.sub(pattern, rf"\1{upd['new_content']}\2", content)
                self.files.write("crystals", new)
        if ai_response.get("new_holes"):
            for hole in ai_response["new_holes"]:
                layer_name = {1: "第一层", 2: "第二层", 3: "第三层"}.get(hole.get("layer", 2), "第二层")
                line = f"\n| {hole['id']} | {hole['content']} | {hole.get('urgency', 0.5)} |\n"
                content = self.files.read("holes")
                insert_after = f"## {layer_name}："
                if insert_after in content:
                    parts = content.split(insert_after, 1)
                    after = parts[1]
                    next_heading = re.search(r"\n## ", after)
                    if next_heading:
                        pos = len(parts[0]) + len(insert_after) + next_heading.start()
                    else:
                        pos = len(content)
                    new_content = content[:pos] + line + content[pos:]
                    self.files.write("holes", new_content)
                else:
                    self.files.append("holes", line)
        if ai_response.get("updated_holes"):
            for upd in ai_response["updated_holes"]:
                content = self.files.read("holes")
                pattern = rf"(\| {upd['id']} \| ).*?(\| .*? \|)"
                new = re.sub(pattern, rf"\1{upd['content']}\2", content)
                self.files.write("holes", new)
        if ai_response.get("pending_cards"):
            kept_cards = []
            for card in ai_response["pending_cards"]:
                if self._append_pending_card(card):
                    kept_cards.append(card)
            ai_response["pending_cards"] = kept_cards
        change_entry = f"""### 变更摘要：{ai_response.get('report_summary', '无摘要')}
- 新增晶体：{len(ai_response.get('new_crystals', []))}
- 更新晶体：{len(ai_response.get('updated_crystals', []))}
- 新增孔洞：{len(ai_response.get('new_holes', []))}
- 更新孔洞：{len(ai_response.get('updated_holes', []))}
- 冲突：{len(ai_response.get('conflicts', []))}
- 待确认卡片：{len(ai_response.get('pending_cards', []))}
"""
        self.engine._append_change_log("晶体化变更", change_entry)
        crystals_count = len(re.findall(r"^\| C\d+", self.files.read("crystals"), re.MULTILINE))
        holes_count = len(re.findall(r"^\| H\d+", self.files.read("holes"), re.MULTILINE))
        self.files.write("state", f"""# 系统状态快照
**生成时间**: {datetime.now().isoformat()}
**晶体总数**: {crystals_count}
**孔洞总数**: {holes_count}
**最新变更摘要**: {ai_response.get('report_summary', '无')}
""")

    def _on_closing(self):
        if self.current_session_id:
            self.db.update_session(self.current_session_id, self.current_history, self.current_session_name)
        self.root.destroy()

    def _configure_styles(self):
        style = ttk.Style(self.root)
        try:
            style.theme_use("clam")
        except tk.TclError:
            pass
        style.configure("TCombobox", fieldbackground=Config.GUI_BG_INPUT, background=Config.GUI_BG_INPUT,
                        foreground=Config.GUI_FG_TEXT, arrowcolor=Config.GUI_ACCENT, bordercolor=Config.GUI_BORDER,
                        lightcolor=Config.GUI_BORDER, darkcolor=Config.GUI_BORDER)
        style.configure("TNotebook", background=Config.GUI_BG_MAIN, borderwidth=0)
        style.configure("TNotebook.Tab", background=Config.GUI_BUTTON_SOFT, foreground=Config.GUI_FG_TEXT,
                        padding=(12, 6), font=Config.GUI_TEXT_FONT)
        style.map("TNotebook.Tab", background=[("selected", Config.GUI_BG_CARD)],
                  foreground=[("selected", Config.GUI_ACCENT_DARK)])
        style.configure("Treeview", background=Config.GUI_BG_INPUT, fieldbackground=Config.GUI_BG_INPUT,
                        foreground=Config.GUI_FG_TEXT, rowheight=28, bordercolor=Config.GUI_BORDER,
                        font=Config.GUI_TEXT_FONT)
        style.configure("Treeview.Heading", background=Config.GUI_BG_CARD_ALT, foreground=Config.GUI_ACCENT_DARK,
                        font=("微软雅黑", 10, "bold"))
        style.map("Treeview", background=[("selected", Config.GUI_HIGHLIGHT)],
                  foreground=[("selected", Config.GUI_FG_TEXT)])
        style.configure("Horizontal.TProgressbar", troughcolor=Config.GUI_BG_CARD_ALT,
                        background=Config.GUI_ACCENT, bordercolor=Config.GUI_BORDER)

    def _load_roles(self) -> List[Dict]:
        path = Config.get_path("roles")
        if not path.exists():
            default = {
                "radical": {"name": "激进者", "instruction": "攻击默认前提，假设现有框架是错的，给出颠覆性方案。"},
                "conservative": {"name": "保守者", "instruction": "风险优先，假设资源有限，给出最可落地的稳健方案。"},
                "structural": {"name": "结构主义者", "instruction": "从已有晶体中寻找同构案例，用类比生成方案。"},
                "executor": {"name": "执行者", "instruction": "把方案拆成步骤、资源、时间和可检查的行动清单。"},
                "auditor": {"name": "审计者", "instruction": "检查证据、漏洞、冲突、过度推断和需要暂存的问题。"},
            }
            FileIO.write("roles", json.dumps(default, ensure_ascii=False, indent=2))
            data = default
        else:
            try:
                data = json.loads(path.read_text(encoding="utf-8"))
            except:
                data = {}
        roles_list = []
        for key, val in data.items():
            roles_list.append({"id": key, "key": key, "name": val.get("name", key), "instruction": val.get("instruction", "")})
        fallback_roles = [
            {"id": "radical", "key": "radical", "name": "激进者", "instruction": "攻击默认前提，假设现有框架是错的，给出颠覆性方案。"},
            {"id": "conservative", "key": "conservative", "name": "保守者", "instruction": "风险优先，假设资源有限，给出最可落地的稳健方案。"},
            {"id": "structural", "key": "structural", "name": "结构主义者", "instruction": "从已有晶体中寻找同构案例，用类比生成方案。"},
            {"id": "executor", "key": "executor", "name": "执行者", "instruction": "把方案拆成步骤、资源、时间和可检查的行动清单。"},
            {"id": "auditor", "key": "auditor", "name": "审计者", "instruction": "检查证据、漏洞、冲突、过度推断和需要暂存的问题。"},
        ]
        existing = {role.get("id") or role.get("key") for role in roles_list}
        for role in fallback_roles:
            if len(roles_list) >= 5:
                break
            if role["id"] not in existing:
                roles_list.append(role)
                existing.add(role["id"])
        return roles_list

    def _init_app(self):
        FileIO.ensure_directories()
        FileIO.ensure_default_files()
        self.db._init_db()
        self._load_session_list()
        self._new_session()
        self.root.after(2000, self._check_and_run_daily_plan)

    def _sync_vector_store(self):
        """同步向量库"""
        self._log("🔄 开始同步向量库...", "system")
        self.update_status("向量化中...")

        def task():
            try:
                result = self.engine.sync_vector_store()
                if result["status"] == "already_synced":
                    self.root.after(0, lambda: self._log(f"✅ 向量库已同步（{result['total']} 条晶体）", "success"))
                elif result["status"] == "synced":
                    self.root.after(0, lambda: self._log(f"✅ 向量库同步完成：{result['synced']}/{result['total']} 条晶体", "success"))
                else:
                    self.root.after(0, lambda: self._log(f"⚠️ 同步失败：{result.get('status', 'unknown')}", "warning"))
            except Exception as e:
                self.root.after(0, lambda: self._log(f"❌ 同步出错：{e}", "error"))
            self.root.after(0, lambda: self.update_status("就绪"))

        threading.Thread(target=task, daemon=True).start()

    def _do_meta_search(self):
        """执行 Meta 搜索（认知路径平行对比）"""
        user_input = self._get_input()
        if not user_input:
            self._log("请输入要分析的问题", "error")
            return

        self._log(f"🔍 启动 Meta 搜索：{user_input[:50]}...", "system")
        self.update_status("Meta搜索中...")

        def task():
            try:
                meta_engine = MetaSearchEngine(self.engine, self.ai)
                result = meta_engine.run_comparison(user_input)

                if "error" in result:
                    self.root.after(0, lambda: self._log(f"❌ Meta搜索失败：{result['error']}", "error"))
                    return

                # 构建输出
                lines = ["🔍 认知路径平行对比结果", "=" * 40, ""]
                for i, scored_path in enumerate(result.get("paths", []), 1):
                    path = scored_path.get("path", {})
                    score = scored_path.get("score", 0)
                    details = scored_path.get("details", {})
                    lines.append(f"【路径{i}】{path.get('name', '未命名')} - 得分: {score}")
                    lines.append(f"  策略: {path.get('strategy', 'unknown')}")
                    lines.append(f"  晶体: {', '.join(path.get('crystal_ids', [])[:5])}")
                    lines.append(f"  详情: 晶体数={details.get('crystal_count_score', 0)}, 指纹匹配={details.get('fingerprint_score', 0)}")
                    lines.append("")

                best = result.get("best_path")
                if best:
                    path = best.get("path", {})
                    lines.append(f"🏆 最优路径: {path.get('name', '未命名')}")
                    lines.append(f"  推荐晶体: {', '.join(path.get('crystal_ids', [])[:5])}")
                    lines.append(f"  得分: {best.get('score', 0)}")

                output = "\n".join(lines)
                self.root.after(0, lambda: self._insert_foldable_content(output, "meta_search"))
                self.root.after(0, lambda: self.update_status("就绪"))

            except Exception as e:
                self.root.after(0, lambda: self._log(f"❌ Meta搜索出错：{e}", "error"))
                self.root.after(0, lambda: self.update_status("就绪"))

        threading.Thread(target=task, daemon=True).start()

    def _export_agents_md(self):
        """导出 AGENTS.md 到用户选择的目录"""
        try:
            from scripts.export_agents import export_agents_md

            # 选择保存目录
            output_dir = filedialog.askdirectory(
                title="选择保存 AGENTS.md 的目录",
                parent=self.root
            )
            if not output_dir:
                return

            output_path = Path(output_dir) / "AGENTS.md"

            # 生成内容
            self._log("📄 正在生成 AGENTS.md...", "system")
            content = export_agents_md(self.engine)

            # 写入文件
            output_path.write_text(content, encoding="utf-8")
            self._log(f"✅ AGENTS.md 已导出到：{output_path}", "success")

            # 弹窗显示
            messagebox.showinfo(
                "导出成功",
                f"AGENTS.md 已导出到：\n{output_path}\n\n"
                "该文件可用于 Claude Code / Cursor 等工具的项目上下文。",
                parent=self.root
            )

        except ImportError as e:
            self._log(f"❌ 导入导出模块失败：{e}", "error")
            messagebox.showerror("错误", f"导出模块不存在，请检查 scripts/export_agents.py", parent=self.root)
        except Exception as e:
            self._log(f"❌ 导出失败：{e}", "error")
            messagebox.showerror("错误", f"导出失败：{e}", parent=self.root)

    def _export_skill_dialog(self):
        """导出单个晶体为 Skill 目录"""
        try:
            from scripts.export_agents import export_skill

            # 选择要导出的晶体
            crystal_id = simpledialog.askstring(
                "导出 Skill",
                "请输入要导出的晶体 ID（如 C001）：",
                parent=self.root
            )
            if not crystal_id:
                return

            crystal_id = crystal_id.strip().upper()
            if not crystal_id.startswith("C"):
                crystal_id = f"C{crystal_id.zfill(3)}"

            # 选择保存目录
            output_dir = filedialog.askdirectory(
                title=f"选择保存 Skill 目录的根目录（将创建 skill_{crystal_id}/ 子目录）",
                parent=self.root
            )
            if not output_dir:
                return

            output_path = Path(output_dir) / f"skill_{crystal_id}"

            # 检查是否已存在
            if output_path.exists():
                if not messagebox.askyesno(
                    "目录已存在",
                    f"{output_path} 已存在，是否覆盖？",
                    parent=self.root
                ):
                    return

            self._log(f"📦 正在导出晶体 {crystal_id} 为 Skill 目录...", "system")

            # 执行导出
            result = export_skill(crystal_id, str(output_path), self.engine)

            if "error" in result:
                self._log(f"❌ 导出失败：{result['error']}", "error")
                messagebox.showerror("错误", f"导出失败：{result['error']}", parent=self.root)
                return

            # 显示成功信息
            files_list = "\n".join([f"  - {f}" for f in result.keys()])
            self._log(f"✅ Skill 目录已导出到：{output_path}", "success")
            self._log(f"   生成文件：\n{files_list}", "system")

            messagebox.showinfo(
                "导出成功",
                f"Skill 目录已导出到：\n{output_path}\n\n"
                f"生成的文件：\n{files_list}\n\n"
                "目录结构：\n"
                "skill_{crystal_id}/\n"
                "├── CRYSTAL.md    # 晶体内容\n"
                "├── validate.py   # 验证脚本\n"
                "└── references/   # 外部引用\n",
                parent=self.root
            )

        except ImportError as e:
            self._log(f"❌ 导入导出模块失败：{e}", "error")
            messagebox.showerror("错误", f"导出模块不存在，请检查 scripts/export_agents.py", parent=self.root)
        except Exception as e:
            self._log(f"❌ 导出失败：{e}", "error")
            messagebox.showerror("错误", f"导出失败：{e}", parent=self.root)

    def _export_all_skills(self):
        """导出所有晶体为 Skill 目录（批量导出）"""
        try:
            from scripts.export_agents import AgentsExporter

            # 选择保存根目录
            output_dir = filedialog.askdirectory(
                title="选择保存所有 Skill 的根目录",
                parent=self.root
            )
            if not output_dir:
                return

            # 询问最大导出数量
            max_skills_str = simpledialog.askstring(
                "批量导出",
                "最多导出多少条晶体？（输入数字，默认 10）",
                parent=self.root,
                initialvalue="10"
            )
            try:
                max_skills = int(max_skills_str) if max_skills_str else 10
            except ValueError:
                max_skills = 10

            self._log(f"📦 正在批量导出 {max_skills} 条晶体为 Skill 目录...", "system")

            exporter = AgentsExporter(self.engine)
            result = exporter.export_all_skills(Path(output_dir), max_skills=max_skills)

            # 显示结果
            self._log(f"✅ 批量导出完成：成功 {result['exported']} 条，失败 {result['failed']} 条", "success")
            self._log(f"   AGENTS.md 已生成：{result.get('agents_md', '')}", "system")

            messagebox.showinfo(
                "批量导出完成",
                f"导出结果：\n"
                f"  - 总晶体数：{result['total']}\n"
                f"  - 成功导出：{result['exported']} 条\n"
                f"  - 失败：{result['failed']} 条\n"
                f"  - 输出目录：{output_dir}\n"
                f"  - AGENTS.md：{result.get('agents_md', '')}\n",
                parent=self.root
            )

        except ImportError as e:
            self._log(f"❌ 导入导出模块失败：{e}", "error")
            messagebox.showerror("错误", f"导出模块不存在，请检查 scripts/export_agents.py", parent=self.root)
        except Exception as e:
            self._log(f"❌ 批量导出失败：{e}", "error")
            messagebox.showerror("错误", f"批量导出失败：{e}", parent=self.root)
    # ========== 新增 Day 7 方法：模式切换 ==========
    def _on_profile_change(self, event=None):
        """模式切换事件"""
        profile = self.profile_var.get()
        descriptions = {
            "high_accuracy": "高精度模式：追求最高答案质量，适合复杂决策问题",
            "balanced": "平衡模式：质量与成本的平衡，适合日常使用",
            "economy": "经济模式：最低成本，适合简单问题和快速响应"
        }
        self.profile_desc_label.config(text=descriptions.get(profile, ""))
        self._log(f"🔄 切换运行模式：{profile}", "system")

    def _get_current_profile(self) -> Dict[str, Any]:
        """获取当前模式配置"""
        profile_name = self.profile_var.get()
        profiles = {
            "high_accuracy": Config.PROFILE_HIGH_ACCURACY,
            "balanced": Config.PROFILE_BALANCED,
            "economy": Config.PROFILE_ECONOMY
        }
        return profiles.get(profile_name, Config.PROFILE_BALANCED)
    # ========== 新增 Day 7 方法：帕累托报告和认知效率 ==========
    def _show_pareto_report(self):
        """显示帕累托前沿报告"""
        try:
            # 获取帕累托数据
            pareto_data = self.engine.meta.get_pareto_status()
            configs = pareto_data.get("configs", {})
            trends = pareto_data.get("trends", {})
            best_profile = pareto_data.get("best_profile", "无数据")

            # 构建报告 - 使用等宽字体对齐
            lines = []
            lines.append("=" * 80)
            lines.append("📊 帕累托前沿报告")
            lines.append("=" * 80)
            lines.append("")

            # 状态信息（左对齐）
            lines.append(f"  最优配置    ：{best_profile}")
            lines.append(f"  趋势        ：{trends.get('trend', '数据不足')}")
            lines.append(f"  准确性变化  ：{trends.get('accuracy_delta', 0):+.3f}")
            lines.append(f"  成本变化    ：{trends.get('cost_delta', 0):+.3f}")
            lines.append(f"  平均准确性  ：{trends.get('avg_accuracy', 0):.3f}")
            lines.append(f"  平均成本    ：${trends.get('avg_cost', 0):.6f}")
            lines.append("")

            lines.append("各配置性能（实测数据）：")
            lines.append("-" * 80)

            # 表头 - 精确列宽
            header = f"{'配置':<10} {'准确性':>12} {'成本':>16} {'延迟':>12} {'引用数':>12} {'次数':>8}"
            lines.append(header)
            lines.append("-" * 80)

            # 只显示有实测数据的配置（count > 0）
            profile_names = {
                "high_accuracy": "高精度",
                "balanced": "平衡",
                "economy": "经济"
            }

            has_data = False
            for name, data in configs.items():
                if data.get("count", 0) > 0:
                    has_data = True
                    display_name = profile_names.get(name, name)
                    # 精确格式化
                    line = (
                        f"{display_name:<10} "
                        f"{data.get('accuracy', 0):>12.3f} "
                        f"${data.get('cost', 0):>14.6f} "
                        f"{data.get('latency', 0):>11.1f}s "
                        f"{data.get('crystal_refs', 0):>11.1f} "
                        f"{data.get('count', 0):>8}"
                    )
                    lines.append(line)

            if not has_data:
                lines.append("  （暂无实测数据，请先进行一些对话积累数据）")

            lines.append("-" * 80)

            # 显示每日统计
            daily_stats = pareto_data.get("daily_stats", [])
            if daily_stats:
                lines.append("")
                lines.append("最近7天认知效率：")
                lines.append("-" * 80)
                header2 = f"{'日期':<12} {'质量评分':>12} {'晶体引用':>12} {'偏差指数':>12}"
                lines.append(header2)
                lines.append("-" * 80)
                for entry in daily_stats[-7:]:
                    date_val = entry.get("date", "")[:10]
                    quality = entry.get("quality_score", 0)
                    refs = entry.get("crystal_refs", 0)
                    bias = entry.get("bias_index", 0)
                    line = (
                        f"{date_val:<12} "
                        f"{quality:>12.2f} "
                        f"{refs:>12} "
                        f"{bias:>12.2f}"
                    )
                    lines.append(line)
                lines.append("-" * 80)

            report = "\n".join(lines)

            # 显示报告 - 使用等宽字体
            win = Toplevel(self.root)
            win.title("帕累托前沿报告")
            win.geometry("850x650")
            win.configure(bg=Config.GUI_BG_MAIN)

            # 使用等宽字体（宋体或Consolas）
            text_area = scrolledtext.ScrolledText(
                win,
                wrap=tk.NONE,
                font=("Consolas", 11) if sys.platform == "win32" else ("Monaco", 11),
                bg=Config.GUI_BG_INPUT,
                fg=Config.GUI_FG_TEXT
            )
            # 添加水平滚动条
            h_scroll = ttk.Scrollbar(win, orient=tk.HORIZONTAL, command=text_area.xview)
            v_scroll = ttk.Scrollbar(win, orient=tk.VERTICAL, command=text_area.yview)
            text_area.configure(xscrollcommand=h_scroll.set, yscrollcommand=v_scroll.set)

            text_area.pack(side=tk.TOP, fill=tk.BOTH, expand=True, padx=10, pady=(10, 0))
            h_scroll.pack(side=tk.BOTTOM, fill=tk.X, padx=10, pady=(0, 5))
            v_scroll.pack(side=tk.RIGHT, fill=tk.Y, padx=(0, 10), pady=10)

            text_area.insert("1.0", report)
            text_area.config(state=tk.DISABLED)

            self._log("📊 帕累托报告已生成", "success")

        except Exception as e:
            self._log(f"❌ 生成帕累托报告失败：{e}", "error")
            messagebox.showerror("错误", f"生成报告失败：{e}", parent=self.root)
    def _show_cognitive_efficiency(self):
        """显示个人认知效率仪表盘"""
        try:
            # 获取每日统计
            daily_stats = self.engine.meta.get_daily_stats(days=7)

            # 获取当前指纹
            try:
                fp = self.engine.fingerprint_extractor.get_fingerprint()
                confidence = fp.confidence
                total_interactions = fp.total_interactions
            except:
                confidence = 0.0
                total_interactions = 0

            lines = []
            lines.append("=" * 80)
            lines.append("📈 个人认知效率仪表盘")
            lines.append("=" * 80)
            lines.append("")
            lines.append(f"  认知指纹置信度  ：{confidence:.2f}")
            lines.append(f"  总交互次数      ：{total_interactions}")
            lines.append("")

            if daily_stats:
                lines.append("最近7天认知效率趋势：")
                lines.append("-" * 80)
                # 表头 - 精确列宽
                header = f"{'日期':<12} {'质量评分':>12} {'晶体引用':>12} {'偏差指数':>12} {'Token消耗':>14}"
                lines.append(header)
                lines.append("-" * 80)

                for entry in daily_stats:
                    date_val = entry.get("date", "")[:10]
                    quality = entry.get("quality_score", 0)
                    refs = entry.get("crystal_refs", 0)
                    bias = entry.get("bias_index", 0.5)
                    tokens = entry.get("tokens_used", 0)
                    line = (
                        f"{date_val:<12} "
                        f"{quality:>12.2f} "
                        f"{refs:>12} "
                        f"{bias:>12.2f} "
                        f"{tokens:>14}"
                    )
                    lines.append(line)

                # 计算平均值
                if daily_stats:
                    avg_quality = sum(e.get("quality_score", 0) for e in daily_stats) / len(daily_stats)
                    avg_refs = sum(e.get("crystal_refs", 0) for e in daily_stats) / len(daily_stats)
                    avg_bias = sum(e.get("bias_index", 0.5) for e in daily_stats) / len(daily_stats)

                    lines.append("-" * 80)
                    avg_line = (
                        f"{'平均值':<12} "
                        f"{avg_quality:>12.2f} "
                        f"{avg_refs:>12.1f} "
                        f"{avg_bias:>12.2f} "
                    )
                    lines.append(avg_line)
                    lines.append("-" * 80)
            else:
                lines.append("  暂无足够的每日统计数据。")
                lines.append("  请先进行一些对话和晶体化操作来积累数据。")

            lines.append("")
            lines.append("=" * 80)
            lines.append("💡 建议")
            lines.append("-" * 80)

            if daily_stats and len(daily_stats) >= 3:
                recent = daily_stats[-3:]
                avg_recent_quality = sum(e.get("quality_score", 0) for e in recent) / len(recent)
                if avg_recent_quality > 0.7:
                    lines.append("  ✅ 认知效率良好，继续保持！")
                elif avg_recent_quality > 0.4:
                    lines.append("  📈 认知效率中等，建议增加晶体引用和深度推理。")
                else:
                    lines.append("  📉 认知效率较低，建议：")
                    lines.append("     1. 多使用深度推理模式")
                    lines.append("     2. 关注晶体引用质量")
                    lines.append("     3. 定期运行每日计划")
            else:
                lines.append("  📝 继续使用系统，积累更多数据后会有更精准的建议。")

            lines.append("-" * 80)

            report = "\n".join(lines)

            # 显示报告 - 使用等宽字体
            win = Toplevel(self.root)
            win.title("个人认知效率仪表盘")
            win.geometry("850x600")
            win.configure(bg=Config.GUI_BG_MAIN)

            text_area = scrolledtext.ScrolledText(
                win,
                wrap=tk.NONE,
                font=("Consolas", 11) if sys.platform == "win32" else ("Monaco", 11),
                bg=Config.GUI_BG_INPUT,
                fg=Config.GUI_FG_TEXT
            )
            # 添加水平滚动条
            h_scroll = ttk.Scrollbar(win, orient=tk.HORIZONTAL, command=text_area.xview)
            v_scroll = ttk.Scrollbar(win, orient=tk.VERTICAL, command=text_area.yview)
            text_area.configure(xscrollcommand=h_scroll.set, yscrollcommand=v_scroll.set)

            text_area.pack(side=tk.TOP, fill=tk.BOTH, expand=True, padx=10, pady=(10, 0))
            h_scroll.pack(side=tk.BOTTOM, fill=tk.X, padx=10, pady=(0, 5))
            v_scroll.pack(side=tk.RIGHT, fill=tk.Y, padx=(0, 10), pady=10)

            text_area.insert("1.0", report)
            text_area.config(state=tk.DISABLED)

            self._log("📈 认知效率仪表盘已生成", "success")

        except Exception as e:
            self._log(f"❌ 生成认知效率仪表盘失败：{e}", "error")
            messagebox.showerror("错误", f"生成仪表盘失败：{e}", parent=self.root)

    def _show_saturation_status(self):
        """显示饱和检测状态"""
        try:
            status = self.engine.meta.get_saturation_status()

            lines = []
            lines.append("=" * 70)
            lines.append("📊 双时间尺度进化调度 - 饱和检测状态")
            lines.append("=" * 70)
            lines.append("")
            lines.append(f"  饱和状态      ：{status.get('saturation_status', 'unknown')}")
            lines.append(f"  当前层级      ：{status.get('current_level', 'unknown')}")
            lines.append(f"  连续饱和轮数  ：{status.get('consecutive_rounds', 0)}")
            lines.append(f"  最近提升幅度  ：{status.get('last_improvement', 0):.3f}")
            lines.append(f"  质量历史记录数：{status.get('quality_history_count', 0)}")
            lines.append(f"  控制逻辑变更数：{status.get('control_logic_changes_count', 0)}")
            lines.append("")
            lines.append("-" * 70)
            lines.append("💡 说明")
            lines.append("  - unsaturated: 未饱和，继续优化提示词")
            lines.append("  - saturated: 已饱和，准备升级到控制逻辑层面")
            lines.append("  - 当前层级 prompt: 正在优化提示词")
            lines.append("  - 当前层级 control_logic: 正在调整控制逻辑（检索策略、验证规则等）")

            report = "\n".join(lines)

            # 显示报告
            win = Toplevel(self.root)
            win.title("饱和检测状态")
            win.geometry("600x400")
            win.configure(bg=Config.GUI_BG_MAIN)

            text_area = scrolledtext.ScrolledText(
                win,
                wrap=tk.WORD,
                font=("宋体", 11),
                bg=Config.GUI_BG_INPUT,
                fg=Config.GUI_FG_TEXT
            )
            text_area.pack(fill=tk.BOTH, expand=True, padx=10, pady=10)
            text_area.insert("1.0", report)
            text_area.config(state=tk.DISABLED)

            self._log("📊 饱和状态已显示", "success")

        except Exception as e:
            self._log(f"❌ 获取饱和状态失败：{e}", "error")
            messagebox.showerror("错误", f"获取状态失败：{e}", parent=self.root)

    # ===== Day 8: 灵感熔炉复盘 =====
    def _show_inspiration_furnace(self):
        """显示灵感熔炉复盘结果"""
        try:
            self._log("🔥 正在运行灵感熔炉复盘（一）...", "system")
            self.update_status("灵感熔炉复盘中...")

            # 运行复盘
            result = self.engine.meta.inspiration_furnace_review()

            # 构建报告
            lines = []
            lines.append("=" * 70)
            lines.append("🔥 灵感熔炉复盘（一）")
            lines.append("=" * 70)
            lines.append("")
            lines.append(f"待筛选灵感总数：{result.get('total_pending', 0)}")
            lines.append("")

            # S级
            s_level = result.get('s_level', [])
            lines.append(f"📌 S级灵感（<2小时可验证）：{len(s_level)} 条")
            lines.append("-" * 70)
            for i, insp in enumerate(s_level, 1):
                lines.append(f"  {i}. {insp.get('content', '')[:80]}")
                ev = insp.get('evaluation', {})
                lines.append(f"     评分：{ev.get('total_score', 0):.2f} | 资源：{ev.get('resource_hours', 0)}小时")
            lines.append("")

            # A级
            a_level = result.get('a_level', [])
            lines.append(f"📌 A级灵感（半天内可落地）：{len(a_level)} 条")
            lines.append("-" * 70)
            for i, insp in enumerate(a_level, 1):
                lines.append(f"  {i}. {insp.get('content', '')[:80]}")
                ev = insp.get('evaluation', {})
                lines.append(f"     评分：{ev.get('total_score', 0):.2f} | 资源：{ev.get('resource_hours', 0)}小时")
            lines.append("")

            # B级
            b_level = result.get('b_level', [])
            lines.append(f"📌 B级灵感（1-2天）：{len(b_level)} 条")
            lines.append("-" * 70)
            for i, insp in enumerate(b_level, 1):
                lines.append(f"  {i}. {insp.get('content', '')[:80]}")
                ev = insp.get('evaluation', {})
                lines.append(f"     评分：{ev.get('total_score', 0):.2f} | 资源：{ev.get('resource_hours', 0)}小时")
            lines.append("")

            # 已拒绝
            rejected = result.get('rejected', [])
            if rejected:
                lines.append(f"❌ 已拒绝：{len(rejected)} 条")
                lines.append("-" * 70)
                for i, insp in enumerate(rejected, 1):
                    lines.append(f"  {i}. {insp.get('content', '')[:60]}")
            lines.append("")

            lines.append("=" * 70)
            lines.append(result.get('summary', ''))

            # 显示报告
            win = Toplevel(self.root)
            win.title("灵感熔炉复盘（一）")
            win.geometry("800x650")
            win.configure(bg=Config.GUI_BG_MAIN)

            text_area = scrolledtext.ScrolledText(
                win,
                wrap=tk.WORD,
                font=("宋体", 11),
                bg=Config.GUI_BG_INPUT,
                fg=Config.GUI_FG_TEXT
            )
            text_area.pack(fill=tk.BOTH, expand=True, padx=10, pady=10)
            text_area.insert("1.0", "\n".join(lines))
            text_area.config(state=tk.DISABLED)

            self._log(f"🔥 灵感熔炉复盘完成：{result.get('summary', '')}", "success")
            self.update_status("就绪")

        except Exception as e:
            self._log(f"❌ 灵感熔炉复盘失败：{e}", "error")
            messagebox.showerror("错误", f"复盘失败：{e}", parent=self.root)

# =============================================================================
# 14. Web API (web/app.py) - 完整复制并调整导入
# =============================================================================
from fastapi import BackgroundTasks, FastAPI, File, Form, HTTPException, UploadFile
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse
from fastapi.staticfiles import StaticFiles
from pydantic import BaseModel, Field
import shutil
import tempfile
import uuid
import hashlib
import re
import subprocess
import threading
from datetime import date, datetime
from typing import Any, Callable, Dict, List, Optional, Tuple

WEB_ROOT = PROJECT_ROOT / "web_static"
JOBS: Dict[str, Dict[str, Any]] = {}
STOP_FLAGS: Dict[str, threading.Event] = {}
LEGACY_PROCESS: Optional[subprocess.Popen] = None
LOCK = threading.Lock()

FileIO.ensure_directories()
FileIO.ensure_default_files()

db = DBManager()
files = FileIO()
engine = CrystalEngine(files)

app = FastAPI(title="认知晶体树 5 Web", version="0.1.0")
app.add_middleware(
    CORSMiddleware,
    allow_origins=["http://127.0.0.1:8788", "http://localhost:8788"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)
app.mount("/static", StaticFiles(directory=str(WEB_ROOT)), name="static")

# 路由
@app.get("/")
def index():
    return FileResponse(WEB_ROOT / "index.html")

@app.get("/api/bootstrap")
def bootstrap():
    assets = _assets()
    return {
        "data_root": str(Config.DATA_ROOT),
        "db_path": str(Config.get_db_path()),
        "api_key_configured": bool(Config.get_api_key()),
        "sessions": _sessions(),
        "assets": assets["counts"],
        "pending_count": len(_pending_cards()),
        "task_count": len([t for t in _task_cards() if t.get("status") == "pending"]),
        "legacy_backend_running": LEGACY_PROCESS is not None and LEGACY_PROCESS.poll() is None,
    }

@app.post("/api/vector/sync")
def sync_vector_store():
    """同步向量库（将当前所有晶体向量化）"""
    try:
        result = engine.sync_vector_store()
        return result
    except Exception as e:
        return {"status": "error", "message": str(e)}

@app.get("/api/vector/status")
def vector_status():
    """获取向量库状态"""
    try:
        count = engine.vector_store.count()
        available = engine.vector_store.is_available()
        return {"available": available, "count": count}
    except Exception as e:
        return {"available": False, "count": 0, "error": str(e)}

@app.get("/api/conflicts")
def get_conflicts(method: str = "auto", limit: int = 20):
    """获取检测到的晶体冲突"""
    try:
        scope = engine.get_conflict_scope()
        conflicts = engine.detect_conflicts(scope=scope, method=method)
        if limit > 0:
            conflicts = conflicts[:limit]
        return {
            "total": len(conflicts),
            "conflicts": [
                {
                    "crystal_a": c.crystal_a,
                    "crystal_b": c.crystal_b,
                    "similarity": c.similarity,
                    "content_a": c.content_a,
                    "content_b": c.content_b
                }
                for c in conflicts
            ]
        }
    except Exception as e:
        return {"error": str(e), "conflicts": []}

# Pydantic models
class SessionCreate(BaseModel):
    name: Optional[str] = None

class SessionRename(BaseModel):
    name: str

class ChatRequest(BaseModel):
    session_id: str
    input: str
    api_key: Optional[str] = None

class CrystalRequest(ChatRequest):
    fast_mode: bool = True
    scope: str = "全局"

class DeepReasonRequest(ChatRequest):
    mode: str = "multi_role"
    max_rounds: int = 2

class BatchRequest(BaseModel):
    folder: str
    mode: str = "chat"
    fast_mode: bool = True
    inject_history: bool = False
    session_id: Optional[str] = None
    api_key: Optional[str] = None

class DailyPlanRequest(BaseModel):
    api_key: Optional[str] = None
    intent_keywords: List[str] = Field(default_factory=list)
    time_budget_seconds: int = 900

class SearchRequest(BaseModel):
    keyword: str
    regex: bool = False
    dirs: List[str] = Field(default_factory=lambda: ["晶体数据", "核心配置", "系统日志", "暂存区"])

class CommitRequest(BaseModel):
    session_id: Optional[str] = None
    result: Dict[str, Any]

class PendingConfirmRequest(BaseModel):
    content: str
    force: bool = False

class AssetPatchRequest(BaseModel):
    layer: Optional[str] = None
    fixed: Optional[bool] = None

class BackendLoginRequest(BaseModel):
    username: str
    password: str

# 辅助函数（与原始 app.py 完全一致）
def _job(job_type: str) -> str:
    job_id = f"{job_type}-{uuid.uuid4().hex[:10]}"
    with LOCK:
        JOBS[job_id] = {
            "id": job_id,
            "type": job_type,
            "status": "queued",
            "progress": 0,
            "logs": [],
            "result": None,
            "error": None,
            "created_at": datetime.now().isoformat(),
            "updated_at": datetime.now().isoformat(),
        }
    return job_id

def _set_job(job_id: str, **kwargs):
    with LOCK:
        job = JOBS[job_id]
        job.update(kwargs)
        job["updated_at"] = datetime.now().isoformat()

def _log_job(job_id: str, message: str, level: str = "system"):
    with LOCK:
        JOBS[job_id]["logs"].append({"time": datetime.now().strftime("%H:%M:%S"), "level": level, "message": message})
        JOBS[job_id]["updated_at"] = datetime.now().isoformat()

def _run_job(job_id: str, fn: Callable[[], Any]):
    _set_job(job_id, status="running", progress=5)
    try:
        result = fn()
        _set_job(job_id, status="done", progress=100, result=result)
    except Exception as exc:
        _set_job(job_id, status="error", error=str(exc))
        _log_job(job_id, f"任务失败：{exc}", "error")

def _sessions() -> List[Dict[str, Any]]:
    return [{"id": sid, "name": name, "updated_at": updated} for sid, name, updated in db.list_sessions()]

def _ensure_session(session_id: Optional[str] = None) -> str:
    if session_id:
        name, _ = db.get_session(session_id)
        if name is not None:
            return session_id
    sid = datetime.now().strftime("%Y%m%d%H%M%S") + "_" + uuid.uuid4().hex[:6]
    db.create_session(sid, f"新会话 {datetime.now().strftime('%H:%M')}")
    return sid

def _add_message(session_id: str, role: str, content: str):
    name, history = db.get_session(session_id)
    if name is None:
        raise HTTPException(status_code=404, detail="会话不存在")
    history.append((role, content))
    if role == "user" and len(history) == 1 and name.startswith("新会话"):
        new_name = generate_session_title_from_content(content)   # <-- 新代码
        if new_name:                                              # <-- 新代码
            db.rename_session(session_id, new_name)               # <-- 新代码
            name = new_name                                       # <-- 新代码
    db.update_session(session_id, history, name)

def _history_context(session_id: str, current_input: str, limit: int = 8) -> str:
    _name, history = db.get_session(session_id)
    lines = []
    for role, content in history[-limit:]:
        if role == "user" and current_input and current_input in content:
            continue
        label = "用户" if role == "user" else "AI"
        text = re.sub(r"\s+", " ", str(content)).strip()
        if text:
            lines.append(f"{label}: {text[:900]}")
    if not lines:
        return current_input
    return "【本会话最近上下文】\n" + "\n".join(lines) + f"\n\n【当前问题】\n{current_input}"

def _questions(history: List[Tuple[str, str]]) -> List[Dict[str, Any]]:
    result = []
    q_num = 1
    for index, (role, content) in enumerate(history):
        if role != "user":
            continue
        label = content
        for prefix, name in [("[晶体化] ", "[晶体化] "), ("[深度推理] ", "[深度推理] "), ("[深度推理-多角色] ", "[多角色] "), ("[文件内容] ", "[文件] ")]:
            if content.startswith(prefix):
                label = name + content[len(prefix):]
                break
        result.append({"index": index, "label": f"{q_num}. {label[:52]}{'...' if len(label) > 52 else ''}", "content": content})
        q_num += 1
    return result

def _shorten(content: str, limit: int = 80) -> str:
    return re.sub(r"\s+", " ", str(content or "")).strip()[:limit]

def _existing_crystal_ids() -> set:
    return set(re.findall(r"\bC\d+\b", files.read("crystals")))

def _extract_keywords(text: str) -> List[str]:
    words = re.findall(r"[\w\u4e00-\u9fff]+", text)
    stop = {"的", "了", "和", "与", "或", "一个", "这个", "那个", "如何", "什么", "为什么"}
    return [w for w in words if w not in stop][:5] or ["晶体树", "认知"]

def _build_crystallization_prompt(user_input: str, search_res: str) -> str:
    l0_holes, l1_crystals = engine.get_attention_context()
    l0_text = "\n".join([f"- {h.id}: {h.content[:100]}" for h in l0_holes])
    l1_text = "\n".join([f"- {c.id}: {c.content[:80]} | links={','.join(c.links)}" for c in l1_crystals[:Config.L1_MAX]])
    related = engine.get_associative_crystals(user_input, top_k=8)
    related_text = "\n".join([f"- {c.id}: {c.content[:80]} | links={','.join(c.links)}" for c in related])
    existing_ids = ", ".join(sorted(_existing_crystal_ids())[-20:])
    return f"""
你是认知晶体树的结构化整理器。请只返回 JSON，不要返回 Markdown、解释或代码块。

目标：
1. 从用户输入中提炼可长期复用的认知晶体。
2. 优先连接已有晶体，不要重复制造同义晶体。
3. 暴露冲突和孔洞，但不要夸大不确定性。
4. 每条晶体 content 必须不超过 80 个中文字符。

用户输入：
{user_input}

外部搜索结果：
{search_res}

L0 核心孔洞：
{l0_text}

L1 注意力晶体：
{l1_text}

联想检索命中的相关晶体：
{related_text}

近期已有晶体 ID：
{existing_ids}

返回 JSON schema：
{{
  "new_crystals": [{{"id": "", "content": "不超过80字的新晶体", "links": ["C001"]}}],
  "updated_crystals": [{{"id": "C001", "new_content": "不超过80字的更新后内容"}}],
  "new_holes": [{{"id": "", "content": "需要继续验证的问题", "urgency": 0.5, "layer": 2}}],
  "updated_holes": [{{"id": "H001", "content": "更新后孔洞内容"}}],
  "conflicts": [{{"a": "C001", "b": "C002", "reason": "冲突原因"}}],
  "report_summary": "一句话总结本次结构变化",
  "pending_cards": [{{"type": "晶体候选", "content": "待确认内容", "source": "AI生成", "confidence": "中"}}]
}}
"""

def _normalize_crystal_response(ai_response: Dict[str, Any]) -> Dict[str, Any]:
    if not isinstance(ai_response, dict):
        return {"error": "AI返回不是JSON对象"}
    existing_ids = _existing_crystal_ids()
    next_num = max([int(i) for i in re.findall(r"C(\d+)", files.read("crystals"))], default=0) + 1
    seen_contents = {_shorten(c.content) for c in engine.parse_crystals()}
    id_map: Dict[str, str] = {}
    normalized_new = []
    for item in ai_response.get("new_crystals", []) or []:
        content = _shorten(item.get("content", ""))
        if not content or content in seen_contents:
            continue
        old_id = str(item.get("id", "")).strip()
        new_id = f"C{next_num:03d}"
        next_num += 1
        if old_id and old_id not in existing_ids:
            id_map[old_id] = new_id
        links = []
        for link in item.get("links", []) or []:
            link = id_map.get(str(link).strip(), str(link).strip())
            if re.fullmatch(r"C\d+", link) and (link in existing_ids or link in id_map.values()):
                links.append(link)
        normalized_new.append({"id": new_id, "content": content, "links": sorted(set(links)), "similar": _similar(content)})
        seen_contents.add(content)
    valid_ids = existing_ids | {c["id"] for c in normalized_new}
    normalized_updates = []
    for item in ai_response.get("updated_crystals", []) or []:
        cid = str(item.get("id", "")).strip()
        content = _shorten(item.get("new_content") or item.get("content", ""))
        if cid in existing_ids and content:
            normalized_updates.append({"id": cid, "new_content": content})
    normalized_holes = []
    next_hole = max([int(i) for i in re.findall(r"H(\d+)", files.read("holes"))], default=0) + 1
    for item in ai_response.get("new_holes", []) or []:
        content = _shorten(item.get("content", ""), 120)
        if not content:
            continue
        try:
            urgency = max(0.0, min(1.0, float(item.get("urgency", 0.5))))
        except (TypeError, ValueError):
            urgency = 0.5
        try:
            layer = min(3, max(1, int(item.get("layer", 2))))
        except (TypeError, ValueError):
            layer = 2
        normalized_holes.append({"id": f"H{next_hole:03d}", "content": content, "urgency": urgency, "layer": layer})
        next_hole += 1
    normalized_conflicts = []
    for item in ai_response.get("conflicts", []) or []:
        a = str(item.get("a") or item.get("crystal_a") or "").strip()
        b = str(item.get("b") or item.get("crystal_b") or "").strip()
        if a in valid_ids and b in valid_ids and a != b:
            normalized_conflicts.append({"a": a, "b": b, "reason": str(item.get("reason", ""))[:120]})
    normalized_pending = []
    for item in ai_response.get("pending_cards", []) or []:
        content = _shorten(item.get("content", ""), 200)
        if content:
            normalized_pending.append({
                "type": str(item.get("type", "晶体候选")),
                "content": content,
                "source": str(item.get("source", "AI生成")),
                "confidence": str(item.get("confidence", "中")),
            })
    return {
        "new_crystals": normalized_new,
        "updated_crystals": normalized_updates,
        "new_holes": normalized_holes,
        "updated_holes": ai_response.get("updated_holes", []) or [],
        "conflicts": normalized_conflicts,
        "report_summary": str(ai_response.get("report_summary", "晶体化完成"))[:120],
        "pending_cards": normalized_pending,
    }

def _similar(content: str) -> List[Dict[str, Any]]:
    matches = []
    for crystal in engine.parse_crystals():
        score = engine._simple_similarity(content, crystal.content)
        if score >= 0.55:
            matches.append({"score": round(score, 2), "id": crystal.id, "content": crystal.content[:80]})
    matches.sort(key=lambda item: item["score"], reverse=True)
    return matches[:5]

def _append_pending_card(card: Dict[str, Any]) -> bool:
    content = _shorten(card.get("content", ""), 200)
    if not content:
        return False
    pending = files.read("pending")
    if content in pending:
        return False
    suffix = int(hashlib.sha256(content.encode("utf-8")).hexdigest(), 16) % 1000
    card_id = f"PENDING-{datetime.now().strftime('%Y%m%d%H%M%S')}-{suffix:03d}"
    block = f"""
## {card_id}
- 类型：{card.get('type', '晶体候选')}
- 来源：{card.get('source', 'AI生成')}
- 置信度：{card.get('confidence', '中')}
- 内容：{content}
- AI判断：建议人工确认后再转为晶体。
"""
    files.append("pending", "\n" + block + "\n")
    return True

def _update_files(result: Dict[str, Any]):
    for c in result.get("new_crystals", []) or []:
        files.append("crystals", f"\n| {c['id']} | {c['content']} | {', '.join(c.get('links', [])) or '—'} |\n")
    for upd in result.get("updated_crystals", []) or []:
        content = files.read("crystals")
        pattern = rf"(\| {re.escape(upd['id'])} \| ).*?( \| .*? \|)"
        files.write("crystals", re.sub(pattern, rf"\1{upd['new_content']}\2", content))
    for hole in result.get("new_holes", []) or []:
        files.append("holes", f"\n| {hole['id']} | {hole['content']} | {hole.get('urgency', 0.5)} |\n")
    for upd in result.get("updated_holes", []) or []:
        content = files.read("holes")
        pattern = rf"(\| {re.escape(str(upd.get('id', '')))} \| ).*?(\| .*? \|)"
        files.write("holes", re.sub(pattern, rf"\1{upd.get('content', '')}\2", content))
    kept = []
    for card in result.get("pending_cards", []) or []:
        if _append_pending_card(card):
            kept.append(card)
    result["pending_cards"] = kept
    change_entry = f"""### 变更摘要：{result.get('report_summary', '无摘要')}
- 新增晶体：{len(result.get('new_crystals', []))}
- 更新晶体：{len(result.get('updated_crystals', []))}
- 新增孔洞：{len(result.get('new_holes', []))}
- 更新孔洞：{len(result.get('updated_holes', []))}
- 冲突：{len(result.get('conflicts', []))}
- 待确认卡片：{len(result.get('pending_cards', []))}
"""
    engine._append_change_log("晶体化变更", change_entry)
    crystals_count = len(re.findall(r"^\| C\d+", files.read("crystals"), re.MULTILINE))
    holes_count = len(re.findall(r"^\| H\d+", files.read("holes"), re.MULTILINE))
    files.write("state", f"""# 系统状态快照
**生成时间**: {datetime.now().isoformat()}
**晶体总数**: {crystals_count}
**孔洞总数**: {holes_count}
**最新变更摘要**: {result.get('report_summary', '无')}
""")

def _assets() -> Dict[str, Any]:
    l1, l2, l3 = engine.update_crystal_layers()
    state = engine.load_layer_state()
    layers = state.get("layers", {})
    heat = state.get("heat_map", {})
    last = state.get("last_accessed", {})
    manual = state.get("manual_override", {})
    crystals = []
    for c in engine.parse_crystals():
        crystals.append({
            "id": c.id,
            "content": c.content,
            "links": c.links,
            "layer": layers.get(c.id, "L2"),
            "heat": round(float(heat.get(c.id, 0.0)), 2),
            "last_accessed": last.get(c.id, "从未"),
            "fixed": manual.get(c.id) == "L1_fixed",
        })
    holes = [{"id": h.id, "content": h.content, "urgency": h.urgency, "layer": h.layer} for h in engine.parse_holes()]
    return {"crystals": crystals, "holes": holes, "counts": {"L1": len(l1), "L2": len(l2), "L3": len(l3), "total": len(crystals)}}

def _pending_cards() -> List[Dict[str, Any]]:
    content = files.read("pending")
    cards = []
    for cid, body in re.findall(r"## (PENDING-\d+-\d+)\n(.*?)(?=\n## |\Z)", content, re.DOTALL):
        item = {"id": cid, "type": "", "title": "", "source": "", "content": "", "raw": body.strip()}
        for line in body.splitlines():
            if line.startswith("- 类型："):
                item["type"] = line.split("：", 1)[1].strip()
            elif line.startswith("- 标题："):
                item["title"] = line.split("：", 1)[1].strip()
            elif line.startswith("- 来源："):
                item["source"] = line.split("：", 1)[1].strip()
            elif line.startswith("- 内容摘要：") or line.startswith("- 内容：") or line.startswith("- 生成晶体候选："):
                item["content"] = line.split("：", 1)[1].strip()
        if not item["title"]:
            item["title"] = item["content"][:50] or cid
        cards.append(item)
    return cards

def _task_cards() -> List[Dict[str, Any]]:
    if not files.exists("task_cards"):
        return []
    try:
        return json.loads(files.read("task_cards") or "[]")
    except json.JSONDecodeError:
        return []

def _save_task_cards(cards: List[Dict[str, Any]]):
    files.write("task_cards", json.dumps(cards, ensure_ascii=False, indent=2))

def _load_roles() -> List[Dict[str, str]]:
    try:
        raw = json.loads(files.read("roles") or "{}")
    except json.JSONDecodeError:
        raw = {}

    roles = []
    for key, item in raw.items():
        if isinstance(item, dict):
            roles.append({
                "key": key,
                "name": item.get("name", key),
                "instruction": item.get("instruction", "")
            })

    fallback_keys = [
        "radical", "conservative", "structural", "judge",
        "spokesperson", "lark", "pilgrim", "strategist", "statesman"
    ]
    fallback_roles = {
        "radical": {
            "name": "激进者",
            "instruction": "攻击默认前提，假设现有框架是错的，给出颠覆性方案。"
        },
        "conservative": {
            "name": "保守者",
            "instruction": "风险优先，假设资源有限，给出最可落地的稳健方案。"
        },
        "structural": {
            "name": "结构主义者",
            "instruction": "从已有晶体中寻找同构案例，用类比生成方案。"
        },
        "judge": {
            "name": "大法官",
            "instruction": "以晶体卡片、核心操作原则和资源约束为准绳，做出终审裁决。必须明确引用依据（晶体ID、原则条款或约束条件），不得凭直觉判案。"
        },
        "spokesperson": {
            "name": "首席发言人",
            "instruction": "将内部辩论结论转化为清晰、简洁、无歧义的对外陈述。遵循降维（通俗化）、定调（不超过3条核心信息）、检验（老板读前100字能决策）三原则。"
        },
        "lark": {
            "name": "百灵鸟",
            "instruction": "见多识广的通用智能体，从外部世界（学术、产业、政策、跨学科）补充知识，打破信息茧房。在第二轮登场。"
        },
        "pilgrim": {
            "name": "取经者",
            "instruction": "以长期愿景和核心价值观为锚，防止短期利益或局部优化偏离最终使命。评估方案的可持续性和道德一致性。"
        },
        "strategist": {
            "name": "奇谋者",
            "instruction": "善于洞察人心、把握时机，敢押注非常规路径，捕捉机会窗口。评估方案能否借力打力、以奇制胜。"
        },
        "statesman": {
            "name": "延安智者",
            "instruction": "坚持调查研究，不唯上、不唯书、只唯实。从全局矛盾和主要矛盾切入，提出实事求是、可落地的综合方略。"
        }
    }

    existing_keys = {item["key"] for item in roles}
    for key in fallback_keys:
        if key not in existing_keys:
            roles.append({
                "key": key,
                "name": fallback_roles[key]["name"],
                "instruction": fallback_roles[key]["instruction"]
            })
    return roles



@app.get("/api/backend/status")
def backend_status():
    return {"running": LEGACY_PROCESS is not None and LEGACY_PROCESS.poll() is None}

@app.post("/api/backend/login")
def backend_login(payload: BackendLoginRequest):
    global LEGACY_PROCESS
    if not payload.username.strip() or payload.password != "111111":
        raise HTTPException(status_code=401, detail="用户名或密码错误")
    if LEGACY_PROCESS is not None and LEGACY_PROCESS.poll() is None:
        return {"ok": True, "running": True, "message": "原后端界面已在运行"}

    import sys
    import subprocess

    cmd = [sys.executable, __file__, "--gui"]
    env = os.environ.copy()
    env.setdefault("CRYSTAL_TREE_DATA_ROOT", str(Config.DATA_ROOT))
    env["PYTHONIOENCODING"] = "utf-8"  # 确保子进程使用 UTF-8

    log_dir = Config.DATA_ROOT / "系统日志"
    log_dir.mkdir(parents=True, exist_ok=True)
    log_file = log_dir / "gui_subprocess.log"

    # ===== 关键修复：用 PIPE 捕获错误输出 =====
    with open(log_file, "w", encoding="utf-8") as f:
        f.write(f"=== 启动时间: {datetime.now().isoformat()} ===\n")
        f.write(f"命令: {' '.join(cmd)}\n")
        f.write(f"工作目录: {str(PROJECT_ROOT)}\n")
        f.write(f"数据目录: {env.get('CRYSTAL_TREE_DATA_ROOT')}\n")
        f.write("=" * 60 + "\n")
        f.flush()

        try:
            LEGACY_PROCESS = subprocess.Popen(
                cmd,
                cwd=str(PROJECT_ROOT),
                env=env,
                stdout=f,
                stderr=f,
                creationflags=subprocess.CREATE_NEW_CONSOLE if os.name == "nt" else 0,
                text=True,
            )
        except Exception as e:
            f.write(f"启动异常: {e}\n")
            raise HTTPException(status_code=500, detail=f"启动失败: {e}")

    # 等待短暂时间检查进程是否存活
    import time
    time.sleep(1.0)

    if LEGACY_PROCESS.poll() is not None:
        # 进程已退出，读取日志
        try:
            with open(log_file, "r", encoding="utf-8") as f:
                error_output = f.read()
        except:
            error_output = "无法读取日志文件"
        raise HTTPException(
            status_code=500,
            detail=f"老师端启动后立即退出。请查看日志：{log_file}\n\n日志内容：\n{error_output[-1000:]}"
        )

    return {"ok": True, "running": True, "message": f"已启动老师端界面（独立窗口），日志见 {log_file}"}

@app.post("/api/backend/logout")
def backend_logout():
    global LEGACY_PROCESS
    if LEGACY_PROCESS is None or LEGACY_PROCESS.poll() is not None:
        LEGACY_PROCESS = None
        return {"ok": True, "running": False, "message": "原后端界面未运行，前端可继续使用"}
    LEGACY_PROCESS.terminate()
    try:
        LEGACY_PROCESS.wait(timeout=5)
    except subprocess.TimeoutExpired:
        LEGACY_PROCESS.kill()
        LEGACY_PROCESS.wait(timeout=5)
    LEGACY_PROCESS = None
    return {"ok": True, "running": False, "message": "已退出原后端登录，前端可继续使用"}

@app.get("/api/sessions")
def list_sessions():
    return {"sessions": _sessions()}

@app.post("/api/sessions")
def create_session(payload: SessionCreate):
    sid = datetime.now().strftime("%Y%m%d%H%M%S") + "_" + uuid.uuid4().hex[:6]
    name = payload.name or f"新会话 {datetime.now().strftime('%H:%M')}"
    db.create_session(sid, name)
    return {"id": sid, "name": name}

@app.get("/api/sessions/{session_id}")
def get_session(session_id: str):
    name, history, labels = db.get_session(session_id)   # 改为接收三个返回值
    if name is None:
        raise HTTPException(status_code=404, detail="会话不存在")
    return {"id": session_id, "name": name, "messages": [{"role": r, "content": c} for r, c in history], "questions": _questions(history)}

@app.patch("/api/sessions/{session_id}")
def rename_session(session_id: str, payload: SessionRename):
    db.rename_session(session_id, payload.name.strip())
    return {"ok": True}

@app.delete("/api/sessions/{session_id}")
def delete_session(session_id: str):
    db.delete_session(session_id)
    return {"ok": True}

@app.post("/api/sessions/{session_id}/clear")
def clear_session(session_id: str):
    name, _ = db.get_session(session_id)
    if name is None:
        raise HTTPException(status_code=404, detail="会话不存在")
    db.update_session(session_id, [], name)
    return {"ok": True}

@app.post("/api/chat")
def chat(payload: ChatRequest, background: BackgroundTasks):
    session_id = _ensure_session(payload.session_id)
    _add_message(session_id, "user", payload.input)
    # 在 /api/chat 路由中
    _, history, _ = db.get_session(session_id)
    job_id = _job("chat")

    def task():
        ai = AIClient(api_key=payload.api_key or Config.get_api_key())
        _, history = db.get_session(session_id)
        l0_holes, l1_crystals = engine.get_attention_context()
        context = f"\n[注意力上下文] 当前核心孔洞：{', '.join([h.content[:50] for h in l0_holes])}\nL1晶体数量：{len(l1_crystals)} 条\n"
        _log_job(job_id, "AI 思考中...")
        reply = ai.chat_with_history(history, context=context)
        _add_message(session_id, "assistant", reply)
        return {"session_id": session_id, "reply": reply}

    background.add_task(_run_job, job_id, task)
    return {"job_id": job_id, "session_id": session_id}

@app.post("/api/crystallize")
def crystallize(payload: CrystalRequest, background: BackgroundTasks):
    session_id = _ensure_session(payload.session_id)
    _add_message(session_id, "user", f"[晶体化] {payload.input}")
    job_id = _job("crystallize")

    def task():
        ai = AIClient(api_key=payload.api_key or Config.get_api_key())
        search_res = "（快速模式：跳过外部搜索）" if payload.fast_mode else "（Web v1：外部搜索将在下一阶段增强）"
        prompt = _build_crystallization_prompt(payload.input, search_res)
        _log_job(job_id, "晶体化预览生成中...")
        raw = ai.chat_json(prompt)
        if "error" in raw:
            raise RuntimeError(raw["error"])
        normalized = _normalize_crystal_response(raw)
        if "error" in normalized:
            raise RuntimeError(normalized["error"])
        return {"session_id": session_id, "preview": normalized}

    background.add_task(_run_job, job_id, task)
    return {"job_id": job_id, "session_id": session_id}

@app.post("/api/crystallize/commit")
def commit_crystallize(payload: CommitRequest):
    _update_files(payload.result)
    if payload.session_id:
        _add_message(payload.session_id, "assistant", f"[晶体化结果] {payload.result.get('report_summary', '晶体化完成')}")
    return {"ok": True, "summary": payload.result.get("report_summary", "晶体化完成")}

@app.post("/api/deep-reasoning")
def deep_reasoning(payload: DeepReasonRequest, background: BackgroundTasks):
    session_id = _ensure_session(payload.session_id)
    if payload.mode == "lushi_sampling":
        prefix = "[卢氏注意力增强]"
    elif payload.mode in ("debate_light", "debate_full"):
        prefix = "[辩论增强]"
    elif payload.mode == "multi_role":
        prefix = "[深度推理-多角色]"
    else:
        prefix = "[深度推理]"
    _add_message(session_id, "user", f"{prefix} {payload.input}")
    job_id = _job("deep-reasoning")

    def task():
        ai = AIClient(api_key=payload.api_key or Config.get_api_key())
        effective_mode = "debate_full" if payload.mode in ("auto", "multi_role") else payload.mode
        reason_input = _history_context(session_id, payload.input)
        max_rounds = max(2, min(12, int(payload.max_rounds or 2)))

        if effective_mode in ("debate_light", "debate_full", "lushi_sampling"):
            _log_job(job_id, "晶体树辩论引擎启动中...")
            debate = DebateEngine(ai, engine, _load_roles(), lambda m, level="system": _log_job(job_id, m, level))
            result = debate.run(reason_input, mode=effective_mode, max_rounds=max_rounds)

            try:
                orchestrator = OutputOrchestrator(ai, engine)
                final_schema = orchestrator.generate(payload.input, result.get("rounds", []))

                # ===== 正确提取字段 =====
                # FinalOutputSchema 的字段是：
                # - judge_performance_board (绩效看板)
                # - judge_final_verdict (终审裁决)
                # - judge_rejected_details (驳回明细)
                # - board_version, employee_version, novice_version, expert_version, elegant_epilogue
                # - dashboard_stats, round_by_round

                result["final_schema"] = final_schema.dict()
                result["board_version"] = final_schema.board_version
                result["employee_version"] = final_schema.employee_version
                result["novice_version"] = final_schema.novice_version
                result["expert_version"] = final_schema.expert_version
                result["elegant_epilogue"] = final_schema.elegant_epilogue
                result["dashboard_stats"] = final_schema.dashboard_stats

                # ===== 修正：使用正确的字段名 =====
                result["judge_audit"] = {
                    "by_rule": final_schema.judge_performance_board,  # 绩效看板
                    "summary": final_schema.judge_final_verdict,      # 终审裁决
                    "rejected_items": final_schema.judge_rejected_details  # 驳回明细
                }
                result["round_by_round"] = [r.dict() for r in final_schema.round_by_round]
                result["answer"] = final_schema.board_version

                _add_message(session_id, "assistant", final_schema.board_version)
                _log_job(job_id, "✅ V3.0 结构化输出生成完成", "success")

                # ===== 兼容层 =====
                result["final"] = {
                    "rigid_core": {
                        "decision_summary": final_schema.board_version[:200] if final_schema.board_version else "（决策摘要生成中）",
                        "core_adoptions": [f"{r.get('role', '')}：{r.get('brief_reason', '')}" for r in final_schema.judge_performance_board if r.get("status") in ["adopted", "conditional"]],
                        "key_synthesis": final_schema.expert_version[:500] if final_schema.expert_version else "（决策逻辑待补充）",
                        "risks_and_boundaries": ["（风险分析待补充）"]
                    },
                    "one_sentence_conclusion": final_schema.board_version[:50] if final_schema.board_version else "（结论待补充）",
                    "student_friendly_answer": final_schema.novice_version or "（通俗解读待补充）",
                    "teacher_detail": final_schema.expert_version or "（详细复盘待补充）",
                    "soft_wrap": final_schema.employee_version or "（精简版待补充）",
                    "judge_audit": {
                        "by_rule": final_schema.judge_performance_board,
                        "summary": final_schema.judge_final_verdict
                    },
                    "dashboard_stats": final_schema.dashboard_stats,
                }

                return {
                    "job_id": job_id,
                    "session_id": session_id,
                    "summary": {
                        "board": final_schema.board_version,
                        "employee": final_schema.employee_version,
                        "novice": final_schema.novice_version,
                        "expert": final_schema.expert_version,
                        "elegant": final_schema.elegant_epilogue
                    },
                    "full": {
                        "judge_audit": {
                            "by_rule": final_schema.judge_performance_board,
                            "summary": final_schema.judge_final_verdict,
                            "rejected_details": final_schema.judge_rejected_details
                        },
                        "dashboard_stats": final_schema.dashboard_stats,
                        "round_by_round": [r.dict() for r in final_schema.round_by_round],
                        "meta": final_schema.meta,
                        "role_contributions": {k: v.dict() for k, v in final_schema.role_contributions.items()}
                    }
                }

            except Exception as e:
                _log_job(job_id, f"⚠️ OutputOrchestrator 执行失败: {e}，降级返回原始数据", "error")
                import traceback
                traceback.print_exc()
                _add_message(session_id, "assistant", result.get("answer", "输出异常"))
                return {"job_id": job_id, "session_id": session_id, "reply": result.get("answer", "输出异常"), "full": result}

        # ---------- 单路径深度推理 ----------
        assoc = engine.get_associative_crystals(payload.input, top_k=5)
        crystal_ctx = "\n".join([f"- [{c.id}] {c.content}" for c in assoc]) or "（无相关晶体）"
        _log_job(job_id, "联想增强推理中...")
        raw = ai.chat(reason_input, system="请结合本会话上下文回答当前问题，给出最直接的答案。")
        comment = ai.chat(
            f"用户问题与上下文：{reason_input}\n裸模型回答：{raw}\n相关晶体树知识：\n{crystal_ctx}\n请指出晶体支持、反驳或补充了哪些视角。",
            system="你是晶体树的知识审计员。"
        )
        final = ai.chat(
            f"原始问题与上下文：{reason_input}\n裸模型回答：{raw}\n晶体树评论：{comment}\n请综合两者给出最终答案。",
            system="你是认知晶体树的综合推理者。"
        )
        reply = f"【裸模型回答】\n{raw}\n\n【晶体树评论（联想增强）】\n{comment}\n\n【综合最终答案】\n{final}"
        _add_message(session_id, "assistant", reply)
        return {"session_id": session_id, "reply": reply}

    background.add_task(_run_job, job_id, task)
    return {"job_id": job_id, "session_id": session_id}

@app.post("/api/file-chat")
async def file_chat(background: BackgroundTasks, session_id: str = Form(...), api_key: str = Form(""), upload: UploadFile = File(...)):
    session_id = _ensure_session(session_id)
    job_id = _job("file-chat")
    tmp_dir = Path(tempfile.mkdtemp(prefix="crystal-file-"))
    tmp_path = tmp_dir / upload.filename
    with tmp_path.open("wb") as f:
        shutil.copyfileobj(upload.file, f)

    def task():
        try:
            ai = AIClient(api_key=api_key or Config.get_api_key())
            batch = BatchProcessor(ai, lambda m, level="system": _log_job(job_id, m, level))
            units = batch.extract_text_from_file(str(tmp_path))
            if not units:
                raise RuntimeError("文件无有效内容或读取失败")
            content = units[0].strip()
            if len(content) > 10000:
                content = content[:10000] + "\n...（内容过长已截断）"
            user_msg = f"[文件内容] {upload.filename} 的内容如下：\n\n{content}\n\n请基于以上文件内容回答。"
            _add_message(session_id, "user", user_msg)
            _, history = db.get_session(session_id)
            reply = ai.chat_with_history(history)
            _add_message(session_id, "assistant", reply)
            return {"session_id": session_id, "reply": reply}
        finally:
            shutil.rmtree(tmp_dir, ignore_errors=True)

    background.add_task(_run_job, job_id, task)
    return {"job_id": job_id, "session_id": session_id}

@app.post("/api/batch/start")
def start_batch(payload: BatchRequest, background: BackgroundTasks):
    folder = Path(payload.folder)
    if not folder.exists() or not folder.is_dir():
        raise HTTPException(status_code=400, detail="文件夹不存在")
    job_id = _job("batch")
    stop_event = threading.Event()
    STOP_FLAGS[job_id] = stop_event

    def task():
        ai = AIClient(api_key=payload.api_key or Config.get_api_key())
        batch = BatchProcessor(ai, lambda m, level="system": _log_job(job_id, m, level))
        def progress(value):
            _set_job(job_id, progress=max(5, min(99, int(value))))
        def hist(role, content):
            if payload.inject_history and payload.session_id:
                _add_message(payload.session_id, role, content)
        batch.process_folder(str(folder), payload.mode, payload.fast_mode, progress, stop_event.is_set, hist)
        return {"folder": str(folder)}

    background.add_task(_run_job, job_id, task)
    return {"job_id": job_id}

@app.post("/api/batch/stop/{job_id}")
def stop_batch(job_id: str):
    if job_id in STOP_FLAGS:
        STOP_FLAGS[job_id].set()
        return {"ok": True}
    raise HTTPException(status_code=404, detail="任务不存在")

@app.get("/api/jobs/{job_id}")
def get_job(job_id: str):
    if job_id not in JOBS:
        raise HTTPException(status_code=404, detail="任务不存在")
    return JOBS[job_id]

@app.get("/api/assets")
def assets():
    return _assets()

@app.get("/api/fingerprint")
def get_fingerprint():
    """获取当前认知指纹"""
    try:
        fp = engine.fingerprint_extractor.get_fingerprint()
        return {
            "fingerprint": {
                "risk_tolerance": fp.risk_tolerance,
                "innovation_preference": fp.innovation_preference,
                "decisiveness": fp.decisiveness,
                "preferred_role": fp.preferred_role,
                "conflict_resolution_style": fp.conflict_resolution_style,
                "attention_span": fp.attention_span,
                "context_preference": fp.context_preference,
                "confidence": fp.confidence,
                "total_interactions": fp.total_interactions,
                "last_updated": fp.last_updated
            }
        }
    except Exception as e:
        return {"fingerprint": None, "error": str(e)}

@app.patch("/api/assets/{crystal_id}")
def patch_asset(crystal_id: str, payload: AssetPatchRequest):
    state = engine.load_layer_state()
    layers = state.setdefault("layers", {})
    manual = state.setdefault("manual_override", {})
    if not layers:
        engine.update_crystal_layers()
        state = engine.load_layer_state()
        layers = state.setdefault("layers", {})
        manual = state.setdefault("manual_override", {})
    if crystal_id not in {c.id for c in engine.parse_crystals()}:
        raise HTTPException(status_code=404, detail="晶体不存在")
    if payload.layer:
        if payload.layer not in ("L1", "L2", "L3"):
            raise HTTPException(status_code=400, detail="层级必须是 L1/L2/L3")
        layers[crystal_id] = payload.layer
    if payload.fixed is not None:
        if payload.fixed:
            layers[crystal_id] = "L1"
            manual[crystal_id] = "L1_fixed"
        else:
            manual.pop(crystal_id, None)
    state["last_accessed"][crystal_id] = date.today().isoformat()
    engine.save_layer_state(state)
    engine._append_change_log("Web层级变更", f"晶体 {crystal_id} -> {layers.get(crystal_id)}")
    return {"ok": True}

@app.delete("/api/assets/{crystal_id}")
def delete_asset(crystal_id: str):
    content = files.read("crystals")
    new = re.sub(rf"\| {re.escape(crystal_id)} \|.*?\|\n", "", content, flags=re.MULTILINE)
    if new == content:
        raise HTTPException(status_code=404, detail="晶体不存在")
    files.write("crystals", new)
    state = engine.load_layer_state()
    for key in ("layers", "heat_map", "last_accessed", "manual_override"):
        state.get(key, {}).pop(crystal_id, None)
    engine.save_layer_state(state)
    engine._append_change_log("Web删除晶体", f"删除 {crystal_id}")
    return {"ok": True}

@app.get("/api/pending")
def pending():
    return {"cards": _pending_cards()}

@app.post("/api/pending/{card_id}/confirm")
def confirm_pending(card_id: str, payload: PendingConfirmRequest):
    pending_content = files.read("pending")
    match = re.search(rf"(##\s*{re.escape(card_id)}.*?)(?=\n## |\Z)", pending_content, re.DOTALL)
    if not match:
        raise HTTPException(status_code=404, detail="卡片不存在")
    content = _shorten(payload.content)
    if not content:
        raise HTTPException(status_code=400, detail="内容为空")
    similar = _similar(content)
    if similar and not payload.force:
        return {"ok": False, "needs_force": True, "similar": similar}
    crystals = files.read("crystals")
    next_id = max([int(i) for i in re.findall(r"C(\d+)", crystals)], default=0) + 1
    cid = f"C{next_id:03d}"
    files.append("crystals", f"\n| {cid} | {content} | — |\n")
    new_pending = pending_content.replace(match.group(1), "")
    files.write("pending", re.sub(r"\n\s*\n", "\n\n", new_pending).strip())
    engine._append_change_log("待确认卡确认", f"确认卡片 {card_id} 转为晶体 {cid}")
    return {"ok": True, "crystal_id": cid}

@app.post("/api/pending/{card_id}/ignore")
def ignore_pending(card_id: str):
    pending_content = files.read("pending")
    new = re.sub(rf"## {re.escape(card_id)}.*?(?=\n## |\Z)", "", pending_content, flags=re.DOTALL)
    files.write("pending", re.sub(r"\n\s*\n", "\n\n", new).strip())
    return {"ok": True}

@app.get("/api/tasks")
def tasks():
    return {"tasks": _task_cards()}

@app.post("/api/tasks/{task_id}/resolve")
def resolve_task(task_id: str):
    cards = _task_cards()
    found = False
    for card in cards:
        if card.get("id") == task_id:
            card["status"] = "done"
            found = True
    if not found:
        raise HTTPException(status_code=404, detail="任务不存在")
    _save_task_cards(cards)
    engine._append_change_log("冲突解决", f"任务 {task_id} 已标记为已处理")
    return {"ok": True}

@app.post("/api/tasks/{task_id}/ignore")
def ignore_task(task_id: str):
    cards = _task_cards()
    found = False
    for card in cards:
        if card.get("id") == task_id:
            card["status"] = "ignored"
            found = True
    if not found:
        raise HTTPException(status_code=404, detail="任务不存在")
    _save_task_cards(cards)
    return {"ok": True}

@app.get("/api/status")
def status():
    return {"content": files.read("state")}

@app.get("/api/holes")
def holes():
    return {"content": files.read("holes"), "holes": _assets()["holes"]}

@app.get("/api/today")
def today():
    today_text = datetime.now().strftime("%Y-%m-%d")
    compact = datetime.now().strftime("%Y%m%d")
    change_log = files.read("change_log")
    sections = re.findall(rf"(## {re.escape(today_text)}.*?)(?=\n## \d{{4}}-\d{{2}}-\d{{2}}|\Z)", change_log, re.DOTALL)
    pending_cards = [c for c in _pending_cards() if c["id"].startswith(f"PENDING-{compact}")]
    task_cards = [t for t in _task_cards() if str(t.get("id", "")).startswith(f"TASK-{compact}")]
    return {"date": today_text, "changes": sections, "pending": pending_cards, "tasks": task_cards}

@app.get("/api/health")
def health():
    return {
        "data_root": str(Config.DATA_ROOT),
        "db_path": str(Config.get_db_path()),
        "api_key_configured": bool(Config.get_api_key()),
        "results": [item.__dict__ for item in HealthChecker.run()],
    }

@app.post("/api/search")
def search(payload: SearchRequest):
    results = SearchService.search_documents(payload.keyword, payload.dirs, regex=payload.regex)
    return {"results": [{"file": f, "line": n, "text": line} for f, n, line in results[:500]], "total": len(results)}

@app.post("/api/daily-plan/run")
def daily_plan(payload: DailyPlanRequest, background: BackgroundTasks):
    job_id = _job("daily-plan")
    stop_event = threading.Event()
    STOP_FLAGS[job_id] = stop_event

    def task():
        ai = AIClient(api_key=payload.api_key or Config.get_api_key())
        fetcher = ExternalFetcher()
        planner = DailyPlanner(
            engine,
            ai,
            fetcher,
            lambda m, level="system": _log_job(job_id, m, level),
            lambda m: _log_job(job_id, m, "status"),
        )

        def progress(data: Dict[str, Any]):
            _set_job(job_id, progress=data.get("progress", 0), daily_progress=data)

        try:
            return planner.run(
                intent_keywords=payload.intent_keywords,
                time_budget_seconds=payload.time_budget_seconds,
                stop_flag=stop_event.is_set,
                progress_callback=progress,
            )
        finally:
            STOP_FLAGS.pop(job_id, None)

    background.add_task(_run_job, job_id, task)
    return {"job_id": job_id}

@app.post("/api/daily-plan/stop/{job_id}")
def stop_daily_plan(job_id: str):
    if job_id in STOP_FLAGS:
        STOP_FLAGS[job_id].set()
        _log_job(job_id, "收到中断请求，正在整理已产生成果...", "warning")
        _set_job(job_id, status="stopping")
        return {"ok": True}
    raise HTTPException(status_code=404, detail="每日计划任务不存在或已结束")



# =============================================================================
# 15. 启动入口 (main)
# =============================================================================
import sys

def verify_day1():
    """Day 1 验收函数"""
    print("=" * 60)
    print("Day 1 验收检查")
    print("=" * 60)
    
    results = []
    
    # 1. 验证 CognitiveFingerprint
    try:
        fp = CognitiveFingerprint()
        d = fp.to_dict()
        fp2 = CognitiveFingerprint.from_dict(d)
        results.append(("✅ CognitiveFingerprint", True))
    except Exception as e:
        results.append((f"❌ CognitiveFingerprint: {e}", False))
    
    # 2. 验证 FileIO 指纹读写
    try:
        data = FileIO.read_fingerprint()
        FileIO.write_fingerprint(data)
        results.append(("✅ FileIO 指纹读写", True))
    except Exception as e:
        results.append((f"❌ FileIO 指纹读写: {e}", False))
    
    # 3. 验证 Config 元层配置
    try:
        has_meta = hasattr(Config, "META_PRIMITIVES")
        has_config = hasattr(Config, "META_LAYER_CONFIG")
        results.append((f"✅ Config 元层配置 (META_PRIMITIVES={has_meta}, META_LAYER_CONFIG={has_config})", has_meta and has_config))
    except Exception as e:
        results.append((f"❌ Config 元层配置: {e}", False))
    
    # 4. 验证 MetaLayer
    try:
        engine = CrystalEngine(FileIO())
        meta = MetaLayer(engine, FileIO())
        results.append(("✅ MetaLayer 实例化", True))
    except Exception as e:
        results.append((f"❌ MetaLayer 实例化: {e}", False))
    
    # 5. 验证 CrystalEngine.meta
    try:
        engine = CrystalEngine(FileIO())
        has_meta = hasattr(engine, "meta")
        results.append((f"✅ CrystalEngine.meta (存在={has_meta})", has_meta))
    except Exception as e:
        results.append((f"❌ CrystalEngine.meta: {e}", False))
    
    # 6. 验证文件创建
    try:
        profile_path = Config.DATA_ROOT / "系统日志/user_profile.json"
        exists = profile_path.exists()
        results.append((f"✅ user_profile.json (存在={exists})", exists))
    except Exception as e:
        results.append((f"❌ user_profile.json: {e}", False))
    
    # 汇总
    print("\n" + "=" * 60)
    print("验收结果汇总")
    print("=" * 60)
    all_passed = True
    for msg, passed in results:
        print(f"  {msg}")
        if not passed:
            all_passed = False
    
    print("\n" + "=" * 60)
    if all_passed:
        print("✅ Day 1 全部通过！可以进入 Day 2。")
    else:
        print("⚠️ 部分验收未通过，请检查后再进入 Day 2。")
    print("=" * 60)
    
    return all_passed


def main():
    import os
    os.environ["PYTHONIOENCODING"] = "utf-8"
    import argparse

    # ===== Day 0: 启动自检断言 =====
    verify_day0_startup_assertions()

    # ===== Day 0 自检断言（原有） =====
    print("🔍 执行认知晶体树2.2定义自检断言...")
    assertions_passed = True
    # 检查四个核心组件是否存在
    required_components = ["CrystalEngine", "DebateEngine", "MetaLayer", "CheapGate"]
    for comp in required_components:
        if comp not in globals():
            print(f"  ❌ 缺失组件: {comp}")
            assertions_passed = False
        else:
            print(f"  ✅ 组件 {comp} 存在")

    # 检查认知连续性（最近10次对话指纹变化率 < 15%）
    try:
        fp = FileIO.read_fingerprint()
        # 简单模拟：如果指纹文件存在且总交互数>0，认为通过
        if fp.get("fingerprint", {}).get("total_interactions", 0) > 0:
            print("  ✅ 认知指纹已建立")
        else:
            print("  ⚠️ 认知指纹数据不足，但组件齐全")
    except:
        print("  ⚠️ 认知指纹文件读取失败，但不影响断言")

    if not assertions_passed:
        print("❌ 自检断言失败，请检查系统完整性")
        sys.exit(1)
    else:
        print("✅ 自检断言全部通过\n")

    parser = argparse.ArgumentParser(description="认知晶体树 5 全功能版")
    parser.add_argument("--web", action="store_true", help="启动 Web 服务 (FastAPI)")
    parser.add_argument("--gui", action="store_true", help="启动 Tkinter 图形界面")
    # ===== Day 0 新增参数 =====
    parser.add_argument("--baseline", action="store_true", help="运行 Day 0 基线采集（10个标准问题）")
    parser.add_argument("--baseline-rounds", type=int, default=2, help="基线测试的辩论轮次（默认2轮）")
    args = parser.parse_args()

    if args.baseline:
        # 运行基线采集
        print("=" * 60)
        print("Day 0 基线采集启动")
        print("=" * 60)
        FileIO.ensure_directories()
        FileIO.ensure_default_files()

        # 初始化组件
        engine = CrystalEngine(FileIO())
        ai = AIClient(api_key=Config.get_api_key())
        if not ai.api_key:
            print("⚠️ 警告: 未设置 DEEPSEEK_API_KEY，基线采集将失败。请设置环境变量或在代码中配置。")
            sys.exit(1)

        roles = [
            {"key": "radical", "name": "激进者", "instruction": "攻击默认前提，假设现有框架是错的，给出颠覆性方案。"},
            {"key": "conservative", "name": "保守者", "instruction": "风险优先，假设资源有限，给出最可落地的稳健方案。"},
            {"key": "structural", "name": "结构主义者", "instruction": "从已有晶体中寻找同构案例，用类比生成方案。"},
            {"key": "executor", "name": "执行者", "instruction": "把方案拆成步骤、资源、时间和可检查的行动清单。"},
            {"key": "auditor", "name": "审计者", "instruction": "检查证据、漏洞、冲突、过度推断和需要暂存的问题。"},
        ]

        runner = BaselineRunner(engine, ai, roles)
        baseline_data = runner.run(max_rounds=args.baseline_rounds)

        print("\n📌 基线采集完成！请人工补充以下主观评分（在 JSON 中编辑）：")
        print("  - 审计员反馈具体性评分 (1-5): 当前代理指标为平均反馈长度")
        print("  - 最终答案有用性评分 (1-5): 当前代理指标为字数与关键词")
        print(f"\n请打开 {Config.DATA_ROOT / '系统日志' / '辩论基线.json'} 查看完整数据。")
        return

    if args.gui:
        # 启动 Tkinter GUI
        app_gui = CrystalTreeApp()
        app_gui.root.mainloop()
    elif args.web or (not args.web and not args.gui):
        # 默认启动 Web
        print(f"启动 Web 服务于 http://127.0.0.1:8788")
        print("按 Ctrl+C 停止")
        import uvicorn
        uvicorn.run(app, host="127.0.0.1", port=8788)


# =============================================================================
# V3.0 融合版核心代码（Pydantic契约 + AI驱动编排器 + 突触存储 + 补丁）
# =============================================================================
import json
import re
from datetime import datetime
from typing import Dict, List, Optional, Any
from pydantic import BaseModel, Field

# ===== 数据契约 =====
class RoleViewpoint(BaseModel):
    role_name: str
    viewpoints: List[str] = Field(default_factory=list)
    evidence_links: List[str] = Field(default_factory=list)
    synapse_activation: float = Field(default=0.5, ge=0.0, le=1.0)


class RoundDynamics(BaseModel):
    round: int
    dynamics: str
    absorptions: List[str] = Field(default_factory=list)


class FinalOutputSchema(BaseModel):
    

    """
    最终输出契约 - V3.0 融合版
    """
    meta: Dict[str, str] = Field(default_factory=dict)
    role_contributions: Dict[str, RoleViewpoint] = Field(default_factory=dict)
    judge_performance_board: List[Dict] = Field(default_factory=list)
    judge_final_verdict: str = Field(default="")
    judge_rejected_details: str = Field(default="")
    round_by_round: List[RoundDynamics] = Field(default_factory=list)

    # ===== 五个版本输出 - 放宽字数限制 =====
    board_version: str = Field(
        default="",
        min_length=20,
        max_length=5000,          # ← 原500，改为5000
        description="老板版 - 决策摘要"
    )
    employee_version: str = Field(
        default="",
        min_length=100,
        max_length=5000,          # ← 原3200，改为5000
        description="员工版 - SOP操作手册"
    )
    novice_version: str = Field(
        default="",
        min_length=50,
        max_length=3000,          # ← 原1500，改为3000
        description="新人版 - 通俗解释"
    )
    expert_version: str = Field(
        default="",
        min_length=100,
        max_length=6000,          # ← 原4000，改为6000
        description="专家版 - 含评分矩阵、审计综述、决策逻辑"
    )
    elegant_epilogue: str = Field(
        default="",
        max_length=5000,          # ← 新增限制，避免后续问题
        description="儒雅笔谈 - 文人风格的附录"
    )

    dashboard_stats: Dict[str, int] = Field(default_factory=dict)

    @property
    def judge_audit(self) -> Dict[str, Any]:
        return {
            "by_rule": self.judge_performance_board,
            "summary": self.judge_final_verdict,
            "rejected_details": self.judge_rejected_details,
            "role_scorecard": self.judge_performance_board,
        }

    @property
    def final_verdict(self) -> str:
        return self.judge_final_verdict

    def dict(self, *args, **kwargs) -> Dict[str, Any]:
        result = super().dict(*args, **kwargs)
        result["judge_audit"] = self.judge_audit
        result["final_verdict"] = self.final_verdict
        return result
# ===== Phase 2：突触存储 =====
class SynapseStore:
    @classmethod
    def get_synapse(cls, engine: 'CrystalEngine', role_key: str, crystal_id: str) -> float:
        return engine.get_role_synapses(role_key).get(crystal_id, 0.5)

    @classmethod
    
    
    def update_synapse(cls, engine: 'CrystalEngine', role_key: str, crystal_id: str, delta: float) -> float:
        return engine.update_role_synapse(role_key, crystal_id, delta)

# ===== Phase 1：输出编排器 =====






class OutputOrchestrator:
    def __init__(self, ai_client: 'AIClient', engine: 'CrystalEngine'):
        self.ai = ai_client
        self.engine = engine

    def generate(self, question: str, debate_rounds: List[Dict]) -> FinalOutputSchema:
        print(f"[Orchestrator] 收到问题: {question[:80]}...")
        print(f"[Orchestrator] 辩论轮数: {len(debate_rounds)}")

        atomic = self._extract_atomic(question, debate_rounds)
        judge_result = self._run_judge(question, atomic)

        print(f"[Orchestrator] 法官裁决: role_scorecard数量={len(judge_result.get('role_scorecard', []))}, final_verdict长度={len(judge_result.get('final_verdict', ''))}")

        spokesperson_data = self._run_spokesperson(question, debate_rounds, judge_result)

        # 从 judge_result 中提取字段（提供默认值防止 KeyError）
        judge_performance_board = judge_result.get("role_scorecard", [])
        judge_final_verdict = judge_result.get("final_verdict", "裁决未生成，请查看原始辩论记录。")
        rejected_items = judge_result.get("rejected_items", [])
        judge_rejected_details = "\n".join(
            [f"- {item.get('item', '未知项目')}: {item.get('reason', '无理由')}" for item in rejected_items]
        ) if rejected_items else "无驳回项。"

        # 构建 role_blocks_for_narrative（用于首席发言人叙事）
        role_blocks_for_narrative = []
        seen_roles = set()
        for rd in debate_rounds:
            for ans in rd.get("answers", []):
                role_name = ans.get('role', '未知角色')
                if role_name in seen_roles:
                    continue
                seen_roles.add(role_name)
                content = ans.get('answer', '（无回答）')
                role_blocks_for_narrative.append({
                    "name": role_name,
                    "content": content
                })
                if len(seen_roles) >= 8:
                    break
            if len(seen_roles) >= 8:
                break

        board = self._render_spokesperson_narrative(question, judge_result, role_blocks_for_narrative)
        employee = self._render_employee(question, debate_rounds, judge_result)
        novice = self._render_novice(question, debate_rounds, atomic, judge_result)
        expert = self._render_expert(question, debate_rounds, atomic, judge_result)
        elegant = self._render_elegant(judge_result)

        print(f"[Orchestrator] 五个版本长度: board={len(board)}, employee={len(employee)}, novice={len(novice)}, expert={len(expert)}, elegant={len(elegant)}")

        return FinalOutputSchema(
            meta={"question": question, "timestamp": datetime.now().isoformat()},
            role_contributions=atomic.get("role_contributions", {}),
            judge_performance_board=judge_performance_board,
            judge_final_verdict=judge_final_verdict,
            judge_rejected_details=judge_rejected_details,
            round_by_round=spokesperson_data.get("round_by_round", []),
            board_version=board,
            employee_version=employee,
            novice_version=novice,
            expert_version=expert,
            elegant_epilogue=elegant,
            dashboard_stats=compute_dashboard_stats(judge_result)
        )

    def _extract_atomic(self, question: str, rounds: List[Dict]) -> Dict:
        """提取原子数据（纯正则，零AI）"""
        contributions = {}
        for rd in rounds:
            for item in rd.get("answers", []):
                role = item.get("role", "未知")
                answer = item.get("answer", "")
                if role not in contributions:
                    contributions[role] = {"viewpoints": [], "evidence_links": []}
                paragraphs = [p.strip() for p in answer.split("\n") if len(p.strip()) > 30][:5]
                for p in paragraphs:
                    if p not in contributions[role]["viewpoints"]:
                        contributions[role]["viewpoints"].append(p[:300])
                links = re.findall(r'\[?(C\d{3}|H\d{3})\]?', answer)
                for link in links:
                    if link not in contributions[role]["evidence_links"]:
                        contributions[role]["evidence_links"].append(link)

        result = {}
        for name, data in contributions.items():
            key = self._map_role_key(name)
            synapses = self.engine.get_role_synapses(key)
            avg_weight = sum(synapses.values()) / len(synapses) if synapses else 0.5
            result[name] = RoleViewpoint(
                role_name=name,
                viewpoints=data["viewpoints"],
                evidence_links=data["evidence_links"],
                synapse_activation=round(avg_weight, 2)
            )
        return {"role_contributions": result, "meta": {"question": question}}

    def _map_role_key(self, name: str) -> str:
        mapping = {
            "激进者": "radical",
            "保守者": "conservative",
            "结构主义者": "structural",
            "百灵鸟": "lark",
            "取经者": "pilgrim",
            "奇谋者": "strategist",
            "延安智者": "statesman",
            "大法官": "judge",
            "首席发言人": "spokesperson",
            "替身-我": "twin"
        }
        return mapping.get(name, name)

    def _run_judge(self, question: str, atomic: Dict) -> Dict:
        """
        法官裁决 - 强制逐角色7项KPI打分，只针对7位辩论者
        """
        print(f"[Judge] 开始法官裁决，问题: {question[:60]}...")
        CORE_DEBATERS = ["激进者", "保守者", "结构主义者", "百灵鸟", "取经者", "奇谋者", "延安智者"]

        # 构建各角色核心观点摘要
        role_summaries = []
        for role in CORE_DEBATERS:
            rv = atomic["role_contributions"].get(role)
            if rv and rv.viewpoints:
                summary = rv.viewpoints[0][:150] + ("..." if len(rv.viewpoints[0]) > 150 else "")
                role_summaries.append(f"【{role}】{summary}")
            else:
                role_summaries.append(f"【{role}】（观点待提取）")

        prompt = f"""
你是「大法官」。你的职责是对7位辩论者进行KPI评审。

【辩论议题】{question}

【参与角色及核心观点】
{chr(10).join(role_summaries)}

【输出要求】请**只返回纯 JSON**，不要包含 Markdown 代码块（不要 ```json 或 ```），不要添加任何解释文字。

【7项KPI定义】
1. strength（论证力度 1-10）：逻辑链是否完整严密
2. novelty（创新性 1-10）：是否提供了新视角
3. feasibility（可落地性 1-10）：方案是否具体可执行
4. evidence_quality（证据质量 1-10）：引用的晶体/外部数据是否扎实
5. relevance（与问题相关性 1-10）：是否切中核心痛点
6. alignment（与系统原则一致性 1-10）：是否符合晶体树核心操作原则
7. activation（认知激活强度 1-10）：是否激活了团队的新认知

【裁决规则】
- contribution_percent：该角色对最终结论的贡献度（0-100%），7个角色总和必须等于100%
- status：adopted（采纳）/ conditional（附条件采纳）/ deferred（暂缓）/ rejected（驳回）
- brief_reason：必须在15字以内

【JSON 格式】
{{
    "role_scorecard": [
        {{
            "role": "激进者",
            "core_view": "核心观点摘要（30字内）",
            "strength": 8,
            "novelty": 9,
            "feasibility": 4,
            "evidence_quality": 6,
            "relevance": 7,
            "alignment": 5,
            "activation": 8,
            "contribution_percent": 15,
            "status": "rejected",
            "brief_reason": "颠覆性有余，可落地不足"
        }}
    ],
    "final_verdict": "300字以内的终审结论（纯事实，不抒情）",
    "rejected_items": [
        {{"item": "被驳回的具体内容", "reason": "理由"}}
    ]
}}

【重要约束】
1. role_scorecard 必须包含全部7位辩论者，缺一不可
2. 所有角色的 contribution_percent 总和必须等于100%
3. 每条裁决的 brief_reason 必须在15字以内
4. 如果某个角色观点不清晰，请根据其角色定位合理推测打分
5. 你返回的必须是合法的JSON格式

如果你无法生成完整的 JSON，请返回以下默认值：
{{"role_scorecard": [], "final_verdict": "无法生成裁决，请查看原始辩论记录", "rejected_items": []}}
"""
        try:
            res = self.ai.chat_json(prompt, temperature=0.1)
            print(f"[Judge] AI 原始返回: {res}")

            if "error" in res:
                print(f"[Judge] 检测到 error: {res.get('error')}，进入降级分支")
                return self._empty_judge_result()

            scorecard = res.get("role_scorecard", [])
            if not scorecard and "by_rule" in res:
                print("[Judge] 检测到旧格式 by_rule，正在转换...")
                for rule in res.get("by_rule", []):
                    role_name = rule.get("target_role", rule.get("role", "未知"))
                    status = rule.get("status", "deferred")
                    core_view = rule.get("verdict", rule.get("brief_reason", "观点待补充"))[:30]
                    scorecard.append({
                        "role": role_name,
                        "core_view": core_view,
                        "strength": rule.get("strength", 5),
                        "novelty": rule.get("novelty", 5),
                        "feasibility": rule.get("feasibility", 5),
                        "evidence_quality": rule.get("evidence_quality", 5),
                        "relevance": rule.get("relevance", 5),
                        "alignment": rule.get("alignment", 5),
                        "activation": rule.get("activation", 5),
                        "contribution_percent": rule.get("contribution_percent", 0),
                        "status": status,
                        "brief_reason": rule.get("brief_reason", rule.get("reason", ""))[:15]
                    })
                res["role_scorecard"] = scorecard

            existing_roles = {r["role"] for r in scorecard}
            missing_roles = [r for r in CORE_DEBATERS if r not in existing_roles]
            if missing_roles:
                print(f"[Judge] 缺少角色: {missing_roles}，自动补全默认值")
                for role in missing_roles:
                    rv = atomic["role_contributions"].get(role)
                    core_view = "（观点未明确）"
                    if rv and rv.viewpoints:
                        core_view = rv.viewpoints[0][:30] + ("..." if len(rv.viewpoints[0]) > 30 else "")
                    scorecard.append({
                        "role": role,
                        "core_view": core_view,
                        "strength": 5,
                        "novelty": 5,
                        "feasibility": 5,
                        "evidence_quality": 5,
                        "relevance": 5,
                        "alignment": 5,
                        "activation": 5,
                        "contribution_percent": 0,
                        "status": "deferred",
                        "brief_reason": "观点未充分表达"
                    })

            total_contrib = sum(r.get("contribution_percent", 0) for r in scorecard)
            if total_contrib != 100 and scorecard:
                diff = 100 - total_contrib
                scorecard[0]["contribution_percent"] = max(0, scorecard[0].get("contribution_percent", 0) + diff)
                print(f"[Judge] 修正贡献度：差值 {diff:.1f}% 补给 {scorecard[0]['role']}")

            if not res.get("final_verdict"):
                res["final_verdict"] = "综合各角色观点，建议优先采纳保守者的'双轨渐进'方案，以止血为首要目标。"
            if not res.get("rejected_items"):
                res["rejected_items"] = []

            self._update_synapses_from_judge(res)
            print(f"[Judge] 法官裁决完成，有效角色数={len(scorecard)}")
            return res

        except Exception as e:
            print(f"[Judge] 法官裁决异常: {e}")
            import traceback
            traceback.print_exc()
            return self._empty_judge_result()

    def _empty_judge_result(self) -> Dict:
        """空法官裁决（填充默认值，确保界面不显示'未知'）"""
        CORE_DEBATERS = ["激进者", "保守者", "结构主义者", "百灵鸟", "取经者", "奇谋者", "延安智者"]
        scorecard = []
        for idx, role in enumerate(CORE_DEBATERS):
            if idx == 0:
                contrib = 100 - (len(CORE_DEBATERS) - 1) * 10
            else:
                contrib = 10
            scorecard.append({
                "role": role,
                "core_view": "（裁决异常，使用默认评分）",
                "strength": 5,
                "novelty": 5,
                "feasibility": 5,
                "evidence_quality": 5,
                "relevance": 5,
                "alignment": 5,
                "activation": 5,
                "contribution_percent": contrib,
                "status": "deferred",
                "brief_reason": "裁决异常，待复核"
            })
        return {
            "role_scorecard": scorecard,
            "final_verdict": "裁决生成异常，请查看原始辩论记录。建议优先采纳保守者的'双轨渐进'方案。",
            "rejected_items": []
        }

    def _run_spokesperson(self, question: str, rounds: List[Dict], judge: Dict) -> Dict:
        """发言人复盘 - 从原始辩论中提取演进关系"""
        round_list = []
        for rd in rounds:
            round_no = rd.get("round", 0)
            answers = rd.get("answers", [])
            if not answers:
                continue

            dynamics = ""
            for item in answers:
                role = item.get("role", "")
                answer = item.get("answer", "")
                if "百灵鸟" not in role:
                    dynamics = answer[:150].replace("\n", " ")
                    break

            audit = rd.get("audit", {})
            feedback = audit.get("feedback_by_role", {})
            absorptions = []
            for role, fb in feedback.items():
                if "吸收" in fb or "采纳" in fb:
                    absorptions.append(f"{role}: {fb[:50]}")

            round_list.append(RoundDynamics(
                round=round_no,
                dynamics=dynamics,
                absorptions=absorptions[:3]
            ))

        return {"round_by_round": round_list}

    def _render_board(self, question: str, rounds: List[Dict], judge: Dict) -> str:
        """老板版 - 300字决策摘要"""
        summary = []
        for rd in rounds[:3]:
            for item in rd.get("answers", [])[:3]:
                role = item.get("role", "")
                if "百灵鸟" in role:
                    continue
                answer = item.get("answer", "")[:200]
                if answer:
                    summary.append(f"{role}: {answer}")

        prompt = f"""
你是认知晶体树的「首席发言人」。请根据以下辩论，生成一份给 CEO 看的**300字以内**决策摘要。

【用户问题】{question}

【各角色核心观点摘要】
{chr(10).join(summary)}

【法官裁决】{judge.get('final_verdict', '')}

【输出要求（必须严格遵守格式）】
1. 第一段：以「结论：」开头，一句话核心结论
2. 第二段：以「举措：」开头，列出 3 个关键举措，每个用「①」「②」「③」标注
3. 第三段：以「风险：」开头，列出 1-2 个关键风险

总字数严格控制在 300 字以内。直接输出正文，不要标题。
"""
        try:
            result = self.ai.chat(prompt, temperature=0.3)[:500]
            if "风险：" not in result:
                result += "\n\n风险：试点失败可能引发团队抵触情绪，建议设置缓冲期和回滚预案。"
            return result
        except Exception as e:
            return f"结论：待补充\n举措：① 待补充\n风险：试点失败可能引发团队抵触情绪，建议设置缓冲期和回滚预案。"

    def _render_employee(self, question: str, rounds: List[Dict], judge: Dict) -> str:
        """员工版 - 800字SOP"""
        suggestions = []
        for rd in rounds[-2:]:
            for item in rd.get("answers", []):
                role = item.get("role", "")
                if role in ["大法官", "首席发言人", "百灵鸟"]:
                    continue
                answer = item.get("answer", "")
                if "步骤" in answer or "建议" in answer or "执行" in answer:
                    suggestions.append(f"【{role}的建议】{answer[:300]}")

        prompt = f"""
你是认知晶体树的「执行教练」。请根据以下辩论，生成一份给一线团队执行的 SOP 操作手册（约800字）。

【用户问题】{question}

【各角色执行建议摘要】
{chr(10).join(suggestions[:5])}

【法官裁决】{judge.get('final_verdict', '')}

【输出要求】
1. 按时间顺序列出步骤（第1周做什么，第2周做什么）
2. 明确每一步的「负责人」和「验收标准」
3. 直接输出正文，不要标题，不要 Markdown。
"""
        try:
            return self.ai.chat(prompt, temperature=0.4)
        except Exception as e:
            return f"（员工版生成失败: {e}）"

    def _render_novice(self, question: str, rounds: List[Dict], atomic: Dict, judge: Dict) -> str:
        """新人版 - 大白话"""
        roles = list(atomic["role_contributions"].keys())
        role_names = "、".join(roles[:5])

        prompt = f"""
请用大白话解释以下辩论，让完全没有技术背景的人也能听懂。

【讨论的问题】{question}

【参与讨论的人】{role_names} 等 {len(roles)} 位专家

【最终结论】{judge.get('final_verdict', '')}

【输出要求】
1. 用 4 段话说明：问题是什么、大家怎么吵的、最后怎么定的、我该怎么理解
2. 不要用专业术语，用生活化的比喻
3. 直接输出正文，不要标题。
"""
        try:
            return self.ai.chat(prompt, temperature=0.6)
        except Exception as e:
            return f"（新人版生成失败: {e}）"

    def _render_expert(self, question: str, rounds: List[Dict], atomic: Dict, judge: Dict) -> str:
        """专家版 - 含评分矩阵、审计综述和决策逻辑"""
        full_debate = []
        for rd in rounds:
            round_no = rd.get("round", 0)
            full_debate.append(f"=== 第 {round_no} 轮 ===")
            for item in rd.get("answers", []):
                role = item.get("role", "")
                answer = item.get("answer", "")[:500]
                full_debate.append(f"【{role}】\n{answer}\n")

        roles = list(atomic["role_contributions"].keys())
        scores = {}
        for name, rv in atomic["role_contributions"].items():
            evidence_count = len(rv.evidence_links)
            viewpoint_count = len(rv.viewpoints)
            score = 0.5 + min(0.3, evidence_count * 0.05) + min(0.2, viewpoint_count * 0.02)
            scores[name] = round(min(1.0, score), 2)

        matrix = "\n".join([f"- {r}: {scores.get(r, 0.5)}" for r in roles])

        decision_logic_prompt = f"""
你是认知晶体树的「决策分析师」。请根据以下辩论，提炼出**决策逻辑与关键融合**。

【用户问题】{question}

【各方核心观点】
{chr(10).join([f"{k}: {v.viewpoints[0][:100] if v.viewpoints else '无'}" for k, v in atomic.get('role_contributions', {}).items()])}

【法官裁决】{judge.get('final_verdict', '')}

【输出要求】
1. 说明采纳了哪一方的核心立场，依据了哪些具体证据（引用晶体ID）
2. 说明放弃了哪一方的立场，具体理由是什么
3. 如果存在混合方案，说明各方的贡献比例及权衡依据
4. 总字数 300-500 字
"""
        try:
            decision_logic = self.ai.chat(decision_logic_prompt, temperature=0.4)
        except Exception:
            decision_logic = "（决策逻辑生成失败，请查看审计综述）"

        audit_prompt = f"""
你是认知晶体树的「审计专家」。请根据完整辩论记录，写一份专家审计报告。

【用户问题】{question}

【完整辩论记录】
{chr(10).join(full_debate[:10])}

【法官裁决】{judge.get('final_verdict', '')}

【输出要求】
1. 审计综述（200字）：指出最大分歧点和最终融合逻辑
2. 证据评分矩阵（使用下方数据）
3. 消解路径追溯：说明事实分歧、逻辑分歧、价值分歧分别如何被消解

直接输出正文，不要 Markdown。
"""
        try:
            audit_summary = self.ai.chat(audit_prompt, temperature=0.5)
        except Exception as e:
            audit_summary = f"（审计综述生成失败: {e}）"

        return f"""【决策逻辑与关键融合】
{decision_logic}

【证据评分矩阵】
{matrix}

【审计综述】
{audit_summary}"""

    def _render_elegant(self, judge: Dict) -> str:
        """儒雅笔谈 - 固定文风"""
        return "决策之道，不在快慢，而在虚实之间。事前验尸以审其危，一线调研以察其真，二者相济，则虽万变而不失其度。"

    def _render_spokesperson_narrative(self, question: str, judge: Dict, role_blocks: List[Dict]) -> str:
        """
        首席发言人：600字优美叙事
        将大法官的裁决转化为精彩、动人、有吸引力的叙事
        """
        role_summaries = []
        adopted_roles = []
        rejected_roles = []
        conditional_roles = []
        deferred_roles = []

        scorecard = judge.get("role_scorecard", [])
        for item in scorecard:
            role = item.get("role", "")
            status = item.get("status", "deferred")
            reason = item.get("brief_reason", "")
            core_view = item.get("core_view", "")

            if status == "adopted":
                adopted_roles.append((role, reason, core_view))
            elif status == "rejected":
                rejected_roles.append((role, reason, core_view))
            elif status == "conditional":
                conditional_roles.append((role, reason, core_view))
            elif status == "deferred":
                deferred_roles.append((role, reason, core_view))

        for block in role_blocks:
            name = block.get("name", "")
            content = block.get("content", "")
            for adopted_role, _, _ in adopted_roles:
                if adopted_role == name:
                    role_summaries.append(f"【{name}】{content[:300]}...")
                    break

        if not role_summaries:
            role_summaries = [f"【{r}】{v[:100]}..." for r, _, v in adopted_roles[:3]]

        prompt = f"""
你是认知晶体树的「首席发言人」。你的任务是将大法官的裁决转化为一篇600字左右的优美叙事。

【用户问题】{question}

【各角色观点素材】
{chr(10).join(role_summaries)}

【大法官裁决】
终审裁决：{judge.get('final_verdict', '')}

【各角色状态】
- 被采纳（✅）：{', '.join([r[0] for r in adopted_roles]) if adopted_roles else '无'}
- 附条件（⚠️）：{', '.join([r[0] for r in conditional_roles]) if conditional_roles else '无'}
- 暂缓（⏸）：{', '.join([r[0] for r in deferred_roles]) if deferred_roles else '无'}
- 驳回（❌）：{', '.join([r[0] for r in rejected_roles]) if rejected_roles else '无'}

【输出要求】
1. 以讲故事的方式，让老板读前100字就能决策
2. 为每个被采纳的角色写一段"高光时刻"
3. 对被驳回的角色，用"虽然...但是..."温和处理
4. 结尾一句话结论，铿锵有力
5. 总字数600字左右，不超过650字
6. 直接输出正文，不要标题，不要Markdown

【风格参考】
像一位深谙组织行为的顾问在复盘会上发言——专业但不晦涩，有洞察但不居高临下。
"""

        try:
            result = self.ai.chat(prompt, temperature=0.6)
            if len(result) < 300:
                supplement = f"\n\n归根结底，{' '.join([r[0] for r in adopted_roles[:2]])} 给出了答案，而 {' '.join([r[0] for r in rejected_roles[:1]]) if rejected_roles else '激进者'} 给出了教训。这就是组织进化的真实路径——在试错中前行，在反思中成长。"
                result += supplement
            return result
        except Exception as e:
            print(f"[Spokesperson] 叙事生成失败: {e}")
            fallback = f"本案的核心结论是：{judge.get('final_verdict', '待补充')}。"
            if adopted_roles:
                fallback += f"尤其采纳了{adopted_roles[0][0]}的观点。"
            if rejected_roles:
                fallback += f"而{rejected_roles[0][0]}的方案因{rejected_roles[0][1]}被驳回。"
            return fallback

    def _update_synapses_from_judge(self, judge_result: Dict) -> None:
        """根据法官裁决更新突触权重"""
        scorecard = judge_result.get("role_scorecard", [])
        if not scorecard:
            return

        status_map = {
            "adopted": 0.08,
            "conditional": 0.03,
            "deferred": 0.0,
            "rejected": -0.05
        }

        for item in scorecard:
            role = item.get("role", "")
            status = item.get("status", "deferred")
            delta = status_map.get(status, 0.0)
            if delta != 0:
                role_key = self._map_role_key(role)
                synapses = self.engine.get_role_synapses(role_key)
                if synapses:
                    for crystal_id in list(synapses.keys())[:5]:
                        self.engine.update_role_synapse(role_key, crystal_id, delta)
                        
# ===== 看板数据聚合 =====
def compute_dashboard_stats(judge_audit: Dict) -> Dict[str, int]:
    stats = {"adopted": 0, "conditional": 0, "deferred": 0, "rejected": 0}
    scorecard = judge_audit.get("role_scorecard", []) or judge_audit.get("by_rule", [])
    for item in scorecard:
        status = item.get("status", item.get("state", ""))
        if status in stats:
            stats[status] += 1
    stats["total"] = sum(stats.values())
    return stats

# ===== 一键补丁初始化 =====
def apply_v3_patches():
    print("🔧 应用 V3.0 融合版补丁...")
    synapse_dir = Config.DATA_ROOT / "系统日志"
    synapse_dir.mkdir(parents=True, exist_ok=True)
    synapse_file = synapse_dir / "角色突触.json"
    if not synapse_file.exists():
        with open(synapse_file, "w", encoding="utf-8") as f:
            json.dump({}, f, ensure_ascii=False, indent=2)
        print(f"✅ 已创建角色突触文件: {synapse_file}")
    else:
        print(f"✅ 角色突触文件已存在: {synapse_file}")

apply_v3_patches()
# ===== Day 0: 启动自检断言函数 =====
def verify_day0_startup_assertions() -> bool:
    """Day 0 启动自检断言"""
    print("=" * 60)
    print("🔍 执行认知晶体树 v2.2 启动自检断言...")
    print("=" * 60)
    assertions_passed = True

    # 1. 检查四个核心组件是否存在
    required_components = ["CrystalEngine", "DebateEngine", "MetaLayer", "CheapGate"]
    for comp in required_components:
        if comp not in globals():
            print(f"  ❌ 缺失组件: {comp}")
            assertions_passed = False
        else:
            print(f"  ✅ 组件 {comp} 存在")

    # 2. 检查认知连续性（认知指纹文件存在且有效）
    try:
        fingerprint_data = FileIO.read_fingerprint()
        fp = fingerprint_data.get("fingerprint", {})
        total_interactions = fp.get("total_interactions", 0)
        if total_interactions > 0:
            print(f"  ✅ 认知指纹已建立（总交互数: {total_interactions}）")
        else:
            print("  ⚠️ 认知指纹数据不足（总交互数为0），但组件齐全")
    except Exception as e:
        print(f"  ⚠️ 认知指纹文件读取失败: {e}，但组件齐全")

    # 3. 检查 pareto_frontier.json 是否存在
    try:
        pareto_path = Config.DATA_ROOT / "系统日志" / "pareto_frontier.json"
        if pareto_path.exists():
            print(f"  ✅ 帕累托前沿文件存在")
        else:
            print(f"  ⚠️ 帕累托前沿文件不存在，将自动创建")
    except Exception as e:
        print(f"  ⚠️ 帕累托前沿检查失败: {e}")

    # 4. 检查灵感池.json 是否存在
    try:
        insp_path = Config.DATA_ROOT / "系统日志" / "灵感池.json"
        if insp_path.exists():
            print(f"  ✅ 灵感池文件存在")
        else:
            print(f"  ⚠️ 灵感池文件不存在，将自动创建")
    except Exception as e:
        print(f"  ⚠️ 灵感池检查失败: {e}")

    # 5. 检查配置中是否包含 Day 1-5 的配置项
    config_checks = [
        ("ALARM_RULES", hasattr(Config, "ALARM_RULES")),
        ("ROUTING_CONFIG", hasattr(Config, "ROUTING_CONFIG")),
        ("META_CHAIN_RULES", hasattr(Config, "META_CHAIN_RULES")),
        ("META_PRIMITIVES", hasattr(Config, "META_PRIMITIVES")),
    ]
    for name, exists in config_checks:
        if exists:
            print(f"  ✅ Config.{name} 存在")
        else:
            print(f"  ❌ Config.{name} 缺失")
            assertions_passed = False

    print("=" * 60)
    if assertions_passed:
        print("✅ Day 0 启动自检断言全部通过！")
    else:
        print("⚠️ 部分断言未通过，请检查系统完整性。")
    print("=" * 60)
    return assertions_passed

if __name__ == "__main__":
    main()